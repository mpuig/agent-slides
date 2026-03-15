from __future__ import annotations

import io
import re
import time
from collections.abc import Sequence
from contextlib import suppress
from copy import deepcopy
from dataclasses import asdict, dataclass
from hashlib import sha256
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation as load_presentation
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_AXIS_CROSSES,
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
    XL_TICK_LABEL_POSITION,
    XL_TICK_MARK,
)
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml.xmlchemy import OxmlElement
from pptx.presentation import Presentation as PptxPresentation
from pptx.util import Inches, Pt

from .errors import SlidesError
from .io import canonicalize_pptx_bytes, write_bytes
from .model import (
    AddAreaChartOp,
    AddBarChartOp,
    AddBubbleChartOp,
    AddComboChartOverlayOp,
    AddDoughnutChartOp,
    AddIconOp,
    AddImageOp,
    AddLineChartOp,
    AddLineShapeOp,
    AddMediaOp,
    AddNotesOp,
    AddOvalOp,
    AddPieChartOp,
    AddRadarChartOp,
    AddRawShapeXmlOp,
    AddRectangleOp,
    AddRoundedRectangleOp,
    AddScatterChartOp,
    AddSlideOp,
    AddTableOp,
    AddTextOp,
    DeleteSlideOp,
    MoveSlideOp,
    Operation,
    OperationBatch,
    OperationEvent,
    OperationReport,
    ReplaceTextOp,
    SetChartAxisOptionsOp,
    SetChartAxisScaleOp,
    SetChartAxisTitlesOp,
    SetChartComboSecondaryMappingOp,
    SetChartDataLabelsOp,
    SetChartDataLabelsStyleOp,
    SetChartLegendOp,
    SetChartPlotStyleOp,
    SetChartSecondaryAxisOp,
    SetChartSeriesErrorBarsOp,
    SetChartSeriesLineStyleOp,
    SetChartSeriesStyleOp,
    SetChartSeriesTrendlineOp,
    SetChartStyleOp,
    SetChartTitleOp,
    SetCorePropertiesOp,
    SetImageCropOp,
    SetLineSeriesMarkerOp,
    SetLineSeriesSmoothOp,
    SetPlaceholderImageOp,
    SetPlaceholderTextOp,
    SetSemanticTextOp,
    SetSlideBackgroundOp,
    SetTitleSubtitleOp,
    UpdateChartDataOp,
    UpdateTableCellOp,
)
from .validator import (
    ValidationIssue,
    ValidationReport,
    validate_package_bytes,
    validate_presentation,
)


@dataclass(slots=True)
class SlideRef:
    index: int


@dataclass(slots=True)
class DeckSummary:
    slide_count: int
    shape_count: int
    text_shape_count: int
    chart_count: int
    table_count: int
    image_count: int
    unresolved_token_count: int

    def to_dict(self) -> dict[str, int]:
        return asdict(self)


class Presentation:
    """Agent-friendly facade around python-pptx with deterministic serialization hooks."""

    _DEFAULT_ICON_DIR = Path(__file__).parent / "assets" / "icons"

    def __init__(self, prs: PptxPresentation, *, icon_dirs: Sequence[str | Path] | None = None):
        self._prs = prs
        self._icon_dirs: list[Path] = []
        self._icon_catalog_cache: dict[str, Path] | None = None
        for raw in icon_dirs or []:
            path = Path(raw).expanduser()
            if path not in self._icon_dirs:
                self._icon_dirs.append(path)
        if self._DEFAULT_ICON_DIR not in self._icon_dirs:
            self._icon_dirs.append(self._DEFAULT_ICON_DIR)

    @classmethod
    def create(
        cls,
        template_path: str | Path | None = None,
        *,
        icon_dirs: Sequence[str | Path] | None = None,
    ) -> Presentation:
        if template_path is None:
            return cls(load_presentation(), icon_dirs=icon_dirs)
        return cls(load_presentation(str(template_path)), icon_dirs=icon_dirs)

    @classmethod
    def open(
        cls, path: str | Path, *, icon_dirs: Sequence[str | Path] | None = None
    ) -> Presentation:
        return cls(load_presentation(str(path)), icon_dirs=icon_dirs)

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    def add_slide(
        self,
        layout_index: int | None = 6,
        layout_name: str | None = None,
        hidden: bool = False,
    ) -> SlideRef:
        if layout_name is not None:
            matching = [
                i
                for i, layout in enumerate(self._prs.slide_layouts)
                if (layout.name or "").strip() == layout_name.strip()
            ]
            if matching:
                layout_index = matching[0]
            elif layout_index is None:
                raise SlidesError(
                    code="INVALID_LAYOUT_NAME",
                    message=f"layout_name '{layout_name}' not found",
                    path="add_slide.layout_name",
                    suggested_fix="Use a layout name from the template slide layouts",
                )

        if layout_index is None:
            layout_index = 6
        if layout_index < 0 or layout_index >= len(self._prs.slide_layouts):
            raise SlidesError(
                code="INVALID_LAYOUT_INDEX",
                message=f"layout_index {layout_index} out of range",
                path="add_slide.layout_index",
                suggested_fix=f"use index in [0, {len(self._prs.slide_layouts) - 1}]",
            )
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[layout_index])
        self._set_slide_meta_uid(slide, f"s{len(self._prs.slides):08d}")
        if hidden:
            slide._element.set("show", "0")
        return SlideRef(index=len(self._prs.slides) - 1)

    def clear_slides(self) -> None:
        for idx in range(len(self._prs.slides) - 1, -1, -1):
            self.delete_slide(idx)

    def add_text(
        self,
        *,
        slide_index: int,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 20,
        bold: bool = False,
        font_name: str | None = None,
        font_color: str | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        # Set font properties on every run so they render reliably
        color = RGBColor.from_string(font_color) if font_color else None
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if font_name is not None:
                run.font.name = font_name
            if color is not None:
                run.font.color.rgb = color
        self._ensure_autofit(shape)

    def _apply_shape_border(
        self,
        shape: Any,
        border_color: str | None,
        border_width: float | None,
    ) -> None:
        ln = shape.line
        if border_color:
            ln.color.rgb = RGBColor.from_string(border_color)
            ln.width = Pt(border_width if border_width is not None else 1.0)
        else:
            ln.fill.background()

    def add_rectangle(
        self,
        *,
        slide_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: str,
        border_color: str | None = None,
        border_width: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
        self._apply_shape_border(shape, border_color, border_width)

    def add_rounded_rectangle(
        self,
        *,
        slide_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: str,
        corner_radius: int = 5000,
        border_color: str | None = None,
        border_width: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
        # Set corner radius via XML attribute
        sp_elem = shape._element
        prst_geom = sp_elem.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom"
        )
        if prst_geom is not None:
            av_lst = prst_geom.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst"
            )
            if av_lst is None:
                av_lst = OxmlElement("a:avLst")
                prst_geom.append(av_lst)
            else:
                av_lst.clear()
            gd = OxmlElement("a:gd")
            gd.set("name", "adj")
            gd.set("fmla", f"val {corner_radius}")
            av_lst.append(gd)
        self._apply_shape_border(shape, border_color, border_width)

    def add_oval(
        self,
        *,
        slide_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: str,
        border_color: str | None = None,
        border_width: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color)
        self._apply_shape_border(shape, border_color, border_width)

    def add_line_shape(
        self,
        *,
        slide_index: int,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: str = "000000",
        line_width: float = 1.0,
    ) -> None:
        slide = self._slide(slide_index)
        connector = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR_TYPE.STRAIGHT
            Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        connector.line.color.rgb = RGBColor.from_string(color)
        connector.line.width = Pt(line_width)

    @staticmethod
    def _resolve_chart_type(
        style: str,
        style_map: dict[str, int],
        chart_kind: str,
    ) -> int:
        chart_type = style_map.get(style)
        if chart_type is None:
            valid = ", ".join(style_map)
            raise SlidesError(
                code=f"INVALID_{chart_kind.upper()}_STYLE",
                message=f"Unsupported {chart_kind} style: {style}",
                path=f"add_{chart_kind}_chart.style",
                suggested_fix=f"Use one of: {valid}",
            )
        return chart_type

    def add_bar_chart(
        self,
        *,
        slide_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        style: str = "clustered",
        orientation: str = "column",
        chart_space_xml: str | None = None,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        if orientation not in {"column", "bar"}:
            raise SlidesError(
                code="INVALID_BAR_ORIENTATION",
                message=f"Unsupported bar orientation: {orientation}",
                path="add_bar_chart.orientation",
                suggested_fix="Use one of: column, bar",
            )
        if orientation == "bar":
            style_map = {
                "clustered": XL_CHART_TYPE.BAR_CLUSTERED,
                "stacked": XL_CHART_TYPE.BAR_STACKED,
                "percent_stacked": XL_CHART_TYPE.BAR_STACKED_100,
            }
        else:
            style_map = {
                "clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "stacked": XL_CHART_TYPE.COLUMN_STACKED,
                "percent_stacked": XL_CHART_TYPE.COLUMN_STACKED_100,
            }
        chart_type = self._resolve_chart_type(style, style_map, "bar")
        self._add_chart(
            slide_index=slide_index,
            chart_type=chart_type,
            categories=categories,
            series=series,
            chart_space_xml=chart_space_xml,
            left=left,
            top=top,
            width=width,
            height=height,
        )

    def add_line_chart(
        self,
        *,
        slide_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        style: str = "line",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "stacked": XL_CHART_TYPE.LINE_STACKED,
            "stacked_markers": XL_CHART_TYPE.LINE_MARKERS_STACKED,
            "percent_stacked": XL_CHART_TYPE.LINE_STACKED_100,
            "percent_stacked_markers": XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
        }, "line")
        self._add_chart(
            slide_index=slide_index, chart_type=chart_type, categories=categories,
            series=series, left=left, top=top, width=width, height=height,
        )

    def add_pie_chart(
        self,
        *,
        slide_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        style: str = "pie",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "pie": XL_CHART_TYPE.PIE,
            "exploded": XL_CHART_TYPE.PIE_EXPLODED,
            "pie_of_pie": XL_CHART_TYPE.PIE_OF_PIE,
            "bar_of_pie": XL_CHART_TYPE.BAR_OF_PIE,
        }, "pie")
        self._add_chart(
            slide_index=slide_index, chart_type=chart_type, categories=categories,
            series=series, left=left, top=top, width=width, height=height,
        )

    def add_area_chart(
        self,
        *,
        slide_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        style: str = "area",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "area": XL_CHART_TYPE.AREA,
            "stacked": XL_CHART_TYPE.AREA_STACKED,
            "percent_stacked": XL_CHART_TYPE.AREA_STACKED_100,
        }, "area")
        self._add_chart(
            slide_index=slide_index, chart_type=chart_type, categories=categories,
            series=series, left=left, top=top, width=width, height=height,
        )

    def add_doughnut_chart(
        self,
        *,
        slide_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        style: str = "doughnut",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "exploded": XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        }, "doughnut")
        self._add_chart(
            slide_index=slide_index, chart_type=chart_type, categories=categories,
            series=series, left=left, top=top, width=width, height=height,
        )

    def add_scatter_chart(
        self,
        *,
        slide_index: int,
        series: list[tuple[str, list[tuple[float, float]]]],
        style: str = "markers",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "markers": XL_CHART_TYPE.XY_SCATTER,
            "line": XL_CHART_TYPE.XY_SCATTER_LINES,
            "line_no_markers": XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
            "smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
            "smooth_no_markers": XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
        }, "scatter")
        self._add_xy_chart(
            slide_index=slide_index, chart_type=chart_type, series=series,
            left=left, top=top, width=width, height=height,
        )

    def add_radar_chart(
        self,
        *,
        slide_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        style: str = "radar",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "radar": XL_CHART_TYPE.RADAR,
            "filled": XL_CHART_TYPE.RADAR_FILLED,
            "markers": XL_CHART_TYPE.RADAR_MARKERS,
        }, "radar")
        self._add_chart(
            slide_index=slide_index, chart_type=chart_type, categories=categories,
            series=series, left=left, top=top, width=width, height=height,
        )

    def add_bubble_chart(
        self,
        *,
        slide_index: int,
        series: list[tuple[str, list[tuple[float, float, float]]]],
        style: str = "bubble",
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        chart_type = self._resolve_chart_type(style, {
            "bubble": XL_CHART_TYPE.BUBBLE,
            "bubble_3d": XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT,
        }, "bubble")
        self._add_bubble_chart(
            slide_index=slide_index, chart_type=chart_type, series=series,
            left=left, top=top, width=width, height=height,
        )

    def add_combo_chart_overlay(
        self,
        *,
        slide_index: int,
        categories: list[str],
        bar_series: list[tuple[str, list[float | None]]],
        line_series: list[tuple[str, list[float | None]]],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        # Practical combo strategy: overlay clustered-bar and line charts in same geometry.
        self.add_bar_chart(
            slide_index=slide_index,
            categories=categories,
            series=bar_series,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        self.add_line_chart(
            slide_index=slide_index,
            categories=categories,
            series=line_series,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        slide = self._slide(slide_index)
        chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
        if len(chart_shapes) < 2:
            return
        line_chart = chart_shapes[-1].chart
        line_chart.has_legend = False
        with suppress(Exception):
            line_chart.category_axis.visible = False
        try:
            line_chart.value_axis.has_major_gridlines = False
            line_chart.value_axis.has_minor_gridlines = False
            self._set_axis_position(line_chart.value_axis._element, "r")
        except Exception:
            pass

    def add_image(
        self,
        *,
        slide_index: int,
        path: str | Path,
        left: float,
        top: float,
        width: float | None = None,
        height: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        image_path = Path(path)
        if not image_path.exists():
            raise SlidesError(
                code="IMAGE_NOT_FOUND",
                message=f"Image path does not exist: {image_path}",
                path="add_image.path",
                suggested_fix="Provide an absolute or valid relative path to an existing image",
            )

        kwargs: dict[str, Any] = {}
        if width is not None:
            kwargs["width"] = Inches(width)
        if height is not None:
            kwargs["height"] = Inches(height)

        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)

    def add_table(
        self,
        *,
        slide_index: int,
        rows: list[list[str]],
        left: float,
        top: float,
        width: float,
        height: float,
        table_xml: str | None = None,
        font_size: int | None = None,
    ) -> None:
        if not rows or not rows[0]:
            raise SlidesError(
                code="INVALID_TABLE",
                message="rows must contain at least one row and one column",
                path="add_table.rows",
                suggested_fix="Pass a 2D list with at least one row and one column",
            )
        col_count = len(rows[0])
        if any(len(r) != col_count for r in rows):
            raise SlidesError(
                code="INVALID_TABLE",
                message="All table rows must have same number of columns",
                path="add_table.rows",
                suggested_fix="Pad or normalize rows so each has equal length",
            )

        slide = self._slide(slide_index)
        table_shape = None
        placeholder = self._find_placeholder_by_geometry(
            slide=slide,
            left=left,
            top=top,
            width=width,
            height=height,
            inserter_name="insert_table",
        )
        if placeholder is not None:
            table_shape = placeholder.insert_table(rows=len(rows), cols=col_count)
        if table_shape is None:
            table_shape = slide.shapes.add_table(
                rows=len(rows),
                cols=col_count,
                left=Inches(left),
                top=Inches(top),
                width=Inches(width),
                height=Inches(height),
            )
        table = table_shape.table
        for r_idx, row in enumerate(rows):
            for c_idx, cell in enumerate(row):
                table.cell(r_idx, c_idx).text = cell
        if font_size is not None:
            for r_idx in range(len(rows)):
                for c_idx in range(len(rows[0])):
                    for paragraph in table.cell(r_idx, c_idx).text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
        if table_xml:
            try:
                graphic_data = table_shape._element.graphic.graphicData
                graphic_data.remove(graphic_data.tbl)
                graphic_data.append(etree.fromstring(table_xml.encode("utf-8")))
            except Exception as exc:  # noqa: BLE001
                raise SlidesError(
                    code="INVALID_TABLE_XML",
                    message=f"Failed to apply table_xml: {exc}",
                    path="add_table.table_xml",
                    suggested_fix="Provide a valid <a:tbl> DrawingML fragment",
                ) from exc

    def add_raw_shape_xml(
        self,
        *,
        slide_index: int,
        shape_xml: str,
        rel_images: list[tuple[str, str]] | None = None,
        rel_parts: list[tuple[str, str, str, str, str]] | None = None,
        rel_external: list[tuple[str, str, str]] | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        try:
            raw = etree.fromstring(shape_xml.encode("utf-8"))
        except Exception as exc:  # noqa: BLE001
            raise SlidesError(
                code="INVALID_RAW_SHAPE_XML",
                message=f"Failed to parse shape_xml: {exc}",
                path="add_raw_shape_xml.shape_xml",
                suggested_fix="Provide a valid p:sp/p:pic/p:grpSp/etc XML fragment",
            ) from exc

        rel_map: dict[str, str] = {}
        for old_rid, image_path_str in rel_images or []:
            image_path = Path(image_path_str)
            if not image_path.exists():
                continue
            try:
                _, new_rid = slide.part.get_or_add_image_part(str(image_path))
            except Exception:  # noqa: BLE001
                continue
            rel_map[old_rid] = new_rid
        for old_rid, reltype, target_ref in rel_external or []:
            try:
                new_rid = slide.part.relate_to(str(target_ref), reltype, is_external=True)
            except Exception:  # noqa: BLE001
                continue
            rel_map[old_rid] = new_rid
        for old_rid, reltype, content_type, source_partname, blob_path_str in rel_parts or []:
            blob_path = Path(blob_path_str)
            if not blob_path.exists():
                continue
            try:
                blob = blob_path.read_bytes()
            except Exception:  # noqa: BLE001
                continue
            part_tmpl = self._partname_template(source_partname)
            try:
                new_partname = slide.part.package.next_partname(part_tmpl)
            except Exception:  # noqa: BLE001
                try:
                    new_partname = PackURI(source_partname)
                except Exception:  # noqa: BLE001
                    continue
            try:
                part = Part.load(new_partname, content_type, slide.part.package, blob)
                new_rid = slide.part.relate_to(part, reltype)
            except Exception:  # noqa: BLE001
                continue
            rel_map[old_rid] = new_rid
        if rel_map:
            for node in raw.iter():
                for attr_name, attr_value in list(node.attrib.items()):
                    if attr_value in rel_map:
                        node.attrib[attr_name] = rel_map[attr_value]

        self._append_raw_element(slide_index=slide_index, raw=raw)

    def _append_raw_element(
        self,
        *,
        slide_index: int,
        raw: etree._Element,
    ) -> None:
        slide = self._slide(slide_index)
        max_id = max((shape.shape_id for shape in slide.shapes), default=1)
        next_id = max_id + 1
        for node in raw.iter():
            if etree.QName(node).localname != "cNvPr":
                continue
            if "id" in node.attrib:
                node.attrib["id"] = str(next_id)
                next_id += 1

        slide.shapes._spTree.append(raw)

    @staticmethod
    def _normalize_icon_key(name: str) -> str:
        return name.lower().replace(" ", "").replace("_", "")

    def _icon_catalog(self) -> dict[str, Path]:
        if self._icon_catalog_cache is not None:
            return self._icon_catalog_cache
        catalog: dict[str, Path] = {}
        for icon_dir in self._icon_dirs:
            if not icon_dir.exists():
                continue
            for p in icon_dir.glob("*.xml"):
                stem = p.stem
                short = stem.split("_", 1)[-1] if "_" in stem else stem
                for key in {self._normalize_icon_key(stem), self._normalize_icon_key(short)}:
                    if key not in catalog:
                        catalog[key] = p
        self._icon_catalog_cache = catalog
        return catalog

    def add_icon(
        self,
        *,
        slide_index: int,
        icon_name: str,
        left: float,
        top: float,
        size: float = 0.75,
        color: str | None = None,
    ) -> None:
        catalog = self._icon_catalog()
        search = self._normalize_icon_key(icon_name)
        icon_path = catalog.get(search)
        if icon_path is None:
            for key, path in catalog.items():
                if search in key:
                    icon_path = path
                    break
        if icon_path is None:
            available = sorted(catalog.keys())
            raise SlidesError(
                code="ICON_NOT_FOUND",
                message=f"Icon '{icon_name}' not found. Available: {', '.join(available)}",
                path="add_icon.icon_name",
                suggested_fix="Use one of the available icon names listed above",
            )

        raw = deepcopy(etree.fromstring(icon_path.read_bytes()))

        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

        grp_sp_pr = raw.find(f"{{{ns_p}}}grpSpPr")
        if grp_sp_pr is None:
            grp_sp_pr = raw.find(f"{{{ns_p}}}spPr")
        xfrm = grp_sp_pr.find(f"{{{ns_a}}}xfrm") if grp_sp_pr is not None else None
        if xfrm is not None:
            off = xfrm.find(f"{{{ns_a}}}off")
            ext = xfrm.find(f"{{{ns_a}}}ext")
            if off is not None and ext is not None:
                orig_cx = int(ext.get("cx", "914400"))
                orig_cy = int(ext.get("cy", "914400"))
                aspect = orig_cx / orig_cy if orig_cy > 0 else 1.0
                size_emu = int(Inches(size))
                if aspect >= 1:
                    cx = size_emu
                    cy = int(size_emu / aspect)
                else:
                    cy = size_emu
                    cx = int(size_emu * aspect)
                off.set("x", str(int(Inches(left))))
                off.set("y", str(int(Inches(top))))
                ext.set("cx", str(cx))
                ext.set("cy", str(cy))

        if color:
            for fill in raw.iter(f"{{{ns_a}}}solidFill"):
                scheme = fill.find(f"{{{ns_a}}}schemeClr")
                if scheme is not None:
                    fill.remove(scheme)
                    srgb = etree.SubElement(fill, f"{{{ns_a}}}srgbClr")
                    srgb.set("val", color)
                else:
                    srgb = fill.find(f"{{{ns_a}}}srgbClr")
                    if srgb is not None:
                        srgb.set("val", color)

        self._append_raw_element(slide_index=slide_index, raw=raw)

    @staticmethod
    def _partname_template(source_partname: str) -> str:
        base = str(source_partname)
        m = re.match(r"^(.*?)(\d+)(\.[^./]+)$", base)
        if m:
            return f"{m.group(1)}%d{m.group(3)}"
        m = re.match(r"^(.*?)(\.[^./]+)$", base)
        if m:
            return f"{m.group(1)}%d{m.group(2)}"
        return f"{base}%d"

    def add_notes(self, *, slide_index: int, text: str) -> None:
        slide = self._slide(slide_index)
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        if notes_frame is None:
            raise SlidesError(
                code="MISSING_NOTES_PLACEHOLDER",
                message="No notes placeholder available on notes slide",
                path="add_notes.slide_index",
                suggested_fix="Use a template with notes body placeholder",
            )
        notes_frame.clear()
        notes_frame.text = text

    def add_media(
        self,
        *,
        slide_index: int,
        path: str | Path,
        left: float,
        top: float,
        width: float,
        height: float,
        mime_type: str = "video/unknown",
        poster_path: str | Path | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        media_path = Path(path)
        if not media_path.exists():
            raise SlidesError(
                code="MEDIA_NOT_FOUND",
                message=f"Media path does not exist: {media_path}",
                path="add_media.path",
                suggested_fix="Provide a valid path to an existing media file",
            )
        poster_arg: str | None = None
        if poster_path is not None:
            poster_file = Path(poster_path)
            if not poster_file.exists():
                raise SlidesError(
                    code="POSTER_NOT_FOUND",
                    message=f"Poster path does not exist: {poster_file}",
                    path="add_media.poster_path",
                    suggested_fix="Provide a valid poster image path or omit poster_path",
                )
            poster_arg = str(poster_file)

        slide.shapes.add_movie(
            str(media_path),
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            poster_frame_image=poster_arg,
            mime_type=mime_type,
        )

    def set_core_properties(
        self,
        *,
        title: str | None = None,
        subject: str | None = None,
        author: str | None = None,
        keywords: str | None = None,
    ) -> None:
        core = self._prs.core_properties
        if title is not None:
            core.title = title
        if subject is not None:
            core.subject = subject
        if author is not None:
            core.author = author
        if keywords is not None:
            core.keywords = keywords

    def update_table_cell(
        self,
        *,
        slide_index: int,
        table_index: int,
        row: int,
        col: int,
        text: str,
        font_size: int | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        table_shapes = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]
        if table_index < 0 or table_index >= len(table_shapes):
            raise SlidesError(
                code="INVALID_TABLE_INDEX",
                message=f"table_index {table_index} out of range",
                path="update_table_cell.table_index",
                suggested_fix=f"use index in [0, {len(table_shapes) - 1}]",
            )
        table = table_shapes[table_index].table
        if row < 0 or row >= len(table.rows):
            raise SlidesError(
                code="INVALID_TABLE_ROW",
                message=f"row {row} out of range",
                path="update_table_cell.row",
                suggested_fix=f"use row in [0, {len(table.rows) - 1}]",
            )
        if col < 0 or col >= len(table.columns):
            raise SlidesError(
                code="INVALID_TABLE_COL",
                message=f"col {col} out of range",
                path="update_table_cell.col",
                suggested_fix=f"use col in [0, {len(table.columns) - 1}]",
            )
        cell = table.cell(row, col)
        cell.text = text
        if font_size is not None:
            from pptx.util import Pt

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)

    def update_chart_data(
        self,
        *,
        slide_index: int,
        chart_index: int,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
    ) -> None:
        slide = self._slide(slide_index)
        chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
        if chart_index < 0 or chart_index >= len(chart_shapes):
            raise SlidesError(
                code="INVALID_CHART_INDEX",
                message=f"chart_index {chart_index} out of range",
                path="update_chart_data.chart_index",
                suggested_fix=f"use index in [0, {len(chart_shapes) - 1}]",
            )
        chart = chart_shapes[chart_index].chart
        normalized_categories, normalized_series = self._normalize_category_series(
            categories=categories,
            series=series,
        )
        data = CategoryChartData()
        data.categories = normalized_categories
        for name, values in normalized_series:
            data.add_series(name, tuple(values))
        chart.replace_data(data)

    def set_slide_background(self, *, slide_index: int, color_hex: str) -> None:
        slide = self._slide(slide_index)
        normalized = color_hex.strip().lstrip("#").upper()
        if len(normalized) != 6 or any(ch not in "0123456789ABCDEF" for ch in normalized):
            raise SlidesError(
                code="INVALID_COLOR_HEX",
                message=f"Invalid color value: {color_hex}",
                path="set_slide_background.color_hex",
                suggested_fix="Use a 6-digit hex color (for example: 00AA55)",
            )
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(normalized)

    def list_placeholders(self, *, slide_index: int) -> list[dict[str, Any]]:
        slide = self._slide(slide_index)
        placeholders: list[dict[str, Any]] = []
        for shape in slide.placeholders:
            phf = shape.placeholder_format
            placeholders.append(
                {
                    "idx": int(phf.idx),
                    "name": shape.name,
                    "type": str(phf.type),
                    "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                }
            )
        return placeholders

    def set_placeholder_text(
        self,
        *,
        slide_index: int,
        placeholder_idx: int,
        text: str,
        text_xml: str | None = None,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        bbox = None
        if left is not None and top is not None and width is not None and height is not None:
            bbox = (left, top, width, height)
        target = None
        for shape in slide.placeholders:
            if int(shape.placeholder_format.idx) == placeholder_idx:
                target = shape
                break
        if target is None and bbox is not None:
            best_shape = None
            best_score = float("inf")
            for shape in slide.placeholders:
                if not getattr(shape, "has_text_frame", False):
                    continue
                score = (
                    abs(shape.left / 914400 - bbox[0])
                    + abs(shape.top / 914400 - bbox[1])
                    + abs(shape.width / 914400 - bbox[2])
                    + abs(shape.height / 914400 - bbox[3])
                )
                if score < best_score:
                    best_shape = shape
                    best_score = score
            if best_shape is not None and best_score <= 2.0:
                target = best_shape
        if target is None:
            if left is not None and top is not None and width is not None and height is not None:
                self.add_text(
                    slide_index=slide_index,
                    text=text,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                )
                return
            available = [
                int(s.placeholder_format.idx) for s in slide.placeholders
            ]
            raise SlidesError(
                code="PLACEHOLDER_NOT_FOUND",
                message=(
                    f"Placeholder idx {placeholder_idx} not found;"
                    f" available indices: {available}"
                ),
                path="set_placeholder_text.placeholder_idx",
                suggested_fix=(
                    "Use set_semantic_text with role"
                    " (title/body/subtitle) instead"
                ),
            )
        if not getattr(target, "has_text_frame", False):
            if left is not None and top is not None and width is not None and height is not None:
                self.add_text(
                    slide_index=slide_index,
                    text=text,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                )
                return
            raise SlidesError(
                code="PLACEHOLDER_NO_TEXT_FRAME",
                message=f"Placeholder idx {placeholder_idx} does not support text",
                path="set_placeholder_text.placeholder_idx",
                suggested_fix="Choose a text-capable placeholder",
            )
        target.text_frame.clear()
        target.text_frame.text = text
        if text_xml:
            try:
                target._element.remove(target._element.txBody)
                target._element.append(etree.fromstring(text_xml.encode("utf-8")))
            except Exception as exc:  # noqa: BLE001
                raise SlidesError(
                    code="INVALID_PLACEHOLDER_TEXT_XML",
                    message=f"Failed to apply text_xml: {exc}",
                    path="set_placeholder_text.text_xml",
                    suggested_fix="Provide a valid <p:txBody> fragment",
                ) from exc

    def set_placeholder_image(
        self,
        *,
        slide_index: int,
        placeholder_idx: int,
        path: str,
        crop_left: float | None = None,
        crop_right: float | None = None,
        crop_top: float | None = None,
        crop_bottom: float | None = None,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        image_path = Path(path)
        if not image_path.exists():
            raise SlidesError(
                code="IMAGE_NOT_FOUND",
                message=f"Image path does not exist: {image_path}",
                path="set_placeholder_image.path",
                suggested_fix="Provide an absolute or valid relative path to an existing image",
            )
        bbox = None
        if left is not None and top is not None and width is not None and height is not None:
            bbox = (left, top, width, height)
        target = None
        for shape in slide.placeholders:
            if int(shape.placeholder_format.idx) == placeholder_idx:
                target = shape
                break
        if target is None and bbox is not None:
            best_shape = None
            best_score = float("inf")
            for shape in slide.placeholders:
                insert_picture = getattr(shape, "insert_picture", None)
                if not callable(insert_picture):
                    continue
                score = (
                    abs(shape.left / 914400 - bbox[0])
                    + abs(shape.top / 914400 - bbox[1])
                    + abs(shape.width / 914400 - bbox[2])
                    + abs(shape.height / 914400 - bbox[3])
                )
                if score < best_score:
                    best_shape = shape
                    best_score = score
            if best_shape is not None and best_score <= 2.0:
                target = best_shape
        if target is None:
            if left is not None and top is not None and width is not None and height is not None:
                self.add_image(
                    slide_index=slide_index,
                    path=image_path,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                )
                picture = slide.shapes[-1]
                if crop_left is not None:
                    picture.crop_left = crop_left
                if crop_right is not None:
                    picture.crop_right = crop_right
                if crop_top is not None:
                    picture.crop_top = crop_top
                if crop_bottom is not None:
                    picture.crop_bottom = crop_bottom
                return
            raise SlidesError(
                code="PLACEHOLDER_NOT_FOUND",
                message=f"Placeholder idx {placeholder_idx} not found",
                path="set_placeholder_image.placeholder_idx",
                suggested_fix="Inspect placeholders and use a valid idx",
            )
        insert_picture = getattr(target, "insert_picture", None)
        if not callable(insert_picture):
            if left is not None and top is not None and width is not None and height is not None:
                self.add_image(
                    slide_index=slide_index,
                    path=image_path,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                )
                picture = slide.shapes[-1]
                if crop_left is not None:
                    picture.crop_left = crop_left
                if crop_right is not None:
                    picture.crop_right = crop_right
                if crop_top is not None:
                    picture.crop_top = crop_top
                if crop_bottom is not None:
                    picture.crop_bottom = crop_bottom
                return
            raise SlidesError(
                code="PLACEHOLDER_NO_IMAGE_SUPPORT",
                message=f"Placeholder idx {placeholder_idx} does not support image insertion",
                path="set_placeholder_image.placeholder_idx",
                suggested_fix="Choose an image-capable placeholder or use add_image",
            )
        picture = insert_picture(str(image_path))
        if left is not None:
            picture.left = Inches(left)
        if top is not None:
            picture.top = Inches(top)
        if width is not None:
            picture.width = Inches(width)
        if height is not None:
            picture.height = Inches(height)
        if crop_left is not None:
            picture.crop_left = crop_left
        if crop_right is not None:
            picture.crop_right = crop_right
        if crop_top is not None:
            picture.crop_top = crop_top
        if crop_bottom is not None:
            picture.crop_bottom = crop_bottom

    def set_title_subtitle(
        self,
        *,
        slide_index: int,
        title: str | None = None,
        subtitle: str | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        title_shape = slide.shapes.title
        if title is not None:
            if title_shape is None:
                raise SlidesError(
                    code="TITLE_PLACEHOLDER_NOT_FOUND",
                    message="Slide has no title placeholder",
                    path="set_title_subtitle.title",
                    suggested_fix="Use a title layout or set placeholder text by idx",
                )
            title_shape.text = title
        if subtitle is not None:
            subtitle_set = False
            for shape in slide.placeholders:
                phf = shape.placeholder_format
                if str(phf.type).endswith("SUBTITLE (4)"):
                    shape.text = subtitle
                    subtitle_set = True
                    break
            if not subtitle_set:
                # fallback: first non-title text-capable placeholder
                for shape in slide.placeholders:
                    phf = shape.placeholder_format
                    if str(phf.type).endswith("TITLE (1)"):
                        continue
                    if getattr(shape, "has_text_frame", False):
                        shape.text_frame.clear()
                        shape.text_frame.text = subtitle
                        subtitle_set = True
                        break
            if not subtitle_set:
                raise SlidesError(
                    code="SUBTITLE_PLACEHOLDER_NOT_FOUND",
                    message="Slide has no subtitle/text placeholder for subtitle",
                    path="set_title_subtitle.subtitle",
                    suggested_fix="Use a layout with subtitle or set text via add_text",
                )

    def set_semantic_text(self, *, slide_index: int, role: str, text: str) -> None:
        slide = self._slide(slide_index)
        if role == "title":
            title_shape = slide.shapes.title
            if title_shape is None:
                raise SlidesError(
                    code="SEMANTIC_ROLE_NOT_FOUND",
                    message=(
                        f"No title placeholder on layout"
                        f" '{slide.slide_layout.name}';"
                        " use add_text with coordinates"
                    ),
                    path="set_semantic_text.role",
                    suggested_fix=(
                        "Use add_text with explicit"
                        " coordinates and font_color"
                    ),
                )
            title_shape.text = text
            return

        role_tokens = {
            "subtitle": ["SUBTITLE (4)"],
            "body": ["BODY (2)", "OBJECT (7)"],
            "footer": ["FOOTER (15)"],
            "date": ["DATE (16)"],
            "slide_number": ["SLIDE_NUMBER (13)"],
        }
        wanted = role_tokens.get(role)
        if wanted is None:
            raise SlidesError(
                code="INVALID_SEMANTIC_ROLE",
                message=f"Unsupported semantic role: {role}",
                path="set_semantic_text.role",
                suggested_fix="Use one of: title, subtitle, body, footer, date, slide_number",
            )
        for shape in slide.placeholders:
            type_name = str(shape.placeholder_format.type)
            if any(type_name.endswith(token) for token in wanted):
                if not getattr(shape, "has_text_frame", False):
                    continue
                shape.text_frame.clear()
                shape.text_frame.text = text
                if role == "body":
                    self._ensure_autofit(shape)
                return

        raise SlidesError(
            code="SEMANTIC_ROLE_NOT_FOUND",
            message=(
                f"No {role} placeholder on layout"
                f" '{slide.slide_layout.name}';"
                " use add_text with coordinates"
            ),
            path="set_semantic_text.role",
            suggested_fix=(
                "Use add_text with explicit"
                " coordinates and font_color"
            ),
        )

    def set_chart_legend(
        self,
        *,
        slide_index: int,
        chart_index: int,
        visible: bool = True,
        position: str = "right",
        include_in_layout: bool | None = None,
        font_size: int | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_legend",
        )
        chart.has_legend = visible
        if not visible:
            return

        legend = chart.legend
        if legend is None:
            raise SlidesError(
                code="CHART_LEGEND_NOT_AVAILABLE",
                message="Chart legend could not be created/accessed",
                path="set_chart_legend.visible",
            )
        position_map = {
            "right": XL_LEGEND_POSITION.RIGHT,
            "left": XL_LEGEND_POSITION.LEFT,
            "top": XL_LEGEND_POSITION.TOP,
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "corner": XL_LEGEND_POSITION.CORNER,
        }
        enum_value = position_map.get(position)
        if enum_value is None:
            raise SlidesError(
                code="INVALID_LEGEND_POSITION",
                message=f"Unsupported legend position: {position}",
                path="set_chart_legend.position",
                suggested_fix="Use one of: right, left, top, bottom, corner",
            )
        legend.position = enum_value
        if include_in_layout is not None:
            legend.include_in_layout = include_in_layout
        if font_size is not None:
            from pptx.util import Pt

            legend.font.size = Pt(font_size)

    def set_chart_style(self, *, slide_index: int, chart_index: int, style_id: int) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_style",
        )
        if style_id < 1 or style_id > 48:
            raise SlidesError(
                code="INVALID_CHART_STYLE",
                message=f"style_id out of supported range: {style_id}",
                path="set_chart_style.style_id",
                suggested_fix="Use a value in [1, 48]",
            )
        chart.chart_style = style_id

    def set_chart_title(
        self,
        *,
        slide_index: int,
        chart_index: int,
        text: str,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_title",
        )
        chart.has_title = True
        chart.chart_title.text_frame.clear()
        chart.chart_title.text_frame.text = text

    def set_chart_axis_titles(
        self,
        *,
        slide_index: int,
        chart_index: int,
        category_title: str | None = None,
        value_title: str | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_axis_titles",
        )
        if category_title is not None:
            try:
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.clear()
                chart.category_axis.axis_title.text_frame.text = category_title
            except Exception as exc:  # noqa: BLE001
                raise SlidesError(
                    code="CHART_AXIS_TITLE_UNSUPPORTED",
                    message=f"Category axis title is not supported on this chart: {exc}",
                    path="set_chart_axis_titles.category_title",
                ) from exc
        if value_title is not None:
            try:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.clear()
                chart.value_axis.axis_title.text_frame.text = value_title
            except Exception as exc:  # noqa: BLE001
                raise SlidesError(
                    code="CHART_AXIS_TITLE_UNSUPPORTED",
                    message=f"Value axis title is not supported on this chart: {exc}",
                    path="set_chart_axis_titles.value_title",
                ) from exc

    def set_chart_axis_options(
        self,
        *,
        slide_index: int,
        chart_index: int,
        axis: str = "value",
        reverse_order: bool | None = None,
        major_tick_mark: str | None = None,
        minor_tick_mark: str | None = None,
        tick_label_position: str | None = None,
        visible: bool | None = None,
        crosses: str | None = None,
        crosses_at: float | None = None,
        font_size: int | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_axis_options",
        )
        if axis == "category":
            target_axis = chart.category_axis
        elif axis == "value":
            target_axis = chart.value_axis
        else:
            raise SlidesError(
                code="INVALID_AXIS_TYPE",
                message=f"Unsupported axis type: {axis}",
                path="set_chart_axis_options.axis",
            )

        if reverse_order is not None:
            target_axis.reverse_order = reverse_order
        if visible is not None:
            target_axis.visible = visible

        tick_mark_map = {
            "none": XL_TICK_MARK.NONE,
            "inside": XL_TICK_MARK.INSIDE,
            "outside": XL_TICK_MARK.OUTSIDE,
            "cross": XL_TICK_MARK.CROSS,
        }
        label_pos_map = {
            "none": XL_TICK_LABEL_POSITION.NONE,
            "low": XL_TICK_LABEL_POSITION.LOW,
            "high": XL_TICK_LABEL_POSITION.HIGH,
            "next_to_axis": XL_TICK_LABEL_POSITION.NEXT_TO_AXIS,
        }
        crosses_map = {
            "automatic": XL_AXIS_CROSSES.AUTOMATIC,
            "minimum": XL_AXIS_CROSSES.MINIMUM,
            "maximum": XL_AXIS_CROSSES.MAXIMUM,
        }

        if major_tick_mark is not None:
            enum = tick_mark_map.get(major_tick_mark)
            if enum is None:
                raise SlidesError(
                    code="INVALID_TICK_MARK",
                    message=f"Unsupported major_tick_mark: {major_tick_mark}",
                    path="set_chart_axis_options.major_tick_mark",
                )
            target_axis.major_tick_mark = enum
        if minor_tick_mark is not None:
            enum = tick_mark_map.get(minor_tick_mark)
            if enum is None:
                raise SlidesError(
                    code="INVALID_TICK_MARK",
                    message=f"Unsupported minor_tick_mark: {minor_tick_mark}",
                    path="set_chart_axis_options.minor_tick_mark",
                )
            target_axis.minor_tick_mark = enum
        if tick_label_position is not None:
            enum = label_pos_map.get(tick_label_position)
            if enum is None:
                raise SlidesError(
                    code="INVALID_TICK_LABEL_POSITION",
                    message=f"Unsupported tick_label_position: {tick_label_position}",
                    path="set_chart_axis_options.tick_label_position",
                )
            target_axis.tick_label_position = enum
        if crosses is not None:
            enum = crosses_map.get(crosses)
            if enum is None:
                raise SlidesError(
                    code="INVALID_AXIS_CROSSES",
                    message=f"Unsupported crosses value: {crosses}",
                    path="set_chart_axis_options.crosses",
                )
            target_axis.crosses = enum
        if crosses_at is not None:
            target_axis.crosses_at = crosses_at
        if font_size is not None:
            from pptx.util import Pt

            target_axis.tick_labels.font.size = Pt(font_size)

    def set_chart_data_labels(
        self,
        *,
        slide_index: int,
        chart_index: int,
        enabled: bool = True,
        show_value: bool | None = None,
        show_category_name: bool | None = None,
        show_series_name: bool | None = None,
        number_format: str | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_data_labels",
        )
        if not chart.plots:
            raise SlidesError(
                code="CHART_NO_PLOT",
                message="Chart has no plot to configure data labels",
                path="set_chart_data_labels.chart_index",
            )
        plot = chart.plots[0]
        plot.has_data_labels = enabled
        if not enabled:
            return
        labels = plot.data_labels
        labels.show_value = enabled if show_value is None else show_value
        if show_category_name is not None:
            labels.show_category_name = show_category_name
        if show_series_name is not None:
            labels.show_series_name = show_series_name
        if number_format is not None:
            labels.number_format = number_format

    def set_chart_data_labels_style(
        self,
        *,
        slide_index: int,
        chart_index: int,
        position: str | None = None,
        show_legend_key: bool | None = None,
        number_format_is_linked: bool | None = None,
        font_size: int | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_data_labels_style",
        )
        if not chart.plots:
            raise SlidesError(
                code="CHART_NO_PLOT",
                message="Chart has no plot to configure data label style",
                path="set_chart_data_labels_style.chart_index",
            )
        plot = chart.plots[0]
        plot.has_data_labels = True
        labels = plot.data_labels
        pos_map = {
            "best_fit": XL_LABEL_POSITION.BEST_FIT,
            "center": XL_LABEL_POSITION.CENTER,
            "inside_base": XL_LABEL_POSITION.INSIDE_BASE,
            "inside_end": XL_LABEL_POSITION.INSIDE_END,
            "outside_end": XL_LABEL_POSITION.OUTSIDE_END,
            "left": XL_LABEL_POSITION.LEFT,
            "right": XL_LABEL_POSITION.RIGHT,
            "above": XL_LABEL_POSITION.ABOVE,
            "below": XL_LABEL_POSITION.BELOW,
        }
        if position is not None:
            enum = pos_map.get(position)
            if enum is None:
                raise SlidesError(
                    code="INVALID_LABEL_POSITION",
                    message=f"Unsupported label position: {position}",
                    path="set_chart_data_labels_style.position",
                )
            labels.position = enum
        if show_legend_key is not None:
            labels.show_legend_key = show_legend_key
        if number_format_is_linked is not None:
            labels.number_format_is_linked = number_format_is_linked
        if font_size is not None:
            if font_size < 6 or font_size > 72:
                raise SlidesError(
                    code="INVALID_FONT_SIZE",
                    message=f"font_size out of range: {font_size}",
                    path="set_chart_data_labels_style.font_size",
                )
            labels.font.size = Pt(font_size)

    def set_chart_axis_scale(
        self,
        *,
        slide_index: int,
        chart_index: int,
        minimum: float | None = None,
        maximum: float | None = None,
        major_unit: float | None = None,
        minor_unit: float | None = None,
        show_major_gridlines: bool | None = None,
        show_minor_gridlines: bool | None = None,
        number_format: str | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_axis_scale",
        )
        try:
            axis = chart.value_axis
        except Exception as exc:  # noqa: BLE001
            raise SlidesError(
                code="CHART_VALUE_AXIS_UNSUPPORTED",
                message=f"Chart has no configurable value axis: {exc}",
                path="set_chart_axis_scale.chart_index",
            ) from exc

        if minimum is not None:
            axis.minimum_scale = minimum
        if maximum is not None:
            axis.maximum_scale = maximum
        if major_unit is not None:
            axis.major_unit = major_unit
        if minor_unit is not None:
            axis.minor_unit = minor_unit
        if show_major_gridlines is not None:
            axis.has_major_gridlines = show_major_gridlines
        if show_minor_gridlines is not None:
            axis.has_minor_gridlines = show_minor_gridlines
        if number_format is not None:
            axis.tick_labels.number_format = number_format
            axis.tick_labels.number_format_is_linked = False

    def set_chart_plot_style(
        self,
        *,
        slide_index: int,
        chart_index: int,
        vary_by_categories: bool | None = None,
        gap_width: int | None = None,
        overlap: int | None = None,
        plot_area_x: float | None = None,
        plot_area_y: float | None = None,
        plot_area_w: float | None = None,
        plot_area_h: float | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_plot_style",
        )
        if not chart.plots:
            raise SlidesError(
                code="CHART_NO_PLOT",
                message="Chart has no plot to configure style",
                path="set_chart_plot_style.chart_index",
            )
        plot = chart.plots[0]
        if vary_by_categories is not None:
            try:
                plot.vary_by_categories = vary_by_categories
            except Exception as exc:  # noqa: BLE001
                raise SlidesError(
                    code="CHART_PLOT_STYLE_UNSUPPORTED",
                    message=f"vary_by_categories unsupported for this plot: {exc}",
                    path="set_chart_plot_style.vary_by_categories",
                ) from exc
        if gap_width is not None:
            if gap_width < 0 or gap_width > 500:
                raise SlidesError(
                    code="INVALID_GAP_WIDTH",
                    message=f"gap_width out of range: {gap_width}",
                    path="set_chart_plot_style.gap_width",
                    suggested_fix="Use a value in [0, 500]",
                )
            plot.gap_width = gap_width
        if overlap is not None:
            if overlap < -100 or overlap > 100:
                raise SlidesError(
                    code="INVALID_OVERLAP",
                    message=f"overlap out of range: {overlap}",
                    path="set_chart_plot_style.overlap",
                    suggested_fix="Use a value in [-100, 100]",
                )
            plot.overlap = overlap
        if any(v is not None for v in (plot_area_x, plot_area_y, plot_area_w, plot_area_h)):
            from lxml import etree as _etree

            nsmap = chart._chartSpace.nsmap
            c_ns = nsmap.get("c", "http://schemas.openxmlformats.org/drawingml/2006/chart")
            plot_area_el = chart._chartSpace.chart.plotArea
            layout = plot_area_el.find(f"{{{c_ns}}}layout")
            if layout is None:
                layout = _etree.SubElement(plot_area_el, f"{{{c_ns}}}layout")
                plot_area_el.insert(0, layout)
            manual = layout.find(f"{{{c_ns}}}manualLayout")
            if manual is None:
                manual = _etree.SubElement(layout, f"{{{c_ns}}}manualLayout")
            pairs = [("x", plot_area_x), ("y", plot_area_y), ("w", plot_area_w), ("h", plot_area_h)]
            for attr, val in pairs:
                if val is not None:
                    el = manual.find(f"{{{c_ns}}}{attr}")
                    if el is None:
                        el = _etree.SubElement(manual, f"{{{c_ns}}}{attr}")
                    el.set("val", str(val))

    def set_chart_series_style(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        fill_color_hex: str | None = None,
        invert_if_negative: bool | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_series_style",
        )
        if not chart.plots:
            raise SlidesError(
                code="CHART_NO_PLOT",
                message="Chart has no plot to configure series style",
                path="set_chart_series_style.chart_index",
            )
        series = chart.plots[0].series
        if series_index < 0 or series_index >= len(series):
            raise SlidesError(
                code="INVALID_SERIES_INDEX",
                message=f"series_index {series_index} out of range",
                path="set_chart_series_style.series_index",
                suggested_fix=f"use index in [0, {len(series) - 1}]",
            )
        target = series[series_index]
        if invert_if_negative is not None:
            target.invert_if_negative = invert_if_negative
        if fill_color_hex is not None:
            normalized = fill_color_hex.strip().lstrip("#").upper()
            if len(normalized) != 6 or any(ch not in "0123456789ABCDEF" for ch in normalized):
                raise SlidesError(
                    code="INVALID_COLOR_HEX",
                    message=f"Invalid color value: {fill_color_hex}",
                    path="set_chart_series_style.fill_color_hex",
                    suggested_fix="Use a 6-digit hex color (for example: 0A4280)",
                )
            fill = target.format.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(normalized)

    def set_chart_series_line_style(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        line_color_hex: str | None = None,
        line_width_pt: float | None = None,
    ) -> None:
        target = self._series(
            slide_index=slide_index,
            chart_index=chart_index,
            series_index=series_index,
            path="set_chart_series_line_style",
        )
        line = target.format.line
        if line_color_hex is not None:
            normalized = line_color_hex.strip().lstrip("#").upper()
            if len(normalized) != 6 or any(ch not in "0123456789ABCDEF" for ch in normalized):
                raise SlidesError(
                    code="INVALID_COLOR_HEX",
                    message=f"Invalid color value: {line_color_hex}",
                    path="set_chart_series_line_style.line_color_hex",
                )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string(normalized)
        if line_width_pt is not None:
            if line_width_pt <= 0 or line_width_pt > 20:
                raise SlidesError(
                    code="INVALID_LINE_WIDTH",
                    message=f"line_width_pt out of range: {line_width_pt}",
                    path="set_chart_series_line_style.line_width_pt",
                )
            line.width = Pt(line_width_pt)

    def set_chart_secondary_axis(
        self,
        *,
        slide_index: int,
        chart_index: int,
        enable: bool = True,
        series_indices: list[int] | None = None,
    ) -> None:
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_secondary_axis",
        )
        if not enable:
            return
        if series_indices:
            raise SlidesError(
                code="CHART_SECONDARY_AXIS_SERIES_MAPPING_UNSUPPORTED",
                message=(
                    "Series-to-secondary-axis mapping is not exposed by python-pptx "
                    "for generic charts"
                ),
                path="set_chart_secondary_axis.series_indices",
                suggested_fix=(
                    "Use a pre-authored template combo chart with secondary mapping "
                    "and update data via update_chart_data"
                ),
            )
        try:
            self._set_axis_position(chart.value_axis._element, "r")
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.has_minor_gridlines = False
        except Exception as exc:  # noqa: BLE001
            raise SlidesError(
                code="CHART_SECONDARY_AXIS_UNSUPPORTED",
                message=f"Unable to configure secondary axis view: {exc}",
                path="set_chart_secondary_axis",
            ) from exc

    def set_line_series_marker(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        style: str = "circle",
        size: int | None = None,
    ) -> None:
        series = self._series(
            slide_index=slide_index,
            chart_index=chart_index,
            series_index=series_index,
            path="set_line_series_marker",
        )
        marker_map = {
            "auto": XL_MARKER_STYLE.AUTOMATIC,
            "none": XL_MARKER_STYLE.NONE,
            "circle": XL_MARKER_STYLE.CIRCLE,
            "dash": XL_MARKER_STYLE.DASH,
            "diamond": XL_MARKER_STYLE.DIAMOND,
            "dot": XL_MARKER_STYLE.DOT,
            "plus": XL_MARKER_STYLE.PLUS,
            "square": XL_MARKER_STYLE.SQUARE,
            "star": XL_MARKER_STYLE.STAR,
            "triangle": XL_MARKER_STYLE.TRIANGLE,
            "x": XL_MARKER_STYLE.X,
        }
        enum = marker_map.get(style)
        if enum is None:
            raise SlidesError(
                code="INVALID_MARKER_STYLE",
                message=f"Unsupported marker style: {style}",
                path="set_line_series_marker.style",
            )
        series.marker.style = enum
        if size is not None:
            if size < 2 or size > 72:
                raise SlidesError(
                    code="INVALID_MARKER_SIZE",
                    message=f"Marker size out of range: {size}",
                    path="set_line_series_marker.size",
                    suggested_fix="Use a value in [2, 72]",
                )
            series.marker.size = size

    def set_line_series_smooth(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        smooth: bool = True,
    ) -> None:
        series = self._series(
            slide_index=slide_index,
            chart_index=chart_index,
            series_index=series_index,
            path="set_line_series_smooth",
        )
        try:
            series.smooth = smooth
        except Exception as exc:  # noqa: BLE001
            raise SlidesError(
                code="LINE_SMOOTH_UNSUPPORTED",
                message=f"Series smooth property unsupported: {exc}",
                path="set_line_series_smooth",
            ) from exc

    def set_chart_series_trendline(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        trend_type: str = "linear",
    ) -> None:
        series = self._series(
            slide_index=slide_index,
            chart_index=chart_index,
            series_index=series_index,
            path="set_chart_series_trendline",
        )
        ser_el = series._element
        for node in list(ser_el):
            if node.tag.endswith("trendline"):
                ser_el.remove(node)
        trend = OxmlElement("c:trendline")
        trend_type_el = OxmlElement("c:trendlineType")
        trend_type_el.set("val", trend_type)
        trend.append(trend_type_el)
        ser_el.append(trend)

    def set_chart_series_error_bars(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        value: float,
        direction: str = "y",
        bar_type: str = "both",
    ) -> None:
        if value <= 0:
            raise SlidesError(
                code="INVALID_ERROR_BAR_VALUE",
                message=f"Error bar value must be > 0, got: {value}",
                path="set_chart_series_error_bars.value",
            )
        series = self._series(
            slide_index=slide_index,
            chart_index=chart_index,
            series_index=series_index,
            path="set_chart_series_error_bars",
        )
        ser_el = series._element
        for node in list(ser_el):
            if node.tag.endswith("errBars"):
                ser_el.remove(node)
        err = OxmlElement("c:errBars")
        err_dir = OxmlElement("c:errDir")
        err_dir.set("val", direction)
        err_type = OxmlElement("c:errBarType")
        err_type.set("val", bar_type)
        err_val = OxmlElement("c:errValType")
        err_val.set("val", "fixedVal")
        err_fixed = OxmlElement("c:val")
        err_fixed.set("val", str(value))
        err.extend([err_dir, err_type, err_val, err_fixed])
        ser_el.append(err)

    def set_chart_combo_secondary_mapping(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_indices: list[int],
    ) -> None:
        if not series_indices:
            raise SlidesError(
                code="EMPTY_SERIES_SELECTION",
                message="series_indices must contain at least one series index",
                path="set_chart_combo_secondary_mapping.series_indices",
            )
        chart = self._chart(
            slide_index=slide_index,
            chart_index=chart_index,
            path="set_chart_combo_secondary_mapping",
        )
        plot_area = chart._chartSpace.chart.plotArea
        bar_chart = None
        line_chart = None
        cat_axes: dict[str, Any] = {}
        val_axes: dict[str, Any] = {}
        for child in plot_area:
            local = self._local_name(child.tag)
            if local == "barChart" and bar_chart is None:
                bar_chart = child
            elif local == "lineChart" and line_chart is None:
                line_chart = child
            elif local == "catAx":
                ax_id = self._axis_id(child)
                if ax_id is not None:
                    cat_axes[ax_id] = child
            elif local == "valAx":
                ax_id = self._axis_id(child)
                if ax_id is not None:
                    val_axes[ax_id] = child

        base_plot = bar_chart if bar_chart is not None else line_chart
        if base_plot is None:
            raise SlidesError(
                code="COMBO_BASE_CHART_UNSUPPORTED",
                message="Combo mapping requires a bar or line chart plot in the chart part",
                path="set_chart_combo_secondary_mapping.chart_index",
            )
        base_name = self._local_name(base_plot.tag)
        target_name = "lineChart" if base_name == "barChart" else "barChart"
        target_plot = line_chart if target_name == "lineChart" else bar_chart

        base_series = [n for n in base_plot if self._local_name(n.tag) == "ser"]
        selected = set(series_indices)
        selected_nodes: list[Any] = []
        for n in base_series:
            idx = self._ser_idx(n)
            if idx is not None and idx in selected:
                selected_nodes.append(n)
        if len(selected_nodes) != len(selected):
            raise SlidesError(
                code="INVALID_SERIES_INDEX",
                message="One or more series_indices were not found in base plot",
                path="set_chart_combo_secondary_mapping.series_indices",
            )
        if len(selected_nodes) >= len(base_series):
            raise SlidesError(
                code="INVALID_SERIES_SPLIT",
                message="At least one series must remain on primary plot",
                path="set_chart_combo_secondary_mapping.series_indices",
            )

        cat_id, val_primary_id = self._plot_axis_pair(base_plot, cat_axes, val_axes)
        if cat_id is None or val_primary_id is None:
            raise SlidesError(
                code="CHART_AXIS_RESOLUTION_FAILED",
                message="Unable to resolve primary chart axis ids",
                path="set_chart_combo_secondary_mapping.chart_index",
            )
        secondary_val_id = self._ensure_secondary_val_axis(
            plot_area=plot_area,
            cat_axes=cat_axes,
            val_axes=val_axes,
            cat_id=cat_id,
            val_primary_id=val_primary_id,
        )

        if target_plot is None:
            target_plot = OxmlElement(f"c:{target_name}")
            grouping = OxmlElement("c:grouping")
            grouping.set("val", "standard" if target_name == "lineChart" else "clustered")
            target_plot.append(grouping)
            ax_cat = OxmlElement("c:axId")
            ax_cat.set("val", cat_id)
            ax_val = OxmlElement("c:axId")
            ax_val.set("val", secondary_val_id)
            target_plot.append(ax_cat)
            target_plot.append(ax_val)
            plot_area.insert(1, target_plot)
        else:
            self._set_plot_axis_pair(target_plot, cat_id, secondary_val_id)

        for ser in selected_nodes:
            base_plot.remove(ser)
            self._insert_series_before_axids(target_plot, ser)

    def set_image_crop(
        self,
        *,
        slide_index: int,
        image_index: int,
        crop_left: float | None = None,
        crop_right: float | None = None,
        crop_top: float | None = None,
        crop_bottom: float | None = None,
    ) -> None:
        slide = self._slide(slide_index)
        pictures = [
            shape
            for shape in slide.shapes
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
        ]
        if image_index < 0 or image_index >= len(pictures):
            raise SlidesError(
                code="INVALID_IMAGE_INDEX",
                message=f"image_index {image_index} out of range",
                path="set_image_crop.image_index",
                suggested_fix=f"use index in [0, {len(pictures) - 1}]",
            )
        pic = pictures[image_index]
        for field, value in (
            ("crop_left", crop_left),
            ("crop_right", crop_right),
            ("crop_top", crop_top),
            ("crop_bottom", crop_bottom),
        ):
            if value is None:
                continue
            if not (-2.0 <= value <= 2.0):
                raise SlidesError(
                    code="INVALID_IMAGE_CROP",
                    message=f"{field} value out of supported range: {value}",
                    path=f"set_image_crop.{field}",
                    suggested_fix="Use a crop value between -2.0 and 2.0",
                )
            setattr(pic, field, value)

    def summarize(self) -> DeckSummary:
        shape_count = 0
        text_shape_count = 0
        chart_count = 0
        table_count = 0
        image_count = 0
        unresolved_token_count = 0

        for slide in self._prs.slides:
            for shape in slide.shapes:
                shape_count += 1
                if getattr(shape, "has_text_frame", False):
                    text_shape_count += 1
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text or ""
                        if "{{" in text and "}}" in text:
                            unresolved_token_count += 1
                if getattr(shape, "has_chart", False):
                    chart_count += 1
                if getattr(shape, "has_table", False):
                    table_count += 1
                if getattr(shape, "shape_type", None) == 13:
                    image_count += 1

        return DeckSummary(
            slide_count=len(self._prs.slides),
            shape_count=shape_count,
            text_shape_count=text_shape_count,
            chart_count=chart_count,
            table_count=table_count,
            image_count=image_count,
            unresolved_token_count=unresolved_token_count,
        )

    def inspect(self) -> dict[str, Any]:
        slide_width_emu = self._prs.slide_width
        slide_height_emu = self._prs.slide_height
        slide_width_in = round(float(slide_width_emu or 0) / 914400, 3)
        slide_height_in = round(float(slide_height_emu or 0) / 914400, 3)
        layout_names = [
            str(getattr(layout, "name", "") or "").strip() for layout in self._prs.slide_layouts
        ]
        slides: list[dict[str, Any]] = []
        for slide_index, slide in enumerate(self._prs.slides):
            slide_uid = self._slide_uid(slide)
            shape_uids = self._shape_uids(slide)
            shape_entries: list[dict[str, Any]] = []
            inferred_title: str | None = None
            for shape_index, shape in enumerate(slide.shapes):
                text = self._shape_text(shape)
                if inferred_title is None and text:
                    inferred_title = text.splitlines()[0][:140]
                left = top = width = height = None
                with suppress(Exception):
                    left = round(float(shape.left) / 914400, 3)
                with suppress(Exception):
                    top = round(float(shape.top) / 914400, 3)
                with suppress(Exception):
                    width = round(float(shape.width) / 914400, 3)
                with suppress(Exception):
                    height = round(float(shape.height) / 914400, 3)
                shape_entries.append(
                    {
                        "shape_index": shape_index,
                        "shape_id": getattr(shape, "shape_id", None),
                        "shape_uid": shape_uids[shape_index],
                        "name": getattr(shape, "name", None),
                        "kind": self._shape_kind(shape),
                        "is_placeholder": self._shape_is_placeholder(shape),
                        "has_fill": self._shape_has_fill(shape),
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                        "text": text,
                        "text_chars": len(text),
                        "font_names": self._shape_font_names(shape),
                        "font_sizes_pt": self._shape_font_sizes_pt(shape),
                        "font_colors_hex": self._shape_font_colors_hex(shape),
                        "run_count": self._shape_run_count(shape),
                        "explicit_font_runs": self._shape_explicit_font_runs(shape),
                        "italic_runs": self._shape_italic_runs(shape),
                        "underline_runs": self._shape_underline_runs(shape),
                        "has_effects": self._shape_has_effects(shape),
                        "has_3d": self._shape_has_3d(shape),
                        **(
                            self._chart_metrics(shape)
                            if getattr(shape, "has_chart", False) else {}
                        ),
                        **(
                            self._table_metrics(shape)
                            if getattr(shape, "has_table", False) else {}
                        ),
                    }
                )
            slides.append(
                {
                    "slide_index": slide_index,
                    "slide_id": f"slide-{slide_index + 1}",
                    "slide_uid": slide_uid,
                    "layout_name": self._slide_layout_name(slide),
                    "title": inferred_title or "",
                    "shape_count": len(shape_entries),
                    "has_transition": self._slide_has_transition(slide),
                    "has_animation": self._slide_has_animation(slide),
                    "placeholders": self._slide_placeholders(slide),
                    "shapes": shape_entries,
                }
            )
        return {
            "summary": self.summarize().to_dict(),
            "deck": {
                "slide_width_in": slide_width_in,
                "slide_height_in": slide_height_in,
                "slide_count": len(self._prs.slides),
                "master_count": len(self._prs.slide_masters),
                "layout_names": sorted({name for name in layout_names if name}),
            },
            "slides": slides,
        }

    def find_text(self, *, query: str, limit: int = 10) -> list[dict[str, Any]]:
        tokens = [t for t in re.findall(r"[A-Za-z0-9]+", query.lower()) if t]
        if not tokens:
            return []
        hits: list[dict[str, Any]] = []
        for slide_index, slide in enumerate(self._prs.slides):
            slide_uid = self._slide_uid(slide)
            shape_uids = self._shape_uids(slide)
            for shape_index, shape in enumerate(slide.shapes):
                text = self._shape_text(shape)
                if not text:
                    continue
                haystack = text.lower()
                score = sum(haystack.count(tok) for tok in tokens)
                if score <= 0:
                    continue
                snippet = self._best_snippet(text=text, token=tokens[0])
                hits.append(
                    {
                        "slide_index": slide_index,
                        "slide_id": f"slide-{slide_index + 1}",
                        "slide_uid": slide_uid,
                        "shape_index": shape_index,
                        "shape_id": getattr(shape, "shape_id", None),
                        "shape_uid": shape_uids[shape_index],
                        "score": score,
                        "snippet": snippet,
                    }
                )
        hits.sort(key=lambda item: (-int(item["score"]), int(item["slide_index"])))
        return hits[: max(1, limit)]

    def replace_text(self, *, slide_index: int, old: str, new: str) -> int:
        slide = self._slide(slide_index)
        replaced = 0
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if old in paragraph.text:
                    paragraph.text = paragraph.text.replace(old, new)
                    replaced += 1
        return replaced

    def semantic_replace_text(
        self,
        *,
        query: str,
        replacement: str,
        slide_index: int | None = None,
        slide_id: str | None = None,
        slide_uid: str | None = None,
        shape_id: int | None = None,
        shape_uid: str | None = None,
    ) -> dict[str, int]:
        flags = re.IGNORECASE
        pattern = re.compile(re.escape(query), flags=flags)
        replaced_paragraphs = 0
        scanned_paragraphs = 0
        resolved_slide_index = self._resolve_slide_index(
            slide_index=slide_index,
            slide_id=slide_id,
            slide_uid=slide_uid,
        )
        target_slides = (
            [self._slide(resolved_slide_index)]
            if resolved_slide_index is not None
            else list(self._prs.slides)
        )
        for slide in target_slides:
            shape_uids = self._shape_uids(slide)
            for shape_index, shape in enumerate(slide.shapes):
                if shape_id is not None and int(getattr(shape, "shape_id", -1)) != shape_id:
                    continue
                if (
                    shape_uid is not None
                    and not self._shape_uid_matches(
                        requested=shape_uid,
                        actual=shape_uids[shape_index],
                        shape=shape,
                    )
                ):
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    scanned_paragraphs += 1
                    src = paragraph.text or ""
                    if not pattern.search(src):
                        continue
                    paragraph.text = pattern.sub(replacement, src)
                    replaced_paragraphs += 1
        return {
            "replaced_paragraphs": replaced_paragraphs,
            "scanned_paragraphs": scanned_paragraphs,
        }

    def transform_slide_to_timeline(
        self,
        *,
        slide_index: int | None = None,
        slide_id: str | None = None,
        slide_uid: str | None = None,
        title: str | None = None,
    ) -> dict[str, Any]:
        resolved_slide_index = self._resolve_slide_index(
            slide_index=slide_index,
            slide_id=slide_id,
            slide_uid=slide_uid,
        )
        if resolved_slide_index is None:
            raise SlidesError(
                code="MISSING_SLIDE_SELECTOR",
                message="Provide slide_index or slide_id",
                path="transform_slide_to_timeline.slide_index",
            )
        slide_index = resolved_slide_index
        slide = self._slide(slide_index)
        extracted: list[str] = []
        for shape in slide.shapes:
            text = self._shape_text(shape)
            if not text:
                continue
            for line in text.splitlines():
                clean = line.strip().lstrip("-").strip()
                if clean:
                    extracted.append(clean)
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.clear()

        dedup: list[str] = []
        seen: set[str] = set()
        for item in extracted:
            key = item.lower()
            if key in seen:
                continue
            seen.add(key)
            dedup.append(item)
        if not dedup:
            dedup = ["Milestone 1", "Milestone 2", "Milestone 3", "Milestone 4"]
        title_text = title or dedup[0]
        milestones = dedup[1:5] or dedup[:4]

        self.add_text(
            slide_index=slide_index,
            text=title_text,
            left=0.6,
            top=0.35,
            width=11.2,
            height=0.9,
            font_size=28,
            bold=True,
        )
        for idx, point in enumerate(milestones):
            row = idx // 2
            col = idx % 2
            self.add_text(
                slide_index=slide_index,
                text=f"{idx + 1}. {point}",
                left=0.8 + (col * 5.6),
                top=1.6 + (row * 2.1),
                width=5.0,
                height=1.6,
                font_size=18,
                bold=False,
            )
        return {"slide_index": slide_index, "milestones": milestones}

    def _resolve_slide_index(
        self,
        *,
        slide_index: int | None,
        slide_id: str | None,
        slide_uid: str | None,
    ) -> int | None:
        resolved_slide_id = self._slide_id_to_index(slide_id) if slide_id is not None else None
        resolved_slide_uid = self._slide_uid_to_index(slide_uid) if slide_uid is not None else None

        if slide_index is not None:
            if resolved_slide_id is not None and resolved_slide_id != slide_index:
                raise SlidesError(
                    code="SLIDE_SELECTOR_CONFLICT",
                    message="Slide selectors conflict",
                    path="slide_selector",
                )
            if resolved_slide_uid is not None and resolved_slide_uid != slide_index:
                raise SlidesError(
                    code="SLIDE_SELECTOR_CONFLICT",
                    message="Slide selectors conflict",
                    path="slide_selector",
                )
            return slide_index

        if resolved_slide_id is not None and resolved_slide_uid is not None:
            if resolved_slide_id != resolved_slide_uid:
                raise SlidesError(
                    code="SLIDE_SELECTOR_CONFLICT",
                    message="Slide selectors conflict",
                    path="slide_selector",
                )
            return resolved_slide_id
        if resolved_slide_id is not None:
            return resolved_slide_id
        if resolved_slide_uid is not None:
            return resolved_slide_uid
        return None

    def _slide_id_to_index(self, slide_id: str) -> int:
        match = re.fullmatch(r"slide-(\d+)", slide_id.strip().lower())
        if not match:
            raise SlidesError(
                code="INVALID_SLIDE_ID",
                message=f"Unsupported slide_id format: {slide_id}",
                path="slide_id",
                suggested_fix="Use identifiers like slide-1, slide-2, ...",
            )
        slide_num = int(match.group(1))
        slide_index = slide_num - 1
        if slide_index < 0 or slide_index >= len(self._prs.slides):
            raise SlidesError(
                code="INVALID_SLIDE_ID",
                message=f"slide_id out of range: {slide_id}",
                path="slide_id",
                suggested_fix=f"use slide-1 to slide-{len(self._prs.slides)}",
            )
        return slide_index

    def _slide_uid_to_index(self, slide_uid: str) -> int:
        token = slide_uid.strip()
        for idx, slide in enumerate(self._prs.slides):
            if self._slide_uid(slide) == token:
                return idx
            partname = str(getattr(getattr(slide, "part", None), "partname", "") or "")
            if partname and partname == token:
                return idx
        raise SlidesError(
            code="INVALID_SLIDE_UID",
            message=f"Unknown slide_uid: {slide_uid}",
            path="slide_uid",
        )

    def _slide_uid(self, slide: Any) -> str:
        meta_uid = self._slide_meta_uid(slide)
        if meta_uid:
            return f"slideu-{meta_uid}"
        partname = str(getattr(getattr(slide, "part", None), "partname", ""))
        if partname:
            match = re.search(r"slide(\d+)\.xml$", partname)
            if match:
                generated = f"s{int(match.group(1)):08d}"
            else:
                generated = f"h{sha256(partname.encode('utf-8')).hexdigest()[:12]}"
        else:
            signatures = [self._shape_signature(shape) for shape in slide.shapes]
            payload = f"shape_count={len(signatures)}|{'|'.join(signatures)}"
            generated = f"h{sha256(payload.encode('utf-8')).hexdigest()[:12]}"
        self._set_slide_meta_uid(slide, generated)
        return f"slideu-{generated}"

    def _slide_meta_uid(self, slide: Any) -> str | None:
        c_sld = self._slide_csld(slide)
        if c_sld is None:
            return None
        name = str(c_sld.get("name", "") or "")
        match = re.search(r"\[slides-slideuid:([A-Za-z0-9_-]{6,64})\]", name)
        if match:
            return match.group(1)
        return None

    def _set_slide_meta_uid(self, slide: Any, uid: str) -> None:
        c_sld = self._slide_csld(slide)
        if c_sld is None:
            return
        name = str(c_sld.get("name", "") or "")
        if re.search(r"\[slides-slideuid:[A-Za-z0-9_-]{6,64}\]", name):
            name = re.sub(
                r"\[slides-slideuid:[A-Za-z0-9_-]{6,64}\]",
                f"[slides-slideuid:{uid}]",
                name,
                count=1,
            )
        elif name.strip():
            name = f"{name} [slides-slideuid:{uid}]"
        else:
            name = f"[slides-slideuid:{uid}]"
        c_sld.set("name", name)

    def _slide_csld(self, slide: Any) -> Any | None:
        with suppress(Exception):
            return slide._element.cSld
        return None

    def _slide_placeholders(self, slide: Any) -> list[dict[str, Any]]:
        items: list[dict[str, Any]] = []
        for placeholder in slide.placeholders:
            left = top = width = height = None
            with suppress(Exception):
                left = round(float(placeholder.left) / 914400, 3)
            with suppress(Exception):
                top = round(float(placeholder.top) / 914400, 3)
            with suppress(Exception):
                width = round(float(placeholder.width) / 914400, 3)
            with suppress(Exception):
                height = round(float(placeholder.height) / 914400, 3)
            with suppress(Exception):
                phf = placeholder.placeholder_format
                items.append(
                    {
                        "idx": int(phf.idx),
                        "type": str(phf.type),
                        "text": self._shape_text(placeholder),
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    }
                )
        items.sort(key=lambda item: int(item.get("idx", 0)))
        return items

    def _shape_uid(self, *, slide_uid: str, shape: Any) -> str:
        meta_uid = self._shape_meta_uid(shape)
        if meta_uid:
            return f"{slide_uid}::uid-{meta_uid}"
        shape_id = getattr(shape, "shape_id", None)
        if isinstance(shape_id, int):
            generated = f"s{shape_id:08x}"
            self._set_shape_meta_uid(shape, generated)
            return f"{slide_uid}::uid-{generated}"
        signature = self._shape_signature(shape)
        digest = sha256(signature.encode("utf-8")).hexdigest()[:8]
        generated = f"h{digest}"
        self._set_shape_meta_uid(shape, generated)
        return f"{slide_uid}::uid-{generated}"

    def _shape_uids(self, slide: Any) -> list[str]:
        slide_uid = self._slide_uid(slide)
        counts: dict[str, int] = {}
        uids: list[str] = []
        for shape in slide.shapes:
            base = self._shape_uid(slide_uid=slide_uid, shape=shape)
            counts[base] = counts.get(base, 0) + 1
            uids.append(f"{base}-{counts[base]}")
        return uids

    def _shape_meta_uid(self, shape: Any) -> str | None:
        c_nv_pr = self._shape_cnvpr(shape)
        if c_nv_pr is None:
            return None
        descr = str(c_nv_pr.get("descr", "") or "")
        match = re.search(r"\[slides-uid:([A-Za-z0-9_-]{6,64})\]", descr)
        if match:
            return match.group(1)
        return None

    def _set_shape_meta_uid(self, shape: Any, uid: str) -> None:
        c_nv_pr = self._shape_cnvpr(shape)
        if c_nv_pr is None:
            return
        descr = str(c_nv_pr.get("descr", "") or "")
        if re.search(r"\[slides-uid:[A-Za-z0-9_-]{6,64}\]", descr):
            descr = re.sub(
                r"\[slides-uid:[A-Za-z0-9_-]{6,64}\]",
                f"[slides-uid:{uid}]",
                descr,
                count=1,
            )
        elif descr.strip():
            descr = f"{descr} [slides-uid:{uid}]"
        else:
            descr = f"[slides-uid:{uid}]"
        c_nv_pr.set("descr", descr)

    def _shape_cnvpr(self, shape: Any) -> Any | None:
        with suppress(Exception):
            matches = shape._element.xpath(".//*[local-name()='cNvPr']")
            if matches:
                return matches[0]
        return None

    def _shape_uid_matches(self, *, requested: str, actual: str, shape: Any) -> bool:
        req = requested.strip()
        act = actual.strip()
        if req == act:
            return True
        if req and "::" not in req and f"::{req}" in act:
            return True
        req_base = re.sub(r"-\d+$", "", req)
        act_base = re.sub(r"-\d+$", "", act)
        if req_base == act_base:
            return True
        legacy_match = re.search(r"::shape-(\d+)", req)
        if legacy_match:
            legacy_id = int(legacy_match.group(1))
            shape_id = getattr(shape, "shape_id", None)
            if isinstance(shape_id, int) and legacy_id == shape_id:
                return True
        digest_match = re.search(r"shapeu-([0-9a-f]{8,12})", req)
        if digest_match:
            digest = digest_match.group(1)
            sig_digest = sha256(self._shape_signature(shape).encode("utf-8")).hexdigest()
            if sig_digest.startswith(digest):
                return True
        return False

    def _shape_has_effects(self, shape: Any) -> bool:
        with suppress(Exception):
            nodes = shape._element.xpath(
                ".//*[local-name()='effectLst' or local-name()='effectDag' or "
                "local-name()='outerShdw' or local-name()='innerShdw' or "
                "local-name()='glow' or local-name()='reflection' or "
                "local-name()='softEdge']"
            )
            if nodes:
                return True
        return self._shape_has_3d(shape)

    def _shape_has_3d(self, shape: Any) -> bool:
        with suppress(Exception):
            nodes = shape._element.xpath(
                ".//*[local-name()='scene3d' or local-name()='sp3d' or "
                "local-name()='bevelT' or local-name()='bevelB' or local-name()='extrusionClr']"
            )
            return bool(nodes)
        return False

    def _slide_has_transition(self, slide: Any) -> bool:
        with suppress(Exception):
            nodes = slide._element.xpath(".//*[local-name()='transition']")
            return bool(nodes)
        return False

    def _slide_has_animation(self, slide: Any) -> bool:
        with suppress(Exception):
            nodes = slide._element.xpath(
                ".//*[local-name()='timing' or local-name()='anim' or "
                "local-name()='animClr' or local-name()='animEffect' or "
                "local-name()='animMotion' or "
                "local-name()='animRot' or local-name()='animScale' or local-name()='cmd']"
            )
            return bool(nodes)
        return False

    def _shape_signature(self, shape: Any) -> str:
        kind = self._shape_kind(shape)
        text = self._shape_text(shape)
        left = top = width = height = ""
        with suppress(Exception):
            left = f"{round(float(shape.left) / 914400, 4)}"
        with suppress(Exception):
            top = f"{round(float(shape.top) / 914400, 4)}"
        with suppress(Exception):
            width = f"{round(float(shape.width) / 914400, 4)}"
        with suppress(Exception):
            height = f"{round(float(shape.height) / 914400, 4)}"
        return f"{kind}|{left}|{top}|{width}|{height}|{text}"

    def _slide_layout_name(self, slide: Any) -> str:
        with suppress(Exception):
            return str(slide.slide_layout.name or "")
        return ""

    def _shape_kind(self, shape: Any) -> str:
        if getattr(shape, "has_chart", False):
            return "chart"
        if getattr(shape, "has_table", False):
            return "table"
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.LINE:
            return "line"
        if getattr(shape, "has_text_frame", False):
            return "text"
        return "other"

    def _shape_is_placeholder(self, shape: Any) -> bool:
        with suppress(Exception):
            return bool(shape.is_placeholder)
        return False

    def _shape_has_fill(self, shape: Any) -> bool:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            return True
        element = getattr(shape, "_element", None)
        xml = getattr(element, "xml", "")
        return any(
            token in xml
            for token in ("<a:solidFill", "<a:gradFill", "<a:pattFill", "<a:blipFill")
        )

    def _shape_text(self, shape: Any) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        text_lines: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                text_lines.append(text)
        return "\n".join(text_lines)

    def _shape_font_names(self, shape: Any) -> list[str]:
        if not getattr(shape, "has_text_frame", False):
            return []
        names: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                name = getattr(getattr(run, "font", None), "name", None)
                if not name:
                    continue
                if name not in names:
                    names.append(name)
        return names

    def _shape_font_sizes_pt(self, shape: Any) -> list[float]:
        if not getattr(shape, "has_text_frame", False):
            return []
        sizes: list[float] = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                size = getattr(getattr(run, "font", None), "size", None)
                if size is None:
                    continue
                try:
                    value = round(float(size.pt), 1)
                except Exception:  # noqa: BLE001
                    continue
                if value not in sizes:
                    sizes.append(value)
        return sizes

    def _shape_font_colors_hex(self, shape: Any) -> list[str]:
        if not getattr(shape, "has_text_frame", False):
            return []
        colors: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                color = getattr(getattr(run, "font", None), "color", None)
                rgb = getattr(color, "rgb", None)
                if rgb is None:
                    continue
                value = str(rgb).upper()
                if value and value not in colors:
                    colors.append(value)
        return colors

    def _shape_run_count(self, shape: Any) -> int:
        if not getattr(shape, "has_text_frame", False):
            return 0
        return sum(len(paragraph.runs) for paragraph in shape.text_frame.paragraphs)

    def _shape_explicit_font_runs(self, shape: Any) -> int:
        if not getattr(shape, "has_text_frame", False):
            return 0
        count = 0
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                name = getattr(getattr(run, "font", None), "name", None)
                if name:
                    count += 1
        return count

    def _shape_italic_runs(self, shape: Any) -> int:
        if not getattr(shape, "has_text_frame", False):
            return 0
        count = 0
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if bool(getattr(getattr(run, "font", None), "italic", False)):
                    count += 1
        return count

    def _shape_underline_runs(self, shape: Any) -> int:
        if not getattr(shape, "has_text_frame", False):
            return 0
        count = 0
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                underline = getattr(getattr(run, "font", None), "underline", None)
                if underline not in (None, False):
                    count += 1
        return count

    def _chart_metrics(self, shape: Any) -> dict[str, Any]:
        """Extract font-size and layout metrics from a chart shape."""
        metrics: dict[str, Any] = {}
        try:
            chart = shape.chart
        except Exception:  # noqa: BLE001
            return metrics
        font_sizes: list[float] = []
        for axis_attr in ("category_axis", "value_axis"):
            try:
                axis = getattr(chart, axis_attr, None)
            except Exception:  # noqa: BLE001
                continue
            if axis is None:
                continue
            tick_labels = getattr(axis, "tick_labels", None)
            if tick_labels is None:
                continue
            font = getattr(tick_labels, "font", None)
            size = getattr(font, "size", None)
            if size is not None:
                with suppress(Exception):
                    font_sizes.append(round(float(size.pt), 1))
        legend = getattr(chart, "legend", None)
        legend_font_size: float | None = None
        if legend is not None:
            font = getattr(legend, "font", None)
            size = getattr(font, "size", None)
            if size is not None:
                with suppress(Exception):
                    legend_font_size = round(float(size.pt), 1)
            include = getattr(legend, "include_in_layout", None)
            if include is False:
                metrics["chart_legend_outside_plot"] = True
        if font_sizes:
            metrics["chart_axis_font_sizes_pt"] = font_sizes
        if legend_font_size is not None:
            metrics["chart_legend_font_size_pt"] = legend_font_size
        return metrics

    def _table_metrics(self, shape: Any) -> dict[str, Any]:
        """Extract cell font-size metrics from a table shape."""
        metrics: dict[str, Any] = {}
        try:
            table = shape.table
        except Exception:  # noqa: BLE001
            return metrics
        font_sizes: list[float] = []
        row_count = len(table.rows)
        col_count = len(table.columns)
        for row in table.rows:
            for cell in row.cells:
                tf = getattr(cell, "text_frame", None)
                if tf is None:
                    continue
                for para in tf.paragraphs:
                    for run in para.runs:
                        size = getattr(getattr(run, "font", None), "size", None)
                        if size is None:
                            continue
                        try:
                            val = round(float(size.pt), 1)
                        except Exception:  # noqa: BLE001
                            continue
                        if val not in font_sizes:
                            font_sizes.append(val)
        metrics["table_row_count"] = row_count
        metrics["table_col_count"] = col_count
        if font_sizes:
            metrics["table_cell_font_sizes_pt"] = font_sizes
        return metrics

    def _best_snippet(self, *, text: str, token: str) -> str:
        lowered = text.lower()
        pos = lowered.find(token.lower())
        if pos < 0:
            return text[:180]
        start = max(0, pos - 60)
        end = min(len(text), pos + 120)
        return text[start:end]

    def delete_slide(self, slide_index: int) -> None:
        slide_id_list = self._prs.slides._sldIdLst
        if slide_index < 0 or slide_index >= len(slide_id_list):
            raise SlidesError(
                code="INVALID_SLIDE_INDEX",
                message=f"slide_index {slide_index} out of range",
                path="delete_slide.slide_index",
                suggested_fix=f"use index in [0, {len(slide_id_list) - 1}]",
            )
        slide_ids = list(slide_id_list)
        slide_id = slide_ids[slide_index]
        rel_id = slide_id.rId
        self._prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)

    def move_slide(self, *, from_index: int, to_index: int) -> None:
        slide_id_list = self._prs.slides._sldIdLst
        slide_count = len(slide_id_list)
        if from_index < 0 or from_index >= slide_count:
            raise SlidesError(
                code="INVALID_SLIDE_INDEX",
                message=f"from_index {from_index} out of range",
                path="move_slide.from_index",
                suggested_fix=f"use index in [0, {slide_count - 1}]",
            )
        if to_index < 0 or to_index >= slide_count:
            raise SlidesError(
                code="INVALID_SLIDE_INDEX",
                message=f"to_index {to_index} out of range",
                path="move_slide.to_index",
                suggested_fix=f"use index in [0, {slide_count - 1}]",
            )
        if from_index == to_index:
            return
        slide_ids = list(slide_id_list)
        moving = slide_ids[from_index]
        slide_id_list.remove(moving)
        slide_id_list.insert(to_index, moving)

    def repair(self) -> ValidationReport:
        """Apply conservative, deterministic repair steps and return post-repair validation."""
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text or ""
                    if "{{" in text and "}}" in text:
                        paragraph.text = text.replace("{{", "").replace("}}", "")
        return self.validate()

    def validate(
        self,
        deep: bool = False,
        *,
        xsd_dir: str | Path | None = None,
        require_xsd: bool = False,
    ) -> ValidationReport:
        base = validate_presentation(self._prs)
        if not deep:
            return base
        pkg = validate_package_bytes(
            self.to_bytes(deterministic=True),
            xsd_dir=xsd_dir,
            require_xsd=require_xsd,
        )
        issues: list[ValidationIssue] = [*base.issues, *pkg.issues]
        return ValidationReport(ok=all(i.severity != "error" for i in issues), issues=issues)

    def to_bytes(self, deterministic: bool = True) -> bytes:
        buf = io.BytesIO()
        self._prs.save(buf)
        raw = buf.getvalue()
        return canonicalize_pptx_bytes(raw) if deterministic else raw

    def save(self, path: str | Path, deterministic: bool = True) -> None:
        write_bytes(path, self.to_bytes(deterministic=deterministic))

    def fingerprint(self) -> str:
        """Return SHA-256 fingerprint of deterministic package bytes."""
        return sha256(self.to_bytes(deterministic=True)).hexdigest()

    def apply_operations(
        self,
        batch: OperationBatch,
        *,
        dry_run: bool = False,
        transactional: bool = True,
    ) -> OperationReport:
        events: list[OperationEvent] = []
        snapshot: bytes | None = None
        applied_count = 0

        if transactional and not dry_run:
            snapshot = self.to_bytes(deterministic=False)

        try:
            for idx, operation in enumerate(batch.operations):
                op_name = operation.op
                if dry_run:
                    events.append(OperationEvent(index=idx, op=op_name, status="planned"))
                    continue

                start = time.perf_counter()
                self._apply_op(operation)
                duration_ms = int((time.perf_counter() - start) * 1000)
                events.append(
                    OperationEvent(
                        index=idx,
                        op=op_name,
                        status="applied",
                        duration_ms=duration_ms,
                    )
                )
                applied_count += 1

            return OperationReport(
                ok=True,
                dry_run=dry_run,
                events=events,
                applied_count=applied_count,
                failed_index=None,
            )
        except Exception as exc:  # noqa: BLE001
            failed_index = len(events)
            events.append(
                OperationEvent(
                    index=failed_index,
                    op=getattr(batch.operations[failed_index], "op", "unknown")
                    if batch.operations
                    else "unknown",
                    status="failed",
                    detail=str(exc),
                )
            )
            if snapshot is not None:
                self._prs = load_presentation(io.BytesIO(snapshot))
            return OperationReport(
                ok=False,
                dry_run=dry_run,
                events=events,
                applied_count=applied_count,
                failed_index=failed_index,
            )

    # Maps each operation type to the method name that handles it.
    # _apply_op() uses model_dump(exclude={"op"}) as **kwargs, so
    # op model field names must match method parameter names exactly.
    _OP_DISPATCH: dict[type, str] = {
        AddSlideOp: "add_slide",
        AddTextOp: "add_text",
        AddBarChartOp: "add_bar_chart",
        AddLineChartOp: "add_line_chart",
        AddPieChartOp: "add_pie_chart",
        AddAreaChartOp: "add_area_chart",
        AddDoughnutChartOp: "add_doughnut_chart",
        AddScatterChartOp: "add_scatter_chart",
        AddRadarChartOp: "add_radar_chart",
        AddBubbleChartOp: "add_bubble_chart",
        AddComboChartOverlayOp: "add_combo_chart_overlay",
        AddImageOp: "add_image",
        AddTableOp: "add_table",
        AddRectangleOp: "add_rectangle",
        AddRoundedRectangleOp: "add_rounded_rectangle",
        AddOvalOp: "add_oval",
        AddLineShapeOp: "add_line_shape",
        AddRawShapeXmlOp: "add_raw_shape_xml",
        AddIconOp: "add_icon",
        ReplaceTextOp: "replace_text",
        DeleteSlideOp: "delete_slide",
        MoveSlideOp: "move_slide",
        AddNotesOp: "add_notes",
        AddMediaOp: "add_media",
        SetCorePropertiesOp: "set_core_properties",
        UpdateTableCellOp: "update_table_cell",
        UpdateChartDataOp: "update_chart_data",
        SetSlideBackgroundOp: "set_slide_background",
        SetPlaceholderTextOp: "set_placeholder_text",
        SetPlaceholderImageOp: "set_placeholder_image",
        SetTitleSubtitleOp: "set_title_subtitle",
        SetSemanticTextOp: "set_semantic_text",
        SetChartLegendOp: "set_chart_legend",
        SetChartStyleOp: "set_chart_style",
        SetImageCropOp: "set_image_crop",
        SetChartTitleOp: "set_chart_title",
        SetChartAxisTitlesOp: "set_chart_axis_titles",
        SetChartAxisOptionsOp: "set_chart_axis_options",
        SetChartDataLabelsOp: "set_chart_data_labels",
        SetChartDataLabelsStyleOp: "set_chart_data_labels_style",
        SetChartAxisScaleOp: "set_chart_axis_scale",
        SetChartPlotStyleOp: "set_chart_plot_style",
        SetChartSeriesStyleOp: "set_chart_series_style",
        SetChartSeriesLineStyleOp: "set_chart_series_line_style",
        SetChartSecondaryAxisOp: "set_chart_secondary_axis",
        SetLineSeriesMarkerOp: "set_line_series_marker",
        SetLineSeriesSmoothOp: "set_line_series_smooth",
        SetChartSeriesTrendlineOp: "set_chart_series_trendline",
        SetChartSeriesErrorBarsOp: "set_chart_series_error_bars",
        SetChartComboSecondaryMappingOp: "set_chart_combo_secondary_mapping",
    }

    def _apply_op(self, operation: Operation) -> None:
        method_name = self._OP_DISPATCH.get(type(operation))
        if method_name is None:
            raise SlidesError(
                code="UNKNOWN_OPERATION",
                message=f"Unsupported operation type: {type(operation).__name__}",
                path="apply_operations.operation",
                suggested_fix="Provide an operation from the supported schema",
            )
        kwargs = operation.model_dump(exclude={"op"})
        getattr(self, method_name)(**kwargs)

    @staticmethod
    def _ensure_autofit(shape: Any) -> None:
        """Enable normAutofit (shrink-on-overflow) on a placeholder's text body."""
        from lxml.etree import SubElement

        txBody = shape._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
        )
        if txBody is None:
            return
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        # Remove any existing autofit elements first
        for tag in ("noAutofit", "spAutoFit", "normAutofit"):
            for child in txBody.findall(f"{{{ns}}}{tag}"):
                txBody.remove(child)
        SubElement(txBody, f"{{{ns}}}normAutofit")

    def _slide(self, slide_index: int) -> Any:
        if slide_index < 0 or slide_index >= len(self._prs.slides):
            raise SlidesError(
                code="INVALID_SLIDE_INDEX",
                message=f"slide_index {slide_index} out of range",
                path="slide_index",
                suggested_fix=f"use index in [0, {len(self._prs.slides) - 1}]",
            )
        return self._prs.slides[slide_index]

    def _chart(self, *, slide_index: int, chart_index: int, path: str) -> Any:
        slide = self._slide(slide_index)
        chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
        if chart_index < 0 or chart_index >= len(chart_shapes):
            raise SlidesError(
                code="INVALID_CHART_INDEX",
                message=f"chart_index {chart_index} out of range",
                path=f"{path}.chart_index",
                suggested_fix=f"use index in [0, {len(chart_shapes) - 1}]",
            )
        return chart_shapes[chart_index].chart

    def _series(
        self,
        *,
        slide_index: int,
        chart_index: int,
        series_index: int,
        path: str,
    ) -> Any:
        chart = self._chart(slide_index=slide_index, chart_index=chart_index, path=path)
        if not chart.plots:
            raise SlidesError(
                code="CHART_NO_PLOT",
                message="Chart has no plot",
                path=f"{path}.chart_index",
            )
        series = chart.plots[0].series
        if series_index < 0 or series_index >= len(series):
            raise SlidesError(
                code="INVALID_SERIES_INDEX",
                message=f"series_index {series_index} out of range",
                path=f"{path}.series_index",
                suggested_fix=f"use index in [0, {len(series) - 1}]",
            )
        return series[series_index]

    def _set_axis_position(self, axis_element: Any, pos: str) -> None:
        for child in axis_element:
            if child.tag.endswith("axPos"):
                child.set("val", pos)
                return

    def _local_name(self, tag: str) -> str:
        if "}" in tag:
            return tag.split("}", 1)[1]
        return tag

    def _axis_id(self, axis_element: Any) -> str | None:
        for child in axis_element:
            if self._local_name(child.tag) == "axId":
                return child.get("val")
        return None

    def _ser_idx(self, ser_element: Any) -> int | None:
        for child in ser_element:
            if self._local_name(child.tag) == "idx":
                val = child.get("val")
                if val is None:
                    return None
                try:
                    return int(val)
                except ValueError:
                    return None
        return None

    def _plot_axis_pair(
        self,
        plot_element: Any,
        cat_axes: dict[str, Any],
        val_axes: dict[str, Any],
    ) -> tuple[str | None, str | None]:
        cat_id = None
        val_id = None
        for child in plot_element:
            if self._local_name(child.tag) != "axId":
                continue
            axis_id = child.get("val")
            if axis_id is None:
                continue
            if axis_id in cat_axes and cat_id is None:
                cat_id = axis_id
            if axis_id in val_axes and val_id is None:
                val_id = axis_id
        return cat_id, val_id

    def _next_axis_id(self, existing: set[str]) -> str:
        ints: set[int] = set()
        for raw in existing:
            try:
                ints.add(int(raw))
            except ValueError:
                continue
        candidate = 1000000000
        if ints:
            candidate = max(max(ints) + 1, candidate)
            while candidate in ints:
                candidate += 1
        return str(candidate)

    def _ensure_secondary_val_axis(
        self,
        *,
        plot_area: Any,
        cat_axes: dict[str, Any],
        val_axes: dict[str, Any],
        cat_id: str,
        val_primary_id: str,
    ) -> str:
        for val_id, val_axis in val_axes.items():
            cross_ax = None
            ax_pos = None
            for child in val_axis:
                local = self._local_name(child.tag)
                if local == "crossAx":
                    cross_ax = child.get("val")
                if local == "axPos":
                    ax_pos = child.get("val")
            if cross_ax == cat_id and ax_pos == "r":
                return val_id

        primary_axis = val_axes.get(val_primary_id)
        if primary_axis is None:
            raise SlidesError(
                code="CHART_AXIS_RESOLUTION_FAILED",
                message="Primary value axis element not found",
                path="set_chart_combo_secondary_mapping.chart_index",
            )
        new_id = self._next_axis_id(set(cat_axes.keys()) | set(val_axes.keys()))
        secondary = deepcopy(primary_axis)
        for child in secondary:
            local = self._local_name(child.tag)
            if local == "axId":
                child.set("val", new_id)
            elif local == "crossAx":
                child.set("val", cat_id)
            elif local == "axPos":
                child.set("val", "r")
        plot_area.append(secondary)
        val_axes[new_id] = secondary
        return new_id

    def _set_plot_axis_pair(self, plot_element: Any, cat_id: str, val_id: str) -> None:
        ax_nodes = [n for n in plot_element if self._local_name(n.tag) == "axId"]
        for node in ax_nodes:
            plot_element.remove(node)
        cat_node = OxmlElement("c:axId")
        cat_node.set("val", cat_id)
        val_node = OxmlElement("c:axId")
        val_node.set("val", val_id)
        plot_element.append(cat_node)
        plot_element.append(val_node)

    def _insert_series_before_axids(self, plot_element: Any, ser_element: Any) -> None:
        children = list(plot_element)
        for i, child in enumerate(children):
            if self._local_name(child.tag) == "axId":
                plot_element.insert(i, ser_element)
                return
        plot_element.append(ser_element)

    def _add_chart(
        self,
        *,
        slide_index: int,
        chart_type: Any,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
        chart_space_xml: str | None = None,
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        slide = self._slide(slide_index)
        normalized_categories, normalized_series = self._normalize_category_series(
            categories=categories,
            series=series,
        )
        data = CategoryChartData()
        data.categories = normalized_categories
        for name, values in normalized_series:
            data.add_series(name, tuple(values))
        placeholder = self._find_placeholder_by_geometry(
            slide=slide,
            left=left,
            top=top,
            width=width,
            height=height,
            inserter_name="insert_chart",
        )
        if placeholder is not None:
            chart = placeholder.insert_chart(chart_type, data).chart
        else:
            chart = slide.shapes.add_chart(
                chart_type,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
                data,
            ).chart
        if chart_space_xml:
            try:
                chart.part._element = etree.fromstring(chart_space_xml.encode("utf-8"))
            except Exception as exc:  # noqa: BLE001
                raise SlidesError(
                    code="INVALID_CHART_XML",
                    message=f"Failed to apply chart_space_xml: {exc}",
                    path="add_bar_chart.chart_space_xml",
                    suggested_fix="Provide a valid c:chartSpace XML fragment",
                ) from exc

    def _find_placeholder_by_geometry(
        self,
        *,
        slide: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        inserter_name: str,
    ) -> Any | None:
        best_shape = None
        best_score = float("inf")
        for shape in slide.placeholders:
            inserter = getattr(shape, inserter_name, None)
            if not callable(inserter):
                continue
            score = (
                abs(shape.left / 914400 - left)
                + abs(shape.top / 914400 - top)
                + abs(shape.width / 914400 - width)
                + abs(shape.height / 914400 - height)
            )
            if score < best_score:
                best_shape = shape
                best_score = score
        if best_shape is not None and best_score <= 2.0:
            return best_shape
        return None

    def _normalize_category_series(
        self,
        *,
        categories: list[str],
        series: list[tuple[str, list[float | None]]],
    ) -> tuple[list[str], list[tuple[str, list[float | None]]]]:
        if not series:
            raise SlidesError(
                code="CHART_NO_SERIES",
                message="Chart data contains no series",
                path="chart.series",
                suggested_fix="Provide at least one series with one or more values",
            )
        max_len = max((len(values) for _, values in series), default=0)
        if max_len == 0:
            raise SlidesError(
                code="CHART_SERIES_EMPTY",
                message="Chart series values are empty",
                path="chart.series",
                suggested_fix="Provide at least one numeric or null value",
            )

        normalized_categories = list(categories)
        if not normalized_categories:
            normalized_categories = [str(i) for i in range(1, max_len + 1)]

        target_len = max(max_len, len(normalized_categories))
        if len(normalized_categories) < target_len:
            start = len(normalized_categories) + 1
            normalized_categories.extend(str(i) for i in range(start, target_len + 1))
        elif len(normalized_categories) > target_len:
            normalized_categories = normalized_categories[:target_len]

        normalized_series: list[tuple[str, list[float | None]]] = []
        for name, values in series:
            vals = list(values)
            if len(vals) < target_len:
                vals.extend([None] * (target_len - len(vals)))
            elif len(vals) > target_len:
                vals = vals[:target_len]
            normalized_series.append((name, vals))
        return normalized_categories, normalized_series

    def _add_xy_chart(
        self,
        *,
        slide_index: int,
        chart_type: Any,
        series: list[tuple[str, list[tuple[float, float]]]],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        slide = self._slide(slide_index)
        data = XyChartData()
        for name, points in series:
            xy_series = data.add_series(name)
            for x_value, y_value in points:
                xy_series.add_data_point(float(x_value), float(y_value))
        slide.shapes.add_chart(
            chart_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            data,
        )

    def _add_bubble_chart(
        self,
        *,
        slide_index: int,
        chart_type: Any,
        series: list[tuple[str, list[tuple[float, float, float]]]],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        slide = self._slide(slide_index)
        data = BubbleChartData()
        for name, points in series:
            bubble_series = data.add_series(name)
            for x_value, y_value, size_value in points:
                bubble_series.add_data_point(
                    float(x_value),
                    float(y_value),
                    float(size_value),
                )
        slide.shapes.add_chart(
            chart_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            data,
        )
