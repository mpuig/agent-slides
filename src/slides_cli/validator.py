from __future__ import annotations

import io
import posixpath
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx.presentation import Presentation

PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


@dataclass(slots=True)
class ValidationIssue:
    code: str
    message: str
    severity: str
    path: str | None = None


@dataclass(slots=True)
class ValidationReport:
    ok: bool
    issues: list[ValidationIssue]


def validate_presentation(prs: Presentation) -> ValidationReport:
    issues: list[ValidationIssue] = []

    # Baseline structural checks; deeper schema/semantic validation can be layered later.
    if prs.slide_width is None or prs.slide_height is None:
        issues.append(
            ValidationIssue(
                code="MISSING_SLIDE_SIZE",
                message="Presentation is missing slide size metadata",
                severity="error",
            )
        )

    if len(prs.slide_layouts) == 0:
        issues.append(
            ValidationIssue(
                code="NO_LAYOUTS",
                message="Presentation has no slide layouts",
                severity="error",
            )
        )

    if len(prs.slides) == 0:
        issues.append(
            ValidationIssue(
                code="NO_SLIDES",
                message="Presentation has no slides",
                severity="warning",
            )
        )

    # Detect unresolved slide relationships and common placeholder tokens.
    slide_id_list = prs.slides._sldIdLst
    seen_slide_ids: set[int] = set()
    for i, slide_id in enumerate(slide_id_list):
        if slide_id.id in seen_slide_ids:
            issues.append(
                ValidationIssue(
                    code="DUPLICATE_SLIDE_ID",
                    message=f"Duplicate slide id detected: {slide_id.id}",
                    severity="error",
                    path=f"slides[{i}]",
                )
            )
        seen_slide_ids.add(slide_id.id)
        try:
            prs.part.related_part(slide_id.rId)
        except KeyError:
            issues.append(
                ValidationIssue(
                    code="BROKEN_SLIDE_REL",
                    message=f"Slide relationship missing: {slide_id.rId}",
                    severity="error",
                    path=f"slides[{i}]",
                )
            )

    for s_idx, slide in enumerate(prs.slides):
        try:
            _ = slide.slide_layout
        except Exception:  # noqa: BLE001
            issues.append(
                ValidationIssue(
                    code="BROKEN_SLIDE_LAYOUT_REL",
                    message="Slide has a broken slide-layout relationship",
                    severity="error",
                    path=f"slides[{s_idx}]",
                )
            )
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                pass
            else:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text or ""
                    if "{{" in text and "}}" in text:
                        issues.append(
                            ValidationIssue(
                                code="UNRESOLVED_TEMPLATE_TOKEN",
                                message="Unresolved template token found in text",
                                severity="warning",
                                path=f"slides[{s_idx}]",
                            )
                        )

            if getattr(shape, "has_chart", False):
                try:
                    chart = shape.chart
                    if chart.plots:
                        plot = chart.plots[0]
                        expected = len(plot.categories)
                        for series in plot.series:
                            if len(series.values) != expected:
                                issues.append(
                                    ValidationIssue(
                                        code="CHART_DATA_LENGTH_MISMATCH",
                                        message="Chart series length does not match category count",
                                        severity="warning",
                                        path=f"slides[{s_idx}]",
                                    )
                                )
                                break
                except Exception:  # noqa: BLE001
                    issues.append(
                        ValidationIssue(
                            code="CHART_VALIDATION_FAILED",
                            message="Unable to validate chart data integrity",
                            severity="warning",
                            path=f"slides[{s_idx}]",
                        )
                    )

            if getattr(shape, "has_table", False):
                try:
                    table = shape.table
                    if len(table.rows) == 0 or len(table.columns) == 0:
                        issues.append(
                            ValidationIssue(
                                code="INVALID_TABLE_SHAPE",
                                message="Table shape has zero rows or columns",
                                severity="error",
                                path=f"slides[{s_idx}]",
                            )
                        )
                    else:
                        expected_cols = len(table.columns)
                        for r_idx, row in enumerate(table.rows):
                            if len(row.cells) != expected_cols:
                                issues.append(
                                    ValidationIssue(
                                        code="TABLE_ROW_WIDTH_MISMATCH",
                                        message="Table row cell count does not match column count",
                                        severity="error",
                                        path=f"slides[{s_idx}].tables[{r_idx}]",
                                    )
                                )
                                break
                except Exception:  # noqa: BLE001
                    issues.append(
                        ValidationIssue(
                            code="TABLE_VALIDATION_FAILED",
                            message="Unable to validate table integrity",
                            severity="warning",
                            path=f"slides[{s_idx}]",
                        )
                    )

    return ValidationReport(ok=all(i.severity != "error" for i in issues), issues=issues)


def _resolve_relationship_target(
    rels_path: str, target: str, names: set[str]
) -> tuple[str | None, str | None]:
    target_no_fragment = target.split("#", 1)[0]
    if not target_no_fragment:
        return None, None

    if target_no_fragment.startswith("/"):
        resolved = target_no_fragment.lstrip("/")
    else:
        if rels_path == "_rels/.rels":
            source_dir = ""
        elif "/_rels/" in rels_path and rels_path.endswith(".rels"):
            parent, rel_file = rels_path.split("/_rels/", 1)
            source_name = rel_file.removesuffix(".rels")
            source_path = posixpath.normpath(posixpath.join(parent, source_name))
            source_dir = posixpath.dirname(source_path)
        else:
            return None, f"Unrecognized relationships location: {rels_path}"
        resolved = posixpath.normpath(posixpath.join(source_dir, target_no_fragment))

    if resolved.startswith("../") or resolved == "..":
        return None, f"Relationship escapes package root: {target}"
    if resolved not in names:
        return None, f"Relationship target part missing: {resolved}"
    return resolved, None


def validate_package_bytes(
    data: bytes,
    *,
    xsd_dir: str | Path | None = None,
    require_xsd: bool = False,
) -> ValidationReport:
    issues: list[ValidationIssue] = []
    required_entries = {
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
        "ppt/_rels/presentation.xml.rels",
    }

    try:
        with zipfile.ZipFile(io.BytesIO(data), "r") as zf:
            names = set(zf.namelist())
            missing = sorted(required_entries - names)
            for entry in missing:
                issues.append(
                    ValidationIssue(
                        code="MISSING_PACKAGE_PART",
                        message=f"Required package part missing: {entry}",
                        severity="error",
                        path=entry,
                    )
                )
            issues.extend(_validate_content_types(zf, names))
            issues.extend(_validate_core_package_semantics(names))
            for rels_path in sorted(name for name in names if name.endswith(".rels")):
                try:
                    root = ET.fromstring(zf.read(rels_path))
                except ET.ParseError:
                    issues.append(
                        ValidationIssue(
                            code="INVALID_RELS_XML",
                            message=f"Relationships part is not valid XML: {rels_path}",
                            severity="error",
                            path=rels_path,
                        )
                    )
                    continue

                if root.tag != f"{{{PACKAGE_REL_NS}}}Relationships":
                    issues.append(
                        ValidationIssue(
                            code="INVALID_RELS_ROOT",
                            message=f"Relationships part has unexpected root tag: {root.tag}",
                            severity="error",
                            path=rels_path,
                        )
                    )
                    continue

                rel_ids: set[str] = set()
                for rel in root.findall(f"{{{PACKAGE_REL_NS}}}Relationship"):
                    rel_id = rel.attrib.get("Id", "")
                    if not rel_id:
                        issues.append(
                            ValidationIssue(
                                code="MISSING_REL_ID",
                                message="Relationship element missing Id attribute",
                                severity="error",
                                path=rels_path,
                            )
                        )
                    elif rel_id in rel_ids:
                        issues.append(
                            ValidationIssue(
                                code="DUPLICATE_REL_ID",
                                message=f"Duplicate relationship Id: {rel_id}",
                                severity="error",
                                path=rels_path,
                            )
                        )
                    rel_ids.add(rel_id)

                    if "Type" not in rel.attrib:
                        issues.append(
                            ValidationIssue(
                                code="MISSING_REL_TYPE",
                                message="Relationship element missing Type attribute",
                                severity="error",
                                path=rels_path,
                            )
                        )
                    mode = rel.attrib.get("TargetMode", "")
                    if mode == "External":
                        continue
                    target = rel.attrib.get("Target", "")
                    _, rel_err = _resolve_relationship_target(rels_path, target, names)
                    if rel_err is not None:
                        issues.append(
                            ValidationIssue(
                                code="BROKEN_REL_TARGET",
                                message=rel_err,
                                severity="error",
                                path=rels_path,
                            )
                        )
            issues.extend(
                _validate_with_xsd(
                    zf,
                    names,
                    xsd_dir=xsd_dir,
                    require_xsd=require_xsd,
                )
            )
    except zipfile.BadZipFile:
        issues.append(
            ValidationIssue(
                code="INVALID_ZIP",
                message="Output is not a valid ZIP archive",
                severity="error",
            )
        )

    return ValidationReport(ok=all(i.severity != "error" for i in issues), issues=issues)


def _validate_content_types(zf: zipfile.ZipFile, names: set[str]) -> list[ValidationIssue]:
    issues: list[ValidationIssue] = []
    if "[Content_Types].xml" not in names:
        return issues

    try:
        root = ET.fromstring(zf.read("[Content_Types].xml"))
    except ET.ParseError:
        issues.append(
            ValidationIssue(
                code="INVALID_CONTENT_TYPES_XML",
                message="[Content_Types].xml is not valid XML",
                severity="error",
                path="[Content_Types].xml",
            )
        )
        return issues

    if root.tag != f"{{{CONTENT_TYPES_NS}}}Types":
        issues.append(
            ValidationIssue(
                code="INVALID_CONTENT_TYPES_ROOT",
                message=f"[Content_Types].xml has unexpected root tag: {root.tag}",
                severity="error",
                path="[Content_Types].xml",
            )
        )
        return issues

    override_types: dict[str, str] = {}
    for node in root.findall(f"{{{CONTENT_TYPES_NS}}}Override"):
        part_name = node.attrib.get("PartName", "")
        content_type = node.attrib.get("ContentType", "")
        if not part_name.startswith("/"):
            continue
        override_types[part_name.lstrip("/")] = content_type

    required_overrides = {
        "ppt/presentation.xml": (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
        ),
        "docProps/core.xml": "application/vnd.openxmlformats-package.core-properties+xml",
    }
    for part, expected in required_overrides.items():
        actual = override_types.get(part)
        if actual is None:
            issues.append(
                ValidationIssue(
                    code="MISSING_CONTENT_TYPE_OVERRIDE",
                    message=f"Missing content-type override for {part}",
                    severity="error",
                    path="[Content_Types].xml",
                )
            )
        elif actual != expected:
            issues.append(
                ValidationIssue(
                    code="CONTENT_TYPE_MISMATCH",
                    message=f"Unexpected content type for {part}: {actual}",
                    severity="error",
                    path="[Content_Types].xml",
                )
            )

    slide_parts = sorted(
        name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")
    )
    for slide in slide_parts:
        if slide not in override_types:
            issues.append(
                ValidationIssue(
                    code="MISSING_SLIDE_CONTENT_TYPE_OVERRIDE",
                    message=f"Slide part missing content-type override: {slide}",
                    severity="warning",
                    path="[Content_Types].xml",
                )
            )
    return issues


def _validate_core_package_semantics(names: set[str]) -> list[ValidationIssue]:
    issues: list[ValidationIssue] = []
    if "docProps/core.xml" not in names:
        issues.append(
            ValidationIssue(
                code="MISSING_CORE_PROPERTIES",
                message="Missing core properties part: docProps/core.xml",
                severity="error",
                path="docProps/core.xml",
            )
        )
    if "docProps/app.xml" not in names:
        issues.append(
            ValidationIssue(
                code="MISSING_APP_PROPERTIES",
                message="Missing extended properties part: docProps/app.xml",
                severity="warning",
                path="docProps/app.xml",
            )
        )
    if "ppt/presProps.xml" not in names:
        issues.append(
            ValidationIssue(
                code="MISSING_PRESENTATION_PROPERTIES",
                message="Missing presentation properties part: ppt/presProps.xml",
                severity="warning",
                path="ppt/presProps.xml",
            )
        )
    if "ppt/viewProps.xml" not in names:
        issues.append(
            ValidationIssue(
                code="MISSING_VIEW_PROPERTIES",
                message="Missing view properties part: ppt/viewProps.xml",
                severity="warning",
                path="ppt/viewProps.xml",
            )
        )
    return issues


def _validate_with_xsd(
    zf: zipfile.ZipFile,
    names: set[str],
    *,
    xsd_dir: str | Path | None,
    require_xsd: bool,
) -> list[ValidationIssue]:
    issues: list[ValidationIssue] = []
    try:
        import lxml.etree as LET
    except ImportError:
        severity = "error" if require_xsd else "warning"
        issues.append(
            ValidationIssue(
                code="XSD_VALIDATOR_UNAVAILABLE",
                message="lxml is required for XSD validation but is not installed",
                severity=severity,
                path=None,
            )
        )
        return issues

    if xsd_dir is None:
        severity = "error" if require_xsd else "warning"
        issues.append(
            ValidationIssue(
                code="XSD_DIR_NOT_CONFIGURED",
                message="XSD validation requested but no schema directory provided",
                severity=severity,
                path=None,
            )
        )
        return issues

    schema_root = Path(xsd_dir)
    if not schema_root.exists() or not schema_root.is_dir():
        issues.append(
            ValidationIssue(
                code="XSD_DIR_INVALID",
                message=f"XSD directory does not exist or is not a directory: {schema_root}",
                severity="error",
                path=str(schema_root),
            )
        )
        return issues

    routes: list[tuple[re.Pattern[str], list[str]]] = [
        (re.compile(r"^ppt/presentation\.xml$"), ["pml-presentation.xsd"]),
        (re.compile(r"^ppt/slides/slide\d+\.xml$"), ["pml-slide.xsd"]),
        (re.compile(r"^ppt/slideLayouts/slideLayout\d+\.xml$"), ["pml-slideLayout.xsd"]),
        (re.compile(r"^ppt/slideMasters/slideMaster\d+\.xml$"), ["pml-slideMaster.xsd"]),
        (re.compile(r"^ppt/notesSlides/notesSlide\d+\.xml$"), ["pml-notesSlide.xsd"]),
        (re.compile(r"^ppt/notesMasters/notesMaster\d+\.xml$"), ["pml-notesMaster.xsd"]),
        (re.compile(r"^ppt/handoutMasters/handoutMaster\d+\.xml$"), ["pml-handoutMaster.xsd"]),
        (re.compile(r"^ppt/theme/theme\d+\.xml$"), ["dml-theme.xsd"]),
        (re.compile(r"^ppt/tableStyles\.xml$"), ["pml-tableStyles.xsd"]),
        (re.compile(r"^ppt/viewProps\.xml$"), ["pml-viewProps.xsd"]),
        (re.compile(r"^ppt/presProps\.xml$"), ["pml-presProps.xsd"]),
        (re.compile(r"^ppt/commentAuthors\.xml$"), ["pml-commentAuthors.xsd"]),
        (re.compile(r"^ppt/comments/comment\d+\.xml$"), ["pml-comments.xsd"]),
        (re.compile(r"^ppt/charts/chart\d+\.xml$"), ["dml-chart.xsd"]),
        (
            re.compile(r"^docProps/core\.xml$"),
            ["opc-coreProperties.xsd", "core-properties.xsd"],
        ),
        (
            re.compile(r"^docProps/app\.xml$"),
            ["opc-extendedProperties.xsd", "extended-properties.xsd"],
        ),
    ]
    schema_cache: dict[str, Any] = {}

    for part in sorted(name for name in names if name.endswith(".xml")):
        schema_candidates: list[str] | None = None
        for pattern, candidates in routes:
            if pattern.match(part):
                schema_candidates = candidates
                break
        if schema_candidates is None:
            severity = "error" if require_xsd else "warning"
            issues.append(
                ValidationIssue(
                    code="XSD_PART_UNROUTED",
                    message=f"No XSD route configured for XML part: {part}",
                    severity=severity,
                    path=part,
                )
            )
            continue

        selected_schema: str | None = None
        for schema_name in schema_candidates:
            if (schema_root / schema_name).exists():
                selected_schema = schema_name
                break
        if selected_schema is None:
            severity = "error" if require_xsd else "warning"
            issues.append(
                ValidationIssue(
                    code="XSD_SCHEMA_NOT_FOUND",
                    message=(
                        f"Missing schema for part {part}; expected one of: "
                        + ", ".join(schema_candidates)
                    ),
                    severity=severity,
                    path=part,
                )
            )
            continue

        schema_path = schema_root / selected_schema
        if selected_schema not in schema_cache:
            try:
                schema_doc = LET.parse(str(schema_path))
                schema_cache[selected_schema] = LET.XMLSchema(schema_doc)
            except Exception as exc:  # noqa: BLE001
                issues.append(
                    ValidationIssue(
                        code="XSD_SCHEMA_LOAD_FAILED",
                        message=f"Failed to load schema {selected_schema}: {exc}",
                        severity="error",
                        path=str(schema_path),
                    )
                )
                continue

        try:
            xml_doc = LET.fromstring(zf.read(part))
        except Exception as exc:  # noqa: BLE001
            issues.append(
                ValidationIssue(
                    code="XSD_XML_PARSE_FAILED",
                    message=f"Failed to parse XML for XSD validation: {exc}",
                    severity="error",
                    path=part,
                )
            )
            continue

        schema = schema_cache[selected_schema]
        is_valid = schema.validate(xml_doc)
        if not is_valid:
            error_log = schema.error_log
            details = "unknown validation error"
            if len(error_log) > 0:
                first = error_log[0]
                line = getattr(first, "line", None)
                column = getattr(first, "column", None)
                domain = getattr(first, "domain_name", "unknown")
                err_type = getattr(first, "type_name", "unknown")
                details = (
                    f"{first.message} "
                    f"[line={line}, column={column}, domain={domain}, type={err_type}]"
                )
            issues.append(
                ValidationIssue(
                    code="XSD_VALIDATION_ERROR",
                    message=f"Schema validation failed for {part}: {details}",
                    severity="error",
                    path=part,
                )
            )
    return issues
