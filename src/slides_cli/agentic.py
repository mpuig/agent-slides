from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, cast

from pptx import Presentation as load_presentation
from pydantic import BaseModel, ConfigDict, Field

from .model import AddMediaOp, OperationBatch, SetPlaceholderImageOp


class StorySpineEntry(BaseModel):
    model_config = ConfigDict(extra="forbid")

    role: str
    archetype_id: str


class DesignProfile(BaseModel):
    model_config = ConfigDict(extra="forbid")

    name: str = "default"
    template_path: str | None = None
    icon_pack_dir: str | None = None
    content_layout_catalog_path: str | None = None
    archetypes_catalog_path: str | None = None
    primary_color_hex: str | None = None
    max_bullets_per_slide: int = 6
    max_text_chars_per_shape: int = 800
    required_sections: list[str] = Field(default_factory=list)
    allow_external_images: bool = True
    allowed_fonts: list[str] = Field(default_factory=list)
    min_font_size_pt: int = 10
    max_font_size_pt: int = 44
    allowed_colors_hex: list[str] = Field(default_factory=list)
    asset_roots: list[str] = Field(default_factory=list)
    min_margin_in: float = 0.15
    max_overlap_ratio: float = 0.2
    enforce_template_layouts: bool = True
    enforce_template_layout_sequence: bool = False
    enforce_template_placeholders: bool = False
    enforce_template_placeholder_text: bool = False
    template_tolerance_in: float = 0.02
    enforce_template_geometry: bool = True
    geometry_tolerance_in: float = 0.35
    template_shape_kind_tolerance: int = 1
    placeholder_tolerance_in: float = 0.35
    enforce_theme_font_inheritance: bool = False
    require_visual_on_content_slides: bool = True
    require_source_for_data_slides: bool = True
    enforce_chart_style_contract: bool = False
    require_chart_axis_titles: bool = False
    enforce_no_effects: bool = False
    enforce_no_transitions: bool = False
    enforce_no_animations: bool = False
    allowed_image_extensions: list[str] = Field(
        default_factory=lambda: [".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"]
    )
    chart_max_axis_font_size_pt: int = 12
    chart_max_legend_font_size_pt: int = 11
    table_max_cell_font_size_pt: int = 16
    default_font_size_pt: int = 20
    slide_margin_left_in: float = 0.69
    slide_margin_top_in: float = 1.35
    slide_margin_right_in: float = 0.69
    slide_margin_bottom_in: float = 0.55
    text_color_light: str = "FFFFFF"
    text_color_dark: str = "333333"
    story_spine: list[StorySpineEntry] | None = None
    default_layout_index: int = 6


class SlidePlan(BaseModel):
    model_config = ConfigDict(extra="forbid")

    slide_number: int
    story_role: str
    archetype_id: str
    action_title: str
    key_points: list[str] = Field(default_factory=list)
    visual_hint: str | None = None
    source_note: str | None = None


class DeckPlan(BaseModel):
    model_config = ConfigDict(extra="forbid")

    deck_title: str
    brief: str
    audience: str | None = None
    objective: str | None = None
    slides: list[SlidePlan]
    assumptions: list[str] = Field(default_factory=list)


class SlidesDocument(BaseModel):
    model_config = ConfigDict(extra="forbid")

    plan: DeckPlan
    ops: OperationBatch | None = None


class RequestEnvelope(BaseModel):
    model_config = ConfigDict(extra="forbid", populate_by_name=True)

    input: str | None = None
    template: str | None = None
    output: str | None = None
    no_save: bool = False
    dry_run: bool = False
    ops_context_only: bool = False
    no_transaction: bool = False
    non_deterministic: bool = False

    ops: OperationBatch | None = None
    slides: SlidesDocument | None = None
    brief_text: str | None = None
    profile: DesignProfile | None = None

    inspect: bool = False
    find_query: str | None = None
    find_limit: int | None = None
    lint_design: bool = False
    assets_verify: bool = False
    qa_run: bool = False

    do_validate: bool = Field(default=False, alias="validate")
    deep_validate: bool = False
    summary: bool = False
    fingerprint: bool = False

    edit_query: str | None = None
    edit_replacement: str | None = None
    edit_slide: int | None = None
    edit_slide_id: str | None = None
    edit_slide_uid: str | None = None
    edit_shape_id: int | None = None
    edit_shape_uid: str | None = None

    transform_slide: int | None = None
    transform_slide_id: str | None = None
    transform_slide_uid: str | None = None
    transform_to: str | None = None

    fields: str | None = None
    ndjson: bool = False
    verbose: bool = False
    compact: bool = False
    page_size: int | None = None
    page_token: str | None = None
    page_all: bool = False


@dataclass(slots=True)
class DesignIssue:
    code: str
    severity: str
    message: str
    slide_index: int | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "code": self.code,
            "severity": self.severity,
            "message": self.message,
            "slide_index": self.slide_index,
        }


@dataclass(slots=True)
class AssetIssue:
    code: str
    severity: str
    message: str
    op_index: int | None = None
    path: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "code": self.code,
            "severity": self.severity,
            "message": self.message,
            "op_index": self.op_index,
            "path": self.path,
        }


@dataclass(slots=True)
class TemplateStyle:
    default_layout: str | None
    use_placeholders: bool
    layout_names: list[str]
    layout_index_by_name: dict[str, int]
    placeholder_roles_by_layout: dict[str, set[str]]
    placeholder_indices_by_layout: dict[str, dict[str, list[int]]]
    placeholder_boxes_by_layout: dict[str, dict[str, list[tuple[float, float, float, float]]]]


@dataclass(slots=True)
class ContentLayoutCatalog:
    ranked_layouts_by_archetype: dict[str, list[str]]


@dataclass(slots=True)
class ArchetypesCatalog:
    visual_hints: dict[str, str]
    action_title_templates: dict[str, str]
    fallback_layout_tokens: dict[str, list[str]]
    default_geometries: dict[str, dict[str, Any]]


def plan_from_brief(
    brief: str,
    *,
    profile: DesignProfile | None = None,
    audience: str | None = None,
    objective: str | None = None,
    slide_count: int | None = None,
) -> DeckPlan:
    profile = profile or DesignProfile()
    archetypes_catalog = _load_archetypes_catalog(profile)
    cleaned = brief.strip()
    if not cleaned:
        cleaned = "Untitled presentation"
    title = _infer_title(cleaned)
    desired_slides = (
        max(4, min(20, int(slide_count)))
        if slide_count is not None
        else _infer_slide_count(cleaned)
    )
    key_points = _extract_key_points(cleaned)
    inferred_audience = audience or _infer_audience(cleaned)
    inferred_objective = objective or _infer_objective(cleaned)
    topic = _infer_topic(cleaned)

    story = _story_spine(desired_slides, profile=profile)
    slides: list[SlidePlan] = []
    for i, (role, archetype_id) in enumerate(story, start=1):
        points = _points_for_slide(role=role, key_points=key_points, ordinal=i, topic=topic)
        slides.append(
            SlidePlan(
                slide_number=i,
                story_role=role,
                archetype_id=archetype_id,
                action_title=_action_title(
                    role=role,
                    deck_title=title,
                    ordinal=i,
                    topic=topic,
                    archetypes_catalog=archetypes_catalog,
                ),
                key_points=points,
                visual_hint=_visual_hint(
                    archetype_id=archetype_id,
                    archetypes_catalog=archetypes_catalog,
                ),
            )
        )

    assumptions = [
        f"Applied default bullet cap of {profile.max_bullets_per_slide}.",
        "Used deterministic archetype mapping because no explicit slide map was provided.",
    ]
    if audience is None and inferred_audience is None:
        assumptions.append("Audience not specified; used generic executive audience.")
    if objective is None and inferred_objective is None:
        assumptions.append("Objective not specified; defaulted to decision-support framing.")

    return DeckPlan(
        deck_title=title,
        brief=cleaned,
        audience=inferred_audience or "Executive stakeholders",
        objective=inferred_objective or "Support a concrete decision with clear next steps.",
        slides=slides,
        assumptions=assumptions,
    )


def compile_plan_to_operations(
    plan: DeckPlan,
    *,
    profile: DesignProfile | None = None,
) -> OperationBatch:
    profile = profile or DesignProfile()
    template_style = _load_template_style(profile)
    content_layout_catalog = _load_content_layout_catalog(profile)
    archetypes_catalog = _load_archetypes_catalog(profile)
    template_layout_catalog = _load_template_layout_catalog(profile)
    operations: list[dict[str, Any]] = []

    for i, slide in enumerate(plan.slides):
        slide_index = i
        layout_name, layout_index = _resolve_layout_hint(
            archetype_id=slide.archetype_id,
            template_style=template_style,
            content_layout_catalog=content_layout_catalog,
            archetypes_catalog=archetypes_catalog,
            default_layout_index=profile.default_layout_index,
        )
        operations.append(
            {
                "op": "add_slide",
                "layout_index": layout_index,
                "layout_name": layout_name,
            }
        )

        if profile.primary_color_hex and slide.archetype_id in {"title_slide", "section_divider"}:
            operations.append(
                {
                    "op": "set_slide_background",
                    "slide_index": slide_index,
                    "color_hex": profile.primary_color_hex,
                }
            )

        layout_geometry = _resolve_layout_geometry(
            layout_name=layout_name,
            archetype_id=slide.archetype_id,
            template_layout_catalog=template_layout_catalog,
            archetypes_catalog=archetypes_catalog,
            profile=profile,
        )

        operations.extend(
            _compile_slide(
                slide_index=slide_index,
                slide=slide,
                plan=plan,
                profile=profile,
                template_style=template_style,
                layout_name=layout_name,
                layout_geometry=layout_geometry,
            )
        )

    return OperationBatch.model_validate({"operations": operations})


def lint_design(
    *,
    deck_index: dict[str, Any],
    profile: DesignProfile,
    template_index: dict[str, Any] | None = None,
    batch: OperationBatch | None = None,
) -> dict[str, Any]:
    issues: list[DesignIssue] = []

    summary = deck_index.get("summary", {})
    if summary.get("unresolved_token_count", 0) > 0:
        issues.append(
            DesignIssue(
                code="UNRESOLVED_TEMPLATE_TOKEN",
                severity="error",
                message="Deck contains unresolved template tokens.",
            )
        )

    slides = deck_index.get("slides", [])
    deck_meta = deck_index.get("deck", {})
    deck_w = float(deck_meta.get("slide_width_in", 0.0) or 0.0)
    deck_h = float(deck_meta.get("slide_height_in", 0.0) or 0.0)
    allowed_fonts = {f.strip().lower() for f in profile.allowed_fonts if f.strip()}
    allowed_colors = {c.strip().upper() for c in profile.allowed_colors_hex if c.strip()}
    template_fonts: set[str] = set()
    if template_index is not None:
        for template_slide in template_index.get("slides", []):
            for template_shape in template_slide.get("shapes", []):
                for font_name in template_shape.get("font_names", []):
                    name = str(font_name).strip().lower()
                    if name:
                        template_fonts.add(name)
    if not allowed_fonts and template_fonts:
        allowed_fonts = template_fonts
    searchable_titles: list[str] = []
    for slide in slides:
        slide_index = slide.get("slide_index")
        title = str(slide.get("title", "")).strip()
        if title:
            searchable_titles.append(title.lower())
        if not title:
            issues.append(
                DesignIssue(
                    code="MISSING_SLIDE_TITLE",
                    severity="warning",
                    message="Slide has no inferred title.",
                    slide_index=slide_index,
                )
            )

        bullets = 0
        boxes: list[tuple[float, float, float, float]] = []
        visual_shape_count = 0
        has_source_line = False
        has_chart_or_table = False
        for shape in slide.get("shapes", []):
            text = str(shape.get("text", ""))
            has_exception = "[design-exception:" in text.lower()
            kind = str(shape.get("kind", "")).strip().lower()
            if kind in {"chart", "table", "image"}:
                visual_shape_count += 1
            if kind in {"chart", "table"}:
                has_chart_or_table = True
            if kind == "chart":
                ax_sizes = shape.get("chart_axis_font_sizes_pt", [])
                if ax_sizes:
                    for ax_size in ax_sizes:
                        if ax_size > profile.chart_max_axis_font_size_pt:
                            issues.append(
                                DesignIssue(
                                    code="CHART_AXIS_FONT_TOO_LARGE",
                                    severity="warning",
                                    message=(
                                        f"Chart axis font {ax_size}pt exceeds max "
                                        f"{profile.chart_max_axis_font_size_pt}pt."
                                    ),
                                    slide_index=slide_index,
                                )
                            )
                else:
                    issues.append(
                        DesignIssue(
                            code="CHART_AXIS_FONT_UNSET",
                            severity="warning",
                            message=(
                                "Chart axis font size not explicitly set; "
                                "PowerPoint defaults may be too large for the layout. "
                                "Use set_chart_axis_options with font_size."
                            ),
                            slide_index=slide_index,
                        )
                    )
                legend_size = shape.get("chart_legend_font_size_pt")
                if legend_size is not None and legend_size > profile.chart_max_legend_font_size_pt:
                        issues.append(
                            DesignIssue(
                                code="CHART_LEGEND_FONT_TOO_LARGE",
                                severity="warning",
                                message=(
                                    f"Chart legend font {legend_size}pt exceeds max "
                                    f"{profile.chart_max_legend_font_size_pt}pt."
                                ),
                                slide_index=slide_index,
                            )
                        )
                if shape.get("chart_legend_outside_plot"):
                    issues.append(
                        DesignIssue(
                            code="CHART_LEGEND_OVERLAP_RISK",
                            severity="warning",
                            message=(
                                "Chart legend is set to overlay the plot area "
                                "(include_in_layout=False); may obscure data."
                            ),
                            slide_index=slide_index,
                        )
                    )
            if kind == "table":
                cell_sizes = shape.get("table_cell_font_sizes_pt", [])
                if cell_sizes:
                    for cell_size in cell_sizes:
                        if cell_size > profile.table_max_cell_font_size_pt:
                            issues.append(
                                DesignIssue(
                                    code="TABLE_CELL_FONT_TOO_LARGE",
                                    severity="warning",
                                    message=(
                                        f"Table cell font {cell_size}pt exceeds max "
                                        f"{profile.table_max_cell_font_size_pt}pt."
                                    ),
                                    slide_index=slide_index,
                                )
                            )
                            break
                else:
                    row_count = shape.get("table_row_count", 0)
                    if row_count > 3:
                        issues.append(
                            DesignIssue(
                                code="TABLE_CELL_FONT_UNSET",
                                severity="warning",
                                message=(
                                    f"Table has {row_count} rows but no explicit cell "
                                    "font size; PowerPoint defaults (~18pt) may be too "
                                    "large. Use update_table_cell with font_size."
                                ),
                                slide_index=slide_index,
                            )
                        )
            if "source:" in text.lower():
                has_source_line = True
            left = shape.get("left")
            top = shape.get("top")
            width = shape.get("width")
            height = shape.get("height")
            for line in text.splitlines():
                stripped = line.strip()
                if _looks_like_bullet(stripped):
                    bullets += 1
            if all(isinstance(v, (int, float)) for v in (left, top, width, height)):
                left_in = float(left)
                top_in = float(top)
                width_in = float(width)
                height_in = float(height)
                boxes.append((left_in, top_in, left_in + width_in, top_in + height_in))
                if not has_exception and deck_w > 0 and deck_h > 0:
                    margin = max(0.0, profile.min_margin_in)
                    if (
                        left_in < margin
                        or top_in < margin
                        or (left_in + width_in) > (deck_w - margin)
                        or (top_in + height_in) > (deck_h - margin)
                    ):
                        issues.append(
                            DesignIssue(
                                code="SHAPE_OUT_OF_BOUNDS",
                                severity="warning",
                                message=(
                                    "Shape exceeds slide bounds/min margin. "
                                    f"margin={profile.min_margin_in}in"
                                ),
                                slide_index=slide_index,
                            )
                        )
            if len(text) > profile.max_text_chars_per_shape and not has_exception:
                issues.append(
                    DesignIssue(
                        code="TEXT_SHAPE_TOO_LONG",
                        severity="warning",
                        message=(
                            "Text length exceeds max_text_chars_per_shape="
                            f"{profile.max_text_chars_per_shape}."
                        ),
                        slide_index=slide_index,
                    )
                )

            if (
                profile.enforce_no_effects
                and bool(shape.get("has_effects", False))
                and not has_exception
            ):
                issues.append(
                    DesignIssue(
                        code="VISUAL_EFFECT_NOT_ALLOWED",
                        severity="warning",
                        message="Shape uses visual effects (shadow/glow/reflection/3D).",
                        slide_index=slide_index,
                    )
                )

            for font_name in shape.get("font_names", []):
                if allowed_fonts and font_name.lower() not in allowed_fonts and not has_exception:
                    issues.append(
                        DesignIssue(
                            code="FONT_NOT_ALLOWED",
                            severity="warning",
                            message=f"Font '{font_name}' is not in allowed_fonts.",
                            slide_index=slide_index,
                        )
                    )
            if (
                profile.enforce_theme_font_inheritance
                and shape.get("explicit_font_runs", 0) > 0
                and not has_exception
            ):
                issues.append(
                    DesignIssue(
                        code="FONT_INHERITANCE_VIOLATION",
                        severity="warning",
                        message=(
                            "Shape contains explicit font overrides; expected theme inheritance."
                        ),
                        slide_index=slide_index,
                    )
                )

            for font_size in shape.get("font_sizes_pt", []):
                if has_exception:
                    continue
                if font_size < profile.min_font_size_pt or font_size > profile.max_font_size_pt:
                    issues.append(
                        DesignIssue(
                            code="FONT_SIZE_OUT_OF_RANGE",
                            severity="warning",
                            message=(
                                f"Font size {font_size}pt out of [{profile.min_font_size_pt}, "
                                f"{profile.max_font_size_pt}] range."
                            ),
                            slide_index=slide_index,
                        )
                    )

            for color in shape.get("font_colors_hex", []):
                if allowed_colors and color.upper() not in allowed_colors and not has_exception:
                    issues.append(
                        DesignIssue(
                            code="COLOR_NOT_ALLOWED",
                            severity="warning",
                            message=f"Color '{color}' is not in allowed_colors_hex.",
                            slide_index=slide_index,
                        )
                    )
            if shape.get("italic_runs", 0) > 0 and not has_exception:
                issues.append(
                    DesignIssue(
                        code="ITALIC_NOT_ALLOWED",
                        severity="warning",
                        message="Italic text detected.",
                        slide_index=slide_index,
                    )
                )
            if shape.get("underline_runs", 0) > 0 and not has_exception:
                issues.append(
                    DesignIssue(
                        code="UNDERLINE_NOT_ALLOWED",
                        severity="warning",
                        message="Underline text detected.",
                        slide_index=slide_index,
                    )
                )
        if profile.enforce_no_transitions and bool(slide.get("has_transition", False)):
            issues.append(
                DesignIssue(
                    code="SLIDE_TRANSITION_NOT_ALLOWED",
                    severity="warning",
                    message="Slide transition detected.",
                    slide_index=slide_index,
                )
            )
        if profile.enforce_no_animations and bool(slide.get("has_animation", False)):
            issues.append(
                DesignIssue(
                    code="SLIDE_ANIMATION_NOT_ALLOWED",
                    severity="warning",
                    message="Slide animation/timing detected.",
                    slide_index=slide_index,
                )
            )
        if bullets > profile.max_bullets_per_slide:
            issues.append(
                DesignIssue(
                    code="BULLET_OVERFLOW",
                    severity="warning",
                    message=(
                        f"Slide has {bullets} bullets, exceeding max_bullets_per_slide="
                        f"{profile.max_bullets_per_slide}."
                    ),
                    slide_index=slide_index,
                )
            )
        if boxes:
            overlaps = _overlap_violations(boxes, profile.max_overlap_ratio)
            if overlaps > 0:
                issues.append(
                    DesignIssue(
                        code="SHAPE_OVERLAP",
                        severity="warning",
                        message=(
                            f"Detected {overlaps} overlap(s) above max_overlap_ratio="
                            f"{profile.max_overlap_ratio}."
                        ),
                        slide_index=slide_index,
                    )
                )
        text_overlap_risks = _text_overlap_risks(slide.get("shapes", []), profile.max_overlap_ratio)
        if text_overlap_risks > 0:
            issues.append(
                DesignIssue(
                    code="TEXT_OVERLAP_RISK",
                    severity="warning",
                    message=(
                        f"Detected {text_overlap_risks} text overlap risk(s) from estimated "
                        "content bounds."
                    ),
                    slide_index=slide_index,
                )
            )
        imbalance = _visual_imbalance_score(slide.get("shapes", []), deck_w=deck_w, deck_h=deck_h)
        if imbalance is not None and imbalance > 0.55:
            issues.append(
                DesignIssue(
                    code="VISUAL_IMBALANCE_RISK",
                    severity="warning",
                    message=(
                        f"Slide appears visually unbalanced (score={imbalance:.2f}); "
                        "content weight is concentrated in one region."
                    ),
                    slide_index=slide_index,
                )
            )
        low_contrast_hits = _low_contrast_text_risks(slide.get("shapes", []))
        if low_contrast_hits > 0:
            issues.append(
                DesignIssue(
                    code="LOW_TEXT_CONTRAST_RISK",
                    severity="warning",
                    message=(
                        f"Detected {low_contrast_hits} low-contrast text color risk(s) "
                        "against light backgrounds."
                    ),
                    slide_index=slide_index,
                )
            )
        if profile.require_visual_on_content_slides and visual_shape_count <= 0:
            issues.append(
                DesignIssue(
                    code="MISSING_VISUAL_ELEMENT",
                    severity="warning",
                    message="Slide has no chart/table/image visual element.",
                    slide_index=slide_index,
                )
            )
        if profile.require_source_for_data_slides and has_chart_or_table and not has_source_line:
            issues.append(
                DesignIssue(
                    code="MISSING_SOURCE_LINE",
                    severity="warning",
                    message="Chart/table slide should include a source line (e.g. 'Source: ...').",
                    slide_index=slide_index,
                )
            )

    full_text = "\n".join(searchable_titles)
    for required in profile.required_sections:
        token = required.strip().lower()
        if token and token not in full_text:
            issues.append(
                DesignIssue(
                    code="MISSING_REQUIRED_SECTION",
                    severity="error",
                    message=f"Required section not found in slide titles: {required}",
                )
            )

    if not profile.allow_external_images:
        image_count = int(summary.get("image_count", 0))
        if image_count > 0:
            issues.append(
                DesignIssue(
                    code="EXTERNAL_IMAGES_BLOCKED",
                    severity="warning",
                    message="Profile disallows external images, but deck contains image shapes.",
                )
            )

    if template_index is not None:
        template_deck = template_index.get("deck", {})
        tw = float(template_deck.get("slide_width_in", 0.0) or 0.0)
        th = float(template_deck.get("slide_height_in", 0.0) or 0.0)
        tm = int(template_deck.get("master_count", 0) or 0)
        dm = int(deck_meta.get("master_count", 0) or 0)
        tol = max(0.0, profile.template_tolerance_in)
        if tw > 0 and th > 0 and (abs(deck_w - tw) > tol or abs(deck_h - th) > tol):
            issues.append(
                DesignIssue(
                    code="TEMPLATE_SIZE_MISMATCH",
                    severity="error",
                    message=(
                        f"Deck page size differs from template size beyond tolerance ({tol}in)."
                    ),
                )
            )
        if tm > 0 and dm > 0 and tm != dm:
            issues.append(
                DesignIssue(
                    code="TEMPLATE_MASTER_MISMATCH",
                    severity="warning",
                    message=(f"Master count differs from template (template={tm}, deck={dm})."),
                )
            )
        if profile.enforce_template_layouts:
            deck_layouts_available = {
                str(name).strip() for name in deck_meta.get("layout_names", []) if str(name).strip()
            }
            template_layouts_available = {
                str(name).strip()
                for name in template_deck.get("layout_names", [])
                if str(name).strip()
            }
            missing_layouts = sorted(deck_layouts_available - template_layouts_available)
            if missing_layouts:
                issues.append(
                    DesignIssue(
                        code="TEMPLATE_LAYOUT_REGISTRY_MISMATCH",
                        severity="warning",
                        message=(
                            "Deck layout registry includes layouts not in template: "
                            + ", ".join(missing_layouts[:10])
                        ),
                    )
                )
            template_layouts = {
                str(slide.get("layout_name", "")).strip()
                for slide in template_index.get("slides", [])
                if str(slide.get("layout_name", "")).strip()
            }
            if template_layouts:
                for slide in slides:
                    layout_name = str(slide.get("layout_name", "")).strip()
                    if not layout_name:
                        continue
                    if layout_name not in template_layouts:
                        issues.append(
                            DesignIssue(
                                code="TEMPLATE_LAYOUT_MISMATCH",
                                severity="warning",
                                message=f"Layout '{layout_name}' not present in template.",
                                slide_index=slide.get("slide_index"),
                            )
                        )
            if profile.enforce_template_layout_sequence:
                template_slides = template_index.get("slides", [])
                for idx in range(min(len(slides), len(template_slides))):
                    deck_layout = str(slides[idx].get("layout_name", "")).strip()
                    template_layout = str(template_slides[idx].get("layout_name", "")).strip()
                    if deck_layout and template_layout and deck_layout != template_layout:
                        issues.append(
                            DesignIssue(
                                code="TEMPLATE_LAYOUT_SEQUENCE_MISMATCH",
                                severity="warning",
                                message=(
                                    "Layout sequence differs at slide "
                                    f"{idx + 1}: "
                                    f"deck='{deck_layout}', template='{template_layout}'."
                                ),
                                slide_index=slides[idx].get("slide_index"),
                            )
                        )
        if profile.enforce_template_geometry:
            tolerance = max(0.0, profile.geometry_tolerance_in)
            for slide in slides:
                layout_name = str(slide.get("layout_name", "")).strip()
                if not layout_name:
                    continue
                template_slide = _first_template_slide_by_layout(template_index, layout_name)
                if template_slide is None:
                    continue
                misses = _geometry_miss_count(
                    slide_shapes=slide.get("shapes", []),
                    template_shapes=template_slide.get("shapes", []),
                    tolerance=tolerance,
                )
                if misses > 0:
                    issues.append(
                        DesignIssue(
                            code="TEMPLATE_GEOMETRY_DRIFT",
                            severity="warning",
                            message=(
                                f"{misses} shape(s) deviate from template geometry "
                                f"beyond {tolerance}in tolerance."
                            ),
                            slide_index=slide.get("slide_index"),
                        )
                    )
                kind_misses = _shape_kind_mismatch_count(
                    slide_shapes=slide.get("shapes", []),
                    template_shapes=template_slide.get("shapes", []),
                    tolerance=max(0, profile.template_shape_kind_tolerance),
                )
                if kind_misses > 0:
                    issues.append(
                        DesignIssue(
                            code="TEMPLATE_SHAPE_KIND_MISMATCH",
                            severity="warning",
                            message=(
                                f"Shape-kind usage differs from template by {kind_misses} "
                                "item(s) beyond tolerance."
                            ),
                            slide_index=slide.get("slide_index"),
                        )
                    )
        if profile.enforce_template_placeholders:
            placeholder_tolerance = max(0.0, profile.placeholder_tolerance_in)
            for slide in slides:
                layout_name = str(slide.get("layout_name", "")).strip()
                if not layout_name:
                    continue
                template_slide = _first_template_slide_by_layout(template_index, layout_name)
                if template_slide is None:
                    continue
                misses = _placeholder_miss_count(
                    slide_placeholders=slide.get("placeholders", []),
                    template_placeholders=template_slide.get("placeholders", []),
                    tolerance=placeholder_tolerance,
                    compare_text=profile.enforce_template_placeholder_text,
                )
                if misses > 0:
                    issues.append(
                        DesignIssue(
                            code="TEMPLATE_PLACEHOLDER_MISMATCH",
                            severity="warning",
                            message=(
                                f"{misses} placeholder(s) differ from template beyond "
                                f"{placeholder_tolerance}in tolerance."
                            ),
                            slide_index=slide.get("slide_index"),
                        )
                    )

    if profile.enforce_chart_style_contract and batch is not None:
        chart_ordinals: dict[int, int] = {}
        added_keys: set[tuple[int, int]] = set()
        styled_keys: set[tuple[int, int]] = set()
        titled_keys: set[tuple[int, int]] = set()
        for op in batch.operations:
            op_name = getattr(op, "op", "")
            if op_name in {
                "add_bar_chart",
                "add_line_chart",
                "add_pie_chart",
                "add_area_chart",
                "add_doughnut_chart",
                "add_scatter_chart",
                "add_radar_chart",
                "add_bubble_chart",
            }:
                slide_index = int(getattr(op, "slide_index", 0))
                ordinal = chart_ordinals.get(slide_index, 0)
                chart_ordinals[slide_index] = ordinal + 1
                added_keys.add((slide_index, ordinal))
                continue
            if op_name.startswith("set_chart_"):
                slide_index = int(getattr(op, "slide_index", 0))
                chart_index = int(getattr(op, "chart_index", 0))
                key = (slide_index, chart_index)
                styled_keys.add(key)
                if op_name == "set_chart_axis_titles":
                    titled_keys.add(key)
                if op_name == "set_chart_axis_scale" and (
                    bool(getattr(op, "show_major_gridlines", False))
                    or bool(getattr(op, "show_minor_gridlines", False))
                ):
                    issues.append(
                        DesignIssue(
                            code="CHART_GRIDLINES_NOT_ALLOWED",
                            severity="warning",
                            message=(
                                "Chart axis gridlines enabled; style contract "
                                "expects minimal chart junk."
                            ),
                            slide_index=slide_index,
                        )
                    )
        for slide_index, chart_index in sorted(added_keys):
            key = (slide_index, chart_index)
            if key not in styled_keys:
                issues.append(
                    DesignIssue(
                        code="CHART_STYLE_UNSPECIFIED",
                        severity="warning",
                        message=(
                            "Chart has no explicit style operations "
                            "(set_chart_*) in operation batch."
                        ),
                        slide_index=slide_index,
                    )
                )
            if profile.require_chart_axis_titles and key not in titled_keys:
                issues.append(
                    DesignIssue(
                        code="CHART_AXIS_TITLES_MISSING",
                        severity="warning",
                        message="Chart missing explicit axis titles operation.",
                        slide_index=slide_index,
                    )
                )

    by_severity: dict[str, int] = {}
    by_code: dict[str, int] = {}
    for issue in issues:
        by_severity[issue.severity] = by_severity.get(issue.severity, 0) + 1
        by_code[issue.code] = by_code.get(issue.code, 0) + 1

    return {
        "ok": not any(issue.severity == "error" for issue in issues),
        "issue_count": len(issues),
        "summary": {
            "by_severity": by_severity,
            "by_code": by_code,
        },
        "issues": [issue.to_dict() for issue in issues],
        "profile": profile.model_dump(),
    }


def _overlap_violations(boxes: list[tuple[float, float, float, float]], max_ratio: float) -> int:
    violations = 0
    threshold = max(0.0, max_ratio)
    for i in range(len(boxes)):
        l1, t1, r1, b1 = boxes[i]
        a1 = max(0.0, (r1 - l1) * (b1 - t1))
        if a1 <= 0:
            continue
        for j in range(i + 1, len(boxes)):
            l2, t2, r2, b2 = boxes[j]
            a2 = max(0.0, (r2 - l2) * (b2 - t2))
            if a2 <= 0:
                continue
            inter_w = min(r1, r2) - max(l1, l2)
            inter_h = min(b1, b2) - max(t1, t2)
            if inter_w <= 0 or inter_h <= 0:
                continue
            inter = inter_w * inter_h
            ratio = inter / min(a1, a2)
            if ratio > threshold:
                violations += 1
    return violations


def _text_overlap_risks(shapes: list[dict[str, Any]], max_ratio: float) -> int:
    boxes: list[tuple[float, float, float, float]] = []
    for shape in shapes:
        text = str(shape.get("text", "")).strip()
        if len(text) < 12:
            continue
        left = shape.get("left")
        top = shape.get("top")
        width = shape.get("width")
        height = shape.get("height")
        if not all(isinstance(v, (int, float)) for v in (left, top, width, height)):
            continue
        left_in = float(cast(int | float, left))
        top_in = float(cast(int | float, top))
        width_in = float(cast(int | float, width))
        height_in = float(cast(int | float, height))
        if width_in <= 0.1 or height_in <= 0.1:
            continue
        lines = max(1, text.count("\n") + 1)
        sizes = [float(s) for s in shape.get("font_sizes_pt", []) if isinstance(s, (int, float))]
        font_pt = max(sizes) if sizes else 18.0
        est_text_height = lines * ((font_pt / 72.0) * 1.35) + 0.2
        eff_h = min(height_in, max(0.25, est_text_height))
        boxes.append((left_in, top_in, left_in + width_in, top_in + eff_h))
    threshold = max(0.35, max_ratio)
    return _overlap_violations(boxes, threshold)


def _visual_imbalance_score(
    shapes: list[dict[str, Any]],
    *,
    deck_w: float,
    deck_h: float,
) -> float | None:
    if deck_w <= 0.0 or deck_h <= 0.0:
        return None
    weighted: list[tuple[float, float, float]] = []
    for shape in shapes:
        left = shape.get("left")
        top = shape.get("top")
        width = shape.get("width")
        height = shape.get("height")
        if not all(isinstance(v, (int, float)) for v in (left, top, width, height)):
            continue
        left_in = float(cast(int | float, left))
        top_in = float(cast(int | float, top))
        width_in = float(cast(int | float, width))
        height_in = float(cast(int | float, height))
        if width_in <= 0.1 or height_in <= 0.1:
            continue
        area = width_in * height_in
        if area <= 0:
            continue
        text_chars = int(shape.get("text_chars", 0) or 0)
        kind = str(shape.get("kind", "")).strip().lower()
        kind_boost = 0.25 if kind in {"chart", "table", "image"} else 0.0
        density_boost = min(0.35, max(0.0, text_chars / 200.0) * 0.35)
        weight = area * (1.0 + density_boost + kind_boost)
        cx = left_in + (width_in / 2.0)
        cy = top_in + (height_in / 2.0)
        weighted.append((cx, cy, weight))
    if len(weighted) < 3:
        return None

    left_weight = sum(w for cx, _, w in weighted if cx < deck_w / 2.0)
    right_weight = sum(w for cx, _, w in weighted if cx >= deck_w / 2.0)
    top_weight = sum(w for _, cy, w in weighted if cy < deck_h / 2.0)
    bottom_weight = sum(w for _, cy, w in weighted if cy >= deck_h / 2.0)
    total_lr = left_weight + right_weight
    total_tb = top_weight + bottom_weight
    if total_lr <= 0 or total_tb <= 0:
        return None
    lr = abs(left_weight - right_weight) / total_lr
    tb = abs(top_weight - bottom_weight) / total_tb
    return max(lr, tb)


def _low_contrast_text_risks(shapes: list[dict[str, Any]]) -> int:
    risks = 0
    seen: set[str] = set()
    for shape in shapes:
        text_chars = int(shape.get("text_chars", 0) or 0)
        if text_chars <= 0:
            continue
        for color in shape.get("font_colors_hex", []):
            if not isinstance(color, str):
                continue
            token = color.strip().upper()
            if len(token) != 6 or token in seen:
                continue
            seen.add(token)
            if _contrast_ratio_hex(token, "FFFFFF") < 3.0:
                risks += 1
    return risks


def _contrast_ratio_hex(a_hex: str, b_hex: str) -> float:
    def _channel(v: str) -> float:
        raw = int(v, 16) / 255.0
        return raw / 12.92 if raw <= 0.04045 else ((raw + 0.055) / 1.055) ** 2.4

    def _luminance(hex_rgb: str) -> float:
        r = _channel(hex_rgb[0:2])
        g = _channel(hex_rgb[2:4])
        b = _channel(hex_rgb[4:6])
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    la = _luminance(a_hex)
    lb = _luminance(b_hex)
    lighter = max(la, lb)
    darker = min(la, lb)
    return (lighter + 0.05) / (darker + 0.05)


def _first_template_slide_by_layout(
    template_index: dict[str, Any],
    layout_name: str,
) -> dict[str, Any] | None:
    token = layout_name.strip().lower()
    for slide in template_index.get("slides", []):
        name = str(slide.get("layout_name", "")).strip().lower()
        if name == token:
            return slide
    return None


def _geometry_miss_count(
    *,
    slide_shapes: list[dict[str, Any]],
    template_shapes: list[dict[str, Any]],
    tolerance: float,
) -> int:
    def box(shape: dict[str, Any]) -> tuple[float, float, float, float] | None:
        values = [shape.get(k) for k in ("left", "top", "width", "height")]
        if not all(isinstance(v, (int, float)) for v in values):
            return None
        numeric_values = cast(list[int | float], values)
        left, top, width, height = [float(v) for v in numeric_values]
        return (left, top, width, height)

    template_by_kind: dict[str, list[tuple[float, float, float, float]]] = {}
    for shape in template_shapes:
        kind = str(shape.get("kind", "other")).strip().lower()
        b = box(shape)
        if b is None:
            continue
        template_by_kind.setdefault(kind, []).append(b)

    misses = 0
    for shape in slide_shapes:
        kind = str(shape.get("kind", "other")).strip().lower()
        candidate = box(shape)
        if candidate is None:
            continue
        references = template_by_kind.get(kind, [])
        if not references:
            continue
        cl, ct, cw, ch = candidate
        best = None
        for rl, rt, rw, rh in references:
            dist = max(abs(cl - rl), abs(ct - rt), abs(cw - rw), abs(ch - rh))
            if best is None or dist < best:
                best = dist
        if best is not None and best > tolerance:
            misses += 1
    return misses


def _shape_kind_mismatch_count(
    *,
    slide_shapes: list[dict[str, Any]],
    template_shapes: list[dict[str, Any]],
    tolerance: int,
) -> int:
    deck_counts: dict[str, int] = {}
    template_counts: dict[str, int] = {}
    for shape in slide_shapes:
        kind = str(shape.get("kind", "other")).strip().lower()
        deck_counts[kind] = deck_counts.get(kind, 0) + 1
    for shape in template_shapes:
        kind = str(shape.get("kind", "other")).strip().lower()
        template_counts[kind] = template_counts.get(kind, 0) + 1
    mismatch = 0
    for kind in sorted(set(deck_counts) | set(template_counts)):
        mismatch += abs(deck_counts.get(kind, 0) - template_counts.get(kind, 0))
    return max(0, mismatch - max(0, tolerance))


def _placeholder_miss_count(
    *,
    slide_placeholders: list[dict[str, Any]],
    template_placeholders: list[dict[str, Any]],
    tolerance: float,
    compare_text: bool,
) -> int:
    template_by_idx: dict[int, dict[str, Any]] = {}
    for item in template_placeholders:
        idx = int(item.get("idx", -1))
        if idx >= 0:
            template_by_idx[idx] = item
    misses = 0
    for item in slide_placeholders:
        idx = int(item.get("idx", -1))
        if idx < 0:
            continue
        target = template_by_idx.get(idx)
        if target is None:
            misses += 1
            continue
        if str(item.get("type", "")).strip() != str(target.get("type", "")).strip():
            misses += 1
            continue
        if compare_text:
            current_text = str(item.get("text", "")).strip()
            expected_text = str(target.get("text", "")).strip()
            if current_text != expected_text:
                misses += 1
                continue
        for key in ("left", "top", "width", "height"):
            current = item.get(key)
            expected = target.get(key)
            if not isinstance(current, (int, float)) or not isinstance(expected, (int, float)):
                continue
            if abs(float(current) - float(expected)) > tolerance:
                misses += 1
                break
    return misses


def verify_assets(
    *,
    profile: DesignProfile,
    batch: OperationBatch | None = None,
    input_path: Path | None = None,
    template_path: Path | None = None,
) -> dict[str, Any]:
    issues: list[AssetIssue] = []
    roots = [Path(p).resolve() for p in profile.asset_roots if p.strip()]
    allowed_ext = {ext.lower() for ext in profile.allowed_image_extensions if ext.strip()}

    def check_path(path_value: str, *, op_index: int | None, field_path: str) -> None:
        if any(ch in path_value for ch in ["\x00", "\x1f"]):
            issues.append(
                AssetIssue(
                    code="ASSET_PATH_CONTROL_CHAR",
                    severity="error",
                    message=f"Control character detected in path: {path_value!r}",
                    op_index=op_index,
                    path=field_path,
                )
            )
            return
        resolved = Path(path_value).resolve()
        if roots and not any(resolved == root or root in resolved.parents for root in roots):
            issues.append(
                AssetIssue(
                    code="ASSET_ROOT_VIOLATION",
                    severity="error",
                    message=f"Asset outside allowed roots: {resolved}",
                    op_index=op_index,
                    path=field_path,
                )
            )
        ext = resolved.suffix.lower()
        if allowed_ext and ext and ext not in allowed_ext:
            issues.append(
                AssetIssue(
                    code="ASSET_EXTENSION_BLOCKED",
                    severity="warning",
                    message=f"Extension {ext} is not in allowed_image_extensions.",
                    op_index=op_index,
                    path=field_path,
                )
            )
        if not resolved.exists():
            issues.append(
                AssetIssue(
                    code="ASSET_PATH_NOT_FOUND",
                    severity="warning",
                    message=f"Asset file does not exist: {resolved}",
                    op_index=op_index,
                    path=field_path,
                )
            )

    if input_path is not None:
        check_path(str(input_path), op_index=None, field_path="input")
    if template_path is not None:
        check_path(str(template_path), op_index=None, field_path="template")

    if batch is not None:
        for index, op in enumerate(batch.operations):
            op_name = getattr(op, "op", "")
            if op_name in {"add_image", "add_media"} and hasattr(op, "path"):
                check_path(str(op.path), op_index=index, field_path=f"{op_name}.path")
            if isinstance(op, AddMediaOp) and op.poster_path:
                check_path(
                    str(op.poster_path),
                    op_index=index,
                    field_path="add_media.poster_path",
                )
            if isinstance(op, SetPlaceholderImageOp):
                check_path(str(op.path), op_index=index, field_path="set_placeholder_image.path")

    by_severity: dict[str, int] = {}
    by_code: dict[str, int] = {}
    for issue in issues:
        by_severity[issue.severity] = by_severity.get(issue.severity, 0) + 1
        by_code[issue.code] = by_code.get(issue.code, 0) + 1

    return {
        "ok": not any(issue.severity == "error" for issue in issues),
        "issue_count": len(issues),
        "summary": {
            "by_severity": by_severity,
            "by_code": by_code,
        },
        "issues": [issue.to_dict() for issue in issues],
        "profile": profile.model_dump(),
    }


def _infer_title(brief: str) -> str:
    topic = _infer_topic(brief)
    if topic:
        return f"{topic}: executive briefing"
    for line in brief.splitlines():
        clean = line.strip().lstrip("#").strip()
        if clean:
            return clean[:120]
    return "Untitled Deck"


def _infer_slide_count(brief: str) -> int:
    match = re.search(r"(\d{1,2})\s+slides?", brief, flags=re.IGNORECASE)
    if not match:
        return 8
    count = int(match.group(1))
    return max(4, min(20, count))


def _infer_audience(brief: str) -> str | None:
    match = re.search(r"audience\s*:\s*(.+)", brief, flags=re.IGNORECASE)
    if match:
        value = match.group(1).strip()
        return value[:120] if value else None
    if re.search(r"\b(board|c-?suite|executive|leadership)\b", brief, flags=re.IGNORECASE):
        return "Executive stakeholders"
    return None


def _infer_objective(brief: str) -> str | None:
    match = re.search(r"(objective|goal|decision)\s*:\s*(.+)", brief, flags=re.IGNORECASE)
    if match:
        value = match.group(2).strip()
        return value[:200] if value else None
    return None


def _extract_key_points(brief: str) -> list[str]:
    points: list[str] = []
    objective = _infer_objective(brief)
    if objective:
        points.extend(_split_into_points(objective))

    for line in brief.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        lowered = stripped.lower()
        if lowered.startswith(
            (
                "create ",
                "build ",
                "make ",
                "generate ",
                "i need ",
                "tone:",
                "audience:",
            )
        ):
            continue
        if ":" in stripped:
            head, tail = stripped.split(":", 1)
            if head.strip().lower() in {"objective", "goal", "decision", "include", "focus"}:
                points.extend(_split_into_points(tail.strip()))
                continue
        if stripped.startswith(("-", "*")):
            points.append(stripped[1:].strip())
    if points:
        return [p for p in points if p][:20]

    sentence_chunks = [s.strip() for s in re.split(r"[.!?]\s+", brief) if s.strip()]
    return sentence_chunks[:20]


def _story_spine(
    slide_count: int,
    *,
    profile: DesignProfile | None = None,
) -> list[tuple[str, str]]:
    if profile and profile.story_spine:
        spine = [(e.role, e.archetype_id) for e in profile.story_spine]
        return spine[:slide_count]

    base: list[tuple[str, str]] = [
        ("title", "title_slide"),
        ("executive_summary", "executive_summary"),
        ("agenda", "agenda"),
    ]
    middle_archetypes = [
        ("context", "content_text"),
        ("analysis", "two_column"),
        ("evidence", "bar_chart"),
        ("trend", "line_chart"),
        ("mix", "pie_chart"),
        ("turning_points", "table"),
        ("options", "matrix_2x2"),
        ("process", "process_flow"),
        ("lessons", "content_bullets"),
        ("context_deep", "content_text"),
    ]
    while len(base) + 2 < slide_count:
        base.append(middle_archetypes[(len(base) - 3) % len(middle_archetypes)])
    base.extend(
        [
            ("recommendation", "content_bullets"),
            ("next_steps", "end_slide"),
        ]
    )
    return base[:slide_count]


def _action_title(
    *,
    role: str,
    deck_title: str,
    ordinal: int,
    topic: str | None = None,
    archetypes_catalog: ArchetypesCatalog | None = None,
) -> str:
    subject = topic or deck_title

    if archetypes_catalog is not None:
        template = archetypes_catalog.action_title_templates.get(role)
        if template:
            try:
                return template.format(subject=subject, deck_title=deck_title, ordinal=ordinal)
            except (KeyError, IndexError):
                pass

    mapping = {
        "title": deck_title,
        "executive_summary": f"{subject}: what matters most",
        "agenda": f"Agenda: how the {subject} story unfolds",
        "context": f"{subject}: historical and geopolitical context",
        "analysis": f"{subject}: opposing strategies and constraints",
        "evidence": f"{subject}: scale and commitment over time",
        "trend": f"{subject}: trend in support and sentiment",
        "mix": f"{subject}: distribution of impact across dimensions",
        "turning_points": f"{subject}: turning points that changed outcomes",
        "options": f"{subject}: implications across policy and strategy",
        "process": f"{subject}: timeline of major phases",
        "lessons": f"{subject}: strategic lessons for leaders today",
        "context_deep": f"{subject}: deeper drivers and constraints",
        "recommendation": f"{subject}: key takeaways to retain",
        "next_steps": "Next steps and ownership",
    }
    return mapping.get(role, f"Slide {ordinal}: key message")


def _points_for_slide(
    *,
    role: str,
    key_points: list[str],
    ordinal: int,
    topic: str | None = None,
) -> list[str]:
    subject = topic or "the topic"
    if role == "title":
        return key_points[:1] or [f"Prepared for decision-making on {subject}."]
    if role == "agenda":
        agenda = [
            "Origins and context",
            "Escalation and turning points",
            "Outcomes and long-term impact",
            "Lessons for current strategy",
        ]
        return agenda
    if role == "executive_summary":
        return key_points[:3] or [
            f"{subject} escalated through multi-actor strategic competition",
            f"Military outcomes and political outcomes diverged during {subject}",
            f"{subject} still informs intervention and risk policy debates",
        ]
    if role == "next_steps":
        return [
            "Confirm decision owner",
            "Align team on milestones",
            "Launch execution cadence",
        ]
    if role == "process":
        return key_points[:4] or [
            f"Early phase of {subject}",
            f"Escalation phase of {subject}",
            "De-escalation/transition phase",
            "Post-conflict legacy phase",
        ]
    if role == "turning_points":
        return key_points[:5] or [
            f"Initial trigger in {subject}",
            "Escalation decision point",
            "Pivotal campaign or offensive",
            "Negotiation/de-escalation milestone",
            "Conflict end-state milestone",
        ]
    if role in {"evidence", "trend", "mix"}:
        return key_points[:5] or [
            "Period 1 20",
            "Period 2 35",
            "Period 3 50",
            "Period 4 40",
            "Period 5 25",
        ]
    start = ((ordinal - 1) * 3) % max(1, len(key_points))
    window = key_points[start : start + 3]
    if window:
        return window
    return [
        f"{subject}: key insight 1",
        f"{subject}: key insight 2",
        f"{subject}: key insight 3",
    ]


def _infer_topic(brief: str) -> str | None:
    text = " ".join(brief.split())
    patterns = [
        r"\babout\s+([^,.!?]+)",
        r"\bon\s+([^,.!?]+)",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            topic = match.group(1).strip()
            topic = re.sub(r"^(the)\s+", "", topic, flags=re.IGNORECASE)
            topic = re.sub(
                r"\s+for\s+(an?\s+)?(executive|board|leadership|general)\s+audience$",
                "",
                topic,
                flags=re.IGNORECASE,
            ).strip()
            if topic:
                return topic[:120]
    for line in brief.splitlines():
        clean = line.strip().lstrip("#").strip()
        if clean:
            return clean[:120]
    return None


def _split_into_points(text: str) -> list[str]:
    chunks = re.split(r"[;,]|\\band\\b", text, flags=re.IGNORECASE)
    points: list[str] = []
    for chunk in chunks:
        cleaned = " ".join(chunk.split()).strip(" .:-")
        cleaned = re.sub(r"^(and|or)\s+", "", cleaned, flags=re.IGNORECASE)
        if cleaned:
            points.append(cleaned)
    return points


def _visual_hint(
    *,
    archetype_id: str,
    archetypes_catalog: ArchetypesCatalog | None = None,
) -> str:
    if archetypes_catalog is not None:
        hint = archetypes_catalog.visual_hints.get(archetype_id)
        if hint:
            return hint

    hints = {
        "title_slide": "Large title with concise subtitle",
        "executive_summary": "Three concise bullets",
        "agenda": "Numbered sections",
        "content_text": "One paragraph narrative",
        "two_column": "Left-right comparison",
        "three_column": "Three equal insight columns",
        "four_column": "Four compact columns",
        "timeline": "Chronological milestones",
        "matrix_2x2": "2x2 framing: impact vs feasibility",
        "bar_chart": "Category comparison chart with labels",
        "line_chart": "Trend chart by period",
        "pie_chart": "Share-of-total chart",
        "big_number": "Key KPI with short context",
        "process_flow": "Step-by-step execution path",
        "icon_grid": "Grid of labeled concepts",
        "table": "Structured table",
        "quote": "Emphasized quote block",
        "big_statement": "Large message on full green background",
        "green_panel_text": "Green anchor panel with structured text on the right",
        "section_divider": "Section transition slide",
        "disclaimer": "Legal/disclaimer text",
        "content_bullets": "Single message with supporting bullets",
        "end_slide": "Close with ownership and next action",
    }
    return hints.get(archetype_id, "Simple text-first layout")


def _compile_slide(
    *,
    slide_index: int,
    slide: SlidePlan,
    plan: DeckPlan,
    profile: DesignProfile | None = None,
    template_style: TemplateStyle | None = None,
    layout_name: str | None = None,
    layout_geometry: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    ops: list[dict[str, Any]] = []
    profile = profile or DesignProfile()
    use_template_placeholders = bool(template_style and template_style.use_placeholders)
    normalized_layout = str(layout_name or "").strip().lower()
    layout_roles: set[str] = set()
    if use_template_placeholders and template_style is not None:
        layout_roles = template_style.placeholder_roles_by_layout.get(normalized_layout, set())
        layout_role_indices = template_style.placeholder_indices_by_layout.get(
            normalized_layout, {}
        )
        layout_role_boxes = template_style.placeholder_boxes_by_layout.get(normalized_layout, {})
    else:
        layout_role_indices = {}
        layout_role_boxes = {}

    # Extract content box from layout_geometry for positioning
    cb = layout_geometry.get("content_box", {}) if layout_geometry else {}
    cb_left = float(cb.get("left", profile.slide_margin_left_in))
    cb_top = float(cb.get("top", profile.slide_margin_top_in))
    default_w = 13.33 - profile.slide_margin_left_in - profile.slide_margin_right_in
    default_h = 7.5 - profile.slide_margin_top_in - profile.slide_margin_bottom_in
    cb_width = float(cb.get("width", default_w))
    cb_height = float(cb.get("height", default_h))

    def _layout_supports(role: str) -> bool:
        return role in layout_roles

    def _first_placeholder_idx(role: str) -> int | None:
        values = layout_role_indices.get(role) or []
        return values[0] if values else None

    def _first_placeholder_box(role: str) -> tuple[float, float, float, float] | None:
        values = layout_role_boxes.get(role) or []
        return values[0] if values else None

    def text_op(
        text: str,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        font: int,
        bold: bool = False,
    ) -> dict[str, Any]:
        return {
            "op": "add_text",
            "slide_index": slide_index,
            "text": text,
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "font_size": font,
            "bold": bold,
        }

    points = slide.key_points or ["TBD"]
    body_idx = _first_placeholder_idx("body")
    footer_idx = _first_placeholder_idx("footer")
    body_box = _first_placeholder_box("body")
    default_layout_norm = (
        str(template_style.default_layout or "").strip().lower()
        if template_style is not None
        else ""
    )
    is_default_layout = bool(
        use_template_placeholders
        and default_layout_norm
        and normalized_layout == default_layout_norm
    )

    def _set_body_placeholder(text: str) -> bool:
        if use_template_placeholders and body_idx is not None:
            ops.append(
                {
                    "op": "set_placeholder_text",
                    "slide_index": slide_index,
                    "placeholder_idx": body_idx,
                    "text": text,
                }
            )
            return True
        return False

    def _clear_body_placeholder_if_needed() -> None:
        if not (is_default_layout and use_template_placeholders and body_idx is not None):
            return
        # Some templates ship non-empty default body text in the default content layout.
        # Clear it when an archetype renders custom content outside the body placeholder.
        ops.append(
            {
                "op": "set_placeholder_text",
                "slide_index": slide_index,
                "placeholder_idx": body_idx,
                "text": " ",
            }
        )

    def _source_op(source_text: str) -> dict[str, Any]:
        if use_template_placeholders and footer_idx is not None:
            return {
                "op": "set_placeholder_text",
                "slide_index": slide_index,
                "placeholder_idx": footer_idx,
                "text": source_text,
            }
        return text_op(
            source_text,
            left=cb_left,
            top=cb_top + cb_height - 0.3,
            width=cb_width,
            height=0.3,
            font=9,
        )

    def _fit_body_box(
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        pad: float = 0.1,
    ) -> tuple[float, float, float, float]:
        if not (use_template_placeholders and body_box is not None):
            return left, top, width, height
        b_left, b_top, b_width, b_height = body_box
        return (
            b_left + pad,
            b_top + pad,
            max(0.5, b_width - (2 * pad)),
            max(0.5, b_height - (2 * pad)),
        )

    if slide.archetype_id == "title_slide":
        subtitle = points[0] if points else f"Brief: {plan.deck_title}"
        title_idx = _first_placeholder_idx("title")
        subtitle_idx = _first_placeholder_idx("subtitle")
        if use_template_placeholders and title_idx is not None:
            ops.append(
                {
                    "op": "set_placeholder_text",
                    "slide_index": slide_index,
                    "placeholder_idx": title_idx,
                    "text": slide.action_title,
                }
            )
            if subtitle_idx is not None:
                ops.append(
                    {
                        "op": "set_placeholder_text",
                        "slide_index": slide_index,
                        "placeholder_idx": subtitle_idx,
                        "text": subtitle,
                    }
                )
        else:
            ops.append(
                text_op(
                    slide.action_title,
                    left=cb_left,
                    top=cb_top,
                    width=cb_width,
                    height=1.1,
                    font=38,
                    bold=True,
                )
            )
            ops.append(
                text_op(
                    subtitle,
                    left=cb_left,
                    top=cb_top + 1.3,
                    width=cb_width,
                    height=1.0,
                    font=profile.default_font_size_pt,
                )
            )
        return ops

    if slide.archetype_id == "big_statement":
        statement = slide.action_title if slide.action_title else points[0]
        subtitle = points[0] if points else ""
        bg_color = profile.primary_color_hex or "#29BA74"
        ops.append(
            {
                "op": "set_slide_background",
                "slide_index": slide_index,
                "color_hex": bg_color,
            }
        )
        title_idx = _first_placeholder_idx("title")
        body_idx = _first_placeholder_idx("body")
        if use_template_placeholders and title_idx is not None:
            ops.append(
                {
                    "op": "set_placeholder_text",
                    "slide_index": slide_index,
                    "placeholder_idx": title_idx,
                    "text": statement,
                }
            )
            if subtitle and body_idx is not None:
                ops.append(
                    {
                        "op": "set_placeholder_text",
                        "slide_index": slide_index,
                        "placeholder_idx": body_idx,
                        "text": subtitle,
                    }
                )
            return ops
        ops.append(
            text_op(
                statement,
                left=cb_left + 0.4,
                top=cb_top + 0.6,
                width=cb_width - 0.2,
                height=1.5,
                font=34,
                bold=True,
            )
        )
        if subtitle:
            ops.append(
                text_op(
                    subtitle,
                    left=cb_left + 0.9,
                    top=cb_top + 2.4,
                    width=cb_width - 1.0,
                    height=0.9,
                    font=18,
                )
            )
        return ops

    title_idx = _first_placeholder_idx("title")
    end_slide_title_in_body = (
        slide.archetype_id == "end_slide"
        and use_template_placeholders
        and title_idx is None
        and body_idx is not None
    )
    suppress_title_placeholder = use_template_placeholders and (
        slide.archetype_id in {"section_divider", "disclaimer"}
        or (slide.archetype_id == "agenda" and "agenda" in normalized_layout)
        or (slide.archetype_id == "end_slide" and "end" in normalized_layout)
    )
    if not end_slide_title_in_body and not suppress_title_placeholder:
        if use_template_placeholders and title_idx is not None:
            ops.append(
                {
                    "op": "set_placeholder_text",
                    "slide_index": slide_index,
                    "placeholder_idx": title_idx,
                    "text": slide.action_title,
                }
            )
        else:
            ops.append(
                text_op(
                    slide.action_title,
                    left=cb_left,
                    top=cb_top - 1.0,
                    width=cb_width,
                    height=0.9,
                    font=28,
                    bold=False,
                )
            )

    if slide.archetype_id == "content_text":
        body_text = "\n\n".join(points[:6]) if points else "TBD"
        if _set_body_placeholder(body_text):
            return ops
        blocks = _text_blocks_from_points(points, block_count=3)
        top = cb_top
        for heading, body in blocks:
            if heading:
                ops.append(
                    text_op(
                        heading,
                        left=cb_left,
                        top=top,
                        width=cb_width,
                        height=0.4,
                        font=16,
                        bold=True,
                    )
                )
                top += 0.45
            if body:
                ops.append(
                    text_op(body, left=cb_left, top=top, width=cb_width, height=0.95, font=14)
                )
                top += 1.05
        return ops

    if slide.archetype_id == "green_panel_text":
        _clear_body_placeholder_if_needed()
        # Right panel starts at midpoint of content box
        panel_left = cb_left + cb_width * 0.38
        panel_width = cb_width * 0.58
        kpis = _kpi_points(points[:3])
        body_points = points[3:] if len(points) > 3 else points
        if kpis:
            line = "   ".join(f"{value} {label}" for value, label in kpis)
            ops.append(
                text_op(
                    line,
                    left=panel_left,
                    top=cb_top,
                    width=panel_width,
                    height=0.8,
                    font=profile.default_font_size_pt,
                    bold=True,
                )
            )
        blocks = _text_blocks_from_points(body_points, block_count=3)
        top = cb_top + 0.9
        for heading, body in blocks:
            if heading:
                ops.append(
                    text_op(
                        heading,
                        left=panel_left,
                        top=top,
                        width=panel_width,
                        height=0.4,
                        font=16,
                        bold=True,
                    )
                )
                top += 0.42
            if body:
                ops.append(
                    text_op(body, left=panel_left, top=top, width=panel_width, height=0.85, font=13)
                )
                top += 0.95
        return ops

    if slide.archetype_id == "two_column":
        _clear_body_placeholder_if_needed()
        left_text = "\n".join(f"- {p}" for p in points[:3])
        right_text = "\n".join(f"- {p}" for p in points[3:6] or points[:3])
        if use_template_placeholders and body_box is not None:
            area_left, area_top, area_width, area_height = _fit_body_box(
                cb_left, cb_top, cb_width, cb_height, pad=0.05
            )
            gutter = 0.25
            col_width = max(2.0, (area_width - gutter) / 2)
            ops.append(
                text_op(
                    left_text,
                    left=area_left,
                    top=area_top,
                    width=col_width,
                    height=area_height,
                    font=16,
                )
            )
            ops.append(
                text_op(
                    right_text,
                    left=area_left + col_width + gutter,
                    top=area_top,
                    width=col_width,
                    height=area_height,
                    font=16,
                )
            )
            return ops
        gutter = 0.25
        col_width = max(2.0, (cb_width - gutter) / 2)
        ops.append(
            text_op(left_text, left=cb_left, top=cb_top, width=col_width, height=cb_height, font=18)
        )
        ops.append(
            text_op(
                right_text,
                left=cb_left + col_width + gutter,
                top=cb_top,
                width=col_width,
                height=cb_height,
                font=18,
            )
        )
        return ops

    if slide.archetype_id == "three_column":
        _clear_body_placeholder_if_needed()
        cols = [points[0:2], points[2:4], points[4:6]]
        headers = ["Priority 1", "Priority 2", "Priority 3"]
        if use_template_placeholders and body_box is not None:
            area_left, area_top, area_width, area_height = _fit_body_box(
                cb_left, cb_top, cb_width, cb_height, pad=0.05
            )
            gutter = 0.2
            col_width = max(1.5, (area_width - (2 * gutter)) / 3)
            for idx, col_points in enumerate(cols):
                block = "\n".join([headers[idx], *[f"- {p}" for p in col_points if p]])
                ops.append(
                    text_op(
                        block,
                        left=area_left + (idx * (col_width + gutter)),
                        top=area_top,
                        width=col_width,
                        height=area_height,
                        font=15,
                    )
                )
            return ops
        gutter = 0.2
        col_width = max(1.5, (cb_width - (2 * gutter)) / 3)
        for idx, col_points in enumerate(cols):
            block = "\n".join([headers[idx], *[f"- {p}" for p in col_points if p]])
            left = cb_left + (idx * (col_width + gutter))
            ops.append(
                text_op(block, left=left, top=cb_top, width=col_width, height=cb_height, font=17)
            )
        return ops

    if slide.archetype_id == "four_column":
        _clear_body_placeholder_if_needed()
        cols = [points[0:2], points[2:4], points[4:6], points[6:8]]
        if use_template_placeholders and body_box is not None:
            area_left, area_top, area_width, area_height = _fit_body_box(
                cb_left, cb_top, cb_width, cb_height, pad=0.05
            )
            gutter = 0.18
            col_width = max(1.2, (area_width - (3 * gutter)) / 4)
            for idx, col_points in enumerate(cols):
                title = f"Col {idx + 1}"
                body = "\n".join(f"- {p}" for p in (col_points or points[:2]))
                ops.append(
                    text_op(
                        f"{title}\n{body}",
                        left=area_left + (idx * (col_width + gutter)),
                        top=area_top,
                        width=col_width,
                        height=area_height,
                        font=13,
                    )
                )
            return ops
        gutter = 0.18
        col_width = max(1.2, (cb_width - (3 * gutter)) / 4)
        for idx, col_points in enumerate(cols):
            title = f"Col {idx + 1}"
            body = "\n".join(f"- {p}" for p in (col_points or points[:2]))
            left = cb_left + (idx * (col_width + gutter))
            ops.append(
                text_op(
                    f"{title}\n{body}",
                    left=left,
                    top=cb_top,
                    width=col_width,
                    height=cb_height,
                    font=14,
                )
            )
        return ops

    if slide.archetype_id == "timeline":
        _clear_body_placeholder_if_needed()
        if use_template_placeholders and body_idx is not None:
            timeline_lines = [f"{idx + 1}. {point}" for idx, point in enumerate(points[:6])]
            if _set_body_placeholder("\n".join(timeline_lines)):
                return ops
        cell_w = cb_width / 2 - 0.15
        cell_h = cb_height / 2 - 0.15
        for idx, point in enumerate(points[:4]):
            row = idx // 2
            col = idx % 2
            left = cb_left + (col * (cell_w + 0.3))
            top = cb_top + (row * (cell_h + 0.3))
            ops.append(
                text_op(
                    f"{idx + 1}. {point}",
                    left=left,
                    top=top,
                    width=cell_w,
                    height=cell_h,
                    font=18,
                )
            )
        return ops

    if slide.archetype_id == "agenda":
        agenda_lines = [f"{i + 1}. {point}" for i, point in enumerate(points[:8])]
        if _set_body_placeholder("\n".join(agenda_lines)):
            return ops
        ops.append(
            text_op(
                "\n".join(agenda_lines),
                left=cb_left,
                top=cb_top,
                width=cb_width,
                height=cb_height,
                font=profile.default_font_size_pt,
            )
        )
        return ops

    if slide.archetype_id == "matrix_2x2":
        _clear_body_placeholder_if_needed()
        if use_template_placeholders and body_idx is not None:
            q = [points[i] if i < len(points) else f"Item {i + 1}" for i in range(4)]
            matrix_text = (
                "High impact / High feasibility\n"
                f"- {q[0]}\n"
                "High impact / Low feasibility\n"
                f"- {q[1]}\n"
                "Low impact / Low feasibility\n"
                f"- {q[2]}\n"
                "Low impact / High feasibility\n"
                f"- {q[3]}"
            )
            if _set_body_placeholder(matrix_text):
                return ops
        q = [points[i] if i < len(points) else f"Item {i + 1}" for i in range(4)]
        half_w = cb_width / 2 - 0.15
        half_h = cb_height / 2 - 0.3
        mid_x = cb_left + cb_width / 2 + 0.15
        ops.append(
            text_op(
                "High Impact / High Feasibility",
                left=mid_x,
                top=cb_top,
                width=half_w,
                height=0.4,
                font=13,
                bold=True,
            )
        )
        ops.append(
            text_op(q[0], left=mid_x, top=cb_top + 0.45, width=half_w, height=half_h, font=17)
        )
        ops.append(
            text_op(q[1], left=cb_left, top=cb_top + 0.45, width=half_w, height=half_h, font=17)
        )
        ops.append(
            text_op(
                q[2], left=cb_left, top=cb_top + half_h + 0.6, width=half_w, height=half_h, font=17
            )
        )
        ops.append(
            text_op(
                q[3], left=mid_x, top=cb_top + half_h + 0.6, width=half_w, height=half_h, font=17
            )
        )
        ops.append(
            text_op(
                "Low Feasibility",
                left=cb_left,
                top=cb_top + cb_height - 0.4,
                width=2.4,
                height=0.4,
                font=11,
            )
        )
        ops.append(
            text_op(
                "High Feasibility",
                left=cb_left + cb_width - 2.4,
                top=cb_top + cb_height - 0.4,
                width=2.4,
                height=0.4,
                font=11,
            )
        )
        return ops

    if slide.archetype_id == "bar_chart":
        _clear_body_placeholder_if_needed()
        chart_left, chart_top, chart_width, chart_height = _fit_body_box(
            cb_left, cb_top, cb_width * 0.65, cb_height - 0.5
        )
        categories, values = _bar_chart_data_from_points(points)
        ops.append(
            {
                "op": "add_bar_chart",
                "slide_index": slide_index,
                "categories": categories,
                "series": [["Value", values]],
                "style": "clustered",
                "orientation": "column",
                "left": chart_left,
                "top": chart_top,
                "width": chart_width,
                "height": chart_height,
            }
        )
        ops.extend(
            _default_chart_style_ops(
                slide_index=slide_index,
                chart_index=0,
                category_title="Category",
                value_title="Value",
            )
        )
        source = slide.source_note or "Source: internal analysis"
        ops.append(_source_op(source))
        return ops

    if slide.archetype_id == "line_chart":
        _clear_body_placeholder_if_needed()
        chart_left, chart_top, chart_width, chart_height = _fit_body_box(
            cb_left, cb_top, cb_width, cb_height - 0.5
        )
        categories, values = _bar_chart_data_from_points(points)
        ops.append(
            {
                "op": "add_line_chart",
                "slide_index": slide_index,
                "categories": categories,
                "series": [["Trend", values]],
                "style": "line_markers",
                "left": chart_left,
                "top": chart_top,
                "width": chart_width,
                "height": chart_height,
            }
        )
        ops.extend(
            _default_chart_style_ops(
                slide_index=slide_index,
                chart_index=0,
                category_title="Period",
                value_title="Value",
            )
        )
        source = slide.source_note or "Source: internal analysis"
        ops.append(_source_op(source))
        return ops

    if slide.archetype_id == "pie_chart":
        _clear_body_placeholder_if_needed()
        chart_left, chart_top, chart_width, chart_height = _fit_body_box(
            cb_left, cb_top, cb_width * 0.45, cb_height
        )
        categories, values = _bar_chart_data_from_points(points)
        ops.append(
            {
                "op": "add_pie_chart",
                "slide_index": slide_index,
                "categories": categories[:6],
                "series": [["Share", values[:6]]],
                "style": "pie",
                "left": chart_left,
                "top": chart_top,
                "width": chart_width,
                "height": chart_height,
            }
        )
        ops.extend(
            [
                {
                    "op": "set_chart_legend",
                    "slide_index": slide_index,
                    "chart_index": 0,
                    "visible": True,
                    "position": "right",
                    "include_in_layout": False,
                },
                {
                    "op": "set_chart_data_labels",
                    "slide_index": slide_index,
                    "chart_index": 0,
                    "enabled": True,
                    "show_value": True,
                    "show_category_name": True,
                },
            ]
        )
        source = slide.source_note or "Source: internal analysis"
        ops.append(_source_op(source))
        return ops

    if slide.archetype_id == "table":
        _clear_body_placeholder_if_needed()
        table_left, table_top, table_width, table_height = _fit_body_box(
            cb_left, cb_top, cb_width, cb_height
        )
        rows = [["Item", "Value"]]
        categories, values = _bar_chart_data_from_points(points)
        for c, v in zip(categories[:8], values[:8], strict=False):
            rows.append([c, f"{v:.1f}"])
        ops.append(
            {
                "op": "add_table",
                "slide_index": slide_index,
                "rows": rows,
                "left": table_left,
                "top": table_top,
                "width": table_width,
                "height": table_height,
            }
        )
        source = slide.source_note or "Source: internal analysis"
        ops.append(_source_op(source))
        return ops

    if slide.archetype_id == "big_number":
        _clear_body_placeholder_if_needed()
        if use_template_placeholders and body_idx is not None:
            metric = _first_number(points)
            label = points[0] if points else "Key Metric"
            if _set_body_placeholder(f"{metric:.1f}\n{label}"):
                return ops
        metric = _first_number(points)
        label = points[0] if points else "Key Metric"
        ops.append(
            text_op(
                f"{metric:.1f}",
                left=cb_left,
                top=cb_top,
                width=cb_width * 0.4,
                height=2.0,
                font=64,
                bold=True,
            )
        )
        ops.append(
            text_op(
                label,
                left=cb_left,
                top=cb_top + 2.2,
                width=cb_width,
                height=1.0,
                font=profile.default_font_size_pt,
            )
        )
        return ops

    if slide.archetype_id == "process_flow":
        _clear_body_placeholder_if_needed()
        if use_template_placeholders and body_idx is not None:
            items = points[:6] or ["Step 1", "Step 2", "Step 3", "Step 4"]
            flow_text = "\n".join(f"{i + 1}. {item}" for i, item in enumerate(items))
            if _set_body_placeholder(flow_text):
                return ops
        items = points[:4] or ["Step 1", "Step 2", "Step 3", "Step 4"]
        step_count = len(items)
        gutter = 0.2
        step_w = max(1.5, (cb_width - (step_count - 1) * gutter) / step_count)
        for i, item in enumerate(items):
            left = cb_left + (i * (step_w + gutter))
            ops.append(
                text_op(
                    f"{i + 1}. {item}",
                    left=left,
                    top=cb_top + 0.5,
                    width=step_w,
                    height=cb_height - 0.5,
                    font=16,
                )
            )
        return ops

    if slide.archetype_id == "icon_grid":
        _clear_body_placeholder_if_needed()
        if use_template_placeholders and body_idx is not None:
            items = points[:6] or [
                "Concept A",
                "Concept B",
                "Concept C",
                "Concept D",
                "Concept E",
                "Concept F",
            ]
            if _set_body_placeholder("\n".join(f"- {item}" for item in items)):
                return ops
        items = points[:6] or [
            "Concept A",
            "Concept B",
            "Concept C",
            "Concept D",
            "Concept E",
            "Concept F",
        ]
        grid_cols = 3
        gutter = 0.3
        cell_w = max(1.5, (cb_width - (grid_cols - 1) * gutter) / grid_cols)
        cell_h = cb_height / 2 - gutter / 2
        for i, item in enumerate(items):
            row = i // grid_cols
            col = i % grid_cols
            left = cb_left + (col * (cell_w + gutter))
            top = cb_top + (row * (cell_h + gutter))
            ops.append(text_op(item, left=left, top=top, width=cell_w, height=cell_h, font=16))
        return ops

    if slide.archetype_id == "quote":
        _clear_body_placeholder_if_needed()
        if use_template_placeholders and body_idx is not None:
            quote = points[0] if points else "Quote text"
            source = points[1] if len(points) > 1 else "Source"
            if _set_body_placeholder(f'"{quote}"\n— {source}'):
                return ops
        quote = points[0] if points else "Quote text"
        source = points[1] if len(points) > 1 else "Source"
        ops.append(
            text_op(
                f'"{quote}"',
                left=cb_left,
                top=cb_top,
                width=cb_width,
                height=cb_height * 0.6,
                font=30,
                bold=True,
            )
        )
        ops.append(
            text_op(
                f"— {source}",
                left=cb_left + cb_width * 0.5,
                top=cb_top + cb_height * 0.65,
                width=cb_width * 0.5,
                height=0.8,
                font=16,
            )
        )
        return ops

    if slide.archetype_id == "section_divider":
        if use_template_placeholders and body_idx is not None:
            ops.append(
                {
                    "op": "set_placeholder_text",
                    "slide_index": slide_index,
                    "placeholder_idx": body_idx,
                    "text": slide.action_title,
                }
            )
            return ops
        if use_template_placeholders and title_idx is not None:
            ops.append(
                {
                    "op": "set_placeholder_text",
                    "slide_index": slide_index,
                    "placeholder_idx": title_idx,
                    "text": slide.action_title,
                }
            )
            return ops
        ops.append(
            text_op(
                slide.action_title,
                left=cb_left,
                top=cb_top + cb_height * 0.3,
                width=cb_width,
                height=1.0,
                font=40,
                bold=True,
            )
        )
        return ops

    if slide.archetype_id == "disclaimer":
        body = "\n".join(f"- {p}" for p in points[:8]) or "Confidential"
        if _set_body_placeholder(body):
            return ops
        ops.append(
            text_op(body, left=cb_left, top=cb_top, width=cb_width, height=cb_height, font=13)
        )
        return ops

    if slide.archetype_id == "end_slide":
        # Closing layouts are template-owned; avoid adding any extra text/content.
        return ops

    body = "\n".join(f"- {p}" for p in points[:6])
    if _set_body_placeholder(body):
        if slide.source_note:
            ops.append(_source_op(f"Source: {slide.source_note}"))
        return ops
    ops.append(
        text_op(
            body,
            left=cb_left,
            top=cb_top,
            width=cb_width,
            height=cb_height,
            font=profile.default_font_size_pt,
        )
    )
    if slide.source_note:
        ops.append(_source_op(f"Source: {slide.source_note}"))
    return ops


def _looks_like_bullet(line: str) -> bool:
    if not line:
        return False
    if line.startswith(("-", "*", "•", "–")):
        return True
    return bool(re.match(r"^\d+[\.)]\s+", line))


def _bar_chart_data_from_points(points: list[str]) -> tuple[list[str], list[float]]:
    categories: list[str] = []
    values: list[float] = []
    for idx, point in enumerate(points[:6], start=1):
        matches = list(re.finditer(r"(-?\d+(?:\.\d+)?)", point))
        if matches:
            value = float(matches[-1].group(1))
            label = point[: matches[-1].start()].strip(" :-")
            if not label and len(matches) > 1:
                label = matches[0].group(1)
            if not label:
                label = f"Item {idx}"
        else:
            value = float(10 + (idx * 5))
            label = point[:24] if point else f"Item {idx}"
        categories.append(label[:28] or f"Item {idx}")
        values.append(value)
    if not categories:
        categories = ["A", "B", "C"]
        values = [10.0, 20.0, 30.0]
    return categories, values


def _default_chart_style_ops(
    *,
    slide_index: int,
    chart_index: int,
    category_title: str,
    value_title: str,
) -> list[dict[str, Any]]:
    return [
        {
            "op": "set_chart_legend",
            "slide_index": slide_index,
            "chart_index": chart_index,
            "visible": True,
            "position": "right",
            "include_in_layout": False,
        },
        {
            "op": "set_chart_axis_titles",
            "slide_index": slide_index,
            "chart_index": chart_index,
            "category_title": category_title,
            "value_title": value_title,
        },
        {
            "op": "set_chart_axis_scale",
            "slide_index": slide_index,
            "chart_index": chart_index,
            "show_major_gridlines": False,
            "show_minor_gridlines": False,
        },
    ]


def _first_number(points: list[str]) -> float:
    for point in points:
        match = re.search(r"(-?\d+(?:\.\d+)?)", point)
        if match:
            return float(match.group(1))
    return 42.0


def _layout_hint_for_archetype(
    archetype_id: str,
    *,
    default_layout_index: int = 6,
) -> tuple[str | None, int]:
    di = default_layout_index
    mapping: dict[str, tuple[str | None, int]] = {
        "title_slide": ("Title Slide", 0),
        "executive_summary": ("Title and Text", 1),
        "agenda": ("Agenda Full Width Overview", 1),
        "content_text": ("Title and Text", 1),
        "content_bullets": ("Title and Text", 1),
        "green_panel_text": ("Green one third", di),
        "big_statement": ("Big statement green", di),
        "section_divider": ("Section header line", di),
        "disclaimer": ("Disclaimer", di),
        "end_slide": ("End", di),
        "quote": ("Quote", di),
        "table": ("Title and Text", di),
    }
    return mapping.get(archetype_id, (None, di))


def _resolve_layout_hint(
    *,
    archetype_id: str,
    template_style: TemplateStyle | None,
    content_layout_catalog: ContentLayoutCatalog | None,
    archetypes_catalog: ArchetypesCatalog | None = None,
    default_layout_index: int = 6,
) -> tuple[str | None, int]:
    if template_style is not None:
        def _index_for(layout_name: str) -> int:
            return template_style.layout_index_by_name.get(layout_name.strip().lower(), 0)

        if content_layout_catalog is not None:
            candidates = content_layout_catalog.ranked_layouts_by_archetype.get(archetype_id, [])
            if candidates:
                normalized_template = {
                    str(name).strip().lower(): str(name).strip()
                    for name in template_style.layout_names
                    if str(name).strip()
                }
                for candidate in candidates:
                    hit = normalized_template.get(candidate.strip().lower())
                    if hit:
                        return hit, _index_for(hit)
        inferred = _infer_template_layout_name(
            archetype_id,
            template_style.layout_names,
            archetypes_catalog=archetypes_catalog,
        )
        if inferred:
            return inferred, _index_for(inferred)
        if template_style.default_layout:
            return template_style.default_layout, _index_for(template_style.default_layout)
    return _layout_hint_for_archetype(archetype_id, default_layout_index=default_layout_index)


def _load_template_style(profile: DesignProfile) -> TemplateStyle | None:
    if not profile.template_path:
        return None
    template_path = Path(profile.template_path).expanduser()
    if not template_path.exists():
        return None
    default_layout: str | None = None
    use_placeholders = True
    layout_names: list[str] = []
    layout_index_by_name: dict[str, int] = {}
    placeholder_roles_by_layout: dict[str, set[str]] = {}
    placeholder_indices_by_layout: dict[str, dict[str, list[int]]] = {}
    placeholder_boxes_by_layout: dict[str, dict[str, list[tuple[float, float, float, float]]]] = {}
    try:
        prs = load_presentation(str(template_path))
        for idx, layout in enumerate(prs.slide_layouts):
            raw_name = str(getattr(layout, "name", "") or "").strip()
            if not raw_name:
                continue
            layout_names.append(raw_name)
            layout_index_by_name[raw_name.strip().lower()] = idx
            roles: set[str] = set()
            role_indices: dict[str, list[int]] = {}
            role_boxes: dict[str, list[tuple[float, float, float, float]]] = {}
            for placeholder in layout.placeholders:
                type_name = str(placeholder.placeholder_format.type)
                idx = int(placeholder.placeholder_format.idx)
                box = (
                    float(placeholder.left.inches),
                    float(placeholder.top.inches),
                    float(placeholder.width.inches),
                    float(placeholder.height.inches),
                )
                if type_name.endswith(("TITLE (1)", "CENTER_TITLE (3)")):
                    roles.add("title")
                    role_indices.setdefault("title", []).append(idx)
                    role_boxes.setdefault("title", []).append(box)
                if type_name.endswith("SUBTITLE (4)"):
                    roles.add("subtitle")
                    role_indices.setdefault("subtitle", []).append(idx)
                    role_boxes.setdefault("subtitle", []).append(box)
                if type_name.endswith(("BODY (2)", "OBJECT (7)")):
                    roles.add("body")
                    role_indices.setdefault("body", []).append(idx)
                    role_boxes.setdefault("body", []).append(box)
                if type_name.endswith("FOOTER (15)"):
                    roles.add("footer")
                    role_indices.setdefault("footer", []).append(idx)
                    role_boxes.setdefault("footer", []).append(box)
                if type_name.endswith("DATE (16)"):
                    roles.add("date")
                    role_indices.setdefault("date", []).append(idx)
                    role_boxes.setdefault("date", []).append(box)
                if type_name.endswith("SLIDE_NUMBER (13)"):
                    roles.add("slide_number")
                    role_indices.setdefault("slide_number", []).append(idx)
                    role_boxes.setdefault("slide_number", []).append(box)
            placeholder_roles_by_layout[raw_name.strip().lower()] = roles
            placeholder_indices_by_layout[raw_name.strip().lower()] = role_indices
            placeholder_boxes_by_layout[raw_name.strip().lower()] = role_boxes
    except Exception:  # noqa: BLE001
        layout_names = []
        layout_index_by_name = {}
        placeholder_roles_by_layout = {}
        placeholder_indices_by_layout = {}
        placeholder_boxes_by_layout = {}
    default_layout = _pick_template_default_layout(layout_names)
    return TemplateStyle(
        default_layout=default_layout,
        use_placeholders=use_placeholders,
        layout_names=layout_names,
        layout_index_by_name=layout_index_by_name,
        placeholder_roles_by_layout=placeholder_roles_by_layout,
        placeholder_indices_by_layout=placeholder_indices_by_layout,
        placeholder_boxes_by_layout=placeholder_boxes_by_layout,
    )


def _load_content_layout_catalog(profile: DesignProfile) -> ContentLayoutCatalog | None:
    if not profile.content_layout_catalog_path:
        return None
    path = Path(profile.content_layout_catalog_path)
    if not path.exists():
        return None
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return None
    matrix = payload.get("layout_content_matrix")
    if not isinstance(matrix, dict):
        return None
    archetype_map_raw = payload.get("archetypes", {})
    archetype_to_content_layout: dict[str, str] = {}
    if isinstance(archetype_map_raw, dict):
        for archetype_id, meta in archetype_map_raw.items():
            if not isinstance(meta, dict):
                continue
            content_layout_id = str(meta.get("content_layout", "")).strip()
            if content_layout_id and content_layout_id in matrix:
                archetype_to_content_layout[str(archetype_id)] = content_layout_id

    ranked_layouts_by_archetype: dict[str, list[str]] = {}

    def _ordered_layouts(layout_map: dict[str, Any]) -> list[str]:
        preferred: list[tuple[str, int]] = []
        allowed: list[tuple[str, int]] = []
        for layout_name, verdict in layout_map.items():
            if not isinstance(layout_name, str):
                continue
            status = ""
            score = 0
            if isinstance(verdict, dict):
                status = str(verdict.get("status", "")).strip().lower()
                raw_score = verdict.get("score", 0)
                if isinstance(raw_score, (int, float)):
                    score = int(raw_score)
            elif isinstance(verdict, str):
                status = verdict.strip().lower()
            if status == "preferred":
                preferred.append((layout_name, score))
            elif status == "allowed":
                allowed.append((layout_name, score))
        preferred.sort(key=lambda item: item[1], reverse=True)
        allowed.sort(key=lambda item: item[1], reverse=True)
        ordered = [name for name, _ in preferred]
        for name, _ in allowed:
            if name not in ordered:
                ordered.append(name)
        return ordered

    for content_layout_id, layout_map in matrix.items():
        if not isinstance(layout_map, dict):
            continue
        ordered = _ordered_layouts(layout_map)
        if ordered:
            ranked_layouts_by_archetype[str(content_layout_id)] = ordered

    for archetype_id, content_layout_id in archetype_to_content_layout.items():
        mapped = ranked_layouts_by_archetype.get(content_layout_id)
        if mapped:
            ranked_layouts_by_archetype[archetype_id] = mapped

    if not ranked_layouts_by_archetype:
        return None
    return ContentLayoutCatalog(ranked_layouts_by_archetype=ranked_layouts_by_archetype)


def _load_archetypes_catalog(profile: DesignProfile) -> ArchetypesCatalog | None:
    if not profile.archetypes_catalog_path:
        return None
    path = Path(profile.archetypes_catalog_path)
    if not path.exists():
        return None
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return None
    archetypes = payload.get("archetypes")
    if not isinstance(archetypes, dict):
        return None
    visual_hints: dict[str, str] = {}
    action_title_templates: dict[str, str] = {}
    fallback_tokens: dict[str, list[str]] = {}
    default_geometries: dict[str, dict[str, Any]] = {}
    for archetype_id, meta in archetypes.items():
        if not isinstance(meta, dict):
            continue
        vh = meta.get("visual_hint")
        if isinstance(vh, str) and vh:
            visual_hints[archetype_id] = vh
        att = meta.get("action_title_template")
        if isinstance(att, str) and att:
            action_title_templates[archetype_id] = att
        ft = meta.get("fallback_layout_tokens")
        if isinstance(ft, list) and ft:
            fallback_tokens[archetype_id] = [str(t) for t in ft]
        dg = meta.get("default_geometry")
        if isinstance(dg, dict):
            default_geometries[archetype_id] = dg
    return ArchetypesCatalog(
        visual_hints=visual_hints,
        action_title_templates=action_title_templates,
        fallback_layout_tokens=fallback_tokens,
        default_geometries=default_geometries,
    )


def _load_template_layout_catalog(profile: DesignProfile) -> dict[str, Any] | None:
    if not profile.template_path:
        return None
    template_path = Path(profile.template_path).expanduser()
    catalog_path = template_path.parent / "template_layout.json"
    if not catalog_path.exists():
        return None
    try:
        return json.loads(catalog_path.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return None


def _resolve_layout_geometry(
    *,
    layout_name: str | None,
    archetype_id: str,
    template_layout_catalog: dict[str, Any] | None,
    archetypes_catalog: ArchetypesCatalog | None,
    profile: DesignProfile,
) -> dict[str, Any] | None:
    if template_layout_catalog is not None and layout_name:
        normalized = layout_name.strip().lower()
        for layout in template_layout_catalog.get("layouts", []):
            if str(layout.get("layout_name", "")).strip().lower() == normalized:
                content_box = layout.get("content_box")
                visual_def = layout.get("visual_definition", {})
                color_zones = visual_def.get("color_zones", [])
                editable_regions = visual_def.get("editable_regions", [])
                slide_w = float(template_layout_catalog.get("slide_width_in", 13.33))
                slide_h = float(template_layout_catalog.get("slide_height_in", 7.5))
                result: dict[str, Any] = {
                    "slide_width": slide_w,
                    "slide_height": slide_h,
                }
                if content_box:
                    result["content_box"] = content_box
                elif editable_regions:
                    result["content_box"] = editable_regions[0]
                if color_zones:
                    result["color_zones"] = color_zones
                return result

    if archetypes_catalog is not None:
        dg = archetypes_catalog.default_geometries.get(archetype_id)
        if dg:
            return {
                "content_box": dg,
                "slide_width": 13.33,
                "slide_height": 7.5,
            }

    return {
        "content_box": {
            "left": profile.slide_margin_left_in,
            "top": profile.slide_margin_top_in,
            "width": 13.33 - profile.slide_margin_left_in - profile.slide_margin_right_in,
            "height": 7.5 - profile.slide_margin_top_in - profile.slide_margin_bottom_in,
        },
        "slide_width": 13.33,
        "slide_height": 7.5,
    }


def _pick_template_default_layout(layout_names: list[str]) -> str | None:
    priorities = ["Title and Text", "Title and Content", "Title Only", "Blank"]
    normalized = {name.lower(): name for name in layout_names}
    for candidate in priorities:
        hit = normalized.get(candidate.lower())
        if hit:
            return hit
    return layout_names[0] if layout_names else None


def _infer_template_layout_name(
    archetype_id: str,
    layout_names: list[str],
    *,
    archetypes_catalog: ArchetypesCatalog | None = None,
) -> str | None:
    if not layout_names:
        return None
    normalized = {name.lower(): name for name in layout_names}

    def _find(tokens: list[str], *, include_d_prefix: bool = False) -> str | None:
        for token in tokens:
            for key, raw in normalized.items():
                if token not in key:
                    continue
                if not include_d_prefix and key.startswith("d. "):
                    continue
                return raw
        return None

    # Try catalog tokens first
    if archetypes_catalog is not None:
        catalog_tokens = archetypes_catalog.fallback_layout_tokens.get(archetype_id, [])
        if catalog_tokens:
            hit = _find(catalog_tokens, include_d_prefix=False)
            if hit:
                return hit
            hit = _find(catalog_tokens, include_d_prefix=True)
            if hit:
                return hit

    candidate_map: dict[str, list[str]] = {
        "title_slide": ["title slide"],
        "executive_summary": ["title and text", "title only"],
        "agenda": [
            "agenda full width overview",
            "agenda two-thirds",
            "agenda two-thirds",
            "agenda table of contents",
            "agenda section header overview",
            "title and text",
        ],
        "content_text": ["title and text", "title only"],
        "content_bullets": ["title and text", "title only"],
        "two_column": ["arrow half", "green half", "title and text"],
        "three_column": ["arrow two third", "green two third", "white one third", "title and text"],
        "four_column": ["four column green"],
        "process_flow": [
            "left arrow",
            "green left arrow",
            "arrow two third",
            "green arrow two third",
            "arrow half",
            "green arrow half",
            "title and text",
        ],
        "icon_grid": ["white one third", "green one third", "title and text"],
        "big_number": ["green highlight", "big statement icon", "title and text"],
        "matrix_2x2": ["special gray", "title and text", "blank"],
        "bar_chart": ["title and text", "special gray", "title only"],
        "line_chart": ["title and text", "special gray", "title only"],
        "pie_chart": ["title and text", "special gray", "title only"],
        "table": ["title and text", "special gray", "title only"],
        "quote": ["quote"],
        "disclaimer": ["disclaimer"],
        "end_slide": ["end", "title and text", "title and content", "title only"],
        "section_divider": ["section header line", "section header box", "section"],
        "green_panel_text": ["green one third", "green"],
        "big_statement": ["big statement"],
    }
    direct_tokens = candidate_map.get(archetype_id, [])
    hit = _find(direct_tokens, include_d_prefix=False)
    if hit:
        return hit
    hit = _find(direct_tokens, include_d_prefix=True)
    if hit:
        return hit
    if archetype_id in {
        "bar_chart",
        "line_chart",
        "pie_chart",
        "table",
        "big_number",
        "two_column",
        "three_column",
        "four_column",
        "timeline",
        "matrix_2x2",
        "process_flow",
        "icon_grid",
        "executive_summary",
    }:
        fallback_tokens = ["title and text", "title and content", "title only", "blank"]
        hit = _find(fallback_tokens, include_d_prefix=False)
        if hit:
            return hit
        hit = _find(fallback_tokens, include_d_prefix=True)
        if hit:
            return hit
    return None


def _text_blocks_from_points(points: list[str], *, block_count: int) -> list[tuple[str, str]]:
    if not points:
        return [("Insight", "TBD")]
    blocks: list[tuple[str, str]] = []
    cursor = 0
    while cursor < len(points) and len(blocks) < block_count:
        raw = points[cursor].strip()
        cursor += 1
        if not raw:
            continue
        if ":" in raw:
            head, body = raw.split(":", 1)
            blocks.append((head.strip(), body.strip()))
            continue
        title = raw[:60]
        body_parts = [p.strip() for p in points[cursor : cursor + 2] if p.strip()]
        cursor += min(2, max(0, len(points) - cursor))
        blocks.append((title, " ".join(body_parts) if body_parts else raw))
    return blocks or [("Insight", points[0])]


def _kpi_points(points: list[str]) -> list[tuple[str, str]]:
    kpis: list[tuple[str, str]] = []
    for point in points:
        match = re.search(r"(-?\d+(?:\.\d+)?\s*[%xX$]?)", point)
        if not match:
            continue
        value = match.group(1).strip()
        label = re.sub(r"(-?\d+(?:\.\d+)?\s*[%xX$]?)", "", point).strip(" :-")
        if not label:
            label = "KPI"
        kpis.append((value, label[:22]))
    return kpis[:3]
