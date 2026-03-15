from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from contextlib import suppress
from dataclasses import asdict
from functools import lru_cache
from pathlib import Path
from typing import Any, NoReturn

from lxml import etree
from pptx import Presentation as load_presentation

from . import CONTRACT_VERSION
from .agentic import (
    DesignProfile,
    RequestEnvelope,
    SlidesDocument,
    compile_plan_to_operations,
    lint_design,
    verify_assets,
)
from .api import Presentation
from .model import OperationBatch

SCHEMA_RESOURCE_PATHS: dict[str, Path] = {
    "template-layout": Path(__file__).resolve().parent / "schemas" / "template_layout.schema.json",
    "content-layout": Path(__file__).resolve().parent / "schemas" / "content_layout.schema.json",
    "archetypes": Path(__file__).resolve().parent / "schemas" / "archetypes.schema.json",
}

FIELD_MASKS_PATH = Path(__file__).resolve().parent / "schemas" / "field_masks.json"
DEFAULT_PAGE_SIZE = 25


class CliError(Exception):
    def __init__(
        self,
        *,
        code: str,
        message: str,
        hint: str | None = None,
        path: str | None = None,
        exit_code: int = 2,
    ) -> None:
        super().__init__(message)
        self.code = code
        self.message = message
        self.hint = hint
        self.path = path
        self.exit_code = exit_code


class CliArgumentParser(argparse.ArgumentParser):
    def error(self, message: str) -> NoReturn:
        raise CliError(
            code="ARGUMENT_ERROR",
            message=message,
            hint="Run `slides docs` for command contracts.",
            exit_code=2,
        )


def _emit_error(exc: CliError) -> int:
    payload: dict[str, Any] = {
        "ok": False,
        "error": {
            "code": exc.code,
            "message": exc.message,
        },
    }
    if exc.path:
        payload["error"]["path"] = exc.path
    if exc.hint:
        payload["error"]["hint"] = exc.hint
    print(json.dumps(payload, indent=2), file=sys.stderr)
    return exc.exit_code


def _warn_deprecated_flag(flag: str, replacement: str) -> None:
    print(f"warning: {flag} is deprecated; use {replacement} instead.", file=sys.stderr)


def _add_global_options(parser: argparse.ArgumentParser) -> None:
    """Add output-formatting flags shared by all subcommands."""
    parser.add_argument(
        "--fields", type=str, default=None,
        help=(
            "Comma-separated fields to include in JSON output. "
            "Dot notation traverses into arrays: "
            "e.g. --fields summary,slides.slide_index,slides.title"
        ),
    )
    parser.add_argument("--ndjson", action="store_true", help="Emit list outputs as NDJSON")
    parser.add_argument(
        "--verbose", action="store_true",
        help="Emit full stdout payloads (default is quiet when writing outputs)",
    )
    parser.add_argument(
        "--compact", action="store_true",
        help="Strip null/empty values from JSON output to reduce size",
    )


def _add_pagination_options(parser: argparse.ArgumentParser) -> None:
    """Add pagination flags for list-emitting subcommands."""
    parser.add_argument(
        "--page-size", type=int, default=None, help="Page size for list outputs",
    )
    parser.add_argument(
        "--page-token", type=str, default=None, help="Page token (integer offset)",
    )
    parser.add_argument("--page-all", action="store_true", help="Emit all pages")


def _add_profile_options(parser: argparse.ArgumentParser) -> None:
    """Add --profile / --profile-json flags."""
    parser.add_argument("--profile", type=Path, default=None, help="Design profile JSON file")
    parser.add_argument(
        "--profile-json", type=str, default=None, help="Inline or @file design profile JSON",
    )


def _add_icon_pack_option(parser: argparse.ArgumentParser) -> None:
    parser.add_argument(
        "--icon-pack-dir",
        type=Path,
        default=None,
        help="Private icon pack directory (searched before built-in icons)",
    )


def _add_ops_context_options(parser: argparse.ArgumentParser) -> None:
    """Add --ops-json / --slides-json for context-only loading (qa, lint)."""
    parser.add_argument(
        "--ops-json", type=str, default=None, help="Inline or @file OperationBatch JSON",
    )
    parser.add_argument(
        "--slides-json", type=str, default=None, help="Inline or @file slides document JSON",
    )


def _build_parser() -> argparse.ArgumentParser:
    parser = CliArgumentParser(
        description="slides CLI",
        usage="slides <command> [options]",
    )
    sub = parser.add_subparsers(dest="command", title="commands")

    # --- extract ---
    p = sub.add_parser("extract", help="Extract template contracts from a .pptx")
    p.add_argument("source", type=Path, help="Template or sample .pptx file")
    p.add_argument("--output-dir", type=Path, required=True, help="Output directory for artifacts")
    p.add_argument(
        "--base-template-out", type=Path, default=None,
        help="Also output clean base template",
    )
    p.add_argument("--template-out", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--content-layout-out", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--archetypes-out", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--resolved-manifest-out", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--slides-manifest-out", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--slide-analysis-out", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--layout-preview-dir", type=Path, default=None, help=argparse.SUPPRESS)
    p.add_argument("--screenshots-dir", type=Path, default=None, help=argparse.SUPPRESS)
    _add_global_options(p)

    # --- render ---
    p = sub.add_parser("render", help="Render a deck from a slides document (plan + ops)")
    p.add_argument("--slides-json", type=str, required=True, help="Slides document JSON or @file")
    p.add_argument("--template", type=Path, default=None, help="Template PPTX path")
    p.add_argument("--output", type=Path, default=None, help="Output PPTX path")
    _add_icon_pack_option(p)
    p.add_argument("--dry-run", action="store_true", help="Plan operations without applying")
    p.add_argument("--no-transaction", action="store_true", help="Disable transactional rollback")
    p.add_argument(
        "--non-deterministic", action="store_true",
        help="Disable deterministic ZIP save",
    )
    p.add_argument(
        "--slides-out", type=Path, default=None,
        help="Write normalized slides document JSON",
    )
    _add_profile_options(p)
    _add_global_options(p)

    # --- apply ---
    p = sub.add_parser("apply", help="Apply an operations patch to a deck")
    p.add_argument(
        "input", type=Path, nargs="?", default=None,
        help="Input PPTX (omit to start blank)",
    )
    p.add_argument("--ops-json", type=str, required=True, help="OperationBatch JSON or @file")
    p.add_argument("--output", type=Path, required=True, help="Output PPTX path")
    _add_icon_pack_option(p)
    p.add_argument("--dry-run", action="store_true", help="Plan operations without applying")
    p.add_argument(
        "--no-transaction", action="store_true",
        help="Disable transactional rollback",
    )
    p.add_argument(
        "--non-deterministic", action="store_true",
        help="Disable deterministic ZIP save",
    )
    _add_profile_options(p)
    _add_global_options(p)

    # --- inspect ---
    p = sub.add_parser("inspect", help="Inspect deck structure as JSON")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--out", type=Path, default=None, help="Write inspect JSON to file")
    p.add_argument("--summary", action="store_true", help="Print deck summary")
    p.add_argument(
        "--placeholders", type=int, default=None,
        metavar="SLIDE_INDEX", help="List placeholders for slide index",
    )
    p.add_argument("--fingerprint", action="store_true", help="Print deterministic fingerprint")
    _add_global_options(p)
    _add_pagination_options(p)

    # --- find ---
    p = sub.add_parser("find", help="Search text in a deck")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--query", type=str, required=True, help="Text to search for")
    p.add_argument("--limit", type=int, default=10, help="Max results (default: 10)")
    p.add_argument("--out", type=Path, default=None, help="Write find JSON to file")
    _add_global_options(p)
    _add_pagination_options(p)

    # --- plan-inspect ---
    p = sub.add_parser(
        "plan-inspect",
        help="Inspect slides document (plan + ops) with compact, paginated output",
    )
    p.add_argument("--slides-json", type=str, required=True, help="Slides document JSON or @file")
    p.add_argument("--content-only", action="store_true", help="Include only content slides")
    p.add_argument("--summary-only", action="store_true", help="Emit only summary payload")
    p.add_argument("--out", type=Path, default=None, help="Write plan inspect JSON to file")
    _add_global_options(p)
    _add_pagination_options(p)

    # --- validate ---
    p = sub.add_parser("validate", help="Run OOXML validation")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--deep", action="store_true", help="Run model + package validation")
    p.add_argument("--xsd-dir", type=Path, default=None, help="Optional OOXML XSD directory")
    p.add_argument("--require-xsd", action="store_true", help="Fail when XSDs are unavailable")
    p.add_argument("--fail-on-warning", action="store_true", help="Exit non-zero on warning/error")
    p.add_argument("--fail-on-error", action="store_true", help="Exit non-zero on error")
    _add_global_options(p)

    # --- lint ---
    p = sub.add_parser("lint", help="Lint deck against design profile")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--out", type=Path, default=None, help="Write lint report JSON to file")
    _add_profile_options(p)
    _add_ops_context_options(p)
    _add_global_options(p)

    # --- qa ---
    p = sub.add_parser("qa", help="Run validate + lint + assets checks")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--out", type=Path, default=None, help="Write QA report JSON to file")
    p.add_argument("--deep", action="store_true", help="Run deep validation")
    p.add_argument("--xsd-dir", type=Path, default=None, help="Optional OOXML XSD directory")
    p.add_argument("--require-xsd", action="store_true", help="Fail when XSDs are unavailable")
    _add_profile_options(p)
    _add_ops_context_options(p)
    _add_global_options(p)

    # --- edit ---
    p = sub.add_parser("edit", help="Find-and-replace text in a deck")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--query", type=str, required=True, help="Semantic text query")
    p.add_argument("--replacement", type=str, required=True, help="Replacement text")
    p.add_argument("--output", type=Path, required=True, help="Output PPTX path")
    p.add_argument("--slide", type=int, default=None, help="Target slide index")
    p.add_argument("--slide-id", type=str, default=None, help="Target slide id (slide-N)")
    p.add_argument("--slide-uid", type=str, default=None, help="Target stable slide UID")
    p.add_argument("--shape-id", type=int, default=None, help="Target shape id")
    p.add_argument("--shape-uid", type=str, default=None, help="Target stable shape UID")
    _add_global_options(p)

    # --- transform ---
    p = sub.add_parser("transform", help="Transform a slide to a different archetype")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--to", type=str, choices=["timeline"], required=True, help="Target archetype")
    p.add_argument("--output", type=Path, required=True, help="Output PPTX path")
    p.add_argument("--slide", type=int, default=None, help="Slide index")
    p.add_argument("--slide-id", type=str, default=None, help="Slide id (slide-N)")
    p.add_argument("--slide-uid", type=str, default=None, help="Stable slide UID")
    _add_global_options(p)

    # --- repair ---
    p = sub.add_parser("repair", help="Apply conservative repair to a deck")
    p.add_argument("input", type=Path, help="Input PPTX path")
    p.add_argument("--output", type=Path, required=True, help="Output PPTX path")
    p.add_argument(
        "--non-deterministic", action="store_true",
        help="Disable deterministic ZIP save",
    )
    _add_global_options(p)

    # --- preflight ---
    p = sub.add_parser("preflight", help="Verify project contracts, profile paths, and deps")
    p.add_argument(
        "--project-dir",
        type=Path,
        default=Path("."),
        help="Project directory containing extracted contracts and design-profile.json",
    )
    p.add_argument(
        "--require-optional-deps",
        action="store_true",
        help="Fail if optional preview/screenshot dependencies are unavailable",
    )
    _add_profile_options(p)
    _add_global_options(p)

    # --- docs ---
    p = sub.add_parser("docs", help="Print self-discovery docs")
    p.add_argument(
        "target", type=str, nargs="?", default="markdown",
        help="json | markdown | method:<id> | schema:<id>",
    )
    _add_global_options(p)

    # --- version ---
    sub.add_parser("version", help="Print contract version")

    return parser


def _build_discovery_contract() -> dict[str, Any]:
    external_schemas: dict[str, Any] = {}
    for name, schema_path in SCHEMA_RESOURCE_PATHS.items():
        if schema_path.exists():
            external_schemas[name] = json.loads(schema_path.read_text(encoding="utf-8"))

    schemas = {
        "operation-batch": OperationBatch.model_json_schema(),
        "ops": OperationBatch.model_json_schema(),
        "design-profile": DesignProfile.model_json_schema(),
        "slides-document": SlidesDocument.model_json_schema(),
        "request-envelope": RequestEnvelope.model_json_schema(),
        **external_schemas,
    }
    response_schemas: dict[str, dict[str, Any]] = {
        "operation-report": {
            "type": "object",
            "properties": {
                "ok": {"type": "boolean"},
                "dry_run": {"type": "boolean"},
                "events": {"type": "array"},
                "applied_count": {"type": "integer"},
                "failed_index": {"type": ["integer", "null"]},
            },
        },
        "inspect-payload": {"type": "object"},
        "find-payload": {"type": "object"},
        "edit-transform-payload": {"type": "object"},
        "qa-payload": {"type": "object"},
        "preflight-payload": {"type": "object"},
    }
    methods = [
        {
            "id": "extract",
            "description": (
                "Extract slide-layout catalog and content-layout compatibility from template."
            ),
            "cli": "slides extract <source.pptx> --output-dir <dir>",
            "mutates_deck": False,
            "inputs": ["source", "output-dir", "base-template-out"],
            "outputs": [
                "template_catalog.json", "content_layout_catalog.json",
                "archetypes.json", "resolved_manifest.json",
                "slides_manifest.json", "slide_analysis.json",
                "base_template.pptx",
            ],
            "request_schema": "request-envelope",
            "response_schema": "inspect-payload",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "render",
            "description": "Render a deck from a slides document (plan + ops).",
            "cli": "slides render --slides-json @slides.json --output out.pptx",
            "mutates_deck": True,
            "inputs": ["slides-json", "template", "profile", "icon-pack-dir", "output"],
            "outputs": ["output", "operation-report"],
            "request_schema": "slides-document",
            "response_schema": "operation-report",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": True,
        },
        {
            "id": "apply",
            "description": "Apply an operations patch to an existing deck.",
            "cli": "slides apply <deck.pptx> --ops-json @ops.json --output out.pptx",
            "mutates_deck": True,
            "inputs": ["input", "ops-json", "icon-pack-dir", "output"],
            "outputs": ["output", "operation-report"],
            "request_schema": "operation-batch",
            "response_schema": "operation-report",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": True,
        },
        {
            "id": "inspect",
            "description": "Inspect deck structure with pagination and field masks.",
            "cli": "slides inspect <deck.pptx>",
            "mutates_deck": False,
            "inputs": ["input", "fields", "page-size", "page-token", "page-all"],
            "outputs": ["out", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "inspect-payload",
            "supports_field_masks": True,
            "supports_pagination": True,
            "supports_dry_run": False,
        },
        {
            "id": "find",
            "description": "Find semantically similar text and object selectors in deck.",
            "cli": "slides find <deck.pptx> --query \"text\"",
            "mutates_deck": False,
            "inputs": ["input", "query", "limit", "fields", "page-size", "page-token"],
            "outputs": ["out", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "find-payload",
            "supports_field_masks": True,
            "supports_pagination": True,
            "supports_dry_run": False,
        },
        {
            "id": "plan-inspect",
            "description": (
                "Inspect generated slides document with content filtering and pagination."
            ),
            "cli": "slides plan-inspect --slides-json @slides.json --content-only",
            "mutates_deck": False,
            "inputs": [
                "slides-json",
                "content-only",
                "summary-only",
                "fields",
                "page-size",
                "page-token",
                "page-all",
            ],
            "outputs": ["out", "stdout"],
            "request_schema": "slides-document",
            "response_schema": "inspect-payload",
            "supports_field_masks": True,
            "supports_pagination": True,
            "supports_dry_run": False,
        },
        {
            "id": "edit",
            "description": "Apply semantic text find-and-replace.",
            "cli": (
                "slides edit <deck.pptx>"
                ' --query "old" --replacement "new"'
                " --output out.pptx"
            ),
            "mutates_deck": True,
            "inputs": [
                "input", "query", "replacement", "output",
                "slide", "slide-id", "slide-uid", "shape-id", "shape-uid",
            ],
            "outputs": ["output", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "edit-transform-payload",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "transform",
            "description": "Transform a slide to a different archetype.",
            "cli": "slides transform <deck.pptx> --to timeline --slide-uid UID --output out.pptx",
            "mutates_deck": True,
            "inputs": ["input", "to", "output", "slide", "slide-id", "slide-uid"],
            "outputs": ["output", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "edit-transform-payload",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "validate",
            "description": "Run OOXML validation on a deck.",
            "cli": "slides validate <deck.pptx>",
            "mutates_deck": False,
            "inputs": ["input", "deep", "xsd-dir", "require-xsd"],
            "outputs": ["stdout"],
            "request_schema": "request-envelope",
            "response_schema": "inspect-payload",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "lint",
            "description": "Lint deck against design profile.",
            "cli": "slides lint <deck.pptx> --profile p.json",
            "mutates_deck": False,
            "inputs": ["input", "profile", "ops-json", "slides-json"],
            "outputs": ["out", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "inspect-payload",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "qa",
            "description": "Run validate + lint + assets checks.",
            "cli": "slides qa <deck.pptx> --profile p.json",
            "mutates_deck": False,
            "inputs": ["input", "profile", "ops-json", "slides-json"],
            "outputs": ["out", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "qa-payload",
            "supports_field_masks": True,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "preflight",
            "description": "Verify extracted project artifacts, profile paths, and optional deps.",
            "cli": "slides preflight --project-dir output/project",
            "mutates_deck": False,
            "inputs": ["project-dir", "profile", "profile-json", "require-optional-deps"],
            "outputs": ["stdout"],
            "request_schema": "request-envelope",
            "response_schema": "preflight-payload",
            "supports_field_masks": True,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
        {
            "id": "repair",
            "description": "Apply conservative repair to a deck.",
            "cli": "slides repair <deck.pptx> --output out.pptx",
            "mutates_deck": True,
            "inputs": ["input", "output"],
            "outputs": ["output", "stdout"],
            "request_schema": "request-envelope",
            "response_schema": "inspect-payload",
            "supports_field_masks": False,
            "supports_pagination": False,
            "supports_dry_run": False,
        },
    ]
    workflows = [
        {
            "id": "create-deck",
            "steps": [
                "slides render --slides-json @slides.json --dry-run",
                "slides render --slides-json @slides.json --output out.pptx",
                "slides qa out.pptx --profile p.json",
            ],
        },
        {
            "id": "edit-deck",
            "steps": [
                "slides inspect deck.pptx (or slides find)",
                "slides edit deck.pptx --query ... --replacement ... --output out.pptx",
                "slides qa out.pptx --profile p.json",
            ],
        },
    ]
    return {
        "name": "slides",
        "contract_version": CONTRACT_VERSION,
        "discovery_version": "2.0",
        "entrypoint": "slides docs [json|markdown|method:<id>|schema:<id>]",
        "commands": [
            "extract", "render", "apply", "inspect", "find",
            "plan-inspect", "validate", "lint", "qa", "edit", "transform",
            "repair", "preflight", "docs", "version",
        ],
        "methods": methods,
        "workflows": workflows,
        "schemas": schemas,
        "response_schemas": response_schemas,
    }


def _method_markdown(method: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append(f"# slides Method: {method['id']}")
    lines.append("")
    lines.append(method["description"])
    lines.append("")
    if "cli" in method:
        lines.append(f"- cli: `{method['cli']}`")
    lines.append(f"- mutates_deck: `{str(method['mutates_deck']).lower()}`")
    request_ref = method.get("request_schema_ref", method.get("request_schema"))
    response_ref = method.get("response_schema_ref", method.get("response_schema"))
    lines.append(f"- request_schema: `{request_ref}`")
    lines.append(f"- response_schema: `{response_ref}`")
    lines.append(f"- supports_field_masks: `{str(method['supports_field_masks']).lower()}`")
    lines.append(f"- supports_pagination: `{str(method['supports_pagination']).lower()}`")
    lines.append(f"- supports_dry_run: `{str(method['supports_dry_run']).lower()}`")
    lines.append(f"- inputs: `{', '.join(method['inputs'])}`")
    lines.append(f"- outputs: `{', '.join(method['outputs'])}`")
    lines.append("")
    return "\n".join(lines)


def _schema_markdown(schema_id: str, schema: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append(f"# slides Schema: {schema_id}")
    lines.append("")
    title = str(schema.get("title", schema_id))
    schema_type = str(schema.get("type", "object"))
    lines.append(f"- title: `{title}`")
    lines.append(f"- type: `{schema_type}`")
    required = schema.get("required", [])
    if isinstance(required, list) and required:
        lines.append(f"- required: `{', '.join(str(item) for item in required)}`")
    properties = schema.get("properties", {})
    if isinstance(properties, dict) and properties:
        lines.append("")
        lines.append("## Top-Level Properties")
        for key in sorted(properties):
            prop = properties[key]
            if isinstance(prop, dict):
                prop_type = prop.get("type", "any")
                if isinstance(prop_type, list):
                    prop_type_text = " | ".join(str(item) for item in prop_type)
                else:
                    prop_type_text = str(prop_type)
                lines.append(f"- `{key}`: `{prop_type_text}`")
            else:
                lines.append(f"- `{key}`")
    lines.append("")
    return "\n".join(lines)


def _discovery_markdown(contract: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append("# slides Discovery")
    lines.append("")
    lines.append(f"- contract_version: `{contract['contract_version']}`")
    lines.append(f"- discovery_version: `{contract['discovery_version']}`")
    lines.append(f"- entrypoint: `{contract['entrypoint']}`")
    lines.append("")
    lines.append("## Commands")
    for cmd in contract.get("commands", []):
        lines.append(f"- `slides {cmd}`")
    lines.append("")
    lines.append("## Methods")
    for method in contract["methods"]:
        lines.append(f"### {method['id']}")
        lines.append(method["description"])
        lines.append("")
        if "cli" in method:
            lines.append(f"- cli: `{method['cli']}`")
        lines.append(f"- mutates_deck: `{str(method['mutates_deck']).lower()}`")
        lines.append(f"- request_schema: `{method['request_schema']}`")
        lines.append(f"- response_schema: `{method['response_schema']}`")
        lines.append(f"- supports_field_masks: `{str(method['supports_field_masks']).lower()}`")
        lines.append(f"- supports_pagination: `{str(method['supports_pagination']).lower()}`")
        lines.append(f"- supports_dry_run: `{str(method['supports_dry_run']).lower()}`")
        lines.append(f"- inputs: `{', '.join(method['inputs'])}`")
        lines.append(f"- outputs: `{', '.join(method['outputs'])}`")
        lines.append("")
    lines.append("## Workflows")
    for wf in contract["workflows"]:
        lines.append(f"### {wf['id']}")
        for step in wf["steps"]:
            lines.append(f"1. `{step}`")
        lines.append("")
    lines.append("## Schemas")
    for name in contract["schemas"]:
        lines.append(f"- `{name}`")
    lines.append("")
    return "\n".join(lines)


def _resolve_method_alias(method_id: str) -> str:
    alias = {
        "extract.template_catalog": "extract",
        "generate.slides": "render",
        "render.deck": "render",
        "inspect.deck": "inspect",
        "find.deck": "find",
        "edit.deck": "edit",
        "qa.deck": "qa",
    }
    return alias.get(method_id, method_id)


def _select_discovery_method(contract: dict[str, Any], method_id: str) -> dict[str, Any] | None:
    resolved = _resolve_method_alias(method_id)
    for method in contract["methods"]:
        if method["id"] == resolved:
            return method
    return None


def _resolve_method_contract(contract: dict[str, Any], method: dict[str, Any]) -> dict[str, Any]:
    request_ref = method["request_schema"]
    response_ref = method["response_schema"]
    request_schema = contract["schemas"].get(request_ref)
    response_schema = contract["schemas"].get(response_ref) or contract["response_schemas"].get(
        response_ref
    )
    enriched = dict(method)
    enriched["request_schema_ref"] = request_ref
    enriched["response_schema_ref"] = response_ref
    enriched["request_schema"] = request_schema
    enriched["response_schema"] = response_schema
    return enriched


def _parse_docs_value(raw: str) -> tuple[str, str | None]:
    if raw == "fields":
        return "markdown", "fields"
    if raw.startswith("fields:"):
        payload_ref = raw[len("fields:"):].strip()
        if not payload_ref:
            raise ValueError("docs fields selector is empty")
        fmt = "markdown"
        if ":" in payload_ref:
            candidate_id, candidate_fmt = payload_ref.rsplit(":", 1)
            if candidate_fmt in {"json", "markdown"}:
                payload_ref = candidate_id.strip()
                fmt = candidate_fmt
        return fmt, f"fields:{payload_ref}"
    if raw.startswith("schema:"):
        schema_ref = raw[len("schema:") :].strip()
        if not schema_ref:
            raise ValueError("docs schema selector is empty")
        schema_format = "json"
        if ":" in schema_ref:
            candidate_id, candidate_fmt = schema_ref.rsplit(":", 1)
            if candidate_fmt in {"json", "markdown"}:
                schema_ref = candidate_id.strip()
                schema_format = candidate_fmt
        if not schema_ref:
            raise ValueError("docs schema selector is empty")
        return schema_format, f"schema:{schema_ref}"
    if raw in {"json", "markdown"}:
        return raw, None
    if not raw.startswith("method:"):
        raise ValueError(
            "docs must be one of: json, markdown, method:<id>, "
            "schema:<id>, fields, fields:<payload>"
        )
    method_ref = raw[len("method:") :].strip()
    if not method_ref:
        raise ValueError("docs method selector is empty")
    method_format = "json"
    if ":" in method_ref:
        candidate_id, candidate_fmt = method_ref.rsplit(":", 1)
        if candidate_fmt in {"json", "markdown"}:
            method_ref = candidate_id.strip()
            method_format = candidate_fmt
    if not method_ref:
        raise ValueError("docs method selector is empty")
    return method_format, method_ref


def _placeholder_role(type_name: str) -> str | None:
    token = type_name.upper()
    if token.endswith(("TITLE (1)", "CENTER_TITLE (3)")):
        return "title"
    if token.endswith("SUBTITLE (4)"):
        return "subtitle"
    if token.endswith(("BODY (2)", "OBJECT (7)")):
        return "body"
    if token.endswith("FOOTER (15)"):
        return "footer"
    if token.endswith("DATE (16)"):
        return "date"
    if token.endswith("SLIDE_NUMBER (13)"):
        return "slide_number"
    return None


def _shape_type_token(shape: Any) -> str:
    raw = str(getattr(shape, "shape_type", "")).strip().upper()
    if "TEXT_BOX" in raw:
        return "text_box"
    if "PICTURE" in raw:
        return "picture"
    if "AUTO_SHAPE" in raw:
        return "auto_shape"
    if "FREEFORM" in raw:
        return "freeform"
    if "LINE" in raw:
        return "line"
    if "TABLE" in raw:
        return "table"
    if "CHART" in raw:
        return "chart"
    if "EMBEDDED_OLE_OBJECT" in raw:
        return "embedded_ole_object"
    return "other"


def _shape_text(shape: Any) -> str:
    if not bool(getattr(shape, "has_text_frame", False)):
        return ""
    return " ".join(str(getattr(shape, "text", "") or "").split()).strip()


def _resolve_theme_colors(template_path: str | Path) -> dict[str, str]:
    """Extract theme color scheme from a PPTX, returning scheme-name → hex RGB."""
    import zipfile

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    scheme: dict[str, str] = {}
    with suppress(Exception), zipfile.ZipFile(template_path) as z:
            for name in z.namelist():
                if "theme1.xml" in name:
                    root = etree.fromstring(z.read(name))
                    clr_el = root.find(f".//{{{ns_a}}}clrScheme")
                    if clr_el is None:
                        break
                    for child in clr_el:
                        tag = child.tag.split("}")[-1]
                        srgb = child.find(f"{{{ns_a}}}srgbClr")
                        sys_clr = child.find(f"{{{ns_a}}}sysClr")
                        if srgb is not None:
                            scheme[tag] = srgb.get("val", "")
                        elif sys_clr is not None:
                            scheme[tag] = sys_clr.get("lastClr", sys_clr.get("val", ""))
                    scheme.setdefault("bg1", scheme.get("lt1", "FFFFFF"))
                    scheme.setdefault("bg2", scheme.get("lt2", "F2F2F2"))
                    scheme.setdefault("tx1", scheme.get("dk1", "000000"))
                    scheme.setdefault("tx2", scheme.get("dk2", "000000"))
                    break
    return scheme


def _is_dark(hex_rgb: str) -> bool:
    """Return True if a hex RGB color needs white text for contrast.

    Uses W3C relative luminance with a threshold that matches typical
    presentation design conventions (white text on saturated colors).
    """
    try:
        r, g, b = int(hex_rgb[:2], 16), int(hex_rgb[2:4], 16), int(hex_rgb[4:6], 16)

        def _linear(c: int) -> float:
            s = c / 255
            return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4

        rel_lum = 0.2126 * _linear(r) + 0.7152 * _linear(g) + 0.0722 * _linear(b)
        return rel_lum < 0.4
    except (ValueError, IndexError):
        return False


def _text_color_for_bg(bg_hex: str) -> str:
    """Choose a contrasting text color (hex RGB) for a given background."""
    return "FFFFFF" if _is_dark(bg_hex) else "333333"


# Default margin values used in layout analysis when no profile is available
_DEFAULT_MARGIN_LEFT = 0.69
_DEFAULT_MARGIN_TOP = 1.35
_DEFAULT_MARGIN_RIGHT = 0.69
_DEFAULT_MARGIN_BOTTOM = 0.55


def _extract_color_zones(
    *,
    layout: Any,
    slide_width_in: float,
    slide_height_in: float,
    theme_colors: dict[str, str],
) -> tuple[str | None, list[dict[str, Any]]]:
    """Detect background color and panel color zones for a layout.

    Returns (layout_bg_hex, color_zones) where each zone has:
      region, left, width, bg_color, text_color
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    def _resolve_color_ref(el: Any) -> str | None:
        """Resolve a schemeClr or srgbClr element to hex RGB."""
        if el is None:
            return None
        scheme = el.find(f"{{{ns_a}}}schemeClr")
        if scheme is not None:
            return theme_colors.get(scheme.get("val", ""))
        srgb = el.find(f"{{{ns_a}}}srgbClr")
        if srgb is not None:
            return srgb.get("val")
        return None

    def _resolve_fill(parent: Any) -> str | None:
        """Resolve solidFill or gradFill (first stop) to hex RGB."""
        if parent is None:
            return None
        solid = parent.find(f"{{{ns_a}}}solidFill")
        if solid is not None:
            return _resolve_color_ref(solid)
        grad = parent.find(f"{{{ns_a}}}gradFill")
        if grad is not None:
            gs_list = grad.find(f"{{{ns_a}}}gsLst")
            if gs_list is not None:
                first_gs = gs_list.find(f"{{{ns_a}}}gs")
                if first_gs is not None:
                    return _resolve_color_ref(first_gs)
        return None

    # 1. Layout background
    layout_bg_hex: str | None = None
    bg_pr = layout._element.find(f".//{{{ns_p}}}bg/{{{ns_p}}}bgPr")
    layout_bg_hex = _resolve_fill(bg_pr)

    # 2. Panel shapes (large decorative shapes with solid fills)
    panels: list[dict[str, Any]] = []
    for shape in layout.shapes:
        with suppress(Exception):
            _ = shape.placeholder_format
            continue
        with suppress(Exception):
            left_in = float(shape.left) / 914400
            width_in = float(shape.width) / 914400
            height_in = float(shape.height) / 914400
            area_ratio = (width_in * height_in) / max(0.1, slide_width_in * slide_height_in)
            if area_ratio < 0.15 or width_in < 1.5 or height_in < 1.5:
                continue
            sp_pr = shape._element.find(f"{{{ns_p}}}spPr")
            if sp_pr is None:
                continue
            color = _resolve_fill(sp_pr)
            if color:
                panels.append(
                    {
                        "left": round(left_in, 2),
                        "width": round(width_in, 2),
                        "bg_color": color,
                    }
                )

    # 3. Build color zones
    zones: list[dict[str, Any]] = []
    if not panels and layout_bg_hex:
        # Whole slide is one zone
        zones.append(
            {
                "region": "full_slide",
                "left": 0,
                "width": round(slide_width_in, 2),
                "bg_color": layout_bg_hex,
                "text_color": _text_color_for_bg(layout_bg_hex),
            }
        )
    elif panels:
        # Sort panels by left position
        panels.sort(key=lambda p: p["left"])
        covered: list[tuple[float, float, str]] = []
        for p in panels:
            covered.append((p["left"], p["left"] + p["width"], p["bg_color"]))

        # Build zones from gaps and panels
        prev_right = 0.0
        for i, (pl, pr, pc) in enumerate(covered):
            # Gap before this panel = background color
            if pl > prev_right + 0.5 and layout_bg_hex:
                zones.append(
                    {
                        "region": f"gap_{i}",
                        "left": round(prev_right, 2),
                        "width": round(pl - prev_right, 2),
                        "bg_color": layout_bg_hex,
                        "text_color": _text_color_for_bg(layout_bg_hex),
                    }
                )
            # Panel zone
            zones.append(
                {
                    "region": f"panel_{i}",
                    "left": round(pl, 2),
                    "width": round(pr - pl, 2),
                    "bg_color": pc,
                    "text_color": _text_color_for_bg(pc),
                }
            )
            prev_right = pr

        # Gap after last panel
        if prev_right < slide_width_in - 0.5 and layout_bg_hex:
            zones.append(
                {
                    "region": "right_area",
                    "left": round(prev_right, 2),
                    "width": round(slide_width_in - prev_right, 2),
                    "bg_color": layout_bg_hex,
                    "text_color": _text_color_for_bg(layout_bg_hex),
                }
            )

    return layout_bg_hex, zones


def _analyze_layout_visual(
    *,
    layout: Any,
    layout_name: str,
    slide_width_in: float,
    slide_height_in: float,
    content_box: dict[str, float] | None,
    theme_colors: dict[str, str] | None = None,
    placeholders: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    fixed_labels: list[dict[str, Any]] = []
    decorative: list[str] = []
    text_boxes: list[dict[str, float]] = []
    side_panels: list[tuple[str, float, float]] = []
    shape_type_counts: dict[str, int] = {}
    slide_area = max(0.1, slide_width_in * slide_height_in)

    for i in range(len(layout.shapes)):
        with suppress(Exception):
            shape = layout.shapes[i]
            token = _shape_type_token(shape)
            shape_type_counts[token] = shape_type_counts.get(token, 0) + 1
            left = float(shape.left.inches)
            top = float(shape.top.inches)
            width = float(shape.width.inches)
            height = float(shape.height.inches)
            right = left + width
            area_ratio = max(0.0, (width * height) / slide_area)

            text = _shape_text(shape)
            if text:
                if text in {"‹#›", "<#>", "#"} or text.lower().startswith("copyright ©"):
                    continue
                fixed_labels.append(
                    {
                        "text": text[:120],
                        "left": round(left, 2),
                        "top": round(top, 2),
                        "width": round(width, 2),
                        "height": round(height, 2),
                    }
                )
                text_boxes.append({"left": left, "top": top, "width": width, "height": height})

            if token in {"picture", "auto_shape", "freeform"} and area_ratio >= 0.18:
                if left <= (slide_width_in * 0.2) and right <= (slide_width_in * 0.6):
                    side_panels.append(("left", left, right))
                elif right >= (slide_width_in * 0.8) and left >= (slide_width_in * 0.4):
                    side_panels.append(("right", left, right))

            if token == "line" and width >= (slide_width_in * 0.6):
                decorative.append("horizontal_divider")
            if token in {"picture", "auto_shape", "freeform"} and area_ratio >= 0.6:
                decorative.append("background_visual")

    editable_regions: list[dict[str, Any]] = []
    if content_box is not None:
        editable_regions.append(
            {
                "name": "content_box",
                "left": round(float(content_box["left"]), 2),
                "top": round(float(content_box["top"]), 2),
                "width": round(float(content_box["width"]), 2),
                "height": round(float(content_box["height"]), 2),
                "source": "placeholder_union",
            }
        )
    else:
        left = _DEFAULT_MARGIN_LEFT
        right = slide_width_in - _DEFAULT_MARGIN_RIGHT
        top = _DEFAULT_MARGIN_TOP
        bottom = slide_height_in - _DEFAULT_MARGIN_BOTTOM
        for side, panel_left, panel_right in side_panels:
            if side == "left":
                left = max(left, panel_right + 0.2)
            else:
                right = min(right, panel_left - 0.2)
        if text_boxes:
            header_candidates = [t["top"] + t["height"] for t in text_boxes if t["top"] < 2.0]
            if header_candidates:
                header_bottom = min(slide_height_in, max(header_candidates))
                top = max(top, header_bottom + 0.15)
        width = max(0.0, right - left)
        height = max(0.0, bottom - top)
        if width >= 2.0 and height >= 1.0:
            editable_regions.append(
                {
                    "name": "suggested_manual_region",
                    "left": round(left, 2),
                    "top": round(top, 2),
                    "width": round(width, 2),
                    "height": round(height, 2),
                    "source": "visual_inference_no_placeholders",
                }
            )

    labels = [item["text"] for item in fixed_labels]
    summary_parts: list[str] = []
    if labels:
        summary_parts.append(f"Fixed labels present: {', '.join(labels[:3])}.")
    if editable_regions:
        summary_parts.append(
            "Add editable content inside suggested regions; avoid fixed labels/decorative areas."
        )
    else:
        summary_parts.append(
            "No clear editable region inferred; treat layout as template-owned visual."
        )

    # Color zone analysis
    layout_bg_hex = None
    color_zones: list[dict[str, Any]] = []
    if theme_colors:
        layout_bg_hex, color_zones = _extract_color_zones(
            layout=layout,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
            theme_colors=theme_colors,
        )

    # Extract title placeholder region for overlap avoidance
    title_region: dict[str, Any] | None = None
    if placeholders:
        for ph in placeholders:
            if ph.get("role") == "title":
                left_v = float(ph["left"])
                top_v = float(ph["top"])
                width_v = float(ph["width"])
                height_v = float(ph["height"])
                if width_v > 0 and height_v > 0:
                    title_region = {
                        "left": round(left_v, 2),
                        "top": round(top_v, 2),
                        "width": round(width_v, 2),
                        "height": round(height_v, 2),
                    }
                    # Determine which color zone the title sits in
                    title_cx = left_v + width_v / 2
                    for zone in color_zones:
                        zl = zone["left"]
                        zw = zone["width"]
                        if zl <= title_cx <= zl + zw:
                            title_region["zone"] = zone.get("region", "")
                            title_region["text_color"] = zone.get(
                                "text_color", _text_color_for_bg(zone.get("bg_color", "FFFFFF"))
                            )
                            break
                break

    # Add editable_area to each color zone (zone area minus title placeholder)
    if color_zones and title_region:
        tr_left = title_region["left"]
        tr_top = title_region["top"]
        tr_bottom = tr_top + title_region["height"]
        for zone in color_zones:
            zl = zone["left"]
            zw = zone["width"]
            zt = _DEFAULT_MARGIN_TOP
            zb = slide_height_in - _DEFAULT_MARGIN_BOTTOM
            # Check if title overlaps this zone
            title_cx = tr_left + title_region["width"] / 2
            if zl <= title_cx <= zl + zw:
                # Title is in this zone — editable area is above and below it
                zone["editable_above"] = {
                    "left": round(zl, 2),
                    "top": round(zt, 2),
                    "width": round(zw, 2),
                    "height": round(max(0, tr_top - zt - 0.1), 2),
                }
                zone["editable_below"] = {
                    "left": round(zl, 2),
                    "top": round(tr_bottom + 0.1, 2),
                    "width": round(zw, 2),
                    "height": round(max(0, zb - tr_bottom - 0.1), 2),
                }
            else:
                # Title not in this zone — full zone is editable
                zone["editable_area"] = {
                    "left": round(zl, 2),
                    "top": round(zt, 2),
                    "width": round(zw, 2),
                    "height": round(max(0, zb - zt), 2),
                }

    result: dict[str, Any] = {
        "fixed_labels": fixed_labels,
        "shape_type_counts": shape_type_counts,
        "decorative_structure": sorted(set(decorative)),
        "editable_regions": editable_regions,
        "usage_summary": " ".join(summary_parts),
    }
    if title_region:
        result["title_region"] = title_region
    if color_zones:
        result["color_zones"] = color_zones
    if layout_bg_hex:
        result["layout_bg"] = layout_bg_hex
    return result


def _generate_layout_preview_images(
    *,
    template_path: Path,
    layout_count: int,
    preview_dir: Path,
) -> dict[int, str]:
    if shutil.which("soffice") is None or shutil.which("pdftoppm") is None:
        return {}

    preview_dir.mkdir(parents=True, exist_ok=True)
    mapping: dict[int, str] = {}
    with tempfile.TemporaryDirectory(prefix="slides-layout-preview-") as tmp:
        tmp_path = Path(tmp)
        def _preview_placeholder_label(type_name: str, idx: int) -> str:
            token = type_name.upper()
            if token.endswith(("TITLE (1)", "CENTER_TITLE (3)")):
                return "[TITLE]"
            if token.endswith("SUBTITLE (4)"):
                return "[SUBTITLE]"
            if token.endswith(("BODY (2)", "OBJECT (7)")):
                return "[BODY]\n- Bullet 1\n- Bullet 2"
            if token.endswith("FOOTER (15)"):
                return "[FOOTER]"
            if token.endswith("DATE (16)"):
                return "[DATE]"
            if token.endswith("SLIDE_NUMBER (13)"):
                return "[SLIDE_NUMBER]"
            return f"[PLACEHOLDER_{idx}]"

        preview_deck = load_presentation(str(template_path))
        for layout_index in range(layout_count):
            with suppress(Exception):
                slide = preview_deck.slides.add_slide(preview_deck.slide_layouts[layout_index])
                for placeholder in slide.placeholders:
                    with suppress(Exception):
                        if not bool(getattr(placeholder, "has_text_frame", False)):
                            continue
                        type_name = str(placeholder.placeholder_format.type)
                        idx = int(placeholder.placeholder_format.idx)
                        raw_existing = str(getattr(placeholder, "text", "") or "").strip()
                        existing = " ".join(raw_existing.split()).strip()
                        existing_lower = existing.lower()
                        is_generic = (
                            not existing
                            or existing in {"‹#›", "<#>", "#"}
                            or "click to edit" in existing_lower
                            or "master title style" in existing_lower
                            or "master subtitle style" in existing_lower
                            or "master text styles" in existing_lower
                            or "click icon to add" in existing_lower
                        )
                        placeholder.text = (
                            _preview_placeholder_label(type_name, idx)
                            if is_generic
                            else raw_existing
                        )
        pptx_path = tmp_path / "layout_preview.pptx"
        preview_deck.save(str(pptx_path))
        cmd_pdf = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp_path),
            str(pptx_path),
        ]
        proc_pdf = subprocess.run(cmd_pdf, capture_output=True, text=True)
        if proc_pdf.returncode != 0:
            return {}
        pdf_path = tmp_path / "layout_preview.pdf"
        if not pdf_path.exists():
            return {}
        prefix = tmp_path / "layout_page"
        cmd_png = ["pdftoppm", "-png", str(pdf_path), str(prefix)]
        proc_png = subprocess.run(cmd_png, capture_output=True, text=True)
        if proc_png.returncode != 0:
            return {}

        for layout_index in range(layout_count):
            page_path = tmp_path / f"layout_page-{layout_index + 1}.png"
            if not page_path.exists():
                continue
            out_path = preview_dir / f"layout_{layout_index:02d}.png"
            shutil.copy2(page_path, out_path)
            mapping[layout_index] = str(out_path)
    return mapping


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_R = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)

_MAX_ICON_AREA_IN2 = 4.0  # icons are small — max ~2"×2"
_MIN_ICON_AREA_IN2 = 0.04  # ignore tiny invisible shapes


def _shape_geom_hash(el: etree._Element) -> str:
    """Hash the geometry paths of a shape/group for deduplication."""
    parts: list[bytes] = []
    for geom in el.iter(f"{{{_NS_A}}}custGeom"):
        parts.append(etree.tostring(geom))
    for prst in el.iter(f"{{{_NS_A}}}prstGeom"):
        parts.append(etree.tostring(prst))
    return hashlib.sha256(b"".join(parts)).hexdigest()[:12]


def _el_dimensions_inches(
    el: etree._Element,
) -> tuple[float, float] | None:
    """Read width/height in inches from the first xfrm found."""
    for xfrm in el.iter(f"{{{_NS_A}}}xfrm"):
        ext = xfrm.find(f"{{{_NS_A}}}ext")
        if ext is not None:
            cx = int(ext.get("cx", "0"))
            cy = int(ext.get("cy", "0"))
            if cx > 0 and cy > 0:
                return cx / 914400, cy / 914400
    return None


def _has_custom_geometry(el: etree._Element) -> bool:
    """Check if element contains at least one custom geometry path."""
    return next(el.iter(f"{{{_NS_A}}}custGeom"), None) is not None


def _has_text_content(el: etree._Element) -> bool:
    """Check if element contains non-whitespace text."""
    return any(
        t_el.text and t_el.text.strip()
        for t_el in el.iter(f"{{{_NS_A}}}t")
    )


def _has_image_fill(el: etree._Element) -> bool:
    """Check if element has an image/blip fill."""
    return next(el.iter(f"{{{_NS_A}}}blip"), None) is not None


def _clean_icon_name(raw_name: str) -> str:
    """Normalize a shape name into a clean icon filename stem."""
    name = re.sub(r"\s+\d+$", "", raw_name.strip())
    name = re.sub(r"[^\w\s-]", "", name)
    name = re.sub(r"[\s-]+", "_", name).strip("_")
    return name or "icon"


def _is_icon_candidate(
    el: etree._Element, tag: str,
) -> bool:
    """Determine if an XML element looks like a vector icon."""
    if _has_text_content(el) or _has_image_fill(el):
        return False

    dims = _el_dimensions_inches(el)
    if dims is None:
        return False
    w, h = dims
    area = w * h
    if area < _MIN_ICON_AREA_IN2 or area > _MAX_ICON_AREA_IN2:
        return False

    aspect = max(w, h) / min(w, h) if min(w, h) > 0 else 999
    if aspect > 3.0:
        return False

    if tag == "grpSp":
        return _has_custom_geometry(el)
    if tag == "sp":
        return _has_custom_geometry(el)
    return False


def _extract_icons(
    source_path: Path, icons_dir: Path,
) -> list[str]:
    """Extract vector icon shapes from a .pptx into standalone XML files.

    Returns list of extracted icon filenames.
    """
    prs = load_presentation(str(source_path))
    seen_hashes: dict[str, str] = {}
    extracted: list[str] = []

    for slide in prs.slides:
        sp_tree = slide.shapes._spTree

        for el in sp_tree:
            tag = etree.QName(el).localname
            if tag not in {"grpSp", "sp"}:
                continue

            if not _is_icon_candidate(el, tag):
                continue

            geom_hash = _shape_geom_hash(el)
            if not geom_hash or geom_hash in seen_hashes:
                continue

            cnv_pr = el.find(f".//{{{_NS_P}}}cNvPr")
            if cnv_pr is None:
                cnv_pr = el.find(f".//{_NS_P}cNvPr")
            raw_name = ""
            if cnv_pr is not None:
                raw_name = cnv_pr.get("name", "")
            clean = _clean_icon_name(raw_name)
            if clean in set(seen_hashes.values()):
                clean = f"{clean}_{geom_hash[:6]}"

            seen_hashes[geom_hash] = clean
            icons_dir.mkdir(parents=True, exist_ok=True)
            out_path = icons_dir / f"{clean}.xml"
            out_path.write_bytes(
                etree.tostring(el, xml_declaration=False, encoding="unicode")
                .encode("utf-8")
            )
            extracted.append(f"{clean}.xml")

    return extracted


def _extract_template_layout_catalog(
    template_path: Path,
    *,
    preview_dir: Path | None = None,
) -> dict[str, Any]:
    def infer_layout_family(layout_name: str) -> str:
        key = layout_name.strip().lower()
        if "layout guide" in key:
            return "guide"
        if "title slide" in key:
            return "cover"
        if "agenda" in key:
            return "agenda"
        if key == "end" or key.endswith(" end"):
            return "closing"
        if "disclaimer" in key:
            return "legal"
        if "section header" in key:
            return "section_break"
        if "special gray" in key:
            return "data"
        if "quote" in key or "big statement" in key:
            return "highlight"
        return "content"

    def infer_layout_policy(layout_name: str, role_counts: dict[str, int]) -> dict[str, Any]:
        key = layout_name.strip().lower()
        notes: list[str] = []
        intent_tags: list[str] = []
        allowed_content_layouts: list[str] = []
        blocked_content_layouts: list[str] = []
        constraints: dict[str, Any] = {}

        if "title slide" in key:
            intent_tags.append("cover")
            allowed_content_layouts = ["title_slide"]
            blocked_content_layouts = ["bar_chart", "line_chart", "pie_chart", "table"]
            constraints["max_bullets"] = 1
            notes.append("Use only for deck opening/cover slides.")
        elif "agenda" in key:
            intent_tags.append("agenda")
            allowed_content_layouts = ["agenda", "executive_summary"]
            blocked_content_layouts = ["four_column", "bar_chart", "line_chart", "pie_chart"]
            constraints["max_bullets"] = 8
            notes.append("Use for agenda/table-of-contents style pages.")
        elif "section header" in key:
            intent_tags.append("section_break")
            allowed_content_layouts = ["section_divider"]
            blocked_content_layouts = ["content_text", "content_bullets", "table"]
            constraints["max_bullets"] = 0
            notes.append("Use only for section transitions.")
        elif "quote" in key:
            intent_tags.append("highlight")
            allowed_content_layouts = ["quote"]
            blocked_content_layouts = ["bar_chart", "line_chart", "pie_chart", "table"]
            constraints["max_bullets"] = 2
            notes.append("Use for one quote and one attribution.")
        elif "disclaimer" in key:
            intent_tags.append("legal")
            allowed_content_layouts = ["disclaimer"]
            blocked_content_layouts = ["bar_chart", "line_chart", "pie_chart"]
            constraints["max_bullets"] = 10
            notes.append("Use only for legal/disclaimer pages.")
        elif key == "end" or key.endswith(" end"):
            intent_tags.append("closing")
            allowed_content_layouts = ["end_slide"]
            blocked_content_layouts = ["bar_chart", "line_chart", "pie_chart", "table"]
            constraints["max_bullets"] = 5
            notes.append("Use for final slide / next steps.")
        elif "big statement" in key:
            intent_tags.append("hero_message")
            allowed_content_layouts = ["big_statement", "big_number"]
            blocked_content_layouts = ["table", "bar_chart", "line_chart", "pie_chart"]
            constraints["max_bullets"] = 1
            notes.append("Use for one big strategic message.")
        elif "arrow" in key or "one third" in key or "half" in key or "two third" in key:
            intent_tags.append("structured_content")
            allowed_content_layouts = [
                "two_column",
                "three_column",
                "four_column",
                "process_flow",
                "icon_grid",
            ]
            constraints["prefer_balanced_columns"] = True
            notes.append("Use for structured multi-block content; balance text across blocks.")
        elif "special gray" in key:
            intent_tags.append("data_visual")
            allowed_content_layouts = [
                "matrix_2x2",
                "bar_chart",
                "line_chart",
                "pie_chart",
                "table",
            ]
            constraints["requires_source_line"] = True
            notes.append("Use for dense analysis or data visuals.")
        elif "layout guide" in key:
            intent_tags.append("guide")
            blocked_content_layouts = [
                "title_slide",
                "executive_summary",
                "agenda",
                "content_text",
                "content_bullets",
                "two_column",
                "three_column",
                "four_column",
                "table",
                "bar_chart",
                "line_chart",
                "pie_chart",
                "big_number",
                "icon_grid",
                "process_flow",
                "matrix_2x2",
                "big_statement",
                "green_panel_text",
                "quote",
                "section_divider",
                "disclaimer",
                "end_slide",
            ]
            notes.append("Do not use for content slides.")
        else:
            intent_tags.append("general_content")
            allowed_content_layouts = [
                "executive_summary",
                "content_text",
                "content_bullets",
                "two_column",
                "three_column",
                "four_column",
                "table",
                "bar_chart",
                "line_chart",
                "pie_chart",
                "big_number",
                "icon_grid",
                "process_flow",
                "matrix_2x2",
                "green_panel_text",
            ]
            notes.append("General-purpose content layout.")

        if role_counts.get("title", 0) == 0:
            constraints["requires_explicit_title_shape"] = True
        if role_counts.get("body", 0) == 0:
            constraints["limited_body_capacity"] = True

        return {
            "intent_tags": intent_tags,
            "allowed_content_layouts": allowed_content_layouts,
            "blocked_content_layouts": blocked_content_layouts,
            "constraints": constraints,
            "agent_notes": notes,
        }

    template_bytes = template_path.read_bytes()
    template_fingerprint = hashlib.sha256(template_bytes).hexdigest()
    prs = load_presentation(str(template_path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide_width_in = float(getattr(slide_width, "inches", 0.0))
    slide_height_in = float(getattr(slide_height, "inches", 0.0))
    slide_area = max(0.1, slide_width_in * slide_height_in)
    theme_colors = _resolve_theme_colors(template_path)
    preview_map: dict[int, str] = {}
    if preview_dir is not None:
        preview_map = _generate_layout_preview_images(
            template_path=template_path,
            layout_count=len(prs.slide_layouts),
            preview_dir=preview_dir,
        )
    layouts: list[dict[str, Any]] = []
    for layout_index, layout in enumerate(prs.slide_layouts):
        name = str(getattr(layout, "name", "") or "").strip() or f"Layout {layout_index}"
        placeholders: list[dict[str, Any]] = []
        body_boxes: list[tuple[float, float, float, float]] = []
        role_counts: dict[str, int] = {
            "title": 0,
            "subtitle": 0,
            "body": 0,
            "footer": 0,
            "date": 0,
            "slide_number": 0,
        }
        content_box: dict[str, float] | None = None
        for placeholder in layout.placeholders:
            type_name = str(placeholder.placeholder_format.type)
            role = _placeholder_role(type_name)
            if role is not None:
                role_counts[role] += 1
            left = float(placeholder.left.inches)
            top = float(placeholder.top.inches)
            width = float(placeholder.width.inches)
            height = float(placeholder.height.inches)
            placeholders.append(
                {
                    "idx": int(placeholder.placeholder_format.idx),
                    "type": type_name,
                    "role": role,
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                }
            )
            if role == "body":
                body_boxes.append((left, top, width, height))
                if content_box is None:
                    content_box = {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    }
                else:
                    right = max(content_box["left"] + content_box["width"], left + width)
                    bottom = max(content_box["top"] + content_box["height"], top + height)
                    content_box["left"] = min(content_box["left"], left)
                    content_box["top"] = min(content_box["top"], top)
                    content_box["width"] = right - content_box["left"]
                    content_box["height"] = bottom - content_box["top"]
        body_area = 0.0
        body_area_ratio = 0.0
        if content_box is not None:
            body_area = max(0.0, content_box["width"] * content_box["height"])
            body_area_ratio = body_area / slide_area
        columns_hint = max(1, len(body_boxes)) if body_boxes else 0
        family = infer_layout_family(name)
        visual_definition = _analyze_layout_visual(
            layout=layout,
            layout_name=name,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
            content_box=content_box,
            theme_colors=theme_colors,
            placeholders=placeholders,
        )
        policy = infer_layout_policy(name, role_counts)
        notes = list(policy.get("agent_notes", []))
        if visual_definition.get("usage_summary"):
            notes.append(str(visual_definition["usage_summary"]))
        editable_regions = visual_definition.get("editable_regions", [])
        if editable_regions and role_counts.get("body", 0) == 0:
            notes.append(
                "No body placeholders: write with add_text/add_chart inside editable_regions."
            )
        policy["agent_notes"] = notes
        layouts.append(
            {
                "layout_id": hashlib.sha1(
                    f"{template_fingerprint}:{layout_index}:{name}".encode()
                ).hexdigest()[:12],
                "layout_index": layout_index,
                "layout_name": name,
                "layout_family": family,
                "placeholders": placeholders,
                "role_counts": role_counts,
                "content_box": content_box,
                "geometry": {
                    "body_placeholder_count": len(body_boxes),
                    "columns_hint": columns_hint,
                    "body_area_ratio": round(body_area_ratio, 4),
                    "body_area_in2": round(body_area, 3),
                },
                "visual_definition": visual_definition,
                "preview_image": preview_map.get(layout_index),
                "policy": policy,
            }
        )
    return {
        "template_path": str(template_path),
        "template_fingerprint": template_fingerprint,
        "slide_width_in": slide_width_in,
        "slide_height_in": slide_height_in,
        "layout_count": len(layouts),
        "layouts": layouts,
        "theme_colors": theme_colors if theme_colors else None,
        "preview_enabled": bool(preview_dir),
        "preview_dir": str(preview_dir) if preview_dir is not None else None,
    }


_ARCHETYPE_VISUAL_HINTS: dict[str, str] = {
    "title_slide": "Large title with concise subtitle",
    "executive_summary": "Three concise bullets",
    "agenda": "Numbered sections",
    "content_text": "One paragraph narrative",
    "two_column": "Left-right comparison",
    "three_column": "Three equal insight columns",
    "four_column": "Four compact columns",
    "timeline": "Chronological milestones",
    "matrix_2x2": "2x2 framing: impact vs feasibility",
    "bar_chart": "Category comparison chart with labels",
    "line_chart": "Trend chart by period",
    "pie_chart": "Share-of-total chart",
    "big_number": "Key KPI with short context",
    "process_flow": "Step-by-step execution path",
    "icon_grid": "Grid of labeled concepts",
    "table": "Structured table",
    "quote": "Emphasized quote block",
    "big_statement": "Large message on full colored background",
    "green_panel_text": "Colored anchor panel with structured text on the right",
    "section_divider": "Section transition slide",
    "disclaimer": "Legal/disclaimer text",
    "content_bullets": "Single message with supporting bullets",
    "end_slide": "Close with ownership and next action",
}

_ARCHETYPE_ACTION_TITLE_TEMPLATES: dict[str, str] = {
    "title": "{deck_title}",
    "executive_summary": "{subject}: what matters most",
    "agenda": "Agenda: how the {subject} story unfolds",
    "context": "{subject}: historical and geopolitical context",
    "analysis": "{subject}: opposing strategies and constraints",
    "evidence": "{subject}: scale and commitment over time",
    "trend": "{subject}: trend in support and sentiment",
    "mix": "{subject}: distribution of impact across dimensions",
    "turning_points": "{subject}: turning points that changed outcomes",
    "options": "{subject}: implications across policy and strategy",
    "process": "{subject}: timeline of major phases",
    "lessons": "{subject}: strategic lessons for leaders today",
    "context_deep": "{subject}: deeper drivers and constraints",
    "recommendation": "{subject}: key takeaways to retain",
    "next_steps": "Next steps and ownership",
}

_ARCHETYPE_FALLBACK_LAYOUT_TOKENS: dict[str, list[str]] = {
    "title_slide": ["title slide"],
    "executive_summary": ["title and text", "title only"],
    "agenda": [
        "agenda full width overview",
        "agenda two-thirds",
        "agenda table of contents",
        "agenda section header overview",
        "title and text",
    ],
    "content_text": ["title and text", "title only"],
    "content_bullets": ["title and text", "title only"],
    "two_column": ["arrow half", "green half", "title and text"],
    "three_column": ["arrow two third", "green two third", "white one third", "title and text"],
    "four_column": ["four column green"],
    "process_flow": [
        "left arrow",
        "green left arrow",
        "arrow two third",
        "green arrow two third",
        "arrow half",
        "green arrow half",
        "title and text",
    ],
    "icon_grid": ["white one third", "green one third", "title and text"],
    "big_number": ["green highlight", "big statement icon", "title and text"],
    "matrix_2x2": ["special gray", "title and text", "blank"],
    "bar_chart": ["title and text", "special gray", "title only"],
    "line_chart": ["title and text", "special gray", "title only"],
    "pie_chart": ["title and text", "special gray", "title only"],
    "table": ["title and text", "special gray", "title only"],
    "quote": ["quote"],
    "disclaimer": ["disclaimer"],
    "end_slide": ["end", "title and text", "title and content", "title only"],
    "section_divider": ["section header line", "section header box", "section"],
    "green_panel_text": ["green one third", "green"],
    "big_statement": ["big statement"],
}

_ARCHETYPE_DEFAULT_GEOMETRIES: dict[str, dict[str, Any]] = {
    "title_slide": {"left": 0.7, "top": 0.9, "width": 11.0, "height": 2.1, "font_size": 38},
    "executive_summary": {"left": 0.9, "top": 1.6, "width": 10.8, "height": 4.8, "font_size": 20},
    "content_text": {"left": 0.9, "top": 1.6, "width": 10.8, "height": 4.8, "font_size": 14},
    "content_bullets": {"left": 1.0, "top": 1.6, "width": 10.5, "height": 4.8, "font_size": 20},
    "two_column": {"left": 0.8, "top": 1.5, "width": 10.6, "height": 4.8, "font_size": 18},
    "three_column": {"left": 0.7, "top": 1.5, "width": 11.0, "height": 4.8, "font_size": 17},
    "four_column": {"left": 0.5, "top": 1.5, "width": 11.8, "height": 4.8, "font_size": 14},
    "bar_chart": {"left": 0.9, "top": 1.7, "width": 7.2, "height": 4.2, "font_size": 20},
    "line_chart": {"left": 0.9, "top": 1.7, "width": 10.2, "height": 4.2, "font_size": 20},
    "pie_chart": {"left": 1.0, "top": 1.6, "width": 5.2, "height": 4.5, "font_size": 20},
    "table": {"left": 0.8, "top": 1.7, "width": 10.2, "height": 4.8, "font_size": 20},
    "big_number": {"left": 1.0, "top": 2.0, "width": 4.5, "height": 3.2, "font_size": 64},
    "process_flow": {"left": 0.8, "top": 2.5, "width": 10.8, "height": 2.0, "font_size": 16},
    "timeline": {"left": 0.8, "top": 1.6, "width": 10.8, "height": 4.2, "font_size": 18},
    "icon_grid": {"left": 0.8, "top": 1.8, "width": 11.0, "height": 4.2, "font_size": 16},
    "matrix_2x2": {"left": 1.0, "top": 1.35, "width": 10.2, "height": 5.0, "font_size": 17},
    "quote": {"left": 1.0, "top": 2.0, "width": 10.0, "height": 3.0, "font_size": 30},
    "big_statement": {"left": 1.1, "top": 2.2, "width": 11.1, "height": 2.4, "font_size": 34},
    "green_panel_text": {"left": 5.0, "top": 1.55, "width": 6.0, "height": 4.5, "font_size": 13},
    "section_divider": {"left": 1.0, "top": 3.0, "width": 10.0, "height": 1.0, "font_size": 40},
    "disclaimer": {"left": 0.8, "top": 1.6, "width": 10.8, "height": 5.5, "font_size": 13},
    "end_slide": {"left": 0.7, "top": 0.9, "width": 11.0, "height": 2.1, "font_size": 38},
    "agenda": {"left": 1.0, "top": 1.6, "width": 10.5, "height": 4.8, "font_size": 20},
}

ARCHETYPE_SPECS: dict[str, dict[str, Any]] = {
    "title_slide": {
        "description": "Opening cover slide with deck title and subtitle.",
        "when_to_use": ["First slide of every deck"],
        "guidelines": [
            "Title: deck name or key theme (not an action title)",
            "Subtitle: date, audience, author, or context line",
        ],
        "template_layout_hint": "Title Slide or cover layout",
        "op_recipe": [
            {"op": "add_slide", "note": "layout_name from content_layout.json cover family"},
            {"op": "set_title_subtitle", "note": "title + subtitle text"},
        ],
    },
    "executive_summary": {
        "description": "Top-level summary stating the deck's key recommendation.",
        "when_to_use": [
            "Early in the deck (slide 2-3) to state the answer up front",
            "Use SCQA: action title = Answer, body = Situation + Complication + supporting points",
        ],
        "guidelines": [
            "Action title states the core recommendation",
            "Body: 3-5 bullet points covering situation, complication, and key arguments",
            "This is the most important slide — reader should get the full message here",
        ],
        "template_layout_hint": "Title and Text or content layout with body placeholder",
        "op_recipe": [
            {"op": "add_slide", "note": "layout_name from content_layout.json"},
            {"op": "set_semantic_text", "note": "role: title, text: action title (the Answer)"},
            {"op": "set_semantic_text", "note": "role: body, text: bullet points. "
             "Fallback: set_placeholder_text or add_text"},
        ],
    },
    "section_divider": {
        "description": "Visual break between major sections of the deck.",
        "when_to_use": [
            "Before each major argument/section (for decks > 8 slides)",
            "Title: section theme (not an action title)",
        ],
        "guidelines": [
            "Keep minimal — title only, no body content",
            "Optional subtitle for section context",
        ],
        "template_layout_hint": "Section Header layout",
        "op_recipe": [
            {"op": "add_slide", "note": "section header layout"},
            {"op": "set_title_subtitle", "note": "title: section name, subtitle: optional context"},
        ],
    },
    "content_bullets": {
        "description": "Slide with action title and bulleted body text.",
        "when_to_use": [
            "Qualitative arguments, lists of findings, recommendations",
            "When the message is text-based, not data-driven",
        ],
        "guidelines": [
            "Action title states the insight",
            "3-6 bullet points, parallel structure",
            "Use bold lead-ins for scanability",
            "Avoid walls of text — if > 6 bullets, split into two slides",
        ],
        "template_layout_hint": "Title and Text layout",
        "op_recipe": [
            {"op": "add_slide", "note": "content layout with body placeholder"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "set_semantic_text", "note": "role: body, text: bulleted text. "
             "Fallback: set_placeholder_text or add_text"},
        ],
    },
    "content_text": {
        "description": "Narrative text slide for longer-form explanations.",
        "when_to_use": [
            "When a narrative paragraph is more appropriate than bullets",
            "Context-setting slides, methodology descriptions",
        ],
        "guidelines": [
            "Action title states the insight",
            "Body: 1-3 short paragraphs, not bullets",
            "Keep concise — if it needs scrolling, split the slide",
        ],
        "template_layout_hint": "Title and Text layout",
        "op_recipe": [
            {"op": "add_slide", "note": "content layout with body placeholder"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "set_semantic_text", "note": "role: body. "
             "Fallback: set_placeholder_text or add_text"},
        ],
    },
    "big_statement": {
        "description": "Large-text emphasis slide for a single powerful message.",
        "when_to_use": [
            "Key takeaway or provocative statement that needs visual weight",
            "Transition moments in the narrative",
        ],
        "guidelines": [
            "One sentence, large font, centered",
            "No bullets, no charts — just the message",
        ],
        "template_layout_hint": "Title Only or Blank layout",
        "op_recipe": [
            {"op": "add_slide", "note": "minimal layout (Title Only or Blank)"},
            {"op": "add_text", "note": "large centered text with explicit coordinates"},
        ],
    },
    "big_number": {
        "description": "Highlight a single metric with large typography.",
        "when_to_use": [
            "When one number tells the story (revenue, growth rate, market size)",
            "Impact slides",
        ],
        "guidelines": [
            "Number: large, bold, prominent",
            "Label: small text below explaining what the number represents",
            "Optional context line or comparison",
        ],
        "template_layout_hint": "Title Only or Blank layout",
        "op_recipe": [
            {"op": "add_slide", "note": "minimal layout"},
            {"op": "set_semantic_text", "note": "role: title (optional slide heading)"},
            {"op": "add_text", "note": "the number, large font, centered"},
            {"op": "add_text", "note": "label/context below the number, smaller font"},
        ],
    },
    "bar_chart": {
        "description": "Bar or column chart comparing values across categories.",
        "when_to_use": [
            "Comparing values across categories (revenue by segment, cost by department)",
            "Showing change over discrete time periods",
            "Ranking items by value",
        ],
        "guidelines": [
            "Sort by value (largest first) unless chronological",
            "Data labels on bars for readability",
            "Legend at bottom for multi-series; hide for single series",
            "Round numbers (use ~$2.5B not $2,487,392,104)",
            "One chart per slide",
        ],
        "template_layout_hint": "Title Only layout, chart as positioned shape",
        "op_recipe": [
            {"op": "add_slide", "note": "layout from content_layout.json"},
            {"op": "set_semantic_text", "note": "role: title. "
             "Fallback: add_text if no title placeholder"},
            {"op": "add_bar_chart", "note": "categories, series, position from layout geometry"},
            {"op": "set_chart_data_labels", "note": "enabled: true (optional)"},
            {"op": "set_chart_legend", "note": "position: bottom, visible for multi-series"},
            {"op": "add_text", "note": "source line at bottom of slide, small font"},
        ],
    },
    "line_chart": {
        "description": "Line chart showing trends over time.",
        "when_to_use": [
            "Time series data (monthly, quarterly, annual trends)",
            "Comparing trajectories across groups",
        ],
        "guidelines": [
            "X-axis: time periods, Y-axis: metric",
            "2-5 series max — more becomes unreadable",
            "Label end points directly when possible",
        ],
        "template_layout_hint": "Title Only layout",
        "op_recipe": [
            {"op": "add_slide", "note": "layout from content_layout.json"},
            {"op": "set_semantic_text", "note": "role: title. "
             "Fallback: add_text if no title placeholder"},
            {"op": "add_line_chart", "note": "categories (time), series, position"},
            {"op": "set_chart_legend", "note": "position: bottom"},
            {"op": "add_text", "note": "source line"},
        ],
    },
    "pie_chart": {
        "description": "Pie or doughnut chart showing composition/share.",
        "when_to_use": [
            "Market share, revenue mix, budget allocation",
            "When parts-of-a-whole is the message",
        ],
        "guidelines": [
            "Max 6 slices — group small ones into 'Other'",
            "Label slices with percentage + category name",
            "Consider doughnut for a cleaner look",
        ],
        "template_layout_hint": "Title Only layout",
        "op_recipe": [
            {"op": "add_slide", "note": "layout from content_layout.json"},
            {"op": "set_semantic_text", "note": "role: title. "
             "Fallback: add_text if no title placeholder"},
            {"op": "add_pie_chart", "note": "categories, single series, position"},
            {"op": "set_chart_data_labels", "note": "show_category_name + show_value"},
            {"op": "add_text", "note": "source line"},
        ],
    },
    "table": {
        "description": "Data table for structured comparisons.",
        "when_to_use": [
            "Detailed comparisons across multiple dimensions",
            "When exact values matter more than visual pattern",
        ],
        "guidelines": [
            "Header row: bold, distinct background",
            "Max 6-8 rows, 4-6 columns — more needs an appendix",
            "Align numbers right, text left",
            "Highlight key cells with color",
        ],
        "template_layout_hint": "Title Only or Title and Text layout",
        "op_recipe": [
            {"op": "add_slide", "note": "layout from content_layout.json"},
            {"op": "set_semantic_text", "note": "role: title. "
             "Fallback: add_text if no title placeholder"},
            {"op": "add_table", "note": "rows (list of row-lists), position, dimensions"},
            {"op": "add_text", "note": "source line"},
        ],
    },
    "two_column": {
        "description": "Side-by-side comparison in two columns.",
        "when_to_use": [
            "Before/after, pros/cons, current vs. proposed",
            "Two parallel arguments or perspectives",
        ],
        "guidelines": [
            "Column headers: bold, descriptive",
            "Parallel content structure across columns",
        ],
        "template_layout_hint": "Two Content layout or Title Only with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "two-column layout or Title Only"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "left column: header + content, explicit left/top/width"},
            {"op": "add_text", "note": "right column: header + content, explicit left/top/width"},
        ],
    },
    "three_column": {
        "description": "Three-column layout for parallel comparisons or categories.",
        "when_to_use": [
            "Three options, workstreams, or categories",
            "When content is naturally parallel across three items",
        ],
        "guidelines": [
            "Equal column widths, parallel structure",
            "Column headers: bold, descriptive",
        ],
        "template_layout_hint": "Title Only with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "three-column layout or Title Only"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "column 1 content with explicit coordinates"},
            {"op": "add_text", "note": "column 2 content with explicit coordinates"},
            {"op": "add_text", "note": "column 3 content with explicit coordinates"},
        ],
    },
    "four_column": {
        "description": "Four-column layout for parallel comparisons.",
        "when_to_use": [
            "Four workstreams, pillars, or phases",
            "Frameworks with four components",
        ],
        "guidelines": [
            "Keep content brief — space is tight",
            "Equal column widths, parallel structure",
        ],
        "template_layout_hint": "Title Only with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "four-column layout or Title Only"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "column 1-4 content, each with explicit coordinates"},
        ],
    },
    "process_flow": {
        "description": "Sequential process or timeline visualization.",
        "when_to_use": [
            "Implementation roadmaps, workflows, step-by-step processes",
            "When order and sequence matter",
        ],
        "guidelines": [
            "3-6 steps — more needs a different format",
            "Horizontal flow (left to right) preferred",
            "Each step: short label + optional description",
        ],
        "template_layout_hint": "Title Only with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "Title Only layout"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "each step as a positioned text box, "
             "arranged horizontally with consistent spacing"},
        ],
    },
    "matrix_2x2": {
        "description": "2x2 matrix for strategic positioning or categorization.",
        "when_to_use": [
            "Strategic frameworks (growth-share, effort-impact, risk-reward)",
            "Categorizing items along two dimensions",
        ],
        "guidelines": [
            "Label both axes clearly",
            "Name each quadrant",
            "Place items as positioned text in the relevant quadrant",
        ],
        "template_layout_hint": "Title Only with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "Title Only layout"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "axis labels (x and y)"},
            {"op": "add_text", "note": "quadrant labels and items, positioned in each quadrant"},
        ],
    },
    "icon_grid": {
        "description": "Grid of icons with labels for capability maps or feature lists.",
        "when_to_use": [
            "Capability overviews, service offerings, feature lists",
            "When visual variety is needed to break up text-heavy sections",
        ],
        "guidelines": [
            "3-6 items in a grid (2x2, 2x3, 3x3)",
            "Each item: icon/emoji + short label + optional description",
            "Consistent spacing and alignment",
        ],
        "template_layout_hint": "Title Only with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "Title Only layout"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "each grid item as positioned text box "
             "with label + description, arranged in grid pattern"},
        ],
    },
    "quote": {
        "description": "Featured quote with attribution.",
        "when_to_use": [
            "Customer testimonials, expert opinions, strategic framing",
            "When a voice from outside the team adds credibility",
        ],
        "guidelines": [
            "Quote text: large, italic or styled",
            "Attribution: smaller, below the quote",
            "Do NOT use set_semantic_text role: title — "
             "quote layouts often lack title placeholders",
        ],
        "template_layout_hint": "Quote layout or Blank layout",
        "op_recipe": [
            {"op": "add_slide", "note": "quote layout or Blank"},
            {"op": "add_text", "note": "quote text with explicit coordinates, large font"},
            {"op": "add_text", "note": "attribution line below, smaller font"},
        ],
    },
    "agenda": {
        "description": "Table of contents listing deck sections.",
        "when_to_use": [
            "After executive summary, before first content section",
            "To set audience expectations for deck structure",
        ],
        "guidelines": [
            "3-6 agenda items",
            "Each item: section number + title + optional description",
            "Action title describes what the agenda covers (not just 'Agenda')",
        ],
        "template_layout_hint": "Title Only layout with positioned shapes",
        "op_recipe": [
            {"op": "add_slide", "note": "Title Only layout"},
            {"op": "set_semantic_text", "note": "role: title"},
            {"op": "add_text", "note": "each agenda item as positioned text, "
             "vertically stacked with consistent spacing"},
        ],
    },
    "disclaimer": {
        "description": "Legal disclaimer or confidentiality notice.",
        "when_to_use": ["End of deck or where legally required"],
        "guidelines": [
            "Small text, standard legal language",
            "Minimal visual styling",
        ],
        "template_layout_hint": "Disclaimer layout or Blank",
        "op_recipe": [
            {"op": "add_slide", "note": "disclaimer layout or Blank"},
            {"op": "add_text", "note": "disclaimer text, small font, full-width"},
        ],
    },
    "end_slide": {
        "description": "Closing slide (thank you, contact info, next steps).",
        "when_to_use": ["Last slide of the deck"],
        "guidelines": [
            "Keep minimal — title or thank-you message",
            "Optional: contact details or next-steps summary",
        ],
        "template_layout_hint": "End layout or Title Slide layout",
        "op_recipe": [
            {"op": "add_slide", "note": "end layout or cover layout"},
            {"op": "set_title_subtitle", "note": "title: closing message, "
             "subtitle: contact or next steps"},
        ],
    },
    "green_panel_text": {
        "description": "Colored panel with key message text.",
        "when_to_use": [
            "Highlighting a critical finding or call to action",
            "Visual break with emphasis",
        ],
        "guidelines": [
            "Panel fills a significant portion of the slide",
            "Text: concise, high-contrast against panel color",
        ],
        "template_layout_hint": "Title Only or Blank layout",
        "op_recipe": [
            {"op": "add_slide", "note": "layout with colored panel or Title Only"},
            {"op": "set_semantic_text", "note": "role: title (if available)"},
            {"op": "add_text", "note": "panel message text with explicit coordinates"},
        ],
    },
}


def _render_slide_screenshots(
    *,
    source_path: Path,
    screenshots_dir: Path,
) -> dict[int, str]:
    if shutil.which("soffice") is None or shutil.which("pdftoppm") is None:
        return {}
    screenshots_dir.mkdir(parents=True, exist_ok=True)
    mapping: dict[int, str] = {}
    with tempfile.TemporaryDirectory(prefix="slides-slide-shots-") as tmp:
        tmp_path = Path(tmp)
        cmd_pdf = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmp_path),
            str(source_path),
        ]
        proc_pdf = subprocess.run(cmd_pdf, capture_output=True, text=True)
        if proc_pdf.returncode != 0:
            return {}
        pdf_candidates = sorted(tmp_path.glob("*.pdf"))
        if not pdf_candidates:
            return {}
        pdf_path = pdf_candidates[0]
        prefix = tmp_path / "slide_page"
        cmd_png = ["pdftoppm", "-png", str(pdf_path), str(prefix)]
        proc_png = subprocess.run(cmd_png, capture_output=True, text=True)
        if proc_png.returncode != 0:
            return {}
        page_files = sorted(
            tmp_path.glob("slide_page-*.png"),
            key=lambda p: int(re.findall(r"-(\d+)\.png$", p.name)[0]),
        )
        for page_index, page in enumerate(page_files):
            out_path = screenshots_dir / f"slide_{page_index + 1:03d}.png"
            shutil.copy2(page, out_path)
            mapping[page_index] = str(out_path)
    return mapping


def _extract_slides_manifest(
    *,
    source_path: Path,
    template_catalog: dict[str, Any],
    screenshot_map: dict[int, str],
) -> dict[str, Any]:
    prs = load_presentation(str(source_path))
    layout_by_name: dict[str, dict[str, Any]] = {}
    for layout in template_catalog.get("layouts", []):
        name = str(layout.get("layout_name", ""))
        if name and name not in layout_by_name:
            layout_by_name[name] = layout

    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides):
        layout_name = str(getattr(slide.slide_layout, "name", "") or "").strip()
        layout_entry = layout_by_name.get(layout_name, {})
        text_items: list[str] = []
        bullet_lines = 0
        chart_count = 0
        table_count = 0
        image_count = 0
        text_shape_count = 0
        for shape in slide.shapes:
            token = _shape_type_token(shape)
            if token == "chart":
                chart_count += 1
            elif token == "table":
                table_count += 1
            elif token == "picture":
                image_count += 1
            if bool(getattr(shape, "has_text_frame", False)):
                text_shape_count += 1
                raw_text = str(getattr(shape, "text", "") or "")
                text = " ".join(raw_text.split()).strip()
                if text:
                    text_items.append(text[:240])
                with suppress(Exception):
                    for para in shape.text_frame.paragraphs:
                        if int(getattr(para, "level", 0) or 0) > 0:
                            bullet_lines += 1
                        para_text = " ".join(str(getattr(para, "text", "") or "").split()).strip()
                        if para_text.startswith(("-", "*", "•")):
                            bullet_lines += 1

        title_text = ""
        with suppress(Exception):
            title_shape = slide.shapes.title
            if title_shape is not None and bool(getattr(title_shape, "has_text_frame", False)):
                title_text = " ".join(str(getattr(title_shape, "text", "") or "").split()).strip()

        slides.append(
            {
                "slide_index": slide_index,
                "slide_number": slide_index + 1,
                "layout_name": layout_name,
                "layout_id": layout_entry.get("layout_id"),
                "screenshot_path": screenshot_map.get(slide_index),
                "title": title_text,
                "visible_text": text_items[:12],
                "stats": {
                    "text_shape_count": text_shape_count,
                    "chart_count": chart_count,
                    "table_count": table_count,
                    "image_count": image_count,
                    "bullet_line_count": bullet_lines,
                    "columns_hint": int(
                        (layout_entry.get("geometry") or {}).get("columns_hint", 0) or 0
                    ),
                },
            }
        )
    return {
        "schema_version": "1.0",
        "source_path": str(source_path),
        "slide_count": len(slides),
        "slides": slides,
    }


def _suggest_archetype_for_slide(slide: dict[str, Any]) -> tuple[str, list[str], float]:
    layout_name = str(slide.get("layout_name", "")).lower()
    title = str(slide.get("title", "")).lower()
    stats = slide.get("stats", {})
    chart_count = int(stats.get("chart_count", 0) or 0)
    table_count = int(stats.get("table_count", 0) or 0)
    bullet_count = int(stats.get("bullet_line_count", 0) or 0)
    columns_hint = int(stats.get("columns_hint", 0) or 0)
    notes: list[str] = []
    confidence = 0.6

    if "title slide" in layout_name:
        notes.append("layout name indicates a cover/title page.")
        return "title_slide", notes, 0.95
    if "agenda" in layout_name:
        notes.append("layout name indicates an agenda page.")
        return "agenda", notes, 0.95
    if "section header" in layout_name:
        notes.append("layout name indicates a section break.")
        return "section_divider", notes, 0.9
    if "disclaimer" in layout_name:
        notes.append("layout name indicates legal/disclaimer content.")
        return "disclaimer", notes, 0.95
    if layout_name == "end" or layout_name.endswith(" end"):
        notes.append("layout name indicates final/closing page.")
        return "end_slide", notes, 0.95
    if "quote" in layout_name or '"' in title:
        notes.append("quote markers detected in layout/title.")
        return "quote", notes, 0.8
    if chart_count > 0:
        if any(token in title for token in ["share", "mix", "split"]):
            notes.append("chart present with share/mix wording.")
            return "pie_chart", notes, 0.8
        if any(token in title for token in ["trend", "over time", "trajectory"]):
            notes.append("chart present with trend wording.")
            return "line_chart", notes, 0.8
        notes.append("chart shape detected.")
        return "bar_chart", notes, 0.75
    if table_count > 0:
        notes.append("table shape detected.")
        return "table", notes, 0.85
    if columns_hint >= 4:
        notes.append("layout columns hint is 4+.")
        return "four_column", notes, 0.75
    if columns_hint == 3:
        notes.append("layout columns hint is 3.")
        return "three_column", notes, 0.75
    if columns_hint == 2:
        notes.append("layout columns hint is 2.")
        return "two_column", notes, 0.75
    if bullet_count >= 4:
        notes.append("high bullet density detected.")
        return "content_bullets", notes, 0.75
    notes.append("defaulting to narrative text content.")
    return "content_text", notes, confidence


def _analyze_slides_manifest(manifest: dict[str, Any]) -> dict[str, Any]:
    analyses: list[dict[str, Any]] = []
    for slide in manifest.get("slides", []):
        archetype, reasons, confidence = _suggest_archetype_for_slide(slide)
        stats = slide.get("stats", {})
        analyses.append(
            {
                "slide_index": slide.get("slide_index"),
                "slide_number": slide.get("slide_number"),
                "layout_name": slide.get("layout_name"),
                "layout_id": slide.get("layout_id"),
                "screenshot_path": slide.get("screenshot_path"),
                "title": slide.get("title"),
                "recommended_archetype": archetype,
                "confidence": round(float(confidence), 2),
                "reasons": reasons,
                "content_observations": {
                    "bullet_density": (
                        "high" if int(stats.get("bullet_line_count", 0) or 0) >= 4 else "normal"
                    ),
                    "has_chart": bool(int(stats.get("chart_count", 0) or 0) > 0),
                    "has_table": bool(int(stats.get("table_count", 0) or 0) > 0),
                    "columns_hint": int(stats.get("columns_hint", 0) or 0),
                },
            }
        )
    return {
        "schema_version": "1.0",
        "source_path": manifest.get("source_path"),
        "slide_count": len(analyses),
        "slides": analyses,
    }


def _build_archetypes_from_analysis(
    *,
    slide_analysis: dict[str, Any],
    template_catalog: dict[str, Any],
    content_catalog: dict[str, Any],
) -> dict[str, Any]:
    by_archetype: dict[str, dict[str, Any]] = {}
    layouts_by_archetype: dict[str, set[str]] = {}
    for row in slide_analysis.get("slides", []):
        archetype = str(row.get("recommended_archetype", "")).strip()
        if not archetype:
            continue
        group = by_archetype.setdefault(
            archetype,
            {
                "sample_slide_numbers": [],
                "sample_titles": [],
                "confidence_scores": [],
                "reason_counts": {},
            },
        )
        group["sample_slide_numbers"].append(int(row.get("slide_number", 0) or 0))
        title = str(row.get("title", "")).strip()
        if title:
            group["sample_titles"].append(title)
        group["confidence_scores"].append(float(row.get("confidence", 0.0) or 0.0))
        for reason in row.get("reasons", []):
            token = str(reason).strip()
            if token:
                group["reason_counts"][token] = group["reason_counts"].get(token, 0) + 1
        layout_name = str(row.get("layout_name", "")).strip()
        if layout_name:
            layouts_by_archetype.setdefault(archetype, set()).add(layout_name)

    policy_layouts: dict[str, set[str]] = {}
    explicit_policy_layouts: dict[str, set[str]] = {}
    for layout in template_catalog.get("layouts", []):
        layout_name = str(layout.get("layout_name", "")).strip()
        allowed = (layout.get("policy") or {}).get("allowed_content_layouts", [])
        if not layout_name or not isinstance(allowed, list):
            continue
        for archetype in allowed:
            token = str(archetype).strip()
            if token:
                explicit_policy_layouts.setdefault(token, set()).add(layout_name)

    matrix = content_catalog.get("layout_content_matrix", {})
    if isinstance(matrix, dict):
        for archetype, per_layout in matrix.items():
            if not isinstance(per_layout, dict):
                continue
            for layout_name, rule in per_layout.items():
                status = str((rule or {}).get("status", "blocked"))
                if status in {"preferred", "allowed"}:
                    policy_layouts.setdefault(str(archetype), set()).add(str(layout_name))

    final_archetypes: dict[str, dict[str, Any]] = {}
    all_ids = sorted(
        set(ARCHETYPE_SPECS.keys())
        | set(by_archetype.keys())
        | set((content_catalog.get("archetypes") or {}).keys())
    )
    for archetype_id in all_ids:
        spec = ARCHETYPE_SPECS.get(archetype_id, {})
        observed = by_archetype.get(archetype_id, {})
        allowed_layout_names = sorted(
            set(layouts_by_archetype.get(archetype_id, set()))
            | set(policy_layouts.get(archetype_id, set()))
            | set(explicit_policy_layouts.get(archetype_id, set()))
        )
        entry: dict[str, Any] = {
            "content_layout": archetype_id,
            "story_roles": (
                (content_catalog.get("archetypes", {}).get(archetype_id, {}) or {}).get(
                    "story_roles", []
                )
            ),
            "description": spec.get("description", ""),
            "goal": (spec.get("when_to_use", [""])[0] if spec.get("when_to_use") else ""),
            "guidelines": spec.get("guidelines", []),
            "prohibitions": [],
            "when_to_use": spec.get("when_to_use", []),
            "op_recipe": spec.get("op_recipe", []),
            "template_layout_hint": spec.get("template_layout_hint", ""),
            "allowed_slide_layouts": allowed_layout_names,
            "sample_slide_numbers": observed.get("sample_slide_numbers", [])[:10],
            "sample_titles": observed.get("sample_titles", [])[:5],
            "avg_confidence": round(
                (
                    sum(observed.get("confidence_scores", []))
                    / max(1, len(observed.get("confidence_scores", [])))
                ),
                2,
            ),
        }
        vh = _ARCHETYPE_VISUAL_HINTS.get(archetype_id)
        if vh:
            entry["visual_hint"] = vh
        att = _ARCHETYPE_ACTION_TITLE_TEMPLATES.get(archetype_id)
        if att:
            entry["action_title_template"] = att
        flt = _ARCHETYPE_FALLBACK_LAYOUT_TOKENS.get(archetype_id)
        if flt:
            entry["fallback_layout_tokens"] = flt
        dg = _ARCHETYPE_DEFAULT_GEOMETRIES.get(archetype_id)
        if dg:
            entry["default_geometry"] = dg
        final_archetypes[archetype_id] = entry
    return {
        "schema_version": "1.3",
        "source_template": template_catalog.get("template_path"),
        "archetypes": final_archetypes,
    }


def _build_resolved_manifest(
    *,
    template_catalog: dict[str, Any],
    content_catalog: dict[str, Any],
    archetypes_catalog: dict[str, Any],
) -> dict[str, Any]:
    """Merge template, content, and archetype catalogs into a single resolved manifest.

    Every value is derived from extracted data — no hardcoded positions, sizes, or colors.
    """
    theme_colors = template_catalog.get("theme_colors") or {}
    slide_w = float(template_catalog.get("slide_width_in", 0))
    slide_h = float(template_catalog.get("slide_height_in", 0))

    # --- 1. Resolve theme from extracted palette ---
    theme = _resolve_manifest_theme(theme_colors, slide_w, slide_h)

    # --- 2. Index layouts by name ---
    layouts_by_name: dict[str, dict[str, Any]] = {}
    for layout in template_catalog.get("layouts", []):
        name = str(layout.get("layout_name", "")).strip()
        if name:
            layouts_by_name[name] = layout

    # --- 3. Compatibility matrix ---
    matrix = content_catalog.get("layout_content_matrix", {})

    # --- 4. Resolve each archetype ---
    resolved_archetypes: dict[str, dict[str, Any]] = {}
    for arch_id, arch_spec in (archetypes_catalog.get("archetypes") or {}).items():
        resolved = _resolve_manifest_archetype(
            archetype_id=arch_id,
            arch_spec=arch_spec,
            matrix=matrix,
            layouts_by_name=layouts_by_name,
            theme=theme,
        )
        resolved_archetypes[arch_id] = resolved

    return {
        "schema_version": "1.0",
        "source_template": template_catalog.get("template_path"),
        "template_fingerprint": template_catalog.get("template_fingerprint"),
        "theme": theme,
        "archetypes": resolved_archetypes,
    }


def _resolve_manifest_theme(
    theme_colors: dict[str, str],
    slide_w: float,
    slide_h: float,
) -> dict[str, Any]:
    """Build a top-level theme object from extracted theme colors."""
    palette = {k: v for k, v in theme_colors.items() if v}

    text_dark = palette.get("dk1", palette.get("tx1", ""))
    text_light = palette.get("lt1", palette.get("bg1", ""))
    primary = palette.get("accent1", "")
    background = palette.get("lt1", palette.get("bg1", ""))

    return {
        "palette": palette,
        "semantic_colors": {
            "text_dark": text_dark,
            "text_light": text_light,
            "primary": primary,
            "background": background,
        },
        "slide_dimensions": {
            "width_in": round(slide_w, 3),
            "height_in": round(slide_h, 3),
        },
    }


def _resolve_manifest_archetype(
    *,
    archetype_id: str,
    arch_spec: dict[str, Any],
    matrix: dict[str, dict[str, dict[str, Any]]],
    layouts_by_name: dict[str, dict[str, Any]],
    theme: dict[str, Any],
) -> dict[str, Any]:
    """Resolve one archetype: merge spec with all compatible layouts."""
    # Collect compatible layouts from the matrix, sorted by score
    arch_matrix = matrix.get(archetype_id, {})
    compatible: list[tuple[str, dict[str, Any]]] = []
    for layout_name, rule in arch_matrix.items():
        status = str(rule.get("status", "blocked"))
        if status in {"preferred", "allowed"}:
            compatible.append((layout_name, rule))
    compatible.sort(key=lambda x: -x[1].get("score", 0))

    # Fallback: use allowed_slide_layouts from arch_spec
    seen = {c[0] for c in compatible}
    for name in arch_spec.get("allowed_slide_layouts", []):
        if name in layouts_by_name and name not in seen:
            compatible.append((name, {"status": "allowed", "score": 0}))
            seen.add(name)

    # Resolve each compatible layout
    resolved_layouts: list[dict[str, Any]] = []
    for layout_name, rule in compatible:
        layout = layouts_by_name.get(layout_name)
        if not layout:
            continue
        resolved_layout = _resolve_manifest_layout(
            layout=layout,
            rule=rule,
            theme=theme,
        )
        resolved_layouts.append(resolved_layout)

    return {
        "description": arch_spec.get("description", ""),
        "story_roles": arch_spec.get("story_roles", []),
        "when_to_use": arch_spec.get("when_to_use", []),
        "guidelines": arch_spec.get("guidelines", []),
        "prohibitions": arch_spec.get("prohibitions", []),
        "op_recipe": arch_spec.get("op_recipe", []),
        "visual_hint": arch_spec.get("visual_hint"),
        "action_title_template": arch_spec.get("action_title_template"),
        "resolved_layouts": resolved_layouts,
    }


def _resolve_manifest_layout(
    *,
    layout: dict[str, Any],
    rule: dict[str, Any],
    theme: dict[str, Any],
) -> dict[str, Any]:
    """Resolve a single layout binding: merge placeholders, zones, and geometry."""
    placeholders = layout.get("placeholders", [])
    role_counts = layout.get("role_counts", {})
    visual_def = layout.get("visual_definition", {})
    color_zones = visual_def.get("color_zones", [])
    title_region = visual_def.get("title_region")
    content_box = layout.get("content_box")
    editable_regions = visual_def.get("editable_regions", [])
    layout_bg = visual_def.get("layout_bg") or theme["semantic_colors"].get("background", "")

    has_title = role_counts.get("title", 0) > 0
    has_subtitle = role_counts.get("subtitle", 0) > 0
    has_body = role_counts.get("body", 0) > 0

    # Detect split-panel: title is in a side zone when there are 2+ zones
    # and the title zone is not "full_slide"
    is_split_panel = False
    if title_region and len(color_zones) >= 2:
        title_zone = title_region.get("zone", "")
        if title_zone and title_zone != "full_slide":
            is_split_panel = True

    # --- Resolve title method from extracted data ---
    title_spec = _resolve_manifest_title(
        has_title=has_title,
        is_split_panel=is_split_panel,
        title_region=title_region,
        color_zones=color_zones,
        placeholders=placeholders,
        layout_bg=layout_bg,
    )

    # --- Resolve subtitle from extracted placeholders ---
    subtitle_spec = None
    if has_subtitle:
        sub_ph = _find_manifest_placeholder(placeholders, "subtitle")
        if sub_ph:
            subtitle_spec = {
                "method": "set_semantic_text",
                "role": "subtitle",
                "placeholder_idx": sub_ph["idx"],
                "geometry": _manifest_ph_geometry(sub_ph),
            }

    # --- Resolve body from extracted placeholders ---
    body_spec = None
    if has_body:
        body_ph = _find_manifest_placeholder(placeholders, "body")
        if body_ph:
            body_spec = {
                "method": "set_semantic_text",
                "role": "body",
                "placeholder_idx": body_ph["idx"],
                "geometry": _manifest_ph_geometry(body_ph),
            }

    # --- Content area from extracted content_box ---
    content_area = None
    if content_box:
        content_area = {
            "left": round(float(content_box["left"]), 2),
            "top": round(float(content_box["top"]), 2),
            "width": round(float(content_box["width"]), 2),
            "height": round(float(content_box["height"]), 2),
        }

    return {
        "layout_name": layout.get("layout_name", ""),
        "layout_index": layout.get("layout_index"),
        "layout_family": layout.get("layout_family", ""),
        "preference": rule.get("status", "allowed"),
        "score": rule.get("score", 0),
        "capabilities": {
            "has_title_ph": has_title,
            "has_subtitle_ph": has_subtitle,
            "has_body_ph": has_body,
            "is_split_panel": is_split_panel,
            "columns_hint": (layout.get("geometry") or {}).get("columns_hint", 0),
        },
        "title": title_spec,
        "subtitle": subtitle_spec,
        "body": body_spec,
        "content_area": content_area,
        "editable_regions": editable_regions,
        "color_zones": color_zones,
        "title_region": title_region,
        "layout_bg": layout_bg,
        "policy": layout.get("policy"),
    }


def _resolve_manifest_title(
    *,
    has_title: bool,
    is_split_panel: bool,
    title_region: dict[str, Any] | None,
    color_zones: list[dict[str, Any]],
    placeholders: list[dict[str, Any]],
    layout_bg: str,
) -> dict[str, Any] | None:
    """Determine how to set the title, using only extracted data."""
    if not has_title:
        # No title placeholder — agent must use add_text; provide layout_bg for contrast
        return {
            "method": "add_text",
            "note": "No title placeholder on this layout"
            " — use add_text with coordinates from editable_regions",
            "font_color": _text_color_for_bg(layout_bg),
        }

    if is_split_panel and title_region:
        # Title placeholder is in a side panel — clear it, place on opposite side
        title_zone_name = title_region.get("zone", "")
        opposite_zone = _find_manifest_opposite_zone(color_zones, title_zone_name)
        if opposite_zone:
            # Use the opposite zone's editable_area for title placement
            editable = opposite_zone.get("editable_area", {})
            zone_text_color = opposite_zone.get(
                "text_color",
                _text_color_for_bg(opposite_zone.get("bg_color", layout_bg)),
            )
            geometry = None
            if editable:
                geometry = {
                    "left": round(float(editable.get("left", 0)), 2),
                    "top": round(float(editable.get("top", 0)), 2),
                    "width": round(float(editable.get("width", 0)), 2),
                    "height": round(float(title_region.get("height", 0.65)), 2),
                }
            return {
                "method": "add_text",
                "note": (
                    f"Title placeholder is in {title_zone_name}"
                    " — clear it and place title on opposite panel"
                ),
                "clear_placeholder": {
                    "op": "set_semantic_text",
                    "role": "title",
                    "text": " ",
                },
                "geometry": geometry,
                "font_color": zone_text_color,
            }

    # Standard title placeholder
    title_ph = _find_manifest_placeholder(placeholders, "title")
    if title_ph:
        return {
            "method": "set_semantic_text",
            "role": "title",
            "placeholder_idx": title_ph["idx"],
            "geometry": _manifest_ph_geometry(title_ph),
        }

    return None


def _find_manifest_placeholder(
    placeholders: list[dict[str, Any]], role: str
) -> dict[str, Any] | None:
    """Find a placeholder by role in the extracted placeholders list."""
    for ph in placeholders:
        if ph.get("role") == role:
            return ph
    return None


def _manifest_ph_geometry(ph: dict[str, Any]) -> dict[str, float]:
    """Extract geometry from an extracted placeholder."""
    return {
        "left": round(float(ph.get("left", 0)), 2),
        "top": round(float(ph.get("top", 0)), 2),
        "width": round(float(ph.get("width", 0)), 2),
        "height": round(float(ph.get("height", 0)), 2),
    }


def _find_manifest_opposite_zone(
    color_zones: list[dict[str, Any]], zone_name: str
) -> dict[str, Any] | None:
    """Find the zone opposite to the title's zone.

    The opposite zone is the one whose region name differs from the
    title zone. When there are exactly 2 zones this is unambiguous.
    For 3+ zones, pick the zone with the largest ``editable_area``.
    """
    candidates = [
        z for z in color_zones if z.get("region", "") != zone_name
    ]
    if not candidates:
        return None
    if len(candidates) == 1:
        return candidates[0]
    # Multiple candidates — prefer the one with editable_area
    with_area = [c for c in candidates if c.get("editable_area")]
    if with_area:
        return max(
            with_area,
            key=lambda z: float(
                z["editable_area"].get("width", 0)
            ) * float(z["editable_area"].get("height", 0)),
        )
    return candidates[0]


def _enrich_template_catalog_with_analysis(
    *,
    template_catalog: dict[str, Any],
    slide_analysis: dict[str, Any],
) -> dict[str, Any]:
    by_layout: dict[str, dict[str, Any]] = {}
    for row in slide_analysis.get("slides", []):
        layout_name = str(row.get("layout_name", "")).strip()
        if not layout_name:
            continue
        bucket = by_layout.setdefault(
            layout_name,
            {
                "archetypes": {},
                "titles": [],
            },
        )
        archetype = str(row.get("recommended_archetype", "")).strip()
        if archetype:
            bucket["archetypes"][archetype] = bucket["archetypes"].get(archetype, 0) + 1
        title = str(row.get("title", "")).strip()
        if title:
            bucket["titles"].append(title)

    for layout in template_catalog.get("layouts", []):
        name = str(layout.get("layout_name", "")).strip()
        bucket = by_layout.get(name, {})
        archetypes = sorted(
            bucket.get("archetypes", {}).keys(),
            key=lambda item: (-bucket["archetypes"][item], item),
        )
        policy = layout.get("policy", {})
        if isinstance(policy, dict):
            policy["allowed_archetypes"] = archetypes
            if bucket.get("titles"):
                policy["example_titles"] = bucket["titles"][:5]
    return template_catalog


def _build_content_layout_catalog(template_catalog: dict[str, Any]) -> dict[str, Any]:
    content_layouts: dict[str, dict[str, Any]] = {
        "title_slide": {
            "required_roles": ["title"],
            "preferred_tokens": ["title slide"],
            "preferred_layout_families": ["cover"],
            "min_body_placeholders": 0,
            "max_body_placeholders": 2,
        },
        "executive_summary": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text"],
            "preferred_layout_families": ["content", "agenda"],
            "min_body_placeholders": 1,
            "max_body_placeholders": 4,
        },
        "agenda": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["agenda"],
            "preferred_layout_families": ["agenda"],
            "min_body_placeholders": 1,
        },
        "content_text": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "content_bullets": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "two_column": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["arrow half", "green half"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "three_column": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["arrow two third", "green two third", "white one third"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "four_column": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["four column"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "table": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text", "special gray"],
            "preferred_layout_families": ["data", "content"],
            "min_body_placeholders": 1,
        },
        "bar_chart": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text", "special gray"],
            "preferred_layout_families": ["data", "content"],
            "min_body_placeholders": 1,
        },
        "line_chart": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text", "special gray"],
            "preferred_layout_families": ["data", "content"],
            "min_body_placeholders": 1,
        },
        "pie_chart": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["title and text", "special gray"],
            "preferred_layout_families": ["data", "content"],
            "min_body_placeholders": 1,
        },
        "big_number": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["green highlight", "big statement icon"],
            "preferred_layout_families": ["highlight", "content"],
            "min_body_placeholders": 1,
        },
        "icon_grid": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["white one third", "green one third"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "process_flow": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["left arrow", "arrow"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "matrix_2x2": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["special gray", "title and text"],
            "preferred_layout_families": ["data", "content"],
            "min_body_placeholders": 1,
        },
        "big_statement": {
            "required_roles": ["title"],
            "preferred_tokens": ["big statement"],
            "preferred_layout_families": ["highlight", "content"],
        },
        "green_panel_text": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["green one third"],
            "preferred_layout_families": ["content"],
            "min_body_placeholders": 1,
        },
        "quote": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["quote"],
            "preferred_layout_families": ["highlight", "content"],
            "min_body_placeholders": 1,
        },
        "section_divider": {
            "required_roles": ["body"],
            "preferred_tokens": ["section header"],
            "preferred_layout_families": ["section_break"],
            "min_body_placeholders": 0,
        },
        "disclaimer": {
            "required_roles": ["body"],
            "preferred_tokens": ["disclaimer"],
            "preferred_layout_families": ["legal"],
            "min_body_placeholders": 0,
        },
        "end_slide": {
            "required_roles": ["title", "body"],
            "preferred_tokens": ["end"],
            "preferred_layout_families": ["closing"],
            "min_body_placeholders": 0,
        },
    }
    archetypes: dict[str, dict[str, Any]] = {
        archetype_id: {
            "content_layout": archetype_id,
            "story_roles": [],
        }
        for archetype_id in content_layouts
    }
    blocked_tokens = {"layout guide"}
    matrix: dict[str, dict[str, dict[str, Any]]] = {}
    layouts = template_catalog.get("layouts", [])
    for content_layout, rules in content_layouts.items():
        matrix[content_layout] = {}
        required_roles = set(rules.get("required_roles", []))
        preferred_tokens = [str(t).lower() for t in rules.get("preferred_tokens", [])]
        preferred_families = {
            str(item).strip().lower()
            for item in rules.get("preferred_layout_families", [])
            if str(item).strip()
        }
        min_body_placeholders = int(rules.get("min_body_placeholders", 0) or 0)
        max_body_placeholders = rules.get("max_body_placeholders")
        max_body_placeholders_int = (
            int(max_body_placeholders)
            if isinstance(max_body_placeholders, (int, float))
            else None
        )
        for layout in layouts:
            layout_name = str(layout.get("layout_name", ""))
            layout_id = str(layout.get("layout_id", ""))
            layout_key = layout_name.lower()
            layout_family = str(layout.get("layout_family", "content")).strip().lower()
            role_counts = layout.get("role_counts", {})
            geometry_raw = layout.get("geometry", {})
            geometry = geometry_raw if isinstance(geometry_raw, dict) else {}
            body_placeholder_count = int(geometry.get("body_placeholder_count", 0) or 0)
            layout_policy_raw = layout.get("policy", {})
            layout_policy = layout_policy_raw if isinstance(layout_policy_raw, dict) else {}
            policy_allowed = {
                str(item).strip()
                for item in layout_policy.get("allowed_content_layouts", [])
                if str(item).strip()
            }
            policy_blocked = {
                str(item).strip()
                for item in layout_policy.get("blocked_content_layouts", [])
                if str(item).strip()
            }
            present_roles = {role for role, count in role_counts.items() if int(count) > 0}
            missing = sorted(required_roles - present_roles)
            status = "allowed"
            reason = "compatible by role requirements"
            score = 0
            matched_rules: list[str] = []
            if any(token in layout_key for token in blocked_tokens):
                status = "blocked"
                reason = "layout is a guide/support layout"
            elif content_layout in policy_blocked:
                status = "blocked"
                reason = "blocked by template layout policy"
            elif policy_allowed and content_layout not in policy_allowed:
                status = "blocked"
                reason = "not in template layout policy allowlist"
            elif preferred_families and layout_family not in preferred_families:
                status = "blocked"
                reason = f"layout family '{layout_family}' not allowed"
            elif missing:
                status = "blocked"
                reason = f"missing required roles: {', '.join(missing)}"
            elif body_placeholder_count < min_body_placeholders:
                status = "blocked"
                reason = (
                    f"body placeholder count {body_placeholder_count} below minimum "
                    f"{min_body_placeholders}"
                )
            elif (
                max_body_placeholders_int is not None
                and body_placeholder_count > max_body_placeholders_int
            ):
                status = "blocked"
                reason = (
                    f"body placeholder count {body_placeholder_count} exceeds maximum "
                    f"{max_body_placeholders_int}"
                )
            elif content_layout in policy_allowed:
                status = "preferred"
                reason = "preferred by template layout policy"
                score += 60
                matched_rules.append("policy_allow")
            elif preferred_tokens and any(token in layout_key for token in preferred_tokens):
                status = "preferred"
                reason = "layout matches preferred token(s)"
                score += 40
                matched_rules.append("token_match")
            if status != "blocked":
                if layout_family in preferred_families:
                    score += 25
                    matched_rules.append("family_match")
                if body_placeholder_count >= min_body_placeholders:
                    score += 10
                    matched_rules.append("geometry_body_count")
                if status == "allowed" and score >= 45:
                    status = "preferred"
                    reason = "high compatibility score"
            matrix[content_layout][layout_name] = {
                "layout_id": layout_id or None,
                "status": status,
                "reason": reason,
                "score": score,
                "layout_family": layout_family,
                "matched_rules": matched_rules,
            }
    return {
        "schema_version": "1.0",
        "template_fingerprint": template_catalog.get("template_fingerprint"),
        "content_layouts": content_layouts,
        "archetypes": archetypes,
        "layout_content_matrix": matrix,
    }


def _load_json_value(raw: str) -> Any:
    text = raw
    if raw.startswith("@"):
        path = Path(raw[1:])
        if not path.exists():
            raise CliError(
                code="ARGUMENT_ERROR",
                message=f"File not found: {path}",
                hint=f"Check the path passed via @{path}. The file must exist before running.",
                exit_code=2,
            )
        text = path.read_text(encoding="utf-8")
    return json.loads(text)


@lru_cache(maxsize=1)
def _load_field_allowlists() -> dict[str, set[str]]:
    if not FIELD_MASKS_PATH.exists():
        raise ValueError(f"Field-mask policy not found: {FIELD_MASKS_PATH}")
    raw = json.loads(FIELD_MASKS_PATH.read_text(encoding="utf-8"))
    if not isinstance(raw, dict):
        raise ValueError("Field-mask policy must be a JSON object")
    allowlists: dict[str, set[str]] = {}
    for payload_name, fields in raw.items():
        if not isinstance(payload_name, str) or not payload_name:
            raise ValueError("Field-mask policy keys must be non-empty strings")
        if not isinstance(fields, list) or not all(isinstance(item, str) for item in fields):
            raise ValueError(
                f"Field-mask policy entry '{payload_name}' must be an array of strings"
            )
        allowlists[payload_name] = set(fields)
    return allowlists


def _validate_no_control_chars(value: str, *, field: str) -> None:
    for ch in value:
        code = ord(ch)
        if code < 32 and ch not in {"\n", "\r", "\t"}:
            raise ValueError(f"{field} contains control characters")


def _validate_slide_id(value: str, *, field: str) -> None:
    _validate_no_control_chars(value, field=field)
    if any(char in value for char in ["?", "#", "%"]):
        raise ValueError(f"{field} contains forbidden characters (?, #, %)")


def _validate_output_path(path: Path | None) -> None:
    if path is None:
        return
    enforce_cwd = os.environ.get("SLIDES_ENFORCE_CWD", "0")
    if enforce_cwd != "1":
        return
    cwd = Path.cwd().resolve()
    resolved = path.resolve()
    if cwd != resolved and cwd not in resolved.parents:
        raise ValueError(f"Output path outside current workspace: {resolved}")


def _resolve_profile_relative_path(raw: str, *, profile_source: Path | None) -> str:
    path = Path(raw).expanduser()
    if path.is_absolute() or profile_source is None:
        return str(path)
    return str((profile_source.parent / path).resolve())


def _parse_field_list(fields: str | None) -> list[str]:
    if not fields:
        return []
    return [f.strip() for f in fields.split(",") if f.strip()]


def _validate_fields_requested(
    fields: str | None, *, allowlist: set[str], payload_name: str
) -> None:
    requested = _parse_field_list(fields)
    if not requested:
        return
    invalid = [field for field in requested if field not in allowlist]
    if invalid:
        allowed = ", ".join(sorted(allowlist))
        raise ValueError(
            f"Unsupported field mask for {payload_name}: {', '.join(invalid)}. Allowed: {allowed}"
        )


def _project_fields(obj: Any, field_tree: dict[str, Any]) -> Any:
    """Project *obj* keeping only fields described by *field_tree*.

    *field_tree* is a nested dict where leaves are ``True`` (include this
    key) and branches are sub-trees.  Lists are traversed transparently
    so ``{"slides": {"title": True}}`` applied to
    ``{"slides": [{"title": "A", "x": 1}]}`` yields
    ``{"slides": [{"title": "A"}]}``.
    """
    if not field_tree:
        return obj
    if isinstance(obj, dict):
        result: dict[str, Any] = {}
        for key, subtree in field_tree.items():
            if key not in obj:
                continue
            if subtree is True:
                result[key] = obj[key]
            else:
                result[key] = _project_fields(obj[key], subtree)
        return result
    if isinstance(obj, list):
        return [_project_fields(item, field_tree) for item in obj]
    return obj


def _build_field_tree(field_paths: list[str]) -> dict[str, Any]:
    """Turn a list of dot-separated paths into a nested dict tree.

    >>> _build_field_tree(["summary.slide_count", "slides.title"])
    {"summary": {"slide_count": True}, "slides": {"title": True}}
    """
    tree: dict[str, Any] = {}
    for raw in field_paths:
        parts = [p for p in raw.split(".") if p]
        if not parts:
            continue
        cur = tree
        for part in parts[:-1]:
            nxt = cur.get(part)
            if nxt is True:
                break
            if not isinstance(nxt, dict):
                nxt = {}
                cur[part] = nxt
            cur = nxt
        else:
            existing = cur.get(parts[-1])
            if not isinstance(existing, dict):
                cur[parts[-1]] = True
    return tree


def _apply_fields(payload: Any, fields: str | None) -> Any:
    requested = _parse_field_list(fields)
    if not requested:
        return payload
    if not isinstance(payload, dict):
        return payload
    tree = _build_field_tree(requested)
    out = _project_fields(payload, tree)
    return out if out else payload


def _offset_batch_slide_indices(batch: OperationBatch, *, offset: int) -> OperationBatch:
    if offset <= 0:
        return batch
    payload = batch.model_dump()
    operations = payload.get("operations")
    if not isinstance(operations, list):
        return batch
    for op in operations:
        if not isinstance(op, dict):
            continue
        value = op.get("slide_index")
        if isinstance(value, int):
            op["slide_index"] = value + offset
    return OperationBatch.model_validate(payload)


def _strip_empty(obj: Any) -> Any:
    """Recursively remove None values and empty collections from JSON-serializable data."""
    if isinstance(obj, dict):
        return {k: _strip_empty(v) for k, v in obj.items() if v is not None and v != [] and v != {}}
    if isinstance(obj, list):
        return [_strip_empty(item) for item in obj]
    return obj


def _write_payload(
    payload: Any,
    *,
    payload_name: str,
    path: Path | None,
    fields: str | None = None,
    ndjson: bool = False,
    ndjson_root: str | None = None,
    quiet: bool = False,
    compact: bool = False,
) -> None:
    allowlist = _load_field_allowlists().get(payload_name)
    if allowlist is not None:
        _validate_fields_requested(fields, allowlist=allowlist, payload_name=payload_name)
    filtered = _apply_fields(payload, fields)
    if compact:
        filtered = _strip_empty(filtered)

    if (
        ndjson
        and ndjson_root
        and isinstance(filtered, dict)
        and isinstance(filtered.get(ndjson_root), list)
    ):
        items = filtered[ndjson_root]
        if compact:
            items = [_strip_empty(item) for item in items]
        lines = [json.dumps(item, ensure_ascii=True) for item in items]
        text = "\n".join(lines) + ("\n" if lines else "")
        if path is not None:
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_text(text, encoding="utf-8")
            if quiet:
                return
        print(text, end="")
        return

    indent = None if compact else 2
    text = json.dumps(filtered, indent=indent)
    if path is not None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(text + "\n", encoding="utf-8")
        if quiet:
            return
    print(text)


def _paginate_items(
    items: list[dict[str, Any]],
    *,
    page_size: int | None,
    page_token: str | None,
) -> tuple[list[dict[str, Any]], str | None]:
    if page_size is None or page_size <= 0:
        return items, None
    offset = 0
    if page_token:
        if not page_token.isdigit():
            raise ValueError("page_token must be an integer offset")
        offset = int(page_token)
    if offset < 0 or offset > len(items):
        raise ValueError("page_token out of range")
    page = items[offset : offset + page_size]
    next_offset = offset + page_size
    next_token = str(next_offset) if next_offset < len(items) else None
    return page, next_token


def _parse_page_offset(page_token: str | None, total: int) -> int:
    if not page_token:
        return 0
    if not page_token.isdigit():
        raise ValueError("page_token must be an integer offset")
    offset = int(page_token)
    if offset < 0 or offset > total:
        raise ValueError("page_token out of range")
    return offset


def _emit_paginated_results(
    items: list[dict[str, Any]],
    *,
    envelope_key: str,
    summary_fields: dict[str, Any],
    payload_name: str,
    out_path: Path | None,
    page_all: bool,
    page_size: int | None,
    page_token: str | None,
    fields: str | None,
    ndjson: bool,
    quiet: bool,
    compact: bool,
) -> None:
    """Emit a paginated list of items, handling page_all + ndjson and single-page modes."""
    if page_all and page_size:
        offset = _parse_page_offset(page_token, len(items))
        pages: list[dict[str, Any]] = []
        while offset < len(items):
            page = items[offset : offset + page_size]
            offset += page_size
            next_token = str(offset) if offset < len(items) else None
            page_payload = {
                **summary_fields,
                envelope_key: page,
                "page_size": page_size,
                "next_page_token": next_token,
            }
            pages.append(page_payload)
        if ndjson:
            rendered = [_strip_empty(p) for p in pages] if compact else pages
            lines = [json.dumps(p, ensure_ascii=True) for p in rendered]
            text = "\n".join(lines) + ("\n" if lines else "")
            if out_path is not None:
                out_path.parent.mkdir(parents=True, exist_ok=True)
                out_path.write_text(text, encoding="utf-8")
            if not (quiet and out_path is not None):
                print(text, end="")
        else:
            _write_payload(
                {"pages": pages},
                payload_name=payload_name,
                path=out_path,
                fields=fields,
                ndjson=False,
                quiet=quiet,
                compact=compact,
            )
    else:
        items_page, next_token = _paginate_items(
            items,
            page_size=page_size,
            page_token=page_token,
        )
        payload = {
            **summary_fields,
            envelope_key: items_page,
            "page_size": page_size,
            "next_page_token": next_token,
        }
        _write_payload(
            payload,
            payload_name=payload_name,
            path=out_path,
            fields=fields,
            ndjson=ndjson,
            ndjson_root=envelope_key,
            quiet=quiet,
            compact=compact,
        )


def _load_template_index(profile: DesignProfile) -> dict[str, Any] | None:
    if not profile.template_path:
        return None
    template = Path(profile.template_path)
    if not template.exists():
        return None
    try:
        pres = Presentation.open(template)
        return pres.inspect()
    except Exception:  # noqa: BLE001
        return None



def _resolve_output_options(args: argparse.Namespace) -> tuple[str | None, bool, bool, bool]:
    """Extract common output-formatting options from parsed args.

    Returns (fields, ndjson, verbose, compact) — the quad used by all write helpers.
    """
    fields = getattr(args, "fields", None)
    ndjson = getattr(args, "ndjson", False)
    verbose = getattr(args, "verbose", False)
    compact = getattr(args, "compact", False) or (not verbose)
    return fields, ndjson, not verbose, compact  # (fields, ndjson, quiet, compact)


def _resolve_profile(args: argparse.Namespace) -> tuple[DesignProfile, Path | None]:
    """Build a DesignProfile from --profile / --profile-json args."""
    profile_payload = None
    profile_source: Path | None = None
    profile_json_raw = getattr(args, "profile_json", None)
    profile_path = getattr(args, "profile", None)
    if profile_json_raw:
        profile_payload = _load_json_value(profile_json_raw)
        if profile_json_raw.startswith("@"):
            profile_source = Path(profile_json_raw[1:]).resolve()
    elif profile_path:
        profile_payload = json.loads(profile_path.read_text(encoding="utf-8"))
        profile_source = profile_path.resolve()
    profile = DesignProfile.model_validate(profile_payload or {})
    updates: dict[str, str] = {}
    if profile.template_path:
        updates["template_path"] = _resolve_profile_relative_path(
            profile.template_path, profile_source=profile_source
        )
    if profile.icon_pack_dir:
        updates["icon_pack_dir"] = _resolve_profile_relative_path(
            profile.icon_pack_dir, profile_source=profile_source
        )
    if profile.content_layout_catalog_path:
        updates["content_layout_catalog_path"] = _resolve_profile_relative_path(
            profile.content_layout_catalog_path, profile_source=profile_source
        )
    if profile.archetypes_catalog_path:
        updates["archetypes_catalog_path"] = _resolve_profile_relative_path(
            profile.archetypes_catalog_path, profile_source=profile_source
        )
    if updates:
        profile = profile.model_copy(update=updates)
    return profile, profile_source


def _preflight_issue(
    *,
    code: str,
    severity: str,
    message: str,
    path: str | None = None,
) -> dict[str, Any]:
    issue: dict[str, Any] = {
        "code": code,
        "severity": severity,
        "message": message,
    }
    if path is not None:
        issue["path"] = path
    return issue


def _check_preflight_path(
    *,
    issues: list[dict[str, Any]],
    code_prefix: str,
    path: Path,
    field_path: str,
    expect_dir: bool = False,
    required: bool = True,
) -> dict[str, Any]:
    exists = path.exists()
    if not exists:
        if required:
            issues.append(
                _preflight_issue(
                    code=f"{code_prefix}_MISSING",
                    severity="error",
                    message=f"Missing required path: {path}",
                    path=field_path,
                )
            )
        return {"ok": not required, "path": str(path), "exists": False}
    is_valid_type = path.is_dir() if expect_dir else path.is_file()
    if not is_valid_type:
        issues.append(
            _preflight_issue(
                code=f"{code_prefix}_TYPE",
                severity="error",
                message=f"Expected {'directory' if expect_dir else 'file'}: {path}",
                path=field_path,
            )
        )
    return {"ok": is_valid_type, "path": str(path), "exists": True}


def _resolve_icon_pack_dirs(args: argparse.Namespace, profile: DesignProfile | None) -> list[Path]:
    resolved: list[Path] = []
    cli_dir = getattr(args, "icon_pack_dir", None)
    if cli_dir is not None:
        resolved.append(Path(cli_dir).expanduser())
    if profile and profile.icon_pack_dir:
        resolved.append(Path(profile.icon_pack_dir).expanduser())
    env_dir = os.environ.get("SLIDES_ICON_PACK_DIR")
    if env_dir:
        resolved.append(Path(env_dir).expanduser())

    deduped: list[Path] = []
    for path in resolved:
        if path not in deduped:
            deduped.append(path)
    return deduped


def _resolve_ops_context(args: argparse.Namespace) -> OperationBatch | None:
    """Load an OperationBatch from --ops-json / --slides-json for context-only use."""
    ops_raw = getattr(args, "ops_json", None)
    slides_raw = getattr(args, "slides_json", None)
    if ops_raw:
        return OperationBatch.model_validate(_load_json_value(ops_raw))
    if slides_raw:
        slides_doc = SlidesDocument.model_validate(_load_json_value(slides_raw))
        if slides_doc.ops:
            return slides_doc.ops
    return None


_STRUCTURAL_STORY_ROLES = {
    "title",
    "title_slide",
    "agenda",
    "section_divider",
    "disclaimer",
    "end_slide",
}


def _is_content_story_role(role: str) -> bool:
    return role.strip().lower() not in _STRUCTURAL_STORY_ROLES


def _plan_slide_row(slide: Any) -> dict[str, Any]:
    action_title = str(slide.action_title or "")
    points = [str(p) for p in slide.key_points]
    text_chars = len(action_title) + sum(len(p) for p in points)
    return {
        "slide_number": int(slide.slide_number),
        "story_role": str(slide.story_role),
        "archetype_id": str(slide.archetype_id),
        "action_title": action_title,
        "key_point_count": len(points),
        "text_chars_estimate": text_chars,
        "has_source_note": bool(slide.source_note),
        "is_content": _is_content_story_role(str(slide.story_role)),
    }


# ── Subcommand handlers ──────────────────────────────────────────────


def _cmd_extract(args: argparse.Namespace) -> int:
    source_path: Path = args.source
    if not source_path.exists():
        raise CliError(
            code="ARGUMENT_ERROR",
            message=f"Source file not found: {source_path}",
            exit_code=2,
        )
    output_dir: Path = args.output_dir
    _, _, _, compact = _resolve_output_options(args)
    verbose = getattr(args, "verbose", False)

    if args.layout_preview_dir is not None:
        _warn_deprecated_flag("--layout-preview-dir", "--output-dir")
    if args.screenshots_dir is not None:
        _warn_deprecated_flag("--screenshots-dir", "--output-dir")

    template_out = args.template_out or (output_dir / "template_catalog.json")
    content_out = args.content_layout_out or (output_dir / "content_layout_catalog.json")
    archetypes_out = args.archetypes_out or (output_dir / "archetypes.json")
    slides_manifest_out = args.slides_manifest_out or (output_dir / "slides_manifest.json")
    slide_analysis_out = args.slide_analysis_out or (output_dir / "slide_analysis.json")
    screenshots_dir = args.screenshots_dir or (output_dir / "slide_screenshots")
    resolved_manifest_out = args.resolved_manifest_out or (output_dir / "resolved_manifest.json")

    template_catalog = _extract_template_layout_catalog(
        source_path, preview_dir=args.layout_preview_dir,
    )
    content_catalog = _build_content_layout_catalog(template_catalog)
    screenshot_map = _render_slide_screenshots(
        source_path=source_path, screenshots_dir=screenshots_dir,
    )
    slides_manifest = _extract_slides_manifest(
        source_path=source_path, template_catalog=template_catalog,
        screenshot_map=screenshot_map,
    )
    slide_analysis = _analyze_slides_manifest(slides_manifest)
    archetypes_payload = _build_archetypes_from_analysis(
        slide_analysis=slide_analysis, template_catalog=template_catalog,
        content_catalog=content_catalog,
    )
    template_catalog = _enrich_template_catalog_with_analysis(
        template_catalog=template_catalog, slide_analysis=slide_analysis,
    )
    resolved_manifest = _build_resolved_manifest(
        template_catalog=template_catalog, content_catalog=content_catalog,
        archetypes_catalog=archetypes_payload,
    )

    for p in [template_out, content_out, archetypes_out, slides_manifest_out,
              slide_analysis_out, resolved_manifest_out]:
        p.parent.mkdir(parents=True, exist_ok=True)
    indent = None if compact else 2
    template_payload = _strip_empty(template_catalog) if compact else template_catalog
    content_payload = _strip_empty(content_catalog) if compact else content_catalog
    archetypes_payload_out = _strip_empty(archetypes_payload) if compact else archetypes_payload
    slides_manifest_payload = _strip_empty(slides_manifest) if compact else slides_manifest
    slide_analysis_payload = _strip_empty(slide_analysis) if compact else slide_analysis
    resolved_manifest_payload = _strip_empty(resolved_manifest) if compact else resolved_manifest
    template_out.write_text(json.dumps(template_payload, indent=indent) + "\n", encoding="utf-8")
    content_out.write_text(json.dumps(content_payload, indent=indent) + "\n", encoding="utf-8")
    archetypes_out.write_text(
        json.dumps(archetypes_payload_out, indent=indent) + "\n",
        encoding="utf-8",
    )
    slides_manifest_out.write_text(
        json.dumps(slides_manifest_payload, indent=indent) + "\n",
        encoding="utf-8",
    )
    slide_analysis_out.write_text(
        json.dumps(slide_analysis_payload, indent=indent) + "\n",
        encoding="utf-8",
    )
    resolved_manifest_out.write_text(
        json.dumps(resolved_manifest_payload, indent=indent) + "\n", encoding="utf-8",
    )

    icons_dir = output_dir / "icons"
    extracted_icons = _extract_icons(source_path, icons_dir)

    base_template_out = args.base_template_out
    if base_template_out is not None:
        base_template_out.parent.mkdir(parents=True, exist_ok=True)
        base_prs = Presentation.open(source_path)
        base_prs.clear_slides()
        base_prs.save(base_template_out)

    report: dict[str, Any] = {
        "ok": True,
        "source": str(source_path),
        "template_out": str(template_out),
        "content_layout_out": str(content_out),
        "archetypes_out": str(archetypes_out),
        "resolved_manifest_out": str(resolved_manifest_out),
        "slides_manifest_out": str(slides_manifest_out),
        "slide_analysis_out": str(slide_analysis_out),
        "screenshots_dir": str(screenshots_dir),
        "screenshot_count": len(screenshot_map),
        "layout_count": template_catalog["layout_count"],
        "content_layout_count": len(content_catalog["content_layouts"]),
        "archetype_count": len(archetypes_payload.get("archetypes", {})),
        "resolved_archetype_count": len(resolved_manifest.get("archetypes", {})),
        "icons_dir": str(icons_dir),
        "icon_count": len(extracted_icons),
    }
    if base_template_out is not None:
        report["base_template_out"] = str(base_template_out)
    if verbose:
        indent = None if compact else 2
        print(json.dumps(report, indent=indent))
    return 0


def _cmd_render(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    profile, _ = _resolve_profile(args)
    icon_pack_dirs = _resolve_icon_pack_dirs(args, profile)

    slides_payload = _load_json_value(args.slides_json)
    slides_doc = SlidesDocument.model_validate(slides_payload)
    compiled_batch = slides_doc.ops or compile_plan_to_operations(slides_doc.plan, profile=profile)

    if args.slides_out:
        normalized = SlidesDocument(plan=slides_doc.plan, ops=compiled_batch)
        _write_payload(
            normalized.model_dump(), payload_name="slides-document",
            path=args.slides_out, fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
        )

    template_path = args.template
    effective_template_path = template_path
    if effective_template_path is None and profile.template_path:
        effective_template_path = Path(profile.template_path)

    pres = Presentation.create(template_path=effective_template_path, icon_dirs=icon_pack_dirs)
    batch = compiled_batch

    base_slide_count = int(pres.summarize().slide_count)
    if base_slide_count > 0:
        batch = _offset_batch_slide_indices(batch, offset=base_slide_count)

    dry_run = args.dry_run
    report = pres.apply_operations(
        batch, dry_run=dry_run, transactional=not args.no_transaction,
    )
    if not quiet:
        indent = None if compact else 2
        print(report.model_dump_json(indent=indent))
    if not report.ok:
        return 2

    if not dry_run:
        output_path = args.output
        if output_path is None:
            raise CliError(
                code="ARGUMENT_ERROR",
                message="--output is required unless --dry-run is used",
                exit_code=2,
            )
        _validate_output_path(output_path)
        pres.save(output_path, deterministic=not args.non_deterministic)
    return 0


def _cmd_apply(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    profile, _ = _resolve_profile(args)
    icon_pack_dirs = _resolve_icon_pack_dirs(args, profile)

    pres = (
        Presentation.open(args.input, icon_dirs=icon_pack_dirs)
        if args.input is not None
        else Presentation.create(icon_dirs=icon_pack_dirs)
    )
    ops_payload = _load_json_value(args.ops_json)
    batch = OperationBatch.model_validate(ops_payload)

    dry_run = args.dry_run
    report = pres.apply_operations(
        batch, dry_run=dry_run, transactional=not args.no_transaction,
    )
    if not quiet:
        indent = None if compact else 2
        print(report.model_dump_json(indent=indent))
    if not report.ok:
        return 2

    if not dry_run:
        _validate_output_path(args.output)
        pres.save(args.output, deterministic=not args.non_deterministic)
    return 0


def _cmd_inspect(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    pres = Presentation.open(args.input)

    if args.summary:
        _write_payload(
            pres.summarize().to_dict(), payload_name="summary",
            path=None, fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
        )
        return 0

    if args.fingerprint:
        _write_payload(
            {"fingerprint": pres.fingerprint()}, payload_name="fingerprint",
            path=None, fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
        )
        return 0

    if args.placeholders is not None:
        _write_payload(
            {
                "slide_index": args.placeholders,
                "placeholders": pres.list_placeholders(slide_index=args.placeholders),
            },
            payload_name="placeholders", path=None, fields=fields,
            ndjson=ndjson, quiet=quiet, compact=compact,
        )
        return 0

    page_size = getattr(args, "page_size", None)
    page_token = getattr(args, "page_token", None)
    page_all = getattr(args, "page_all", False)
    if page_size is None and not page_all:
        page_size = DEFAULT_PAGE_SIZE
    if page_size is not None and int(page_size) <= 0:
        raise ValueError("page_size must be > 0")

    inspect_payload = pres.inspect()
    _emit_paginated_results(
        list(inspect_payload["slides"]),
        envelope_key="slides",
        summary_fields={"summary": inspect_payload["summary"]},
        payload_name="inspect",
        out_path=args.out,
        page_all=page_all, page_size=page_size, page_token=page_token,
        fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
    )
    return 0


def _cmd_find(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    _validate_no_control_chars(args.query, field="query")
    pres = Presentation.open(args.input)

    page_size = getattr(args, "page_size", None)
    page_token = getattr(args, "page_token", None)
    page_all = getattr(args, "page_all", False)
    if page_size is None and not page_all:
        page_size = DEFAULT_PAGE_SIZE
    if page_size is not None and int(page_size) <= 0:
        raise ValueError("page_size must be > 0")

    all_results = pres.find_text(query=args.query, limit=max(1, args.limit))
    _emit_paginated_results(
        all_results, envelope_key="results",
        summary_fields={"query": args.query},
        payload_name="find", out_path=args.out,
        page_all=page_all, page_size=page_size, page_token=page_token,
        fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
    )
    return 0


def _cmd_plan_inspect(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    page_size = getattr(args, "page_size", None)
    page_token = getattr(args, "page_token", None)
    page_all = getattr(args, "page_all", False)
    if page_size is None and not page_all:
        page_size = DEFAULT_PAGE_SIZE
    if page_size is not None and int(page_size) <= 0:
        raise ValueError("page_size must be > 0")

    slides_payload = _load_json_value(args.slides_json)
    slides_doc = SlidesDocument.model_validate(slides_payload)
    rows = [_plan_slide_row(slide) for slide in slides_doc.plan.slides]
    if args.content_only:
        rows = [row for row in rows if bool(row.get("is_content", False))]

    roles = sorted({str(row["story_role"]) for row in rows})
    archetypes = sorted({str(row["archetype_id"]) for row in rows})
    content_count = sum(1 for row in rows if bool(row.get("is_content", False)))
    summary = {
        "slide_count": len(rows),
        "content_slide_count": content_count,
        "structural_slide_count": len(rows) - content_count,
        "story_roles": roles,
        "archetypes": archetypes,
        "ops_count": len(slides_doc.ops.operations) if slides_doc.ops is not None else 0,
    }
    if args.summary_only:
        _write_payload(
            {"summary": summary},
            payload_name="plan-inspect",
            path=args.out,
            fields=fields,
            ndjson=ndjson,
            quiet=quiet,
            compact=compact,
        )
        return 0

    _emit_paginated_results(
        rows,
        envelope_key="slides",
        summary_fields={"summary": summary},
        payload_name="plan-inspect",
        out_path=args.out,
        page_all=page_all,
        page_size=page_size,
        page_token=page_token,
        fields=fields,
        ndjson=ndjson,
        quiet=quiet,
        compact=compact,
    )
    return 0


def _cmd_validate(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    pres = Presentation.open(args.input)

    validation = pres.validate(
        deep=args.deep, xsd_dir=args.xsd_dir, require_xsd=args.require_xsd,
    )
    severities = [issue.severity for issue in validation.issues]
    _write_payload(
        {"ok": validation.ok, "issues": [asdict(issue) for issue in validation.issues]},
        payload_name="validate", path=None, fields=fields,
        ndjson=ndjson, quiet=quiet, compact=compact,
    )
    if args.fail_on_error and "error" in severities:
        return 3
    if args.fail_on_warning and any(level in ("warning", "error") for level in severities):
        return 4
    return 0


def _cmd_lint(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    profile, _ = _resolve_profile(args)
    pres = Presentation.open(args.input)
    batch = _resolve_ops_context(args)
    template_index = _load_template_index(profile)

    report = lint_design(
        deck_index=pres.inspect(), profile=profile,
        template_index=template_index, batch=batch,
    )
    _write_payload(
        report, payload_name="lint", path=args.out,
        fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
    )
    if not report["ok"]:
        return 5
    return 0


def _cmd_qa(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    profile, _ = _resolve_profile(args)
    pres = Presentation.open(args.input)
    batch = _resolve_ops_context(args)
    template_index = _load_template_index(profile)
    template_path = Path(profile.template_path) if profile.template_path else None

    validation = pres.validate(
        deep=getattr(args, "deep", False),
        xsd_dir=getattr(args, "xsd_dir", None),
        require_xsd=getattr(args, "require_xsd", False),
    )
    lint_report = lint_design(
        deck_index=pres.inspect(), profile=profile,
        template_index=template_index, batch=batch,
    )
    assets_report = verify_assets(
        profile=profile, batch=batch,
        input_path=args.input, template_path=template_path,
    )
    qa_report = {
        "ok": validation.ok and lint_report["ok"] and assets_report["ok"],
        "checks": {
            "validate": {"ok": validation.ok, "issue_count": len(validation.issues)},
            "lint": {"ok": lint_report["ok"], "issue_count": lint_report["issue_count"]},
            "assets": {"ok": assets_report["ok"], "issue_count": assets_report["issue_count"]},
        },
        "summary": pres.summarize().to_dict(),
    }
    _write_payload(
        qa_report, payload_name="qa", path=args.out,
        fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
    )
    if not qa_report["ok"]:
        return 9
    return 0


def _cmd_edit(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    _validate_no_control_chars(args.query, field="query")
    _validate_no_control_chars(args.replacement, field="replacement")
    if args.slide_uid:
        _validate_no_control_chars(args.slide_uid, field="slide_uid")
    if args.shape_uid:
        _validate_no_control_chars(args.shape_uid, field="shape_uid")
    if args.slide_id:
        _validate_slide_id(args.slide_id, field="slide_id")

    if sum(1 for s in [args.slide, args.slide_id, args.slide_uid] if s is not None) > 1:
        raise CliError(
            code="ARGUMENT_ERROR",
            message="Use only one of --slide, --slide-id, --slide-uid",
            exit_code=2,
        )
    if sum(1 for s in [args.shape_id, args.shape_uid] if s is not None) > 1:
        raise CliError(
            code="ARGUMENT_ERROR",
            message="Use only one of --shape-id, --shape-uid",
            exit_code=2,
        )

    pres = Presentation.open(args.input)
    edit_report = pres.semantic_replace_text(
        query=args.query, replacement=args.replacement,
        slide_index=args.slide, slide_id=args.slide_id, slide_uid=args.slide_uid,
        shape_id=args.shape_id, shape_uid=args.shape_uid,
    )
    _write_payload(
        {"edit": edit_report}, payload_name="edit", path=None,
        fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
    )
    if edit_report["replaced_paragraphs"] == 0:
        return 6

    _validate_output_path(args.output)
    pres.save(args.output)
    return 0


def _cmd_transform(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    if args.slide_uid:
        _validate_no_control_chars(args.slide_uid, field="slide_uid")
    if args.slide_id:
        _validate_slide_id(args.slide_id, field="slide_id")

    if sum(1 for s in [args.slide, args.slide_id, args.slide_uid] if s is not None) != 1:
        raise CliError(
            code="ARGUMENT_ERROR",
            message="Provide exactly one of --slide, --slide-id, --slide-uid",
            exit_code=2,
        )

    pres = Presentation.open(args.input)
    transform_report = pres.transform_slide_to_timeline(
        slide_index=args.slide, slide_id=args.slide_id,
        slide_uid=args.slide_uid, title=None,
    )
    _write_payload(
        {"transform": transform_report}, payload_name="transform", path=None,
        fields=fields, ndjson=ndjson, quiet=quiet, compact=compact,
    )
    _validate_output_path(args.output)
    pres.save(args.output)
    return 0


def _cmd_repair(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    pres = Presentation.open(args.input)
    validation = pres.repair()
    _write_payload(
        {
            "repaired": True, "ok": validation.ok,
            "issues": [asdict(issue) for issue in validation.issues],
        },
        payload_name="repair", path=None, fields=fields,
        ndjson=ndjson, quiet=quiet, compact=compact,
    )
    _validate_output_path(args.output)
    pres.save(args.output, deterministic=not args.non_deterministic)
    return 0


def _cmd_preflight(args: argparse.Namespace) -> int:
    fields, ndjson, quiet, compact = _resolve_output_options(args)
    project_dir = args.project_dir.expanduser().resolve()
    issues: list[dict[str, Any]] = []
    artifact_checks: dict[str, dict[str, Any]] = {}
    dependency_checks: dict[str, dict[str, Any]] = {}

    project_ok = project_dir.exists() and project_dir.is_dir()
    if not project_ok:
        issues.append(
            _preflight_issue(
                code="PROJECT_DIR_INVALID",
                severity="error",
                message=f"Project directory not found: {project_dir}",
                path="project_dir",
            )
        )

    profile: DesignProfile | None = None
    profile_source: Path | None = None
    profile_arg = getattr(args, "profile", None)
    profile_json_arg = getattr(args, "profile_json", None)
    if profile_arg is None and profile_json_arg is None and project_ok:
        default_profile = project_dir / "design-profile.json"
        if default_profile.exists():
            profile_arg = default_profile

    if isinstance(profile_arg, Path) and not profile_arg.exists():
        issues.append(
            _preflight_issue(
                code="PROFILE_MISSING",
                severity="error",
                message=f"Design profile not found: {profile_arg}",
                path="profile",
            )
        )
    else:
        profile_args = argparse.Namespace(**vars(args))
        profile_args.profile = profile_arg
        profile_args.profile_json = profile_json_arg
        profile, profile_source = _resolve_profile(profile_args)

    if project_ok:
        artifact_specs = {
            "resolved_manifest": ["resolved_manifest.json"],
            "template_catalog": ["template_catalog.json", "template_layout.json"],
            "content_layout_catalog": ["content_layout_catalog.json", "content_layout.json"],
            "archetypes_catalog": ["archetypes.json"],
        }
        for artifact_name, candidates in artifact_specs.items():
            chosen = next(
                (
                    project_dir / candidate
                    for candidate in candidates
                    if (project_dir / candidate).exists()
                ),
                project_dir / candidates[0],
            )
            artifact_checks[artifact_name] = _check_preflight_path(
                issues=issues,
                code_prefix=artifact_name.upper(),
                path=chosen,
                field_path=f"project_dir/{chosen.name}",
            )

    profile_checks: dict[str, Any] = {
        "loaded": profile_source is not None or bool(profile_json_arg),
        "source": str(profile_source) if profile_source is not None else None,
    }
    if profile is not None:
        if not profile_checks["loaded"]:
            issues.append(
                _preflight_issue(
                    code="PROFILE_NOT_PROVIDED",
                    severity="warning",
                    message="No design profile was provided or found in the project directory.",
                    path="profile",
                )
            )

        if profile.template_path:
            profile_checks["template_path"] = _check_preflight_path(
                issues=issues,
                code_prefix="TEMPLATE_PATH",
                path=Path(profile.template_path).expanduser(),
                field_path="profile.template_path",
            )
        else:
            issues.append(
                _preflight_issue(
                    code="TEMPLATE_PATH_MISSING",
                    severity="error",
                    message="Design profile is missing template_path.",
                    path="profile.template_path",
                )
            )

        if profile.content_layout_catalog_path:
            profile_checks["content_layout_catalog_path"] = _check_preflight_path(
                issues=issues,
                code_prefix="CONTENT_LAYOUT_CATALOG_PATH",
                path=Path(profile.content_layout_catalog_path).expanduser(),
                field_path="profile.content_layout_catalog_path",
            )
        else:
            issues.append(
                _preflight_issue(
                    code="CONTENT_LAYOUT_CATALOG_PATH_MISSING",
                    severity="warning",
                    message="Design profile is missing content_layout_catalog_path.",
                    path="profile.content_layout_catalog_path",
                )
            )

        if profile.archetypes_catalog_path:
            profile_checks["archetypes_catalog_path"] = _check_preflight_path(
                issues=issues,
                code_prefix="ARCHETYPES_CATALOG_PATH",
                path=Path(profile.archetypes_catalog_path).expanduser(),
                field_path="profile.archetypes_catalog_path",
            )

        if profile.icon_pack_dir:
            icon_dir = Path(profile.icon_pack_dir).expanduser()
            icon_check = _check_preflight_path(
                issues=issues,
                code_prefix="ICON_PACK_DIR",
                path=icon_dir,
                field_path="profile.icon_pack_dir",
                expect_dir=True,
            )
            if icon_dir.exists() and icon_dir.is_dir():
                icon_count = len(list(icon_dir.glob("*.xml")))
                icon_check["xml_count"] = icon_count
                if icon_count == 0:
                    issues.append(
                        _preflight_issue(
                            code="ICON_PACK_EMPTY",
                            severity="warning",
                            message=f"Icon pack directory contains no .xml files: {icon_dir}",
                            path="profile.icon_pack_dir",
                        )
                    )
            profile_checks["icon_pack_dir"] = icon_check

        if profile.asset_roots:
            roots: list[dict[str, Any]] = []
            for idx, raw_root in enumerate(profile.asset_roots):
                roots.append(
                    _check_preflight_path(
                        issues=issues,
                        code_prefix="ASSET_ROOT",
                        path=Path(raw_root).expanduser(),
                        field_path=f"profile.asset_roots[{idx}]",
                        expect_dir=True,
                    )
                )
            profile_checks["asset_roots"] = roots

        asset_report = verify_assets(
            profile=profile,
            template_path=(
                Path(profile.template_path).expanduser() if profile.template_path else None
            ),
        )
        issues.extend(asset_report["issues"])

    optional_tools = {
        "soffice": "Required for PDF conversion during layout preview and screenshot extraction",
        "pdftoppm": "Required for PNG rendering during layout preview and screenshot extraction",
    }
    optional_severity = "error" if args.require_optional_deps else "warning"
    for tool_name, description in optional_tools.items():
        available = shutil.which(tool_name) is not None
        dependency_checks[tool_name] = {
            "available": available,
            "description": description,
        }
        if not available:
            issues.append(
                _preflight_issue(
                    code="OPTIONAL_DEPENDENCY_MISSING",
                    severity=optional_severity,
                    message=f"{tool_name} is unavailable. {description}",
                    path=f"dependencies.{tool_name}",
                )
            )

    by_severity: dict[str, int] = {}
    by_code: dict[str, int] = {}
    for issue in issues:
        severity = str(issue["severity"])
        code = str(issue["code"])
        by_severity[severity] = by_severity.get(severity, 0) + 1
        by_code[code] = by_code.get(code, 0) + 1

    report = {
        "ok": not any(issue["severity"] == "error" for issue in issues),
        "project_dir": str(project_dir),
        "profile": profile_checks,
        "checks": {
            "project_dir": {"ok": project_ok, "path": str(project_dir)},
            "artifacts": artifact_checks,
            "optional_dependencies": dependency_checks,
        },
        "issue_count": len(issues),
        "summary": {
            "by_severity": by_severity,
            "by_code": by_code,
        },
        "issues": issues,
    }
    _write_payload(
        report,
        payload_name="preflight",
        path=None,
        fields=fields,
        ndjson=ndjson,
        quiet=quiet,
        compact=compact,
    )
    if not report["ok"]:
        return 5
    return 0


def _cmd_docs(args: argparse.Namespace) -> int:
    compact = getattr(args, "compact", False)
    contract = _build_discovery_contract()
    indent = None if compact else 2
    try:
        docs_format, method_id = _parse_docs_value(args.target)
    except ValueError as exc:
        raise CliError(code="ARGUMENT_ERROR", message=str(exc), exit_code=2) from exc
    if method_id is None:
        if docs_format == "json":
            print(json.dumps(contract, indent=indent))
        else:
            print(_discovery_markdown(contract))
        return 0
    if method_id == "fields":
        allowlists = _load_field_allowlists()
        if docs_format == "json":
            print(json.dumps(
                {k: sorted(v) for k, v in allowlists.items()}, indent=indent,
            ))
        else:
            for name in sorted(allowlists):
                print(f"\n## {name}")
                for f in sorted(allowlists[name]):
                    print(f"  {f}")
        return 0
    if method_id.startswith("fields:"):
        payload_name = method_id[len("fields:"):].strip()
        allowlists = _load_field_allowlists()
        allowlist = allowlists.get(payload_name)
        if allowlist is None:
            available = ", ".join(sorted(allowlists))
            raise CliError(
                code="ARGUMENT_ERROR",
                message=f"Unknown payload: {payload_name}. Available: {available}",
                exit_code=2,
            )
        if docs_format == "json":
            print(json.dumps(sorted(allowlist), indent=indent))
        else:
            print(f"## {payload_name}")
            for f in sorted(allowlist):
                print(f"  {f}")
        return 0
    if method_id.startswith("schema:"):
        schema_id = method_id[len("schema:"):].strip()
        schema = contract["schemas"].get(schema_id)
        if schema is None:
            raise CliError(
                code="ARGUMENT_ERROR",
                message=f"Unknown schema: {schema_id}",
                exit_code=2,
            )
        if docs_format == "json":
            print(json.dumps({"id": schema_id, "schema": schema}, indent=indent))
        else:
            print(_schema_markdown(schema_id, schema))
        return 0
    method = _select_discovery_method(contract, method_id)
    if method is None:
        raise CliError(
            code="ARGUMENT_ERROR",
            message=f"Unknown method: {method_id}",
            exit_code=2,
        )
    method_payload = _resolve_method_contract(contract, method)
    if docs_format == "json":
        print(json.dumps(method_payload, indent=indent))
    else:
        print(_method_markdown(method_payload))
    return 0


def _cmd_version(args: argparse.Namespace) -> int:
    compact = getattr(args, "compact", False)
    indent = None if compact else 2
    print(json.dumps({"contract_version": CONTRACT_VERSION}, indent=indent))
    return 0


_COMMAND_DISPATCH: dict[str, Any] = {
    "extract": _cmd_extract,
    "render": _cmd_render,
    "apply": _cmd_apply,
    "inspect": _cmd_inspect,
    "find": _cmd_find,
    "plan-inspect": _cmd_plan_inspect,
    "validate": _cmd_validate,
    "lint": _cmd_lint,
    "qa": _cmd_qa,
    "edit": _cmd_edit,
    "transform": _cmd_transform,
    "repair": _cmd_repair,
    "preflight": _cmd_preflight,
    "docs": _cmd_docs,
    "version": _cmd_version,
}


def _main(parser: argparse.ArgumentParser) -> int:
    args = parser.parse_args()
    command = args.command
    if command is None:
        parser.print_help()
        return 0
    handler = _COMMAND_DISPATCH.get(command)
    if handler is None:
        parser.error(f"Unknown command: {command}")
    return handler(args)


def main() -> int:
    parser = _build_parser()
    try:
        return _main(parser)
    except CliError as exc:
        return _emit_error(exc)
    except ValueError as exc:
        return _emit_error(
            CliError(
                code="INPUT_VALIDATION_ERROR",
                message=str(exc),
                hint="Check request payload/schema via `slides docs json`.",
                exit_code=2,
            )
        )
    except Exception as exc:  # noqa: BLE001
        return _emit_error(
            CliError(
                code="RUNTIME_ERROR",
                message=str(exc),
                hint="Inspect inputs and retry. Use `slides docs method:<id>` for contracts.",
                exit_code=10,
            )
        )


if __name__ == "__main__":
    raise SystemExit(main())
