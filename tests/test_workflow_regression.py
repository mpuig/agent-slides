from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation

from slides_cli.cli import main


def _run_cli(monkeypatch: Any, argv: list[str]) -> int:
    monkeypatch.setattr("sys.argv", argv)
    return main()


def _build_template(path: Path) -> None:
    prs = PptxPresentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "Neutral Template"
    title.placeholders[1].text = "Regression fixture"

    content = prs.slides.add_slide(prs.slide_layouts[1])
    content.shapes.title.text = "Evidence slide"
    content.placeholders[1].text = "Bullet one\nBullet two"
    prs.save(str(path))


def test_extract_render_qa_workflow_regression(tmp_path: Path, monkeypatch: Any) -> None:
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)

    project_dir = tmp_path / "project"
    base_template_path = project_dir / "base_template.pptx"
    assert _run_cli(
        monkeypatch,
        [
            "slides",
            "extract",
            str(template_path),
            "--output-dir",
            str(project_dir),
            "--base-template-out",
            str(base_template_path),
            "--compact",
        ],
    ) == 0

    design_profile = {
        "name": "workflow-regression",
        "template_path": "base_template.pptx",
        "content_layout_catalog_path": "content_layout_catalog.json",
        "archetypes_catalog_path": "archetypes.json",
        "asset_roots": [str(project_dir.resolve())],
        "default_font_size_pt": 14,
    }
    profile_path = project_dir / "design-profile.json"
    profile_path.write_text(json.dumps(design_profile), encoding="utf-8")

    slides_document = {
        "plan": {
            "deck_title": "Quarterly review",
            "brief": "Quarterly review for leadership",
            "slides": [
                {
                    "slide_number": 1,
                    "story_role": "title",
                    "archetype_id": "title_slide",
                    "action_title": "Quarterly review shows steady progress",
                    "key_points": ["Prepared for leadership review"],
                },
                {
                    "slide_number": 2,
                    "story_role": "evidence",
                    "archetype_id": "bar_chart",
                    "action_title": "Revenue increased across the last four quarters",
                    "key_points": ["Q1 10", "Q2 20", "Q3 30", "Q4 40"],
                    "source_note": "Test dataset",
                },
            ],
        }
    }
    slides_path = project_dir / "slides.json"
    slides_path.write_text(json.dumps(slides_document), encoding="utf-8")

    assert _run_cli(
        monkeypatch,
        [
            "slides",
            "preflight",
            "--project-dir",
            str(project_dir),
            "--profile",
            str(profile_path),
            "--compact",
        ],
    ) == 0

    output_path = project_dir / "output.pptx"
    assert _run_cli(
        monkeypatch,
        [
            "slides",
            "render",
            "--slides-json",
            f"@{slides_path}",
            "--profile",
            str(profile_path),
            "--output",
            str(output_path),
            "--compact",
        ],
    ) == 0

    lint_path = project_dir / "lint.json"
    assert _run_cli(
        monkeypatch,
        [
            "slides",
            "lint",
            str(output_path),
            "--profile",
            str(profile_path),
            "--slides-json",
            f"@{slides_path}",
            "--out",
            str(lint_path),
            "--compact",
        ],
    ) == 0

    qa_path = project_dir / "qa.json"
    assert _run_cli(
        monkeypatch,
        [
            "slides",
            "qa",
            str(output_path),
            "--profile",
            str(profile_path),
            "--slides-json",
            f"@{slides_path}",
            "--out",
            str(qa_path),
            "--compact",
        ],
    ) == 0

    assert output_path.exists()
    lint_report = json.loads(lint_path.read_text(encoding="utf-8"))
    qa_report = json.loads(qa_path.read_text(encoding="utf-8"))

    assert lint_report["issue_count"] == 0
    assert lint_report["summary"]["by_slide_class"] == {"data": 1, "structural": 1}
    assert qa_report["ok"] is True
    assert qa_report["checks"]["validate"]["ok"] is True
    assert qa_report["checks"]["lint"]["ok"] is True
    assert qa_report["checks"]["assets"]["ok"] is True
