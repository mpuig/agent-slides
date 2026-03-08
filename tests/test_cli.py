from __future__ import annotations

import json
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from slides_cli import cli as cli_module
from slides_cli.agentic import (
    DeckPlan,
    DesignProfile,
    SlidePlan,
    TemplateStyle,
    _compile_slide,
    compile_plan_to_operations,
)
from slides_cli.cli import main

_CUSTOM_ICON_XML = """\
<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="2" name="custom_icon"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="0" y="0"/>
      <a:ext cx="914400" cy="914400"/>
    </a:xfrm>
    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="2F6BFF"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>
"""


def test_cli_apply_ops(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "CLI works",
                "left": 0.5,
                "top": 0.5,
                "width": 3,
                "height": 1,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "apply",
            "--ops-json",
            f"@{ops_path}",
            "--output",
            str(output),
        ],
    )
    code = main()
    assert code == 0
    assert output.exists()


def test_cli_apply_add_icon_with_icon_pack_dir(tmp_path: Path, monkeypatch: Any) -> None:
    icon_dir = tmp_path / "icons"
    icon_dir.mkdir()
    (icon_dir / "custom_icon.xml").write_text(_CUSTOM_ICON_XML, encoding="utf-8")
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_icon",
                "slide_index": 0,
                "icon_name": "custom_icon",
                "left": 1.0,
                "top": 1.0,
                "size": 0.8,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "apply",
            "--ops-json",
            f"@{ops_path}",
            "--icon-pack-dir",
            str(icon_dir),
            "--output",
            str(output),
        ],
    )
    assert main() == 0
    assert output.exists()


def test_cli_apply_add_icon_uses_builtin_generic_pack(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_icon",
                "slide_index": 0,
                "icon_name": "generic_circle",
                "left": 1.0,
                "top": 1.0,
                "size": 0.8,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(output)],
    )
    assert main() == 0
    assert output.exists()


def test_cli_apply_add_icon_with_env_icon_pack_dir(tmp_path: Path, monkeypatch: Any) -> None:
    icon_dir = tmp_path / "icons"
    icon_dir.mkdir()
    (icon_dir / "custom_icon.xml").write_text(_CUSTOM_ICON_XML, encoding="utf-8")
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_icon",
                "slide_index": 0,
                "icon_name": "custom_icon",
                "left": 1.0,
                "top": 1.0,
                "size": 0.8,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"
    monkeypatch.setenv("SLIDES_ICON_PACK_DIR", str(icon_dir))

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(output)],
    )
    assert main() == 0
    assert output.exists()


def test_cli_summary_and_validate(tmp_path: Path, monkeypatch: Any) -> None:
    output = tmp_path / "out.pptx"
    Presentation().save(str(output))
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "validate", str(output)],
    )
    code = main()
    assert code == 0


def test_cli_agent_io_default_suppresses_apply_report(
    tmp_path: Path, monkeypatch: Any, capsys: Any
) -> None:
    ops = {"operations": [{"op": "add_slide", "layout_index": 6}]}
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(output)],
    )
    assert main() == 0
    assert capsys.readouterr().out == ""


def test_cli_agent_io_verbose_emits_apply_report(
    tmp_path: Path, monkeypatch: Any, capsys: Any
) -> None:
    ops = {"operations": [{"op": "add_slide", "layout_index": 6}]}
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(output), "--verbose"],
    )
    assert main() == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["ok"] is True


def test_cli_docs_replaces_schema_mode(monkeypatch: Any) -> None:
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "docs",
            "json",
        ],
    )
    code = main()
    assert code == 0


def test_cli_docs_method_render_json(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "docs",
            "method:render",
        ],
    )
    code = main()
    assert code == 0
    out = capsys.readouterr().out
    payload = json.loads(out)
    assert payload["id"] == "render"
    assert payload["supports_dry_run"] is True


def test_cli_docs_schema_template_layout(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "docs",
            "schema:template-layout",
        ],
    )
    code = main()
    assert code == 0
    out = capsys.readouterr().out
    payload = json.loads(out)
    assert payload["id"] == "template-layout"
    assert payload["schema"]["title"] == "TemplateLayoutCatalog"


def test_cli_docs_schema_markdown(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "docs",
            "schema:content-layout:markdown",
        ],
    )
    code = main()
    assert code == 0
    out = capsys.readouterr().out
    assert "# slides Schema: content-layout" in out


def test_cli_extract_layout_catalog_writes_three_contracts(
    tmp_path: Path, monkeypatch: Any
) -> None:
    template = tmp_path / "template.pptx"
    Presentation().save(str(template))
    template_out = tmp_path / "template_layout.json"
    content_out = tmp_path / "content_layout.json"
    archetypes_out = tmp_path / "archetypes.json"

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "extract",
            str(template),
            "--output-dir",
            str(tmp_path),
            "--template-out",
            str(template_out),
            "--content-layout-out",
            str(content_out),
            "--archetypes-out",
            str(archetypes_out),
        ],
    )
    code = main()
    assert code == 0
    assert template_out.exists()
    assert content_out.exists()
    assert archetypes_out.exists()
    payload = json.loads(archetypes_out.read_text(encoding="utf-8"))
    assert payload["schema_version"] in {"1.0", "1.1", "1.2", "1.3"}
    assert "archetypes" in payload


def test_cli_extract_layout_catalog_from_input_generates_manifest_and_analysis(
    tmp_path: Path, monkeypatch: Any
) -> None:
    source = tmp_path / "source.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Agenda"
    prs.save(str(source))
    template_out = tmp_path / "template_layout.json"
    content_out = tmp_path / "content_layout.json"
    archetypes_out = tmp_path / "archetypes.json"
    manifest_out = tmp_path / "slides_manifest.json"
    analysis_out = tmp_path / "slide_analysis.json"

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "extract",
            str(source),
            "--output-dir",
            str(tmp_path),
            "--template-out",
            str(template_out),
            "--content-layout-out",
            str(content_out),
            "--archetypes-out",
            str(archetypes_out),
            "--slides-manifest-out",
            str(manifest_out),
            "--slide-analysis-out",
            str(analysis_out),
        ],
    )
    code = main()
    assert code == 0
    assert template_out.exists()
    assert content_out.exists()
    assert archetypes_out.exists()
    assert manifest_out.exists()
    assert analysis_out.exists()
    manifest_payload = json.loads(manifest_out.read_text(encoding="utf-8"))
    assert manifest_payload["slide_count"] >= 1
    analysis_payload = json.loads(analysis_out.read_text(encoding="utf-8"))
    assert analysis_payload["slide_count"] >= 1


def test_cli_extract_layout_catalog_output_dir_defaults(tmp_path: Path, monkeypatch: Any) -> None:
    template = tmp_path / "template.pptx"
    Presentation().save(str(template))
    out_dir = tmp_path / "extract-out"

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "extract",
            str(template),
            "--output-dir",
            str(out_dir),
        ],
    )
    code = main()
    assert code == 0
    assert (out_dir / "template_catalog.json").exists()
    assert (out_dir / "content_layout_catalog.json").exists()
    assert (out_dir / "archetypes.json").exists()
    assert (out_dir / "resolved_manifest.json").exists()
    assert (out_dir / "slides_manifest.json").exists()
    assert (out_dir / "slide_analysis.json").exists()


def test_cli_fail_on_warning(tmp_path: Path, monkeypatch: Any) -> None:
    output = tmp_path / "out.pptx"
    Presentation().save(str(output))
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "validate", str(output), "--fail-on-warning"],
    )
    code = main()
    assert code == 4


def test_cli_list_placeholders(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {"operations": [{"op": "add_slide", "layout_index": 0}]}
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(output)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "inspect", str(output), "--placeholders", "0"],
    )
    code = main()
    assert code == 0


def test_cli_render_and_lint(tmp_path: Path, monkeypatch: Any) -> None:
    profile = tmp_path / "profile.json"
    profile.write_text(
        json.dumps(
            {
                "name": "test",
                "max_bullets_per_slide": 8,
                "required_sections": ["executive summary", "recommendation"],
            }
        ),
        encoding="utf-8",
    )
    plan = {
        "plan": {
            "deck_title": "AI Adoption Plan",
            "brief": "Improve productivity, reduce cycle time",
            "slides": [
                {
                    "slide_number": 1,
                    "story_role": "context",
                    "archetype_id": "content_text",
                    "action_title": "AI adoption accelerates productivity",
                    "key_points": ["Improve productivity", "Reduce cycle time"],
                }
            ],
            "assumptions": [],
        }
    }
    slides = tmp_path / "slides.json"
    slides.write_text(json.dumps(plan), encoding="utf-8")
    output = tmp_path / "out.pptx"
    lint = tmp_path / "lint.json"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "render", "--slides-json", f"@{slides}", "--profile", str(profile),
         "--output", str(output)],
    )
    assert main() == 0
    assert output.exists()

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "lint", str(output), "--profile", str(profile), "--out", str(lint)],
    )
    code = main()
    assert code in (0, 5)
    assert lint.exists()
    payload = json.loads(lint.read_text(encoding="utf-8"))
    assert "issues" in payload


def test_cli_inspect_and_find(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Pricing model and margin outlook",
                "left": 0.5,
                "top": 0.5,
                "width": 5.5,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    output = tmp_path / "out.pptx"
    inspect_path = tmp_path / "inspect.json"
    find_path = tmp_path / "find.json"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(output)],
    )
    assert main() == 0
    assert output.exists()

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "inspect",
            str(output),
            "--out",
            str(inspect_path),
        ],
    )
    assert main() == 0
    inspect_payload = json.loads(inspect_path.read_text(encoding="utf-8"))
    assert inspect_payload["summary"]["slide_count"] >= 1

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "find",
            str(output),
            "--query",
            "margin",
            "--out",
            str(find_path),
        ],
    )
    assert main() == 0
    find_payload = json.loads(find_path.read_text(encoding="utf-8"))
    assert find_payload["results"]


def test_cli_plan_inspect_content_only_paged(tmp_path: Path, monkeypatch: Any) -> None:
    slides_doc = {
        "plan": {
            "deck_title": "Plan",
            "brief": "Brief",
            "slides": [
                {
                    "slide_number": 1,
                    "story_role": "title_slide",
                    "archetype_id": "title_slide",
                    "action_title": "Title",
                    "key_points": [],
                },
                {
                    "slide_number": 2,
                    "story_role": "analysis",
                    "archetype_id": "two_column",
                    "action_title": "Analysis message",
                    "key_points": ["a", "b"],
                },
            ],
            "assumptions": [],
        },
        "ops": {"operations": []},
    }
    slides_path = tmp_path / "slides.json"
    out = tmp_path / "plan_inspect.json"
    slides_path.write_text(json.dumps(slides_doc), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "plan-inspect",
            "--slides-json",
            f"@{slides_path}",
            "--content-only",
            "--page-size",
            "1",
            "--out",
            str(out),
        ],
    )
    assert main() == 0
    payload = json.loads(out.read_text(encoding="utf-8"))
    assert payload["summary"]["slide_count"] == 1
    assert payload["summary"]["content_slide_count"] == 1
    assert payload["slides"][0]["story_role"] == "analysis"


def test_cli_edit_and_transform(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Roadmap draft with legacy plan",
                "left": 0.5,
                "top": 0.5,
                "width": 7.0,
                "height": 1.0,
            },
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "- Phase 1\\n- Phase 2\\n- Phase 3",
                "left": 0.8,
                "top": 1.6,
                "width": 8.0,
                "height": 2.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    deck = tmp_path / "deck.pptx"
    edited = tmp_path / "edited.pptx"
    transformed = tmp_path / "timeline.pptx"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0
    assert deck.exists()

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "edit",
            str(deck),
            "--query",
            "legacy plan",
            "--replacement",
            "target-state plan",
            "--output",
            str(edited),
        ],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "find", str(edited), "--query", "Recovered selector"],
    )
    assert main() == 0
    assert edited.exists()

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "transform",
            str(edited),
            "--slide",
            "0",
            "--to",
            "timeline",
            "--output",
            str(transformed),
        ],
    )
    assert main() == 0
    assert transformed.exists()


def test_cli_selector_edit_by_slide_id_and_shape_id(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Pricing model with legacy assumptions",
                "left": 0.5,
                "top": 0.5,
                "width": 8.0,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    deck = tmp_path / "deck.pptx"
    index_path = tmp_path / "index.json"
    edited = tmp_path / "edited.pptx"
    find_path = tmp_path / "find.json"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "inspect",
            str(deck),
            "--out",
            str(index_path),
        ],
    )
    assert main() == 0
    index_payload = json.loads(index_path.read_text(encoding="utf-8"))
    first_slide = index_payload["slides"][0]
    first_shape = first_slide["shapes"][0]

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "edit",
            str(deck),
            "--query",
            "legacy assumptions",
            "--replacement",
            "current assumptions",
            "--slide-id",
            first_slide["slide_id"],
            "--shape-id",
            str(first_shape["shape_id"]),
            "--output",
            str(edited),
        ],
    )
    assert main() == 0
    assert edited.exists()

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "find",
            str(edited),
            "--query",
            "current assumptions",
            "--out",
            str(find_path),
        ],
    )
    assert main() == 0
    find_payload = json.loads(find_path.read_text(encoding="utf-8"))
    assert find_payload["results"]


def test_cli_selector_edit_by_stable_uids(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Legacy assumption in UID path",
                "left": 0.5,
                "top": 0.5,
                "width": 8.0,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    inspect_path = tmp_path / "inspect.json"
    edited = tmp_path / "edited.pptx"
    find_path = tmp_path / "find.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "inspect",
            str(deck),
            "--out",
            str(inspect_path),
        ],
    )
    assert main() == 0
    payload = json.loads(inspect_path.read_text(encoding="utf-8"))
    slide_uid = payload["slides"][0]["slide_uid"]
    shape_uid = payload["slides"][0]["shapes"][0]["shape_uid"]

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "edit",
            str(deck),
            "--query",
            "Legacy assumption",
            "--replacement",
            "Current assumption",
            "--slide-uid",
            slide_uid,
            "--shape-uid",
            shape_uid,
            "--output",
            str(edited),
        ],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "inspect",
            str(edited),
            "--out",
            str(inspect_path),
        ],
    )
    assert main() == 0
    edited_payload = json.loads(inspect_path.read_text(encoding="utf-8"))
    assert edited_payload["slides"][0]["shapes"][0]["shape_uid"] == shape_uid

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "find",
            str(edited),
            "--query",
            "Current assumption",
            "--out",
            str(find_path),
        ],
    )
    assert main() == 0
    assert json.loads(find_path.read_text(encoding="utf-8"))["results"]
    with zipfile.ZipFile(edited, "r") as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml")
    assert b"slides-uid" in slide_xml
    assert b"slides-slideuid" in slide_xml


def test_compile_plan_includes_advanced_archetypes() -> None:
    plan = DeckPlan.model_validate({
        "deck_title": "Portfolio Strategy",
        "brief": "Retail 20, Enterprise 35, Public 15",
        "slides": [
            {
                "slide_number": 1,
                "story_role": "evidence",
                "archetype_id": "bar_chart",
                "action_title": "Revenue by segment",
                "key_points": ["Retail 20", "Enterprise 35", "Public 15"],
            }
        ],
        "assumptions": [],
    })
    ops = compile_plan_to_operations(plan).model_dump()
    op_names = [op["op"] for op in ops["operations"]]
    assert "add_bar_chart" in op_names


def test_cli_slides_document_with_archetype_library(tmp_path: Path, monkeypatch: Any) -> None:
    plan = {
        "deck_title": "Archetype Test",
        "brief": "test",
        "slides": [
            {
                "slide_number": 1,
                "story_role": "a",
                "archetype_id": "table",
                "action_title": "Table",
                "key_points": ["A 1", "B 2"],
            },
            {
                "slide_number": 2,
                "story_role": "b",
                "archetype_id": "line_chart",
                "action_title": "Line",
                "key_points": ["Jan 10", "Feb 20"],
            },
            {
                "slide_number": 3,
                "story_role": "c",
                "archetype_id": "pie_chart",
                "action_title": "Pie",
                "key_points": ["X 30", "Y 70"],
            },
            {
                "slide_number": 4,
                "story_role": "d",
                "archetype_id": "process_flow",
                "action_title": "Flow",
                "key_points": ["Plan", "Build", "Run"],
            },
            {
                "slide_number": 5,
                "story_role": "e",
                "archetype_id": "big_number",
                "action_title": "KPI",
                "key_points": ["Growth 28"],
            },
        ],
        "assumptions": [],
    }
    slides_path = tmp_path / "slides.json"
    slides_path.write_text(json.dumps({"plan": plan}), encoding="utf-8")
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "render",
            "--slides-json",
            f"@{slides_path}",
            "--dry-run",
        ],
    )
    assert main() == 0
    compiled = compile_plan_to_operations(DeckPlan.model_validate(plan)).model_dump()
    op_names = [
        op["op"] for op in compiled["operations"]
    ]
    assert "add_table" in op_names
    assert "add_line_chart" in op_names
    assert "add_pie_chart" in op_names


def test_cli_slides_document_with_extended_archetypes(tmp_path: Path, monkeypatch: Any) -> None:
    plan = {
        "deck_title": "Extended Archetypes",
        "brief": "test",
        "slides": [
            {
                "slide_number": 1,
                "story_role": "a",
                "archetype_id": "content_text",
                "action_title": "Narrative",
                "key_points": ["Why now: market shifted", "What changed", "What we do next"],
            },
            {
                "slide_number": 2,
                "story_role": "b",
                "archetype_id": "green_panel_text",
                "action_title": "Principles",
                "key_points": ["2.8x ROI", "$80M Investment", "24 mo Timeline", "Governance first"],
            },
            {
                "slide_number": 3,
                "story_role": "c",
                "archetype_id": "big_statement",
                "action_title": "We should accelerate now",
                "key_points": ["Delay destroys value"],
            },
        ],
        "assumptions": [],
    }
    slides_path = tmp_path / "slides.json"
    slides_path.write_text(json.dumps({"plan": plan}), encoding="utf-8")
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "render",
            "--slides-json",
            f"@{slides_path}",
            "--dry-run",
        ],
    )
    assert main() == 0
    compiled = compile_plan_to_operations(DeckPlan.model_validate(plan)).model_dump()
    operations = compiled["operations"]
    op_names = [op["op"] for op in operations]
    assert "set_slide_background" in op_names
    slide_ops = [op for op in operations if op["op"] == "add_slide"]
    assert any(op.get("layout_name") == "Green one third" for op in slide_ops)
    assert any(op.get("layout_name") == "Big statement green" for op in slide_ops)


def test_cli_slides_document_adds_chart_style_ops(tmp_path: Path, monkeypatch: Any) -> None:
    plan = {
        "deck_title": "Chart Ops",
        "brief": "test",
        "slides": [
            {
                "slide_number": 1,
                "story_role": "evidence",
                "archetype_id": "bar_chart",
                "action_title": "Category performance",
                "key_points": ["A 10", "B 20", "C 30"],
            }
        ],
        "assumptions": [],
    }
    slides_path = tmp_path / "slides.json"
    slides_path.write_text(json.dumps({"plan": plan}), encoding="utf-8")
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "render",
            "--slides-json",
            f"@{slides_path}",
            "--dry-run",
        ],
    )
    assert main() == 0
    compiled = compile_plan_to_operations(DeckPlan.model_validate(plan)).model_dump()
    op_names = [
        op["op"] for op in compiled["operations"]
    ]
    assert "set_chart_axis_titles" in op_names
    assert "set_chart_axis_scale" in op_names


def test_cli_template_without_json_uses_template_defaults(
    tmp_path: Path, monkeypatch: Any
) -> None:
    template_dir = tmp_path / "tpl-no-json"
    template_dir.mkdir(parents=True, exist_ok=True)
    template_pptx = template_dir / "template.pptx"
    prs = Presentation()
    prs.save(str(template_pptx))
    profile = tmp_path / "profile-no-json.json"
    profile.write_text(json.dumps({"template_path": str(template_pptx)}), encoding="utf-8")
    plan = {
        "deck_title": "Template No JSON",
        "brief": "test",
        "slides": [
            {
                "slide_number": 1,
                "story_role": "context",
                "archetype_id": "content_text",
                "action_title": "Template default title",
                "key_points": ["Point A", "Point B", "Point C"],
            }
        ],
        "assumptions": [],
    }
    slides_path = tmp_path / "slides-no-json.json"
    slides_path.write_text(json.dumps({"plan": plan}), encoding="utf-8")
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "render",
            "--slides-json",
            f"@{slides_path}",
            "--profile",
            str(profile),
            "--dry-run",
        ],
    )
    assert main() == 0
    profile_model = DesignProfile.model_validate({"template_path": str(template_pptx)})
    operations = compile_plan_to_operations(
        DeckPlan.model_validate(plan),
        profile=profile_model,
    ).model_dump()["operations"]
    assert operations[0]["op"] == "add_slide"
    assert "title" in operations[0]["layout_name"].lower()
    assert any(
        op["op"] in {"set_semantic_text", "set_placeholder_text"}
        for op in operations
    )


def test_cli_profile_relative_template_path_resolves_from_profile_location(
    tmp_path: Path, monkeypatch: Any
) -> None:
    workspace = tmp_path / "workspace"
    profile_dir = workspace / "profiles"
    profile_dir.mkdir(parents=True, exist_ok=True)
    template = profile_dir / "template.pptx"
    Presentation().save(str(template))

    profile = profile_dir / "profile.json"
    profile.write_text(json.dumps({"template_path": "template.pptx"}), encoding="utf-8")

    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "relative path test",
                "left": 0.5,
                "top": 0.5,
                "width": 3.5,
                "height": 1.0,
            },
        ]
    }
    ops_path = workspace / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    out_path = workspace / "out.pptx"

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "apply",
            "--ops-json",
            f"@{ops_path}",
            "--profile",
            str(profile),
            "--output",
            str(out_path),
        ],
    )
    assert main() == 0
    assert out_path.exists()


def test_compile_plan_uses_template_layout_index_from_catalog(tmp_path: Path) -> None:
    template_pptx = tmp_path / "template.pptx"
    prs = Presentation()
    prs.save(str(template_pptx))
    reloaded = Presentation(str(template_pptx))
    layout_names = [str(getattr(layout, "name", "") or "") for layout in reloaded.slide_layouts]
    target_name = "Title and Content"
    assert target_name in layout_names
    target_index = layout_names.index(target_name)

    catalog = {
        "schema_version": "1.0",
        "layout_content_matrix": {
            "content_text": {
                target_name: {"status": "preferred"},
            }
        },
    }
    catalog_path = tmp_path / "content-layout.json"
    catalog_path.write_text(json.dumps(catalog), encoding="utf-8")

    plan = DeckPlan.model_validate(
        {
            "deck_title": "Template Catalog Index",
            "brief": "test",
            "slides": [
                {
                    "slide_number": 1,
                    "story_role": "context",
                    "archetype_id": "content_text",
                    "action_title": "Use preferred layout by name with matching index",
                    "key_points": ["A", "B"],
                }
            ],
            "assumptions": [],
        }
    )
    profile = DesignProfile.model_validate(
        {
            "template_path": str(template_pptx),
            "content_layout_catalog_path": str(catalog_path),
        }
    )

    ops = compile_plan_to_operations(plan, profile=profile).model_dump()["operations"]
    add_slide = ops[0]
    assert add_slide["op"] == "add_slide"
    assert add_slide["layout_name"] == target_name
    assert add_slide["layout_index"] == target_index


def test_cli_lint_chart_style_contract(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_bar_chart",
                "slide_index": 0,
                "categories": ["A", "B"],
                "series": [["Value", [10, 20]]],
                "left": 1.0,
                "top": 1.4,
                "width": 7.0,
                "height": 4.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    profile = tmp_path / "profile.json"
    deck = tmp_path / "deck.pptx"
    lint = tmp_path / "lint.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    profile.write_text(
        json.dumps({"name": "chart-contract", "enforce_chart_style_contract": True}),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--ops-json",
            f"@{ops_path}",
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "CHART_STYLE_UNSPECIFIED" for issue in issues)


def test_compile_slide_agenda_suppresses_title_on_agenda_layout() -> None:
    plan = DeckPlan.model_validate(
        {"deck_title": "Deck", "brief": "b", "slides": [], "assumptions": []}
    )
    slide = SlidePlan.model_validate(
        {
            "slide_number": 1,
            "story_role": "agenda",
            "archetype_id": "agenda",
            "action_title": "Agenda: should not overwrite",
            "key_points": ["A", "B", "C"],
        }
    )
    style = TemplateStyle(
        default_layout="Title and Text",
        use_placeholders=True,
        layout_names=["Agenda Full Width Overview"],
        layout_index_by_name={"agenda full width overview": 3},
        placeholder_roles_by_layout={"agenda full width overview": {"title", "body"}},
        placeholder_indices_by_layout={"agenda full width overview": {"title": [0], "body": [1]}},
        placeholder_boxes_by_layout={},
    )
    ops = _compile_slide(
        slide_index=0,
        slide=slide,
        plan=plan,
        template_style=style,
        layout_name="Agenda Full Width Overview",
    )
    assert not any(
        op.get("op") == "set_placeholder_text" and op.get("placeholder_idx") == 0 for op in ops
    )
    assert any(
        op.get("op") == "set_placeholder_text" and op.get("placeholder_idx") == 1 for op in ops
    )


def test_compile_slide_clears_default_body_placeholder_for_custom_layout() -> None:
    plan = DeckPlan.model_validate(
        {"deck_title": "Deck", "brief": "b", "slides": [], "assumptions": []}
    )
    slide = SlidePlan.model_validate(
        {
            "slide_number": 1,
            "story_role": "evidence",
            "archetype_id": "bar_chart",
            "action_title": "Evidence title",
            "key_points": ["A 10", "B 20", "C 30"],
        }
    )
    style = TemplateStyle(
        default_layout="Title and Text",
        use_placeholders=True,
        layout_names=["Title and Text"],
        layout_index_by_name={"title and text": 2},
        placeholder_roles_by_layout={"title and text": {"title", "body"}},
        placeholder_indices_by_layout={"title and text": {"title": [0], "body": [1]}},
        placeholder_boxes_by_layout={"title and text": {"body": [(1.0, 1.5, 8.0, 4.5)]}},
    )
    ops = _compile_slide(
        slide_index=0,
        slide=slide,
        plan=plan,
        template_style=style,
        layout_name="Title and Text",
    )
    assert any(
        op.get("op") == "set_placeholder_text"
        and op.get("placeholder_idx") == 1
        and op.get("text") == " "
        for op in ops
    )
    assert any(op.get("op") == "add_bar_chart" for op in ops)


def test_compile_slide_end_slide_adds_no_extra_content() -> None:
    plan = DeckPlan.model_validate(
        {"deck_title": "Deck", "brief": "b", "slides": [], "assumptions": []}
    )
    slide = SlidePlan.model_validate(
        {
            "slide_number": 1,
            "story_role": "next_steps",
            "archetype_id": "end_slide",
            "action_title": "Thank you",
            "key_points": ["Step 1", "Step 2"],
        }
    )
    style = TemplateStyle(
        default_layout="Title and Text",
        use_placeholders=True,
        layout_names=["End"],
        layout_index_by_name={"end": 9},
        placeholder_roles_by_layout={"end": {"title", "body"}},
        placeholder_indices_by_layout={"end": {"title": [0], "body": [1]}},
        placeholder_boxes_by_layout={},
    )
    ops = _compile_slide(
        slide_index=0,
        slide=slide,
        plan=plan,
        template_style=style,
        layout_name="End",
    )
    assert ops == []


def test_compile_slide_end_layout_without_placeholders_adds_no_text() -> None:
    plan = DeckPlan.model_validate(
        {"deck_title": "Deck", "brief": "b", "slides": [], "assumptions": []}
    )
    slide = SlidePlan.model_validate(
        {
            "slide_number": 1,
            "story_role": "next_steps",
            "archetype_id": "end_slide",
            "action_title": "Thank you",
            "key_points": ["Step 1", "Step 2"],
        }
    )
    style = TemplateStyle(
        default_layout="Title and Text",
        use_placeholders=True,
        layout_names=["End"],
        layout_index_by_name={"end": 9},
        placeholder_roles_by_layout={"end": set()},
        placeholder_indices_by_layout={"end": {}},
        placeholder_boxes_by_layout={"end": {}},
    )
    ops = _compile_slide(
        slide_index=0,
        slide=slide,
        plan=plan,
        template_style=style,
        layout_name="End",
    )
    assert ops == []


def test_cli_edit_shape_uid_legacy_selector_recovery(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Legacy selector text",
                "left": 0.5,
                "top": 0.5,
                "width": 8.0,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    inspect = tmp_path / "inspect.json"
    edited = tmp_path / "edited.pptx"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "inspect", str(deck), "--out", str(inspect)],
    )
    assert main() == 0
    payload = json.loads(inspect.read_text(encoding="utf-8"))
    slide_uid = payload["slides"][0]["slide_uid"]
    shape_id = payload["slides"][0]["shapes"][0]["shape_id"]
    legacy_selector = f"{slide_uid}::shape-{shape_id}-1"
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "edit",
            str(deck),
            "--query",
            "Legacy selector",
            "--replacement",
            "Recovered selector",
            "--slide-uid",
            slide_uid,
            "--shape-uid",
            legacy_selector,
            "--output",
            str(edited),
        ],
    )
    assert main() == 0

def test_cli_lint_geometry_tolerance(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Near edge",
                "left": 0.0,
                "top": 0.0,
                "width": 4.0,
                "height": 1.0,
            },
        ]
    }
    profile = {
        "name": "strict-geometry",
        "min_margin_in": 0.5,
    }
    ops_path = tmp_path / "ops.json"
    profile_path = tmp_path / "profile.json"
    deck = tmp_path / "deck.pptx"
    lint_path = tmp_path / "lint.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile_path),
            "--out",
            str(lint_path),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint_path.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "SHAPE_OUT_OF_BOUNDS" for issue in issues)


def test_cli_lint_theme_font_inheritance(tmp_path: Path, monkeypatch: Any) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4.0), Inches(1.0))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Manual override font"
    run.font.name = "Arial"
    deck = tmp_path / "font_override.pptx"
    prs.save(str(deck))

    profile = tmp_path / "profile.json"
    lint = tmp_path / "lint.json"
    profile.write_text(
        json.dumps({"name": "font-theme", "enforce_theme_font_inheritance": True}),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "FONT_INHERITANCE_VIOLATION" for issue in issues)


def test_cli_lint_template_size_conformance(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Template check",
                "left": 0.5,
                "top": 0.5,
                "width": 4.0,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    template = tmp_path / "template.pptx"
    profile = tmp_path / "profile.json"
    lint_path = tmp_path / "lint.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    prs.save(str(template))

    profile.write_text(
        json.dumps(
            {
                "name": "template-conformance",
                "template_path": str(template),
                "enforce_template_layouts": True,
                "template_tolerance_in": 0.01,
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint_path),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint_path.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "TEMPLATE_SIZE_MISMATCH" for issue in issues)


def test_cli_lint_template_geometry_conformance(tmp_path: Path, monkeypatch: Any) -> None:
    template_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Baseline box",
                "left": 0.8,
                "top": 1.2,
                "width": 4.0,
                "height": 1.0,
            },
        ]
    }
    deck_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Shifted box",
                "left": 5.2,
                "top": 4.3,
                "width": 4.0,
                "height": 1.0,
            },
        ]
    }
    template_ops_path = tmp_path / "template_ops.json"
    deck_ops_path = tmp_path / "deck_ops.json"
    template = tmp_path / "template.pptx"
    deck = tmp_path / "deck.pptx"
    profile = tmp_path / "profile.json"
    lint = tmp_path / "lint.json"
    template_ops_path.write_text(json.dumps(template_ops), encoding="utf-8")
    deck_ops_path.write_text(json.dumps(deck_ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{template_ops_path}", "--output", str(template)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{deck_ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    profile.write_text(
        json.dumps(
            {
                "name": "template-geometry",
                "template_path": str(template),
                "enforce_template_geometry": True,
                "geometry_tolerance_in": 0.1,
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "TEMPLATE_GEOMETRY_DRIFT" for issue in issues)


def test_cli_lint_template_shape_kind_mismatch(tmp_path: Path, monkeypatch: Any) -> None:
    template_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Template text",
                "left": 0.8,
                "top": 1.2,
                "width": 4.0,
                "height": 1.0,
            },
        ]
    }
    deck_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_bar_chart",
                "slide_index": 0,
                "categories": ["A", "B"],
                "series": [["Value", [10, 20]]],
                "left": 0.8,
                "top": 1.2,
                "width": 4.0,
                "height": 2.5,
            },
        ]
    }
    template_ops_path = tmp_path / "template_ops.json"
    deck_ops_path = tmp_path / "deck_ops.json"
    template = tmp_path / "template_kind.pptx"
    deck = tmp_path / "deck_kind.pptx"
    profile = tmp_path / "profile_kind.json"
    lint = tmp_path / "lint_kind.json"
    template_ops_path.write_text(json.dumps(template_ops), encoding="utf-8")
    deck_ops_path.write_text(json.dumps(deck_ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{template_ops_path}", "--output", str(template)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{deck_ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    profile.write_text(
        json.dumps(
            {
                "name": "template-kind",
                "template_path": str(template),
                "enforce_template_geometry": True,
                "template_shape_kind_tolerance": 0,
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "TEMPLATE_SHAPE_KIND_MISMATCH" for issue in issues)


def test_cli_lint_template_placeholder_mismatch(tmp_path: Path, monkeypatch: Any) -> None:
    template_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 0},
            {
                "op": "set_placeholder_text",
                "slide_index": 0,
                "placeholder_idx": 0,
                "text": "Baseline placeholder",
            },
        ]
    }
    deck_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 0},
            {
                "op": "set_placeholder_text",
                "slide_index": 0,
                "placeholder_idx": 0,
                "text": "Moved placeholder",
                "left": 2.8,
                "top": 2.6,
                "width": 6.0,
                "height": 1.0,
            },
        ]
    }
    template_ops_path = tmp_path / "template_ph_ops.json"
    deck_ops_path = tmp_path / "deck_ph_ops.json"
    template = tmp_path / "template_ph.pptx"
    deck = tmp_path / "deck_ph.pptx"
    profile = tmp_path / "profile_ph.json"
    lint = tmp_path / "lint_ph.json"
    template_ops_path.write_text(json.dumps(template_ops), encoding="utf-8")
    deck_ops_path.write_text(json.dumps(deck_ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{template_ops_path}", "--output", str(template)],
    )
    assert main() == 0
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{deck_ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    profile.write_text(
        json.dumps(
            {
                "name": "template-placeholder",
                "template_path": str(template),
                "enforce_template_placeholders": True,
                "enforce_template_placeholder_text": True,
                "placeholder_tolerance_in": 0.1,
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "TEMPLATE_PLACEHOLDER_MISMATCH" for issue in issues)


def test_cli_lint_no_transitions_and_animations(tmp_path: Path, monkeypatch: Any) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide._element.append(OxmlElement("p:transition"))
    slide._element.append(OxmlElement("p:timing"))
    deck = tmp_path / "motion.pptx"
    profile = tmp_path / "profile_motion.json"
    lint = tmp_path / "lint_motion.json"
    prs.save(str(deck))
    profile.write_text(
        json.dumps(
            {
                "name": "motion-contract",
                "enforce_no_transitions": True,
                "enforce_no_animations": True,
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "SLIDE_TRANSITION_NOT_ALLOWED" for issue in issues)
    assert any(issue["code"] == "SLIDE_ANIMATION_NOT_ALLOWED" for issue in issues)


def test_cli_lint_no_effects(tmp_path: Path, monkeypatch: Any) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.0), Inches(1.0))
    shape.text = "Effect sample"
    sp_pr = getattr(shape._element, "spPr", None)
    if sp_pr is None:
        sp_pr = OxmlElement("p:spPr")
        shape._element.append(sp_pr)
    effect_lst = OxmlElement("a:effectLst")
    effect_lst.append(OxmlElement("a:outerShdw"))
    sp_pr.append(effect_lst)
    deck = tmp_path / "effects.pptx"
    profile = tmp_path / "profile_effects.json"
    lint = tmp_path / "lint_effects.json"
    prs.save(str(deck))
    profile.write_text(
        json.dumps({"name": "effects-contract", "enforce_no_effects": True}),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    assert any(issue["code"] == "VISUAL_EFFECT_NOT_ALLOWED" for issue in issues)


def test_cli_lint_visual_basics(tmp_path: Path, monkeypatch: Any) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box1 = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(4.2), Inches(1.2))
    box1.text_frame.clear()
    p1 = box1.text_frame.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "Revenue bridge and margin drivers"

    box2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(4.0), Inches(1.1))
    box2.text = "Overlapping commentary on the same region"

    box3 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(4.5), Inches(1.0))
    box3.text_frame.clear()
    p3 = box3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Low contrast line"
    run3.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

    deck = tmp_path / "visual_basics.pptx"
    profile = tmp_path / "profile_visual_basics.json"
    lint = tmp_path / "lint_visual_basics.json"
    prs.save(str(deck))
    profile.write_text(json.dumps({"name": "visual-basics"}), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "lint",
            str(deck),
            "--profile",
            str(profile),
            "--out",
            str(lint),
        ],
    )
    code = main()
    assert code in (0, 5)
    issues = json.loads(lint.read_text(encoding="utf-8"))["issues"]
    codes = {issue["code"] for issue in issues}
    assert "TEXT_OVERLAP_RISK" in codes
    assert "VISUAL_IMBALANCE_RISK" in codes
    assert "LOW_TEXT_CONTRAST_RISK" in codes


def test_cli_apply_inline_ops_json(tmp_path: Path, monkeypatch: Any) -> None:
    output = tmp_path / "out.pptx"
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Inline JSON execution",
                "left": 0.5,
                "top": 0.5,
                "width": 5.0,
                "height": 1.0,
            },
        ]
    }
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", json.dumps(ops), "--output", str(output)],
    )
    assert main() == 0
    assert output.exists()


def test_cli_docs_json(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr("sys.argv", ["slides", "docs", "json"])
    assert main() == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["name"] == "slides"
    assert "commands" in payload
    assert "methods" in payload
    assert "schemas" in payload
    assert "slides-document" in payload["schemas"]
    assert "ops" in payload["schemas"]


def test_cli_docs_markdown(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr("sys.argv", ["slides", "docs", "markdown"])
    assert main() == 0
    out = capsys.readouterr().out
    assert "# slides Discovery" in out
    assert "## Methods" in out
    assert "### render" in out


def test_cli_docs_method_json_alias(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr("sys.argv", ["slides", "docs", "method:render"])
    assert main() == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["id"] == "render"
    assert payload["request_schema_ref"] == "slides-document"
    assert payload["response_schema_ref"] == "operation-report"
    assert payload["request_schema"]["type"] == "object"
    assert payload["response_schema"]["type"] == "object"
    assert payload["supports_dry_run"] is True
    assert payload["supports_field_masks"] is False
    assert payload["supports_pagination"] is False


def test_cli_docs_method_markdown(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr("sys.argv", ["slides", "docs", "method:inspect:markdown"])
    assert main() == 0
    out = capsys.readouterr().out
    assert "# slides Method: inspect" in out
    assert "supports_field_masks: `true`" in out
    assert "supports_pagination: `true`" in out


def test_cli_docs_invalid_selector_returns_json_error(monkeypatch: Any, capsys: Any) -> None:
    monkeypatch.setattr("sys.argv", ["slides", "docs", "method:unknown"])
    assert main() == 2
    payload = json.loads(capsys.readouterr().err)
    assert payload["ok"] is False
    assert payload["error"]["code"] == "ARGUMENT_ERROR"
    assert "Unknown method" in payload["error"]["message"]


def test_cli_ndjson_output(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "Margin improvement",
                "left": 0.5,
                "top": 0.5,
                "width": 5.0,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    deck = tmp_path / "deck.pptx"
    find_ndjson = tmp_path / "find.ndjson"

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "find",
            str(deck),
            "--query",
            "margin",
            "--out",
            str(find_ndjson),
            "--ndjson",
        ],
    )
    assert main() == 0
    lines = [line for line in find_ndjson.read_text(encoding="utf-8").splitlines() if line]
    assert lines
    first = json.loads(lines[0])
    assert "slide_index" in first


def test_cli_rejects_invalid_slide_id(tmp_path: Path, monkeypatch: Any, capsys: Any) -> None:
    output = tmp_path / "out.pptx"
    Presentation().save(str(output))
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "edit",
            str(output),
            "--query",
            "x",
            "--replacement",
            "y",
            "--slide-id",
            "slide-1?fields=name",
            "--output",
            str(output),
        ],
    )
    assert main() == 2
    payload = json.loads(capsys.readouterr().err)
    assert payload["ok"] is False
    assert payload["error"]["code"] == "INPUT_VALIDATION_ERROR"
    assert "forbidden characters" in payload["error"]["message"]


def test_cli_rejects_invalid_field_mask(tmp_path: Path, monkeypatch: Any, capsys: Any) -> None:
    out = tmp_path / "out.pptx"
    Presentation().save(str(out))
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "inspect", str(out), "--summary", "--fields", "does_not_exist"],
    )
    assert main() == 2
    payload = json.loads(capsys.readouterr().err)
    assert payload["ok"] is False
    assert payload["error"]["code"] == "INPUT_VALIDATION_ERROR"
    assert "Unsupported field mask for summary" in payload["error"]["message"]


def test_cli_rejects_malformed_field_mask_policy(
    tmp_path: Path, monkeypatch: Any, capsys: Any
) -> None:
    out = tmp_path / "out.pptx"
    Presentation().save(str(out))
    bad_policy = tmp_path / "bad_field_masks.json"
    bad_policy.write_text(json.dumps({"summary": "not-an-array"}), encoding="utf-8")
    monkeypatch.setattr(cli_module, "FIELD_MASKS_PATH", bad_policy)
    cli_module._load_field_allowlists.cache_clear()
    monkeypatch.setattr(
        "sys.argv",
        ["slides", "inspect", str(out), "--summary", "--fields", "slide_count"],
    )
    assert main() == 2
    payload = json.loads(capsys.readouterr().err)
    assert payload["ok"] is False
    assert payload["error"]["code"] == "INPUT_VALIDATION_ERROR"
    assert "Field-mask policy entry 'summary' must be an array of strings" in payload["error"][
        "message"
    ]
    cli_module._load_field_allowlists.cache_clear()


def test_cli_inspect_pagination(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {"op": "add_slide", "layout_index": 6},
            {"op": "add_slide", "layout_index": 6},
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    inspect_json = tmp_path / "inspect.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "inspect",
            str(deck),
            "--page-size",
            "2",
            "--out",
            str(inspect_json),
        ],
    )
    assert main() == 0
    payload = json.loads(inspect_json.read_text(encoding="utf-8"))
    assert len(payload["slides"]) == 2
    assert payload["next_page_token"] == "2"


def test_cli_inspect_default_pagination(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {"op": "add_slide", "layout_index": 6},
            {"op": "add_slide", "layout_index": 6},
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    inspect_json = tmp_path / "inspect.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    monkeypatch.setattr(cli_module, "DEFAULT_PAGE_SIZE", 2)

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "inspect", str(deck), "--out", str(inspect_json)],
    )
    assert main() == 0
    payload = json.loads(inspect_json.read_text(encoding="utf-8"))
    assert len(payload["slides"]) == 2
    assert payload["next_page_token"] == "2"
    assert payload["page_size"] == 2


def test_cli_find_page_all_ndjson(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "margin one",
                "left": 0.5,
                "top": 0.5,
                "width": 4,
                "height": 1,
            },
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 1,
                "text": "margin two",
                "left": 0.5,
                "top": 0.5,
                "width": 4,
                "height": 1,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    find_out = tmp_path / "find.ndjson"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "find",
            str(deck),
            "--query",
            "margin",
            "--limit",
            "50",
            "--page-size",
            "1",
            "--page-all",
            "--ndjson",
            "--out",
            str(find_out),
        ],
    )
    assert main() == 0
    lines = [line for line in find_out.read_text(encoding="utf-8").splitlines() if line]
    assert len(lines) >= 2
    first_page = json.loads(lines[0])
    assert "results" in first_page
    assert first_page["page_size"] == 1


def test_cli_find_default_pagination(tmp_path: Path, monkeypatch: Any) -> None:
    ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "margin one",
                "left": 0.5,
                "top": 0.5,
                "width": 4,
                "height": 1,
            },
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 1,
                "text": "margin two",
                "left": 0.5,
                "top": 0.5,
                "width": 4,
                "height": 1,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    deck = tmp_path / "deck.pptx"
    find_out = tmp_path / "find.json"
    ops_path.write_text(json.dumps(ops), encoding="utf-8")
    monkeypatch.setattr(cli_module, "DEFAULT_PAGE_SIZE", 1)

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "find",
            str(deck),
            "--query",
            "margin",
            "--limit",
            "50",
            "--out",
            str(find_out),
        ],
    )
    assert main() == 0
    payload = json.loads(find_out.read_text(encoding="utf-8"))
    assert len(payload["results"]) == 1
    assert payload["next_page_token"] == "1"
    assert payload["page_size"] == 1


def test_cli_assets_verify_and_qa_run(tmp_path: Path, monkeypatch: Any) -> None:
    image = tmp_path / "logo.png"
    image.write_bytes(b"not-a-real-image-but-existing-path")
    profile = tmp_path / "profile.json"
    profile.write_text(
        json.dumps(
            {
                "name": "runtime-contract",
                "asset_roots": [str(tmp_path)],
                "allowed_image_extensions": [".png"],
            }
        ),
        encoding="utf-8",
    )
    build_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_text",
                "slide_index": 0,
                "text": "QA gate deck",
                "left": 0.5,
                "top": 0.5,
                "width": 4.0,
                "height": 1.0,
            },
        ]
    }
    asset_ops = {
        "operations": [
            {"op": "add_slide", "layout_index": 6},
            {
                "op": "add_image",
                "slide_index": 0,
                "path": str(image),
                "left": 0.5,
                "top": 0.5,
                "width": 1.0,
                "height": 1.0,
            },
        ]
    }
    ops_path = tmp_path / "ops.json"
    build_ops_path = tmp_path / "build_ops.json"
    deck = tmp_path / "deck.pptx"
    assets = tmp_path / "assets.json"
    qa = tmp_path / "qa.json"
    ops_path.write_text(json.dumps(asset_ops), encoding="utf-8")
    build_ops_path.write_text(json.dumps(build_ops), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["slides", "apply", "--ops-json", f"@{build_ops_path}", "--output", str(deck)],
    )
    assert main() == 0

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "qa",
            str(deck),
            "--ops-json",
            f"@{ops_path}",
            "--profile",
            str(profile),
            "--out",
            str(assets),
        ],
    )
    assert main() == 0
    assets_payload = json.loads(assets.read_text(encoding="utf-8"))
    assert assets_payload["ok"] is True

    monkeypatch.setattr(
        "sys.argv",
        [
            "slides",
            "qa",
            str(deck),
            "--ops-json",
            f"@{ops_path}",
            "--profile",
            str(profile),
            "--out",
            str(qa),
        ],
    )
    assert main() == 0
    qa_payload = json.loads(qa.read_text(encoding="utf-8"))
    assert qa_payload["checks"]["assets"]["ok"] is True


def _make_pptx_with_icon_shape(path: Path, name: str = "TestIcon") -> None:
    """Create a .pptx with a freeform (custGeom) shape that looks like an icon."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Use python-pptx freeform builder to create a custGeom shape
    builder = slide.shapes.build_freeform(Inches(1), Inches(1))
    builder.add_line_segments(
        [(Inches(1.75), Inches(1)), (Inches(1.75), Inches(1.75)),
         (Inches(1), Inches(1.75)), (Inches(1), Inches(1))],
    )
    shape = builder.convert_to_shape()
    shape.name = name
    prs.save(str(path))


def test_extract_icons_from_pptx(tmp_path: Path, monkeypatch: Any) -> None:
    """Extract should find custGeom shapes and write them as icon XML."""
    source = tmp_path / "template.pptx"
    _make_pptx_with_icon_shape(source)

    output_dir = tmp_path / "extracted"
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides", "extract", str(source),
            "--output-dir", str(output_dir),
            "--verbose",
        ],
    )
    assert main() == 0

    icons_dir = output_dir / "icons"
    assert icons_dir.exists()
    xml_files = list(icons_dir.glob("*.xml"))
    assert len(xml_files) == 1
    assert xml_files[0].stem == "TestIcon"

    content = xml_files[0].read_text(encoding="utf-8")
    assert "custGeom" in content


def test_extract_icons_deduplicates(
    tmp_path: Path, monkeypatch: Any,
) -> None:
    """Same icon on two slides should produce only one XML file."""
    prs = Presentation()
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder = slide.shapes.build_freeform(Inches(1), Inches(1))
        builder.add_line_segments(
            [(Inches(1.75), Inches(1)), (Inches(1.75), Inches(1.75)),
             (Inches(1), Inches(1.75)), (Inches(1), Inches(1))],
        )
        shape = builder.convert_to_shape()
        shape.name = "DupeIcon"

    source = tmp_path / "dupes.pptx"
    prs.save(str(source))

    output_dir = tmp_path / "extracted"
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides", "extract", str(source),
            "--output-dir", str(output_dir),
            "--verbose",
        ],
    )
    assert main() == 0

    icons_dir = output_dir / "icons"
    xml_files = list(icons_dir.glob("*.xml"))
    assert len(xml_files) == 1


def test_extract_icons_skips_text_shapes(
    tmp_path: Path, monkeypatch: Any,
) -> None:
    """Shapes with text content should not be extracted as icons."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    builder = slide.shapes.build_freeform(Inches(1), Inches(1))
    builder.add_line_segments(
        [(Inches(1.75), Inches(1)), (Inches(1.75), Inches(1.75)),
         (Inches(1), Inches(1.75)), (Inches(1), Inches(1))],
    )
    shape = builder.convert_to_shape()
    shape.name = "NotAnIcon"
    shape.text_frame.paragraphs[0].text = "Hello World"

    source = tmp_path / "textshape.pptx"
    prs.save(str(source))

    output_dir = tmp_path / "extracted"
    monkeypatch.setattr(
        "sys.argv",
        [
            "slides", "extract", str(source),
            "--output-dir", str(output_dir),
            "--verbose",
        ],
    )
    assert main() == 0

    icons_dir = output_dir / "icons"
    if icons_dir.exists():
        assert len(list(icons_dir.glob("*.xml"))) == 0
